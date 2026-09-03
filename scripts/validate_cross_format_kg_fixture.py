#!/usr/bin/env python3
"""Validate the static cross-format KG contract and its frozen corpus manifest."""

from __future__ import annotations

import argparse
import hashlib
import json
import os
import re
import sys
import unicodedata
from datetime import date, datetime
from pathlib import Path, PurePosixPath
from typing import Any


SCHEMA_VERSION = "0.1"
DATASET_ID = "cross-format-kg-v0.1"
FIXTURE_SPEC = "fixture-spec.json"
BASELINE_CASES = "cases.jsonl"
EXPECTED_GRAPH = "gold/expected-graph.jsonl"
QA_CASES = "gold/qa-cases.jsonl"
CORPUS_MANIFEST = "corpus-manifest.json"
CORPUS_ROOT = "corpus"
BUILD_INPUT_GLOB = "corpus/**"
REQUIRED_SOURCE_PATHS = (
    "project-orion/01_ORION-27_案件定義書.docx",
    "project-orion/02_ORION-27_担当履歴.xlsx",
    "project-orion/03_ORION-27_体制計画_v1.pptx",
    "project-orion/04_ORION-27_体制計画_v2.pptx",
    "project-orion/05_社員ID対応表_署名済.pdf",
)
REQUIRED_FORMATS = {
    REQUIRED_SOURCE_PATHS[0]: "docx",
    REQUIRED_SOURCE_PATHS[1]: "xlsx",
    REQUIRED_SOURCE_PATHS[2]: "pptx",
    REQUIRED_SOURCE_PATHS[3]: "pptx",
    REQUIRED_SOURCE_PATHS[4]: "pdf",
}
XLSX_HEADERS = (
    "Record ID",
    "Project ID",
    "Work ID",
    "Role",
    "Assignee ID",
    "Valid From",
    "Valid To",
    "Status",
)
REQUIRED_SEMANTIC_RELATION_TYPES = {
    "ASSIGNED_TO",
    "IDENTIFIES_PERSON",
    "SUPERSEDES",
    "CONTRADICTS",
}
SHA256 = re.compile(r"^[0-9a-f]{64}$")


class FixtureContractError(ValueError):
    """Raised when a fixture contract or corpus snapshot is unsafe or invalid."""


def _require(condition: bool, message: str) -> None:
    if not condition:
        raise FixtureContractError(message)


def _object_without_duplicate_keys(pairs: list[tuple[str, Any]]) -> dict[str, Any]:
    result: dict[str, Any] = {}
    for key, value in pairs:
        if key in result:
            raise FixtureContractError(f"duplicate JSON key: {key}")
        result[key] = value
    return result


def _decode_json(text: str, label: str) -> Any:
    try:
        return json.loads(text, object_pairs_hook=_object_without_duplicate_keys)
    except json.JSONDecodeError as exc:
        raise FixtureContractError(
            f"{label}: invalid JSON at line {exc.lineno}, column {exc.colno}"
        ) from exc


def load_json(path: Path, label: str) -> dict[str, Any]:
    _require(path.is_file(), f"{label}: file is missing: {path}")
    value = _decode_json(path.read_text(encoding="utf-8"), label)
    _require(isinstance(value, dict), f"{label}: top level must be an object")
    return value


def load_jsonl(path: Path, label: str) -> list[dict[str, Any]]:
    _require(path.is_file(), f"{label}: file is missing: {path}")
    text = path.read_text(encoding="utf-8")
    _require(bool(text), f"{label}: JSONL must not be empty")
    records: list[dict[str, Any]] = []
    for line_number, line in enumerate(text.splitlines(), start=1):
        _require(bool(line.strip()), f"{label}:{line_number}: blank JSONL line")
        value = _decode_json(line, f"{label}:{line_number}")
        _require(
            isinstance(value, dict),
            f"{label}:{line_number}: JSONL record must be an object",
        )
        records.append(value)
    _require(bool(records), f"{label}: JSONL must contain at least one record")
    return records


def _canonical_json(value: Any) -> str:
    return json.dumps(
        value, ensure_ascii=False, sort_keys=True, separators=(",", ":")
    )


def _record_sha256(value: Any) -> str:
    return hashlib.sha256(_canonical_json(value).encode("utf-8")).hexdigest()


def _file_sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _relative_path(value: Any, label: str) -> str:
    _require(isinstance(value, str) and bool(value), f"{label}: non-empty path required")
    _require("\\" not in value, f"{label}: POSIX separators are required")
    normalized = unicodedata.normalize("NFC", value)
    _require(value == normalized, f"{label}: path must be NFC-normalized")
    path = PurePosixPath(value)
    _require(not path.is_absolute(), f"{label}: absolute path is forbidden")
    _require(
        bool(path.parts) and all(part not in {"", ".", ".."} for part in path.parts),
        f"{label}: unsafe relative path",
    )
    return path.as_posix()


def _string_list(value: Any, label: str, *, allow_empty: bool = False) -> list[str]:
    _require(isinstance(value, list), f"{label}: array required")
    _require(allow_empty or bool(value), f"{label}: non-empty array required")
    result: list[str] = []
    for index, item in enumerate(value):
        _require(
            isinstance(item, str) and bool(item),
            f"{label}[{index}]: non-empty string required",
        )
        result.append(item)
    _require(len(result) == len(set(result)), f"{label}: duplicate values")
    return result


def _positive_integer(value: Any, label: str) -> int:
    _require(
        isinstance(value, int) and not isinstance(value, bool) and value >= 1,
        f"{label}: positive integer required",
    )
    return value


def _validate_source_reference(
    value: Any,
    label: str,
    declared_paths: set[str],
) -> str:
    _require(isinstance(value, dict), f"{label}: object required")
    _require(
        set(value) == {"path", "locator", "selector"},
        f"{label}: expected only path, locator, selector",
    )
    path = _relative_path(value["path"], f"{label}.path")
    _require(path in declared_paths, f"{label}: undeclared corpus path: {path}")
    locator = value["locator"]
    _require(isinstance(locator, dict) and bool(locator), f"{label}.locator: non-empty object required")
    selector = value["selector"]
    _require(isinstance(selector, dict), f"{label}.selector: object required")
    _require(
        set(selector) == {"kind", "value"}
        and selector.get("kind") == "exact_phrase"
        and isinstance(selector.get("value"), str)
        and bool(selector["value"]),
        f"{label}.selector: exact_phrase with a non-empty value required",
    )
    return path


def _contains_forbidden_evidence_id(value: Any) -> bool:
    if isinstance(value, dict):
        if any(key in {"evidence_id", "generated_evidence_id"} for key in value):
            return True
        return any(_contains_forbidden_evidence_id(item) for item in value.values())
    if isinstance(value, list):
        return any(_contains_forbidden_evidence_id(item) for item in value)
    return False


def _validate_fixture_spec(spec: dict[str, Any]) -> set[str]:
    _require(spec.get("schema_version") == SCHEMA_VERSION, "fixture-spec: schema_version")
    _require(
        spec.get("record_type") == "cross_format_kg_fixture_spec",
        "fixture-spec: record_type",
    )
    _require(spec.get("dataset_id") == DATASET_ID, "fixture-spec: dataset_id")

    boundary = spec.get("input_boundary")
    _require(isinstance(boundary, dict), "fixture-spec.input_boundary: object required")
    _require(
        boundary.get("build_and_index_inputs") == [BUILD_INPUT_GLOB],
        "build/index input boundary must be exactly corpus/**",
    )
    evaluator_only = _string_list(
        boundary.get("evaluator_only_inputs"),
        "fixture-spec.input_boundary.evaluator_only_inputs",
    )
    _require(
        set(evaluator_only)
        == {
            FIXTURE_SPEC,
            BASELINE_CASES,
            EXPECTED_GRAPH,
            QA_CASES,
            CORPUS_MANIFEST,
        },
        "evaluator-only inputs must contain spec, manifest, cases, and both gold files",
    )
    forbidden_build = _string_list(
        boundary.get("forbidden_build_inputs"),
        "fixture-spec.input_boundary.forbidden_build_inputs",
    )
    _require(
        FIXTURE_SPEC in forbidden_build
        and CORPUS_MANIFEST in forbidden_build
        and BASELINE_CASES in forbidden_build
        and "gold/**" in forbidden_build,
        "build boundary must explicitly forbid spec, manifest, cases, and gold/**",
    )
    for key in (
        "gold_visible_to_extractor",
        "gold_visible_to_graph_builder",
        "gold_visible_to_answerer",
    ):
        _require(boundary.get(key) is False, f"fixture-spec.input_boundary.{key} must be false")

    corpus = spec.get("corpus_contract")
    _require(isinstance(corpus, dict), "fixture-spec.corpus_contract: object required")
    _require(corpus.get("root") == CORPUS_ROOT, "corpus root must be corpus")
    required_files = corpus.get("required_files")
    _require(isinstance(required_files, list), "corpus required_files: array required")
    _require(len(required_files) == 5, "corpus must declare exactly five source files")
    declared: list[str] = []
    for index, record in enumerate(required_files):
        label = f"fixture-spec.corpus_contract.required_files[{index}]"
        _require(isinstance(record, dict), f"{label}: object required")
        path = _relative_path(record.get("path"), f"{label}.path")
        declared.append(path)
        _require(
            record.get("format") == REQUIRED_FORMATS.get(path),
            f"{label}: fixed source format mismatch",
        )
        _require(
            isinstance(record.get("role"), str) and bool(record["role"]),
            f"{label}.role: non-empty string required",
        )
    _require(
        tuple(declared) == REQUIRED_SOURCE_PATHS,
        "corpus required_files must match the five fixed source paths in order",
    )

    graph = spec.get("graph_contract")
    qa = spec.get("qa_contract")
    _require(isinstance(graph, dict) and graph.get("gold") == EXPECTED_GRAPH, "graph gold path mismatch")
    _require(isinstance(qa, dict) and qa.get("gold") == QA_CASES, "QA gold path mismatch")
    for label, value in (("graph gold", graph["gold"]), ("QA gold", qa["gold"])):
        path = _relative_path(value, label)
        _require(
            PurePosixPath(path).parts[0] != CORPUS_ROOT,
            f"{label} must remain outside corpus",
        )
    return set(declared)


def _validate_baseline_cases(
    records: list[dict[str, Any]], declared_paths: set[str]
) -> None:
    seen_ids: set[str] = set()
    covered_sources: set[str] = set()
    for index, record in enumerate(records):
        label = f"cases.jsonl[{index}]"
        _require(record.get("schema_version") == SCHEMA_VERSION, f"{label}: schema_version")
        _require(record.get("record_type") == "general_memory_shadow_eval_case", f"{label}: record_type")
        case_id = record.get("eval_case_id")
        _require(isinstance(case_id, str) and bool(case_id), f"{label}: eval_case_id")
        _require(case_id not in seen_ids, f"{label}: duplicate eval_case_id: {case_id}")
        seen_ids.add(case_id)
        _require(isinstance(record.get("query"), str) and bool(record["query"]), f"{label}: query")
        sources = _string_list(record.get("relevant_sources"), f"{label}.relevant_sources")
        for source in sources:
            normalized = _relative_path(source, f"{label}.relevant_sources")
            _require(normalized in declared_paths, f"{label}: undeclared relevant source: {normalized}")
            covered_sources.add(normalized)
        provenance = record.get("provenance")
        _require(isinstance(provenance, dict), f"{label}.provenance: object required")
        for source_index, location in enumerate(
            _string_list(provenance.get("source_locations"), f"{label}.provenance.source_locations")
        ):
            _require(location.startswith("corpus/"), f"{label}: source location must begin corpus/")
            source_path = location[len("corpus/"):].split("#", 1)[0]
            source_path = _relative_path(source_path, f"{label}.source_locations[{source_index}]")
            _require(source_path in declared_paths, f"{label}: undeclared provenance source: {source_path}")
    _require(covered_sources == declared_paths, "baseline cases must collectively cover all five sources")


def _validate_expected_graph(
    records: list[dict[str, Any]], declared_paths: set[str]
) -> set[str]:
    edge_keys: set[str] = set()
    relation_types: set[str] = set()
    for index, record in enumerate(records):
        label = f"gold/expected-graph.jsonl[{index}]"
        _require(record.get("schema_version") == SCHEMA_VERSION, f"{label}: schema_version")
        _require(
            record.get("record_type") == "cross_format_kg_expected_edge",
            f"{label}: record_type",
        )
        edge_key = record.get("gold_edge_key")
        _require(isinstance(edge_key, str) and bool(edge_key), f"{label}: gold_edge_key")
        _require(edge_key not in edge_keys, f"{label}: duplicate gold_edge_key: {edge_key}")
        edge_keys.add(edge_key)
        _require(record.get("relation_class") == "semantic", f"{label}: relation_class")
        relation_type = record.get("relation_type")
        _require(isinstance(relation_type, str) and bool(relation_type), f"{label}: relation_type")
        relation_types.add(relation_type)
        _require(record.get("expected_status") == "verified", f"{label}: expected_status")
        _require(
            isinstance(record.get("from"), dict)
            and set(record["from"]) == {"canonical_key", "node_type"},
            f"{label}.from: canonical node reference required",
        )
        _require(
            isinstance(record.get("to"), dict)
            and set(record["to"]) == {"canonical_key", "node_type"},
            f"{label}.to: canonical node reference required",
        )
        references = record.get("source_references")
        _require(isinstance(references, list) and bool(references), f"{label}: source_references")
        for source_index, reference in enumerate(references):
            _validate_source_reference(
                reference,
                f"{label}.source_references[{source_index}]",
                declared_paths,
            )
        _require(
            not _contains_forbidden_evidence_id(record),
            f"{label}: generated Evidence IDs are forbidden in gold",
        )
    missing_types = REQUIRED_SEMANTIC_RELATION_TYPES - relation_types
    _require(not missing_types, f"gold graph missing semantic relation types: {sorted(missing_types)}")
    return edge_keys


def _validate_qa_cases(
    records: list[dict[str, Any]],
    declared_paths: set[str],
    gold_edge_keys: set[str],
) -> tuple[int, int]:
    seen_ids: set[str] = set()
    accepted_count = 0
    hold_count = 0
    reference_time_hold = False
    for index, record in enumerate(records):
        label = f"gold/qa-cases.jsonl[{index}]"
        _require(record.get("schema_version") == SCHEMA_VERSION, f"{label}: schema_version")
        _require(record.get("record_type") == "cross_format_kg_qa_case", f"{label}: record_type")
        case_id = record.get("qa_case_id")
        _require(isinstance(case_id, str) and bool(case_id), f"{label}: qa_case_id")
        _require(case_id not in seen_ids, f"{label}: duplicate qa_case_id: {case_id}")
        seen_ids.add(case_id)
        _require(isinstance(record.get("question"), str) and bool(record["question"]), f"{label}: question")

        expected = record.get("expected")
        graph = record.get("graph_requirements")
        _require(isinstance(expected, dict), f"{label}.expected: object required")
        _require(isinstance(graph, dict), f"{label}.graph_requirements: object required")
        decision = expected.get("decision")
        _require(decision in {"ACCEPTED", "HOLD"}, f"{label}: unsupported decision")
        _require(graph.get("cross_document") is True, f"{label}: cross_document must be true")
        _require(
            graph.get("all_used_semantic_edges_must_have_status") == "verified",
            f"{label}: used semantic edges must be verified",
        )
        minimum_documents = _positive_integer(
            graph.get("minimum_distinct_visited_documents"),
            f"{label}.minimum_distinct_visited_documents",
        )
        _require(minimum_documents >= 2, f"{label}: at least two distinct documents required")
        minimum_edges = _positive_integer(
            graph.get("minimum_verified_semantic_edges"),
            f"{label}.minimum_verified_semantic_edges",
        )
        required_documents = _string_list(
            graph.get("required_visited_documents"),
            f"{label}.required_visited_documents",
        )
        normalized_documents = {
            _relative_path(path, f"{label}.required_visited_documents")
            for path in required_documents
        }
        _require(
            normalized_documents <= declared_paths,
            f"{label}: required visited document is undeclared",
        )
        _require(
            len(normalized_documents) >= minimum_documents,
            f"{label}: required visited documents do not meet minimum",
        )
        required_edges = _string_list(
            graph.get("required_gold_edge_keys"),
            f"{label}.required_gold_edge_keys",
        )
        unknown_edges = set(required_edges) - gold_edge_keys
        _require(not unknown_edges, f"{label}: unknown required gold edges: {sorted(unknown_edges)}")
        _require(
            len(required_edges) >= minimum_edges,
            f"{label}: required gold edges do not meet minimum",
        )

        references = record.get("source_references")
        _require(isinstance(references, list) and bool(references), f"{label}: source_references")
        reference_paths = {
            _validate_source_reference(
                reference,
                f"{label}.source_references[{source_index}]",
                declared_paths,
            )
            for source_index, reference in enumerate(references)
        }
        _require(
            normalized_documents <= reference_paths,
            f"{label}: each required visited document needs a source reference",
        )

        provenance = record.get("provenance")
        _require(isinstance(provenance, dict), f"{label}.provenance: object required")
        for key in (
            "gold_available_to_answerer",
            "gold_available_to_extractor",
            "gold_available_to_graph_builder",
        ):
            _require(provenance.get(key) is False, f"{label}.provenance.{key} must be false")
        _require(
            provenance.get("question_released_after_snapshot_freeze") is True,
            f"{label}: question must be released after snapshot freeze",
        )

        ablation = graph.get("edge_ablation")
        _require(isinstance(ablation, dict), f"{label}.edge_ablation: object required")
        if decision == "ACCEPTED":
            accepted_count += 1
            _require(ablation.get("required") is True, f"{label}: accepted case requires ablation")
            _require(
                ablation.get("remove_each_required_edge_one_at_a_time") is True,
                f"{label}: accepted case must ablate each required edge",
            )
            _require(
                ablation.get("expected_decision") == "HOLD",
                f"{label}: accepted case ablation must produce HOLD",
            )
        else:
            hold_count += 1
            if expected.get("reason_code") == "reference_time_required":
                _require(record.get("reference_time") is None, f"{label}: reference time must be absent")
                reference_time_hold = True

    _require(accepted_count >= 1, "QA gold must contain at least one ACCEPTED case")
    _require(hold_count >= 1, "QA gold must contain at least one HOLD case")
    _require(reference_time_hold, "QA gold must contain the reference_time_required HOLD case")
    return accepted_count, hold_count


def validate_contract(dataset_root: Path) -> dict[str, Any]:
    dataset_root = dataset_root.resolve()
    _require(dataset_root.is_dir(), f"dataset directory is missing: {dataset_root}")
    spec = load_json(dataset_root / FIXTURE_SPEC, FIXTURE_SPEC)
    declared_paths = _validate_fixture_spec(spec)
    baseline = load_jsonl(dataset_root / BASELINE_CASES, BASELINE_CASES)
    graph = load_jsonl(dataset_root / EXPECTED_GRAPH, EXPECTED_GRAPH)
    qa = load_jsonl(dataset_root / QA_CASES, QA_CASES)
    _validate_baseline_cases(baseline, declared_paths)
    gold_edge_keys = _validate_expected_graph(graph, declared_paths)
    accepted_count, hold_count = _validate_qa_cases(
        qa, declared_paths, gold_edge_keys
    )
    return {
        "status": "OK",
        "dataset_id": DATASET_ID,
        "declared_source_count": len(declared_paths),
        "baseline_case_count": len(baseline),
        "gold_edge_count": len(graph),
        "qa_case_count": len(qa),
        "accepted_case_count": accepted_count,
        "hold_case_count": hold_count,
        "build_input_boundary": BUILD_INPUT_GLOB,
    }


def _corpus_files(dataset_root: Path) -> dict[str, Path]:
    corpus = (dataset_root / CORPUS_ROOT).resolve()
    _require(corpus.is_dir(), f"corpus directory is missing: {corpus}")
    files: dict[str, Path] = {}
    for path in corpus.rglob("*"):
        if path.is_symlink():
            raise FixtureContractError(f"corpus symlink is forbidden: {path}")
        if not path.is_file():
            continue
        relative = unicodedata.normalize("NFC", path.relative_to(corpus).as_posix())
        _require(relative not in files, f"NFC-normalized corpus path collision: {relative}")
        files[relative] = path
    return files


def _docx_container_text(container: Any) -> list[str]:
    parts = [paragraph.text for paragraph in getattr(container, "paragraphs", ())]
    for table in getattr(container, "tables", ()):
        for row in table.rows:
            parts.extend(cell.text for cell in row.cells)
    return parts


def _extract_docx_text(path: Path) -> str:
    try:
        from docx import Document
    except ImportError as exc:  # pragma: no cover - depends on runtime packaging
        raise FixtureContractError(
            "DOCX content validation requires python-docx"
        ) from exc
    try:
        document = Document(path)
        parts = _docx_container_text(document)
        for section in document.sections:
            for container in (section.header, section.footer):
                parts.extend(_docx_container_text(container))
        return "\n".join(part for part in parts if part)
    except FixtureContractError:
        raise
    except Exception as exc:
        raise FixtureContractError(f"cannot extract DOCX content: {path}") from exc


def _pptx_shape_text(shape: Any) -> list[str]:
    parts: list[str] = []
    if getattr(shape, "has_text_frame", False):
        text = shape.text_frame.text
        if text:
            parts.append(text)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            parts.extend(cell.text for cell in row.cells if cell.text)
    child_shapes = getattr(shape, "shapes", None)
    if child_shapes is not None:
        for child in child_shapes:
            parts.extend(_pptx_shape_text(child))
    return parts


def _extract_pptx_text(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:  # pragma: no cover - depends on runtime packaging
        raise FixtureContractError(
            "PPTX content validation requires python-pptx"
        ) from exc
    try:
        presentation = Presentation(path)
        parts: list[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                parts.extend(_pptx_shape_text(shape))
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text
                if notes_text:
                    parts.append(notes_text)
        return "\n".join(parts)
    except FixtureContractError:
        raise
    except Exception as exc:
        raise FixtureContractError(f"cannot extract PPTX content: {path}") from exc


def _extract_pdf_text(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as exc:  # pragma: no cover - depends on runtime packaging
        raise FixtureContractError("PDF content validation requires pypdf") from exc
    try:
        reader = PdfReader(path)
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    except Exception as exc:
        raise FixtureContractError(f"cannot extract PDF content: {path}") from exc


def _normalized_excel_value(value: Any) -> Any:
    if isinstance(value, datetime):
        if value.hour == value.minute == value.second == value.microsecond == 0:
            return value.date().isoformat()
        return value.isoformat()
    if isinstance(value, date):
        return value.isoformat()
    return value


def _extract_and_validate_xlsx(
    path: Path,
    contract: dict[str, Any],
    label: str,
) -> tuple[str, dict[str, Any]]:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:  # pragma: no cover - depends on runtime packaging
        raise FixtureContractError(
            "XLSX content validation requires openpyxl"
        ) from exc
    try:
        # Artifact-tool generated workbooks may omit worksheet dimension
        # metadata. Normal mode calculates the actual A1:H3 extent reliably.
        workbook = load_workbook(path, read_only=False, data_only=False)
    except Exception as exc:
        raise FixtureContractError(f"cannot extract XLSX content: {path}") from exc
    try:
        sheet_name = contract.get("sheet")
        _require(
            isinstance(sheet_name, str) and bool(sheet_name),
            f"{label}.sheet: non-empty string required",
        )
        _require(
            workbook.sheetnames == [sheet_name],
            f"{label}: workbook sheets must be exactly [{sheet_name!r}]",
        )
        worksheet = workbook[sheet_name]
        _require(
            worksheet.max_row == 3 and worksheet.max_column == 8,
            f"{label}: populated XLSX range must be exactly A1:H3",
        )
        cells = [
            list(row)
            for row in worksheet.iter_rows(
                min_row=1, max_row=3, min_col=1, max_col=8
            )
        ]
        header = tuple(cell.value for cell in cells[0])
        _require(header == XLSX_HEADERS, f"{label}: A1:H1 header mismatch")

        expected_rows = contract.get("must_contain_rows")
        _require(
            isinstance(expected_rows, list) and len(expected_rows) == 2,
            f"{label}.must_contain_rows: exactly two rows required",
        )
        for index, row in enumerate(expected_rows):
            _require(
                isinstance(row, list) and len(row) == 8,
                f"{label}.must_contain_rows[{index}]: eight cells required",
            )
        normalized_rows = [
            [_normalized_excel_value(cell.value) for cell in row]
            for row in cells[1:]
        ]
        _require(
            normalized_rows == expected_rows,
            f"{label}: normalized A2:H3 rows do not match must_contain_rows",
        )

        for row_index, (actual_row, expected_row) in enumerate(
            zip(cells[1:], expected_rows), start=2
        ):
            for column_index in (5, 6):
                cell = actual_row[column_index]
                expected = expected_row[column_index]
                coordinate = cell.coordinate
                if expected is None:
                    _require(
                        cell.value is None,
                        f"{label}: {coordinate} must be a true blank/null",
                    )
                    continue
                _require(
                    isinstance(expected, str)
                    and re.fullmatch(r"\d{4}-\d{2}-\d{2}", expected) is not None,
                    f"{label}: {coordinate} expected date must be ISO YYYY-MM-DD",
                )
                _require(
                    isinstance(cell.value, (date, datetime)) and cell.is_date,
                    f"{label}: {coordinate} must be an XLSX typed date, not text",
                )

        all_values: list[str] = []
        for sheet in workbook.worksheets:
            for row in sheet.iter_rows():
                for cell in row:
                    normalized = _normalized_excel_value(cell.value)
                    if normalized is not None:
                        all_values.append(str(normalized))
        return "\n".join(all_values), {
            "sheet": sheet_name,
            "validated_range": "A1:H3",
            "typed_date_cells": ["F2", "G2", "F3"],
            "blank_date_cells": ["G3"],
        }
    finally:
        workbook.close()


def _validate_exact_phrases(
    text: str,
    contract: dict[str, Any],
    label: str,
) -> None:
    is_xlsx = contract.get("format") == "xlsx"
    _require(
        is_xlsx or "must_contain_exact" in contract,
        f"{label}.must_contain_exact: field is required",
    )
    _require(
        "must_not_contain_exact" in contract,
        f"{label}.must_not_contain_exact: field is required",
    )
    must_contain = _string_list(
        contract.get("must_contain_exact", []),
        f"{label}.must_contain_exact",
        allow_empty=is_xlsx,
    )
    must_not_contain = _string_list(
        contract["must_not_contain_exact"],
        f"{label}.must_not_contain_exact",
    )
    missing = [phrase for phrase in must_contain if phrase not in text]
    forbidden = [phrase for phrase in must_not_contain if phrase in text]
    _require(not missing, f"{label}: missing required exact phrases: {missing}")
    _require(
        not forbidden,
        f"{label}: forbidden exact phrases were present: {forbidden}",
    )


def validate_corpus(dataset_root: Path) -> dict[str, Any]:
    """Validate the five real artifacts against their content contract."""
    dataset_root = dataset_root.resolve()
    contract_result = validate_contract(dataset_root)
    spec = load_json(dataset_root / FIXTURE_SPEC, FIXTURE_SPEC)
    contracts = {
        record["path"]: record
        for record in spec["corpus_contract"]["required_files"]
    }
    files = _corpus_files(dataset_root)
    expected_paths = set(REQUIRED_SOURCE_PATHS)
    actual_paths = set(files)
    _require(
        actual_paths == expected_paths,
        "corpus file set mismatch: "
        f"missing={sorted(expected_paths - actual_paths)} "
        f"extra={sorted(actual_paths - expected_paths)}",
    )

    results: list[dict[str, Any]] = []
    for relative in REQUIRED_SOURCE_PATHS:
        path = files[relative]
        size = path.stat().st_size
        _require(size > 0, f"corpus file is empty: {relative}")
        contract = contracts[relative]
        label = f"fixture-spec corpus contract for {relative}"
        artifact_format = contract["format"]
        details: dict[str, Any] = {}
        if artifact_format == "docx":
            text = _extract_docx_text(path)
        elif artifact_format == "xlsx":
            text, details = _extract_and_validate_xlsx(path, contract, label)
        elif artifact_format == "pptx":
            text = _extract_pptx_text(path)
        elif artifact_format == "pdf":
            text = _extract_pdf_text(path)
        else:  # protected by the static fixture contract
            raise FixtureContractError(
                f"{label}: unsupported artifact format: {artifact_format}"
            )
        _require(bool(text), f"{label}: extracted content is empty")
        _validate_exact_phrases(text, contract, label)
        results.append({
            "path": relative,
            "format": artifact_format,
            "size_bytes": size,
            "extracted_character_count": len(text),
            "content_contract": "PASS",
            **details,
        })
    return {
        **contract_result,
        "corpus_content_validation": "PASS",
        "corpus_file_count": len(results),
        "files": results,
    }


def build_manifest_record(dataset_root: Path) -> dict[str, Any]:
    validate_contract(dataset_root)
    files = _corpus_files(dataset_root.resolve())
    expected_paths = set(REQUIRED_SOURCE_PATHS)
    actual_paths = set(files)
    _require(
        actual_paths == expected_paths,
        "corpus file set mismatch: "
        f"missing={sorted(expected_paths - actual_paths)} "
        f"extra={sorted(actual_paths - expected_paths)}",
    )
    records: list[dict[str, Any]] = []
    for relative in REQUIRED_SOURCE_PATHS:
        path = files[relative]
        size = path.stat().st_size
        _require(size > 0, f"corpus file is empty: {relative}")
        records.append({
            "path": relative,
            "format": REQUIRED_FORMATS[relative],
            "size_bytes": size,
            "sha256": _file_sha256(path),
        })
    return {
        "schema_version": SCHEMA_VERSION,
        "record_type": "cross_format_kg_corpus_manifest",
        "dataset_id": DATASET_ID,
        "corpus_root": CORPUS_ROOT,
        "file_count": len(records),
        "files": records,
        "source_set_sha256": _record_sha256(records),
    }


def _ensure_manifest_outside_corpus(dataset_root: Path, manifest_path: Path) -> None:
    corpus = (dataset_root.resolve() / CORPUS_ROOT).resolve()
    resolved = manifest_path.resolve()
    _require(
        resolved != corpus and corpus not in resolved.parents,
        "corpus manifest must be written outside corpus/",
    )


def write_manifest(
    dataset_root: Path,
    manifest_path: Path,
    *,
    overwrite: bool = False,
) -> dict[str, Any]:
    _ensure_manifest_outside_corpus(dataset_root, manifest_path)
    _require(overwrite or not manifest_path.exists(), f"refusing to overwrite manifest: {manifest_path}")
    record = build_manifest_record(dataset_root)
    manifest_path.parent.mkdir(parents=True, exist_ok=True)
    temporary = manifest_path.with_name(f".{manifest_path.name}.building")
    try:
        temporary.write_text(
            json.dumps(record, ensure_ascii=False, sort_keys=True, indent=2) + "\n",
            encoding="utf-8",
        )
        os.replace(temporary, manifest_path)
    finally:
        temporary.unlink(missing_ok=True)
    return record


def validate_manifest(dataset_root: Path, manifest_path: Path) -> dict[str, Any]:
    _ensure_manifest_outside_corpus(dataset_root, manifest_path)
    actual = load_json(manifest_path, "corpus manifest")
    expected = build_manifest_record(dataset_root)
    _require(actual == expected, "corpus manifest does not match current sha256/size snapshot")
    for index, record in enumerate(actual["files"]):
        _require(SHA256.fullmatch(record["sha256"]) is not None, f"manifest.files[{index}].sha256")
    return {
        "status": "OK",
        "manifest": str(manifest_path.resolve()),
        "file_count": actual["file_count"],
        "source_set_sha256": actual["source_set_sha256"],
    }


def parse_args() -> argparse.Namespace:
    repository = Path(__file__).resolve().parents[1]
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--dataset",
        type=Path,
        default=repository / "evaluation" / DATASET_ID,
    )
    action = parser.add_mutually_exclusive_group()
    action.add_argument("--validate-corpus", action="store_true")
    action.add_argument("--write-manifest", type=Path)
    action.add_argument("--validate-manifest", type=Path)
    parser.add_argument("--overwrite", action="store_true")
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    try:
        _require(
            not args.overwrite or args.write_manifest is not None,
            "--overwrite is valid only with --write-manifest",
        )
        if args.write_manifest is not None:
            result = write_manifest(
                args.dataset, args.write_manifest, overwrite=args.overwrite
            )
        elif args.validate_manifest is not None:
            result = validate_manifest(args.dataset, args.validate_manifest)
        elif args.validate_corpus:
            result = validate_corpus(args.dataset)
        else:
            result = validate_contract(args.dataset)
    except FixtureContractError as exc:
        print(f"ERROR: {exc}", file=sys.stderr)
        return 2
    print(json.dumps(result, ensure_ascii=False, sort_keys=True, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
