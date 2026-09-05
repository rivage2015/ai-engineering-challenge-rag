#!/usr/bin/env python3
"""Generate a small, question-independent intermediate-record sample.

This is a diagnostic probe, not the production answer pipeline.  It accepts
arbitrary input paths and applies the same suffix-based dispatch to every file.
"""

from __future__ import annotations

import argparse
import base64
import codecs
import csv
import hashlib
import io
import json
import mimetypes
import os
import posixpath
import re
import shutil
import stat
import tempfile
import time
import unicodedata
import urllib.parse
import zipfile
from datetime import date, datetime, timezone
from pathlib import Path, PurePosixPath
from typing import Any, Callable
from xml.etree import ElementTree

from evidence_text_chunking import (
    MAX_QUESTION_EVIDENCE_CHARS,
    exact_text_chunks,
)


SCHEMA_VERSION = "0.1"
EXTRACTOR = "intermediate-record-probe"
EXTRACTOR_VERSION = "0.7.1"
FORMULA_CACHED_VALUE_STATUS = "stored_in_file_not_recalculated"

PLAIN_TEXT_SUFFIXES = {
    ".md", ".txt", ".py", ".toml", ".yaml", ".yml", ".rst", ".sql", ".sh", ".command",
}
IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp"}
OCR_QUALITY_BY_AGREEMENT = {
    "independent_agreement": "high",
    "same_engine_agreement": "provisional",
    "provisional_single_pass": "provisional",
    "display_transform_unresolved": "provisional",
}
PROVISIONAL_OCR_MARKER = "[暫定読取]"
OCR_BBOX_COORDINATE_SYSTEMS = {
    "raw_raster_top_left_normalized_1000",
    "display_oriented_top_left_normalized_1000",
    "source_orientation_1_top_left_normalized_1000",
}
OCR_ENGINE_BY_PASS = {
    "apple_vision_primary": "apple_vision",
    "apple_vision_literal": "apple_vision",
    "apple_vision_fast_sparse": "apple_vision",
    "paddleocr_primary": "paddleocr",
    "tesseract_psm3": "tesseract",
    "tesseract_psm6": "tesseract",
    "tesseract_psm11": "tesseract",
}
CODE_SUFFIXES = {".py", ".toml", ".yaml", ".yml", ".sql", ".sh", ".command"}
DIRECT_TEXT_SUFFIXES = PLAIN_TEXT_SUFFIXES | {".csv", ".tsv", ".json", ".xml", ".ipynb"}
MAX_DIRECT_TEXT_BYTES = 64 * 1024 * 1024
STREAM_TEXT_READ_CHARS = 64 * 1024
MAX_STREAM_TEXT_READ_BLOCKS = 4096
TEXT_ENCODING_SNIFF_BYTES = 64 * 1024
TEXT_ENCODING_SCAN_BYTES = 1024 * 1024
MAX_OOXML_ARCHIVE_BYTES = 256 * 1024 * 1024
MAX_OOXML_ZIP_ENTRIES = 10_000
MAX_OOXML_MEMBER_BYTES = 128 * 1024 * 1024
MAX_OOXML_TOTAL_BYTES = 512 * 1024 * 1024
MAX_OOXML_COMPRESSION_RATIO = 500.0
MAX_EMBEDDED_VISUALS_PER_DOCUMENT = 128
MAX_EMBEDDED_VISUAL_BYTES_PER_DOCUMENT = 256 * 1024 * 1024
MAX_EMBEDDED_VISUAL_SECONDS_PER_DOCUMENT = 900.0
VISUAL_OBSERVATION_MODES = {
    "immediate",
    "deferred_per_document",
    "suppressed",
}
MAX_DEFERRED_VISUAL_TASKS = 128
MAX_DEFERRED_VISUAL_SPOOL_BYTES = 256 * 1024 * 1024
OOXML_FORBIDDEN_XML = (b"<!DOCTYPE", b"<!ENTITY")
OOXML_DOCUMENT_RELATIONSHIP_NAMESPACES = {
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "http://purl.oclc.org/ooxml/officeDocument/relationships",
}
OOXML_PACKAGE_RELATIONSHIP_NAMESPACES = {
    "http://schemas.openxmlformats.org/package/2006/relationships",
    "http://purl.oclc.org/ooxml/package/relationships",
}
OOXML_PRESENTATION_NAMESPACES = {
    "http://schemas.openxmlformats.org/presentationml/2006/main",
    "http://purl.oclc.org/ooxml/presentationml/main",
}
OOXML_SPREADSHEET_NAMESPACES = {
    "http://schemas.openxmlformats.org/spreadsheetml/2006/main",
    "http://purl.oclc.org/ooxml/spreadsheetml/main",
}
OOXML_WORDPROCESSING_NAMESPACES = {
    "http://schemas.openxmlformats.org/wordprocessingml/2006/main",
    "http://purl.oclc.org/ooxml/wordprocessingml/main",
}
OOXML_DRAWING_NAMESPACES = {
    "http://schemas.openxmlformats.org/drawingml/2006/main",
    "http://purl.oclc.org/ooxml/drawingml/main",
}
OOXML_DRAWING_TABLE_URIS = {
    "http://schemas.openxmlformats.org/drawingml/2006/table",
    "http://purl.oclc.org/ooxml/drawingml/table",
}
OOXML_CHART_NAMESPACES = {
    "http://schemas.openxmlformats.org/drawingml/2006/chart",
    "http://purl.oclc.org/ooxml/drawingml/chart",
}
OOXML_DIAGRAM_NAMESPACES = {
    "http://schemas.openxmlformats.org/drawingml/2006/diagram",
    "http://purl.oclc.org/ooxml/drawingml/diagram",
}
OOXML_SPREADSHEET_DRAWING_NAMESPACES = {
    "http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing",
    "http://purl.oclc.org/ooxml/drawingml/spreadsheetDrawing",
}
OOXML_RELATIONSHIP_CARRIER_NAMESPACES = (
    OOXML_PRESENTATION_NAMESPACES
    | OOXML_SPREADSHEET_NAMESPACES
    | OOXML_WORDPROCESSING_NAMESPACES
    | OOXML_DRAWING_NAMESPACES
    | OOXML_CHART_NAMESPACES
    | OOXML_DIAGRAM_NAMESPACES
    | OOXML_SPREADSHEET_DRAWING_NAMESPACES
    | {"urn:schemas-microsoft-com:vml", "urn:schemas-microsoft-com:office:office"}
)
OOXML_RELATIONSHIP_CARRIERS: dict[tuple[str, str], frozenset[str]] = {}
for _namespace in OOXML_DRAWING_NAMESPACES:
    OOXML_RELATIONSHIP_CARRIERS.update({
        (_namespace, "blip"): frozenset({"embed", "link"}),
        (_namespace, "hlinkClick"): frozenset({"id"}),
        (_namespace, "hlinkHover"): frozenset({"id"}),
    })
for _namespace in OOXML_CHART_NAMESPACES:
    OOXML_RELATIONSHIP_CARRIERS[(_namespace, "chart")] = frozenset({"id"})
for _namespace in OOXML_DIAGRAM_NAMESPACES:
    OOXML_RELATIONSHIP_CARRIERS[(_namespace, "relIds")] = frozenset(
        {"dm", "lo", "qs", "cs"}
    )
for _namespace in OOXML_PRESENTATION_NAMESPACES:
    OOXML_RELATIONSHIP_CARRIERS.update({
        (_namespace, "sldId"): frozenset({"id"}),
        (_namespace, "oleObj"): frozenset({"id"}),
        (_namespace, "contentPart"): frozenset({"id"}),
    })
for _namespace in OOXML_SPREADSHEET_NAMESPACES:
    OOXML_RELATIONSHIP_CARRIERS.update({
        (_namespace, "sheet"): frozenset({"id"}),
        (_namespace, "drawing"): frozenset({"id"}),
        (_namespace, "legacyDrawing"): frozenset({"id"}),
        (_namespace, "legacyDrawingHF"): frozenset({"id"}),
        (_namespace, "oleObject"): frozenset({"id"}),
    })
for _namespace in OOXML_WORDPROCESSING_NAMESPACES:
    OOXML_RELATIONSHIP_CARRIERS.update({
        (_namespace, "headerReference"): frozenset({"id"}),
        (_namespace, "footerReference"): frozenset({"id"}),
        (_namespace, "altChunk"): frozenset({"id"}),
        (_namespace, "subDoc"): frozenset({"id"}),
    })
OOXML_RELATIONSHIP_CARRIERS.update({
    ("urn:schemas-microsoft-com:vml", "imagedata"): frozenset(
        {"id", "href"}
    ),
    ("urn:schemas-microsoft-com:office:office", "OLEObject"): frozenset(
        {"id"}
    ),
})
DOCX_REQUIRED_OOXML_MEMBERS = frozenset({
    "[Content_Types].xml",
    "_rels/.rels",
    "word/document.xml",
})
PPTX_REQUIRED_OOXML_MEMBERS = frozenset({
    "[Content_Types].xml",
    "_rels/.rels",
    "ppt/presentation.xml",
    "ppt/_rels/presentation.xml.rels",
})
XML_DECLARATION_ENCODING = re.compile(
    r"<\?xml\s+[^>]{0,512}?\bencoding\s*=\s*(['\"])([A-Za-z][A-Za-z0-9._-]*)\1",
    re.IGNORECASE,
)
DATA_URI_PATTERN = re.compile(
    r"data:(image/[A-Za-z0-9.+-]+);base64,([A-Za-z0-9+/=\r\n\t ]+)",
    re.IGNORECASE,
)
NOTEBOOK_ATTACHMENT_PATTERN = re.compile(
    r"attachment:([^)\s\"'>]+)",
    re.IGNORECASE,
)
CONTROL_PATTERN = re.compile(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]")
DATE_TOKEN_PATTERN = re.compile(r"(?<!\d)(20\d{6})(?!\d)")
ALIAS_DATE_PATTERN = re.compile(r"([A-Za-z]{2,})[-_]?((?:20)\d{6})")


def canonical_json(value: Any) -> str:
    return json.dumps(value, ensure_ascii=False, sort_keys=True, separators=(",", ":"))


def ocr_engine(pass_name: object) -> str:
    if not isinstance(pass_name, str):
        raise ValueError("OCR provenance pass name is missing")
    try:
        return OCR_ENGINE_BY_PASS[pass_name]
    except KeyError as exc:
        raise ValueError(f"unsupported OCR provenance pass: {pass_name!r}") from exc


def ocr_independence_group(pass_name: object) -> str:
    return ocr_engine(pass_name)


def _ocr_match_text(value: object) -> str:
    if not isinstance(value, str) or not value.strip():
        raise ValueError("OCR supporter raw text is missing")
    return unicodedata.normalize("NFC", value).strip()


def _ocr_bbox(value: object) -> list[int]:
    if (
        not isinstance(value, list)
        or len(value) != 4
        or any(isinstance(item, bool) or not isinstance(item, int) for item in value)
        or value[0] < 0
        or value[1] < 0
        or value[2] <= 0
        or value[3] <= 0
        or value[0] + value[2] > 1000
        or value[1] + value[3] > 1000
    ):
        raise ValueError("OCR supporter bbox is invalid")
    return list(value)


def _ocr_overlap(first: list[int], second: list[int]) -> float:
    left = max(first[0], second[0])
    top = max(first[1], second[1])
    right = min(first[0] + first[2], second[0] + second[2])
    bottom = min(first[1] + first[3], second[1] + second[3])
    if right <= left or bottom <= top:
        return 0.0
    intersection = (right - left) * (bottom - top)
    smaller = min(first[2] * first[3], second[2] * second[3])
    return intersection / smaller if smaller else 0.0


def validate_ocr_supporters(
    line: dict[str, Any],
    provenance: dict[str, Any],
    *,
    primary_pass: str,
    primary_engine: str,
    primary_group: str,
    audit_pass: str | None,
    audit_engine: str | None,
    audit_group: str | None,
    agreement_type: str,
) -> None:
    """Recompute line agreement from the two raw engine observations."""
    supporters = provenance.get("supporters")
    expected_count = 1 if audit_pass is None else 2
    if not isinstance(supporters, list) or len(supporters) != expected_count:
        raise ValueError("OCR raw supporters are missing or incomplete")
    expected = [
        (
            primary_pass,
            primary_engine,
            primary_group,
            provenance.get("primary_line_id"),
            provenance.get("primary_bbox_coordinate_system"),
            line.get("primary_confidence"),
        )
    ]
    if audit_pass is not None:
        expected.append((
            audit_pass,
            audit_engine,
            audit_group,
            provenance.get("audit_line_id"),
            provenance.get("audit_bbox_coordinate_system"),
            line.get("audit_confidence"),
        ))
    boxes: list[list[int]] = []
    line_text = _ocr_match_text(line.get("text"))
    for supporter, contract in zip(supporters, expected):
        if not isinstance(supporter, dict):
            raise ValueError("OCR supporter must be an object")
        pass_name, engine, group, line_id, coordinate_system, confidence = contract
        if (
            supporter.get("pass") != pass_name
            or supporter.get("engine") != engine
            or supporter.get("independence_group") != group
            or supporter.get("line_id") != line_id
            or supporter.get("bbox_coordinate_system") != coordinate_system
        ):
            raise ValueError("OCR supporter identity disagrees with provenance")
        if _ocr_match_text(supporter.get("raw_text")) != line_text:
            raise ValueError("OCR supporter text does not reproduce the consensus")
        actual_confidence = supporter.get("confidence")
        if (
            isinstance(actual_confidence, bool)
            or not isinstance(actual_confidence, (int, float))
            or confidence is None
            or float(actual_confidence) != float(confidence)
        ):
            raise ValueError("OCR supporter confidence disagrees with provenance")
        boxes.append(_ocr_bbox(supporter.get("bbox")))
    result_bbox = _ocr_bbox(line.get("bbox"))
    union = [
        min(box[0] for box in boxes),
        min(box[1] for box in boxes),
        max(box[0] + box[2] for box in boxes),
        max(box[1] + box[3] for box in boxes),
    ]
    union[2] -= union[0]
    union[3] -= union[1]
    if result_bbox != union:
        raise ValueError("OCR consensus bbox does not reproduce the supporter union")
    claimed_overlap = line.get("overlap")
    if expected_count == 1:
        if claimed_overlap != 0.0 or agreement_type != "provisional_single_pass":
            raise ValueError("single-pass OCR supporter contract is invalid")
        return
    recomputed_overlap = _ocr_overlap(boxes[0], boxes[1])
    if (
        isinstance(claimed_overlap, bool)
        or not isinstance(claimed_overlap, (int, float))
        or abs(float(claimed_overlap) - round(recomputed_overlap, 6)) > 0.000001
        or recomputed_overlap < 0.5
    ):
        raise ValueError("OCR supporter overlap does not reproduce the consensus")


def digest_bytes(value: bytes) -> str:
    return hashlib.sha256(value).hexdigest()


def digest_file(path: Path, chunk_size: int = 1024 * 1024) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        while chunk := handle.read(chunk_size):
            digest.update(chunk)
    return digest.hexdigest()


def digest_value(value: Any) -> str:
    return digest_bytes(canonical_json(value).encode("utf-8"))


def stable_id(prefix: str, value: Any) -> str:
    return f"{prefix}_{digest_value(value)[:32]}"


def normalize_text(value: str) -> str:
    """Normalize transport noise without changing names, numbers, symbols, or negation."""
    normalized = unicodedata.normalize("NFC", value).replace("\r\n", "\n").replace("\r", "\n")
    normalized = CONTROL_PATTERN.sub("", normalized)
    lines = [re.sub(r"[\t\u00a0 ]+", " ", line).rstrip() for line in normalized.split("\n")]
    compact: list[str] = []
    blank = False
    for line in lines:
        if line:
            compact.append(line)
            blank = False
        elif compact and not blank:
            compact.append("")
            blank = True
    while compact and not compact[-1]:
        compact.pop()
    return "\n".join(compact)


def detect_text_encoding(raw: bytes, *, partial_sample: bool = False) -> str:
    """Detect a common native-text encoding from bounded or complete bytes."""
    def decodes(encoding: str) -> bool:
        try:
            decoder = codecs.getincrementaldecoder(encoding)(errors="strict")
            decoder.decode(raw, final=not partial_sample)
            return True
        except UnicodeDecodeError:
            return False

    if raw.startswith((b"\xff\xfe", b"\xfe\xff")):
        return "utf-16"
    if raw and raw.count(b"\x00") / len(raw) > 0.1:
        odd_nuls = raw[1::2].count(0)
        even_nuls = raw[0::2].count(0)
        likely_utf16 = "utf-16-le" if odd_nuls >= even_nuls else "utf-16-be"
        if decodes(likely_utf16):
            return likely_utf16
    for encoding in ("utf-8-sig", "cp932"):
        if decodes(encoding):
            return encoding
    return "utf-8-replacement"


def detect_text_file_encoding(path: Path) -> str:
    """Select a text encoding after a bounded-memory full-file validation.

    An ASCII prefix cannot distinguish UTF-8 from CP932. Choosing from only a
    prefix would corrupt Japanese that appears later in a large file, so every
    candidate decoder is advanced over the complete byte stream while keeping
    only decoder state in memory.
    """
    with path.open("rb") as handle:
        sample = handle.read(TEXT_ENCODING_SNIFF_BYTES)
    if sample.startswith((b"\xff\xfe", b"\xfe\xff")):
        candidates = ["utf-16"]
    elif sample and sample.count(b"\x00") / len(sample) > 0.1:
        odd_nuls = sample[1::2].count(0)
        even_nuls = sample[0::2].count(0)
        likely_utf16 = "utf-16-le" if odd_nuls >= even_nuls else "utf-16-be"
        candidates = [likely_utf16, "utf-8-sig", "cp932"]
    else:
        candidates = ["utf-8-sig", "cp932"]

    decoders = {
        encoding: codecs.getincrementaldecoder(encoding)(errors="strict")
        for encoding in candidates
    }
    with path.open("rb") as handle:
        while block := handle.read(TEXT_ENCODING_SCAN_BYTES):
            for encoding in list(decoders):
                try:
                    decoders[encoding].decode(block, final=False)
                except UnicodeDecodeError:
                    del decoders[encoding]
            if not decoders:
                return "utf-8-replacement"
    for encoding in list(decoders):
        try:
            decoders[encoding].decode(b"", final=True)
        except UnicodeDecodeError:
            del decoders[encoding]
    return next((encoding for encoding in candidates if encoding in decoders), "utf-8-replacement")


def read_text(path: Path) -> tuple[str, str]:
    """Read bounded native text while reporting the selected encoding.

    ``Probe.extract`` routes larger text-like files to the streaming reader
    before this helper is called. Keeping this helper simple preserves exact
    decoding for ordinary files without exposing the adaptive pipeline to an
    unbounded allocation.
    """
    if path.stat().st_size > MAX_DIRECT_TEXT_BYTES:
        raise ValueError("direct_text_resource_limit")
    raw = path.read_bytes()
    encoding = detect_text_encoding(raw)
    if encoding == "utf-8-replacement":
        return raw.decode("utf-8", errors="replace"), encoding
    return raw.decode(encoding), encoding


def validate_ooxml_archive(
    source: str | Path | io.BytesIO,
    *,
    required_members: frozenset[str],
) -> None:
    """Fail closed before parsing OOXML with the standard library.

    The fallback must not turn a ZIP bomb, encrypted member, duplicate/path
    ambiguity, or XML entity payload into an extraction-time resource attack.
    Validation reads XML parts only after their declared sizes and compression
    ratios have passed bounded checks.
    """
    if isinstance(source, io.BytesIO):
        archive_bytes = source.getbuffer().nbytes
        source.seek(0)
    else:
        archive_bytes = Path(source).stat().st_size
    if not 0 < archive_bytes <= MAX_OOXML_ARCHIVE_BYTES:
        raise ValueError("ooxml_archive_resource_limit")

    try:
        with zipfile.ZipFile(source) as archive:
            infos = archive.infolist()
            if not 1 <= len(infos) <= MAX_OOXML_ZIP_ENTRIES:
                raise ValueError("ooxml_archive_invalid")
            seen: set[str] = set()
            logical_paths: set[str] = set()
            total = 0
            for info in infos:
                name = info.filename
                directory = info.is_dir()
                logical_name = name[:-1] if directory else name
                pure = PurePosixPath(logical_name)
                if (
                    not logical_name
                    or "\\" in name
                    or pure.is_absolute()
                    or any(part in {"", ".", ".."} for part in pure.parts)
                    or posixpath.normpath(logical_name) != logical_name
                    or name in seen
                    or logical_name in logical_paths
                    or info.flag_bits & 0x1
                ):
                    raise ValueError("ooxml_archive_invalid")
                seen.add(name)
                logical_paths.add(logical_name)
                if directory:
                    if info.file_size != 0:
                        raise ValueError("ooxml_archive_invalid")
                    continue
                if not 0 <= info.file_size <= MAX_OOXML_MEMBER_BYTES:
                    raise ValueError("ooxml_archive_resource_limit")
                if info.file_size and info.compress_size == 0:
                    raise ValueError("ooxml_archive_invalid")
                if (
                    info.compress_size
                    and info.file_size / info.compress_size > MAX_OOXML_COMPRESSION_RATIO
                ):
                    raise ValueError("ooxml_archive_resource_limit")
                total += info.file_size
                if total > MAX_OOXML_TOTAL_BYTES:
                    raise ValueError("ooxml_archive_resource_limit")
                if name.casefold().endswith((".xml", ".rels")):
                    validate_xml_bytes(archive.read(info))
            if not required_members.issubset(seen):
                raise ValueError("ooxml_archive_invalid")
    finally:
        if isinstance(source, io.BytesIO):
            source.seek(0)


def xml_storage_encoding(value: bytes) -> tuple[str, int]:
    """Return the XML byte encoding selected by BOM/signature sniffing."""
    prefix = value[:4]
    if prefix.startswith(b"\x00\x00\xfe\xff"):
        return "utf-32-be", 4
    if prefix.startswith(b"\xff\xfe\x00\x00"):
        return "utf-32-le", 4
    if prefix.startswith(b"\xfe\xff"):
        return "utf-16-be", 2
    if prefix.startswith(b"\xff\xfe"):
        return "utf-16-le", 2
    if prefix.startswith(b"\xef\xbb\xbf"):
        return "utf-8", 3
    if prefix == b"\x00\x00\x00<":
        return "utf-32-be", 0
    if prefix == b"<\x00\x00\x00":
        return "utf-32-le", 0
    if prefix == b"\x00<\x00?":
        return "utf-16-be", 0
    if prefix == b"<\x00?\x00":
        return "utf-16-le", 0
    return "ascii-compatible", 0


def normalized_xml_encoding(value: str) -> str:
    return re.sub(r"[^a-z0-9]", "", value.casefold())


def validate_xml_bytes(value: bytes) -> None:
    """Reject DTD/entity declarations before ElementTree sees any XML bytes.

    OOXML normally uses UTF-8, but XML also permits BOM/signature-selected
    UTF-16 and UTF-32.  Searching only ASCII bytes leaves those encodings as a
    bypass, so wide encodings are scanned with their exact code-unit layout.
    The XML declaration, when present, must agree with the detected family.
    """
    encoding, bom_bytes = xml_storage_encoding(value)
    if encoding == "ascii-compatible":
        if b"\x00" in value[:512]:
            raise ValueError("ooxml_xml_encoding_invalid")
        upper = value.upper()
        if any(marker in upper for marker in OOXML_FORBIDDEN_XML):
            raise ValueError("ooxml_xml_unsafe")
        prolog = value[:2048].decode("ascii", errors="ignore")
        declaration = XML_DECLARATION_ENCODING.search(prolog)
        if declaration:
            declared = normalized_xml_encoding(declaration.group(2))
            if declared in {"utf16", "utf16le", "utf16be", "utf32", "utf32le", "utf32be"}:
                raise ValueError("ooxml_xml_encoding_invalid")
        return

    if encoding == "utf-8":
        if any(marker in value.upper() for marker in OOXML_FORBIDDEN_XML):
            raise ValueError("ooxml_xml_unsafe")
    elif any(marker.decode("ascii").encode(encoding) in value for marker in OOXML_FORBIDDEN_XML):
        raise ValueError("ooxml_xml_unsafe")
    try:
        prolog = value[bom_bytes:2048].decode(encoding, errors="strict")
    except UnicodeDecodeError as exc:
        raise ValueError("ooxml_xml_encoding_invalid") from exc
    declaration = XML_DECLARATION_ENCODING.search(prolog)
    if not declaration:
        return
    declared = normalized_xml_encoding(declaration.group(2))
    family = (
        "utf32" if encoding.startswith("utf-32")
        else "utf16" if encoding.startswith("utf-16")
        else "utf8"
    )
    compatible = {family, normalized_xml_encoding(encoding)}
    if declared not in compatible:
        raise ValueError("ooxml_xml_encoding_invalid")


def _ooxml_relationship_source(member: str) -> str | None:
    """Map an OOXML ``*.rels`` member to the part it describes."""
    parts = PurePosixPath(member).parts
    if len(parts) < 2 or parts[-2] != "_rels" or not parts[-1].endswith(".rels"):
        return None
    if parts == ("_rels", ".rels"):
        return ""
    source_name = parts[-1][:-5]
    if not source_name:
        return None
    prefix = parts[:-2]
    return PurePosixPath(*prefix, source_name).as_posix()


def _resolve_ooxml_target(source_part: str, target: str) -> str | None:
    """Resolve an internal OOXML relationship target without filesystem access."""
    if not target or "\\" in target or "\x00" in target:
        return None
    if target.startswith("/"):
        candidate = target[1:]
    else:
        candidate = posixpath.join(posixpath.dirname(source_part), target)
    normalized = posixpath.normpath(candidate)
    pure = PurePosixPath(normalized)
    if (
        not normalized
        or normalized == "."
        or pure.is_absolute()
        or normalized == ".."
        or normalized.startswith("../")
        or any(part in {"", ".", ".."} for part in pure.parts)
    ):
        return None
    return normalized


def _xml_name(tag: str) -> tuple[str | None, str]:
    if tag.startswith("{") and "}" in tag:
        namespace, local_name = tag[1:].split("}", 1)
        return namespace, local_name
    return None, tag


def _ooxml_relationship_kind(value: str) -> str | None:
    for namespace in OOXML_DOCUMENT_RELATIONSHIP_NAMESPACES:
        prefix = namespace + "/"
        if value.startswith(prefix):
            kind = value[len(prefix):]
            if kind and "/" not in kind:
                return kind
    return None


def _ooxml_part_root_matches(
    archive: zipfile.ZipFile,
    source_part: str,
    namespaces: set[str],
    local_names: set[str],
) -> bool:
    if source_part not in archive.namelist():
        return False
    raw = archive.read(source_part)
    validate_xml_bytes(raw)
    try:
        root = ElementTree.fromstring(raw)
    except ElementTree.ParseError as exc:
        raise ValueError("ooxml_source_xml_invalid") from exc
    namespace, local_name = _xml_name(root.tag)
    return namespace in namespaces and local_name in local_names


def _word_visual_source_part_allowed(
    archive: zipfile.ZipFile,
    source_part: str,
) -> bool:
    if source_part == "word/document.xml":
        expected_roots = {"document"}
    elif re.fullmatch(r"word/header[0-9]+\.xml", source_part):
        expected_roots = {"hdr"}
    elif re.fullmatch(r"word/footer[0-9]+\.xml", source_part):
        expected_roots = {"ftr"}
    else:
        return False
    return _ooxml_part_root_matches(
        archive,
        source_part,
        OOXML_WORDPROCESSING_NAMESPACES,
        expected_roots,
    )


def _ooxml_relationships(
    archive: zipfile.ZipFile,
) -> dict[str, list[dict[str, str]]]:
    relationships: dict[str, list[dict[str, str]]] = {}
    for member in sorted(archive.namelist()):
        source_part = _ooxml_relationship_source(member)
        if source_part is None:
            continue
        raw = archive.read(member)
        validate_xml_bytes(raw)
        try:
            root = ElementTree.fromstring(raw)
        except ElementTree.ParseError as exc:
            raise ValueError("ooxml_relationship_xml_invalid") from exc
        root_namespace, root_name = _xml_name(root.tag)
        if (
            root_namespace not in OOXML_PACKAGE_RELATIONSHIP_NAMESPACES
            or root_name != "Relationships"
        ):
            raise ValueError("ooxml_relationship_root_invalid")
        rows: list[dict[str, str]] = []
        seen_ids: set[str] = set()
        for item in root:
            item_namespace, item_name = _xml_name(item.tag)
            if (
                item_namespace not in OOXML_PACKAGE_RELATIONSHIP_NAMESPACES
                or item_name != "Relationship"
            ):
                continue
            relationship_id = item.attrib.get("Id", "")
            target = item.attrib.get("Target", "")
            relationship_type = item.attrib.get("Type", "")
            target_mode = item.attrib.get("TargetMode", "")
            if (
                not relationship_id
                or relationship_id in seen_ids
                or not target
                or not relationship_type
                or target_mode.casefold() not in {"", "internal", "external"}
            ):
                raise ValueError("ooxml_relationship_invalid")
            seen_ids.add(relationship_id)
            rows.append({
                "id": relationship_id,
                "target": target,
                "type": relationship_type,
                "target_mode": target_mode,
            })
        relationships[source_part] = rows
    return relationships


def _ooxml_used_relationship_ids(
    archive: zipfile.ZipFile,
    source_part: str,
) -> dict[str, int]:
    if source_part == "":
        return {
            item["id"]: 1
            for item in _ooxml_relationships(archive).get("", [])
        }
    if (
        source_part not in archive.namelist()
        or not source_part.casefold().endswith((".xml", ".vml"))
    ):
        return {}
    raw = archive.read(source_part)
    validate_xml_bytes(raw)
    try:
        root = ElementTree.fromstring(raw)
    except ElementTree.ParseError as exc:
        raise ValueError("ooxml_source_xml_invalid") from exc
    counts: dict[str, int] = {}
    for relationship_use in _ooxml_relationship_uses_from_root(root):
        relationship_id = relationship_use["relationship_id"]
        counts[relationship_id] = counts.get(relationship_id, 0) + 1
    return counts


def _ooxml_relationship_uses_from_root(
    root: ElementTree.Element,
) -> list[dict[str, str]]:
    """Return relationship-id occurrences with a conservative visual role.

    OOXML relationship ids are ordinary attribute values.  Keeping the XML
    ancestor path lets the visual projector distinguish a PowerPoint picture
    from a background or a shape fill without treating every package image as
    displayed content.
    """
    uses: list[dict[str, str]] = []

    def visit(item: ElementTree.Element, ancestors: tuple[str, ...]) -> None:
        item_namespace, local = _xml_name(item.tag)
        path = (*ancestors, local)
        lowered = {value.casefold() for value in path}
        if "bg" in lowered or "background" in lowered:
            usage_kind = "background"
        elif "pic" in lowered:
            usage_kind = "picture"
        elif "blipfill" in lowered:
            usage_kind = "shape_fill"
        elif local.casefold() == "imagedata":
            usage_kind = "picture"
        else:
            usage_kind = "other"
        allowed_attribute_names = OOXML_RELATIONSHIP_CARRIERS.get(
            (str(item_namespace), local),
            frozenset(),
        )
        for key, value in item.attrib.items():
            namespace = (
                key[1:].split("}", 1)[0]
                if key.startswith("{") and "}" in key else None
            )
            attribute_name = _xml_name(key)[1]
            if (
                item_namespace in OOXML_RELATIONSHIP_CARRIER_NAMESPACES
                and attribute_name in allowed_attribute_names
                and
                namespace in OOXML_DOCUMENT_RELATIONSHIP_NAMESPACES
                and isinstance(value, str)
                and value
            ):
                uses.append({
                    "relationship_id": value,
                    "usage_kind": usage_kind,
                })
        for child_item in item:
            visit(child_item, path)

    visit(root, ())
    return uses


def _ooxml_relationship_uses(
    archive: zipfile.ZipFile,
    source_part: str,
) -> list[dict[str, str]]:
    if source_part == "" or source_part not in archive.namelist():
        return []
    if not source_part.casefold().endswith((".xml", ".vml")):
        return []
    raw = archive.read(source_part)
    validate_xml_bytes(raw)
    try:
        root = ElementTree.fromstring(raw)
    except ElementTree.ParseError as exc:
        raise ValueError("ooxml_source_xml_invalid") from exc
    return _ooxml_relationship_uses_from_root(root)


def _direct_xml_children(
    parent: ElementTree.Element,
    namespaces: set[str],
    local_name: str,
) -> list[ElementTree.Element]:
    return [
        child
        for child in parent
        if _xml_name(child.tag)[0] in namespaces
        and _xml_name(child.tag)[1] == local_name
    ]


def _ooxml_xml_root(
    archive: zipfile.ZipFile,
    member: str,
    *,
    namespaces: set[str],
    local_names: set[str],
) -> ElementTree.Element:
    """Read one validated OOXML part and require its exact root contract."""
    if member not in archive.namelist():
        raise ValueError("ooxml_required_part_missing")
    raw = archive.read(member)
    validate_xml_bytes(raw)
    try:
        root = ElementTree.fromstring(raw)
    except ElementTree.ParseError as exc:
        raise ValueError("ooxml_source_xml_invalid") from exc
    namespace, local_name = _xml_name(root.tag)
    if namespace not in namespaces or local_name not in local_names:
        raise ValueError("ooxml_source_root_invalid")
    return root


def _require_ooxml_office_document_binding(
    archive: zipfile.ZipFile,
    relationships: dict[str, list[dict[str, str]]],
    expected_member: str,
) -> None:
    """Require the package root to bind exactly one expected main part."""
    candidates = [
        row
        for row in relationships.get("", [])
        if _ooxml_relationship_kind(row["type"]) == "officeDocument"
    ]
    if len(candidates) != 1 or candidates[0]["target_mode"].casefold() == "external":
        raise ValueError("ooxml_office_document_binding_invalid")
    target = _resolve_ooxml_target("", candidates[0]["target"])
    if target != expected_member or target not in archive.namelist():
        raise ValueError("ooxml_office_document_binding_invalid")


def _xml_attribute(
    item: ElementTree.Element,
    namespaces: set[str],
    local_name: str,
) -> str | None:
    values = [
        value
        for key, value in item.attrib.items()
        if _xml_name(key)[0] in namespaces
        and _xml_name(key)[1] == local_name
        and isinstance(value, str)
    ]
    if len(values) > 1:
        raise ValueError("ooxml_attribute_ambiguous")
    return values[0] if values else None


def _word_paragraph_text(paragraph: ElementTree.Element) -> str:
    """Read visible Word text from canonical run/container paths only."""
    paragraph_namespace, paragraph_name = _xml_name(paragraph.tag)
    if (
        paragraph_namespace not in OOXML_WORDPROCESSING_NAMESPACES
        or paragraph_name != "p"
    ):
        raise ValueError("ooxml_word_paragraph_invalid")
    fragments: list[str] = []
    containers = {
        "customXml", "fldSimple", "hyperlink", "ins", "moveTo", "sdt",
        "sdtContent", "smartTag",
    }

    def visit(container: ElementTree.Element) -> None:
        for child_item in container:
            namespace, local_name = _xml_name(child_item.tag)
            if namespace not in OOXML_WORDPROCESSING_NAMESPACES:
                continue
            if local_name == "r":
                for run_item in child_item:
                    run_namespace, run_name = _xml_name(run_item.tag)
                    if run_namespace not in OOXML_WORDPROCESSING_NAMESPACES:
                        continue
                    if run_name == "t" and run_item.text:
                        fragments.append(run_item.text)
                    elif run_name == "tab":
                        fragments.append("\t")
                    elif run_name in {"br", "cr"}:
                        fragments.append("\n")
            elif local_name in containers:
                visit(child_item)

    visit(paragraph)
    return "".join(fragments)


def _word_paragraph_style_id(paragraph: ElementTree.Element) -> str | None:
    properties = _direct_xml_children(
        paragraph, OOXML_WORDPROCESSING_NAMESPACES, "pPr"
    )
    if len(properties) > 1:
        raise ValueError("ooxml_word_paragraph_properties_ambiguous")
    if not properties:
        return None
    styles = _direct_xml_children(
        properties[0], OOXML_WORDPROCESSING_NAMESPACES, "pStyle"
    )
    if len(styles) > 1:
        raise ValueError("ooxml_word_paragraph_style_ambiguous")
    return (
        _xml_attribute(styles[0], OOXML_WORDPROCESSING_NAMESPACES, "val")
        if styles else None
    )


def _word_paragraph_has_outline_level(paragraph: ElementTree.Element) -> bool:
    properties = _direct_xml_children(
        paragraph, OOXML_WORDPROCESSING_NAMESPACES, "pPr"
    )
    if not properties:
        return False
    outline = _direct_xml_children(
        properties[0], OOXML_WORDPROCESSING_NAMESPACES, "outlineLvl"
    )
    if len(outline) > 1:
        raise ValueError("ooxml_word_outline_level_ambiguous")
    if not outline:
        return False
    value = _xml_attribute(
        outline[0], OOXML_WORDPROCESSING_NAMESPACES, "val"
    )
    if value is None:
        return True
    try:
        return 0 <= int(value) <= 8
    except ValueError:
        return False


def _word_style_catalog(
    archive: zipfile.ZipFile,
    relationships: dict[str, list[dict[str, str]]],
) -> dict[str, dict[str, str | None]]:
    """Return relationship-bound paragraph style names and inheritance."""
    style_rows = [
        row for row in relationships.get("word/document.xml", [])
        if _ooxml_relationship_kind(row["type"]) == "styles"
        and row["target_mode"].casefold() != "external"
    ]
    if len(style_rows) > 1:
        raise ValueError("ooxml_word_styles_binding_ambiguous")
    if not style_rows:
        return {}
    member = _resolve_ooxml_target(
        "word/document.xml", style_rows[0]["target"]
    )
    if member is None:
        raise ValueError("ooxml_word_styles_binding_invalid")
    root = _ooxml_xml_root(
        archive,
        member,
        namespaces=OOXML_WORDPROCESSING_NAMESPACES,
        local_names={"styles"},
    )
    styles: dict[str, dict[str, str | None]] = {}
    for item in _direct_xml_children(
        root, OOXML_WORDPROCESSING_NAMESPACES, "style"
    ):
        style_type = _xml_attribute(
            item, OOXML_WORDPROCESSING_NAMESPACES, "type"
        )
        style_id = _xml_attribute(
            item, OOXML_WORDPROCESSING_NAMESPACES, "styleId"
        )
        if style_type != "paragraph" or not style_id:
            continue
        if style_id in styles:
            raise ValueError("ooxml_word_style_id_ambiguous")
        names = _direct_xml_children(
            item, OOXML_WORDPROCESSING_NAMESPACES, "name"
        )
        based_on = _direct_xml_children(
            item, OOXML_WORDPROCESSING_NAMESPACES, "basedOn"
        )
        if len(names) > 1 or len(based_on) > 1:
            raise ValueError("ooxml_word_style_definition_ambiguous")
        styles[style_id] = {
            "name": (
                _xml_attribute(
                    names[0], OOXML_WORDPROCESSING_NAMESPACES, "val"
                ) if names else None
            ),
            "based_on": (
                _xml_attribute(
                    based_on[0], OOXML_WORDPROCESSING_NAMESPACES, "val"
                ) if based_on else None
            ),
        }
    return styles


def _word_style_is_heading(
    style_id: str | None,
    styles: dict[str, dict[str, str | None]],
) -> bool:
    seen: set[str] = set()
    current = style_id
    while current and current not in seen:
        seen.add(current)
        row = styles.get(current, {})
        candidates = [current, row.get("name")]
        if any(
            isinstance(value, str)
            and re.match(r"^heading(?:\s|[-_])*[0-9]*$", value.strip(), re.I)
            for value in candidates
        ):
            return True
        based_on = row.get("based_on")
        current = based_on if isinstance(based_on, str) else None
    return False


def _word_block_children(
    container: ElementTree.Element,
) -> list[ElementTree.Element]:
    """Return body/cell blocks, descending only through canonical SDT content."""
    blocks: list[ElementTree.Element] = []
    for item in container:
        namespace, local_name = _xml_name(item.tag)
        if namespace not in OOXML_WORDPROCESSING_NAMESPACES:
            continue
        if local_name in {"p", "tbl"}:
            blocks.append(item)
        elif local_name in {"sdt", "sdtContent", "customXml"}:
            blocks.extend(_word_block_children(item))
    return blocks


def _drawing_paragraph_text(paragraph: ElementTree.Element) -> str:
    fragments: list[str] = []
    for item in paragraph:
        namespace, local_name = _xml_name(item.tag)
        if namespace not in OOXML_DRAWING_NAMESPACES:
            continue
        if local_name in {"r", "fld"}:
            texts = _direct_xml_children(
                item, OOXML_DRAWING_NAMESPACES, "t"
            )
            if len(texts) > 1:
                raise ValueError("ooxml_drawing_run_text_ambiguous")
            if texts and texts[0].text:
                fragments.append(texts[0].text)
        elif local_name == "br":
            fragments.append("\n")
    return "".join(fragments)


def _drawing_text_body_text(body: ElementTree.Element) -> str:
    paragraphs = _direct_xml_children(
        body, OOXML_DRAWING_NAMESPACES, "p"
    )
    return "\n".join(_drawing_paragraph_text(item) for item in paragraphs)


def _pptx_shape_tree(root: ElementTree.Element) -> ElementTree.Element | None:
    root_namespace, root_name = _xml_name(root.tag)
    if (
        root_namespace not in OOXML_PRESENTATION_NAMESPACES
        or root_name != "sld"
    ):
        raise ValueError("ooxml_slide_root_invalid")
    common = _direct_xml_children(
        root, OOXML_PRESENTATION_NAMESPACES, "cSld"
    )
    if len(common) > 1:
        raise ValueError("ooxml_slide_content_ambiguous")
    if not common:
        return None
    trees = _direct_xml_children(
        common[0], OOXML_PRESENTATION_NAMESPACES, "spTree"
    )
    if len(trees) > 1:
        raise ValueError("ooxml_slide_shape_tree_ambiguous")
    return trees[0] if trees else None


def _pptx_shape_identity(
    shape: ElementTree.Element,
) -> tuple[str | None, str | None]:
    _, local_name = _xml_name(shape.tag)
    containers = {
        "sp": "nvSpPr",
        "pic": "nvPicPr",
        "graphicFrame": "nvGraphicFramePr",
        "cxnSp": "nvCxnSpPr",
        "grpSp": "nvGrpSpPr",
    }
    container_name = containers.get(local_name)
    if container_name is None:
        return None, None
    candidates = _direct_xml_children(
        shape, OOXML_PRESENTATION_NAMESPACES, container_name
    )
    if len(candidates) > 1:
        raise ValueError("ooxml_slide_shape_identity_ambiguous")
    if not candidates:
        return None, None
    nonvisual = _direct_xml_children(
        candidates[0], OOXML_PRESENTATION_NAMESPACES, "cNvPr"
    )
    if len(nonvisual) > 1:
        raise ValueError("ooxml_slide_shape_identity_ambiguous")
    if not nonvisual:
        return None, None
    return nonvisual[0].attrib.get("id"), nonvisual[0].attrib.get("name")


def _pptx_shape_geometry(
    shape: ElementTree.Element,
    *,
    nested_in_group: bool,
) -> tuple[dict[str, Any] | None, bool]:
    _, local_name = _xml_name(shape.tag)
    if local_name == "graphicFrame":
        transforms = _direct_xml_children(
            shape, OOXML_PRESENTATION_NAMESPACES, "xfrm"
        )
    else:
        property_name = "grpSpPr" if local_name == "grpSp" else "spPr"
        properties = _direct_xml_children(
            shape, OOXML_PRESENTATION_NAMESPACES, property_name
        )
        if len(properties) > 1:
            raise ValueError("ooxml_slide_shape_geometry_ambiguous")
        transforms = (
            _direct_xml_children(
                properties[0], OOXML_DRAWING_NAMESPACES, "xfrm"
            ) if properties else []
        )
    if len(transforms) > 1:
        raise ValueError("ooxml_slide_shape_geometry_ambiguous")
    if not transforms:
        return None, False
    offsets = _direct_xml_children(
        transforms[0], OOXML_DRAWING_NAMESPACES, "off"
    )
    extents = _direct_xml_children(
        transforms[0], OOXML_DRAWING_NAMESPACES, "ext"
    )
    if len(offsets) > 1 or len(extents) > 1:
        raise ValueError("ooxml_slide_shape_geometry_ambiguous")
    if not offsets or not extents:
        return None, False
    try:
        x = int(offsets[0].attrib["x"])
        y = int(offsets[0].attrib["y"])
        width = int(extents[0].attrib["cx"])
        height = int(extents[0].attrib["cy"])
    except (KeyError, ValueError) as exc:
        raise ValueError("ooxml_slide_shape_geometry_invalid") from exc
    if width < 0 or height < 0:
        raise ValueError("ooxml_slide_shape_geometry_invalid")
    return {
        "coordinate_space": "slide" if not nested_in_group else "other",
        "unit": "emu",
        "x": x,
        "y": y,
        "width": width,
        "height": height,
    }, nested_in_group


def _pptx_shape_text(shape: ElementTree.Element) -> str:
    bodies = _direct_xml_children(
        shape, OOXML_PRESENTATION_NAMESPACES, "txBody"
    )
    if len(bodies) > 1:
        raise ValueError("ooxml_slide_shape_text_ambiguous")
    return _drawing_text_body_text(bodies[0]) if bodies else ""


def _pptx_table(
    frame: ElementTree.Element,
) -> ElementTree.Element | None:
    graphics = _direct_xml_children(
        frame, OOXML_DRAWING_NAMESPACES, "graphic"
    )
    if len(graphics) > 1:
        raise ValueError("ooxml_graphic_carrier_ambiguous")
    if not graphics:
        return None
    data = _direct_xml_children(
        graphics[0], OOXML_DRAWING_NAMESPACES, "graphicData"
    )
    if len(data) > 1:
        raise ValueError("ooxml_graphic_data_ambiguous")
    if not data or data[0].attrib.get("uri") not in OOXML_DRAWING_TABLE_URIS:
        return None
    tables = _direct_xml_children(
        data[0], OOXML_DRAWING_NAMESPACES, "tbl"
    )
    if len(tables) > 1:
        raise ValueError("ooxml_drawing_table_ambiguous")
    return tables[0] if tables else None


def _pptx_table_rows(
    table: ElementTree.Element,
) -> list[list[str]]:
    values: list[list[str]] = []
    for row in _direct_xml_children(
        table, OOXML_DRAWING_NAMESPACES, "tr"
    ):
        cells: list[str] = []
        for cell in _direct_xml_children(
            row, OOXML_DRAWING_NAMESPACES, "tc"
        ):
            bodies = _direct_xml_children(
                cell, OOXML_DRAWING_NAMESPACES, "txBody"
            )
            if len(bodies) > 1:
                raise ValueError("ooxml_drawing_table_cell_ambiguous")
            cells.append(
                _drawing_text_body_text(bodies[0]) if bodies else ""
            )
        values.append(cells)
    return values


def _relationship_attribute(
    item: ElementTree.Element,
    local_name: str,
) -> str | None:
    values = [
        value
        for key, value in item.attrib.items()
        if _xml_name(key)[0] in OOXML_DOCUMENT_RELATIONSHIP_NAMESPACES
        and _xml_name(key)[1] == local_name
        and isinstance(value, str)
        and value
    ]
    if len(values) > 1:
        raise ValueError("ooxml_relationship_attribute_ambiguous")
    return values[0] if values else None


def _pptx_slide_context(
    archive: zipfile.ZipFile,
    relationships: dict[str, list[dict[str, str]]],
) -> dict[str, tuple[int, ...]]:
    """Map direct slide parts to their presentation order.

    Layout/master visibility depends on inheritance, suppression flags, and
    background overrides.  Those parts are deliberately not expanded across
    slides until a display resolver exists.
    """
    presentation_part = "ppt/presentation.xml"
    if presentation_part not in archive.namelist():
        return {}
    try:
        root = ElementTree.fromstring(archive.read(presentation_part))
    except ElementTree.ParseError as exc:
        raise ValueError("ooxml_presentation_xml_invalid") from exc
    root_namespace, root_name = _xml_name(root.tag)
    if (
        root_namespace not in OOXML_PRESENTATION_NAMESPACES
        or root_name != "presentation"
    ):
        raise ValueError("ooxml_presentation_root_invalid")
    presentation_targets = {
        row["id"]: _resolve_ooxml_target(presentation_part, row["target"])
        for row in relationships.get(presentation_part, [])
        if row["target_mode"].casefold() != "external"
        and _ooxml_relationship_kind(row["type"]) == "slide"
    }
    slide_parts: list[str] = []
    slide_lists = _direct_xml_children(
        root, OOXML_PRESENTATION_NAMESPACES, "sldIdLst"
    )
    if len(slide_lists) > 1:
        raise ValueError("ooxml_slide_list_ambiguous")
    seen_slide_ids: set[str] = set()
    seen_relationship_ids: set[str] = set()
    seen_targets: set[str] = set()
    if slide_lists:
        for item in _direct_xml_children(
            slide_lists[0], OOXML_PRESENTATION_NAMESPACES, "sldId"
        ):
            slide_id = item.attrib.get("id")
            relationship_id = _relationship_attribute(item, "id")
            if (
                not slide_id
                or slide_id in seen_slide_ids
                or not relationship_id
                or relationship_id in seen_relationship_ids
            ):
                raise ValueError("ooxml_slide_binding_ambiguous")
            target = presentation_targets.get(relationship_id)
            if (
                target is None
                or not target.startswith("ppt/slides/")
                or target in seen_targets
            ):
                raise ValueError("ooxml_slide_binding_invalid")
            if not _ooxml_part_root_matches(
                archive,
                target,
                OOXML_PRESENTATION_NAMESPACES,
                {"sld"},
            ):
                raise ValueError("ooxml_slide_root_invalid")
            seen_slide_ids.add(slide_id)
            seen_relationship_ids.add(relationship_id)
            seen_targets.add(target)
            slide_parts.append(target)

    return {
        slide_part: (slide_number,)
        for slide_number, slide_part in enumerate(slide_parts, 1)
    }


def _count_unresolved_pptx_inherited_media(
    archive: zipfile.ZipFile,
) -> int:
    """Count used layout/master images without expanding them across slides."""
    relationships = _ooxml_relationships(archive)
    count = 0
    for source_part, rows in relationships.items():
        if not source_part.startswith(
            ("ppt/slideLayouts/", "ppt/slideMasters/")
        ):
            continue
        used = _ooxml_used_relationship_ids(archive, source_part)
        for row in rows:
            if (
                row["target_mode"].casefold() != "external"
                and _ooxml_relationship_kind(row["type"]) == "image"
                and row["id"] in used
            ):
                target = _resolve_ooxml_target(source_part, row["target"])
                if target is not None and target.startswith("ppt/media/"):
                    count += used[row["id"]]
    return count


def _xlsx_sheet_context(
    archive: zipfile.ZipFile,
    relationships: dict[str, list[dict[str, str]]],
) -> tuple[dict[str, str], dict[str, str]]:
    """Return worksheet names and drawing-to-worksheet bindings."""
    sheet_names: dict[str, str] = {}
    drawing_sheets: dict[str, str] = {}
    workbook_part = "xl/workbook.xml"
    if workbook_part not in archive.namelist():
        return sheet_names, drawing_sheets
    rel_targets = {
        row["id"]: _resolve_ooxml_target(workbook_part, row["target"])
        for row in relationships.get(workbook_part, [])
        if not row["target_mode"].casefold() == "external"
        and _ooxml_relationship_kind(row["type"]) == "worksheet"
    }
    try:
        workbook_root = ElementTree.fromstring(archive.read(workbook_part))
    except ElementTree.ParseError as exc:
        raise ValueError("ooxml_workbook_xml_invalid") from exc
    workbook_namespace, workbook_name = _xml_name(workbook_root.tag)
    if (
        workbook_namespace not in OOXML_SPREADSHEET_NAMESPACES
        or workbook_name != "workbook"
    ):
        raise ValueError("ooxml_workbook_root_invalid")
    sheet_lists = _direct_xml_children(
        workbook_root, OOXML_SPREADSHEET_NAMESPACES, "sheets"
    )
    if len(sheet_lists) > 1:
        raise ValueError("ooxml_sheet_list_ambiguous")
    seen_sheet_ids: set[str] = set()
    seen_relationship_ids: set[str] = set()
    seen_names: set[str] = set()
    if sheet_lists:
        for sheet in _direct_xml_children(
            sheet_lists[0], OOXML_SPREADSHEET_NAMESPACES, "sheet"
        ):
            relationship_id = _relationship_attribute(sheet, "id")
            name = sheet.attrib.get("name")
            sheet_id = sheet.attrib.get("sheetId")
            normalized_name = name.casefold() if name else ""
            if (
                not relationship_id
                or relationship_id in seen_relationship_ids
                or not name
                or normalized_name in seen_names
                or not sheet_id
                or sheet_id in seen_sheet_ids
            ):
                raise ValueError("ooxml_sheet_binding_ambiguous")
            target = rel_targets.get(relationship_id)
            if target is None or target in sheet_names:
                raise ValueError("ooxml_sheet_binding_invalid")
            seen_relationship_ids.add(relationship_id)
            seen_names.add(normalized_name)
            seen_sheet_ids.add(sheet_id)
            sheet_names[target] = name
    for sheet_part, sheet_name in sheet_names.items():
        if not _ooxml_part_root_matches(
            archive,
            sheet_part,
            OOXML_SPREADSHEET_NAMESPACES,
            {"worksheet"},
        ):
            raise ValueError("ooxml_worksheet_root_invalid")
        try:
            sheet_root = ElementTree.fromstring(archive.read(sheet_part))
        except ElementTree.ParseError as exc:
            raise ValueError("ooxml_worksheet_xml_invalid") from exc
        drawing_uses = _direct_xml_children(
            sheet_root, OOXML_SPREADSHEET_NAMESPACES, "drawing"
        )
        if len(drawing_uses) > 1:
            raise ValueError("ooxml_worksheet_drawing_ambiguous")
        drawing_relationship_ids = {
            relationship_id
            for drawing in drawing_uses
            if (relationship_id := _relationship_attribute(drawing, "id"))
        }
        if len(drawing_relationship_ids) != len(drawing_uses):
            raise ValueError("ooxml_worksheet_drawing_invalid")
        for row in relationships.get(sheet_part, []):
            if (
                row["id"] not in drawing_relationship_ids
                or row["target_mode"].casefold() == "external"
            ):
                continue
            target = _resolve_ooxml_target(sheet_part, row["target"])
            if target and _ooxml_relationship_kind(row["type"]) == "drawing":
                previous = drawing_sheets.get(target)
                if previous is not None and previous != sheet_name:
                    raise ValueError("ooxml_drawing_sheet_binding_ambiguous")
                drawing_sheets[target] = sheet_name
                drawing_relationship_ids.remove(row["id"])
        if drawing_relationship_ids:
            raise ValueError("ooxml_worksheet_drawing_invalid")
    return sheet_names, drawing_sheets


def _xlsx_column_name(index: int) -> str:
    value = index + 1
    result = ""
    while value:
        value, remainder = divmod(value - 1, 26)
        result = chr(65 + remainder) + result
    return result


def _xlsx_drawing_cells(
    archive: zipfile.ZipFile,
    drawing_part: str,
    relationship_id: str,
) -> list[str | None]:
    if drawing_part not in archive.namelist():
        return []
    try:
        root = ElementTree.fromstring(archive.read(drawing_part))
    except ElementTree.ParseError as exc:
        raise ValueError("ooxml_drawing_xml_invalid") from exc
    root_namespace, root_name = _xml_name(root.tag)
    if (
        root_namespace not in OOXML_SPREADSHEET_DRAWING_NAMESPACES
        or root_name != "wsDr"
    ):
        raise ValueError("ooxml_drawing_root_invalid")
    cells: list[str | None] = []
    for anchor in root:
        anchor_namespace, anchor_name = _xml_name(anchor.tag)
        if (
            anchor_namespace not in OOXML_SPREADSHEET_DRAWING_NAMESPACES
            or anchor_name
            not in {"oneCellAnchor", "twoCellAnchor", "absoluteAnchor"}
        ):
            continue
        relationship_occurrences = 0
        for item in anchor.iter():
            for key, value in item.attrib.items():
                namespace = (
                    key[1:].split("}", 1)[0]
                    if key.startswith("{") and "}" in key else None
                )
                if (
                    namespace in OOXML_DOCUMENT_RELATIONSHIP_NAMESPACES
                    and value == relationship_id
                ):
                    relationship_occurrences += 1
        if relationship_occurrences == 0:
            continue
        origin = next(
            (
                item for item in anchor
                if _xml_name(item.tag)[0]
                in OOXML_SPREADSHEET_DRAWING_NAMESPACES
                and _xml_name(item.tag)[1] == "from"
            ),
            None,
        )
        row_value: int | None = None
        column_value: int | None = None
        if origin is not None:
            for item in origin:
                item_namespace, local = _xml_name(item.tag)
                if item_namespace not in OOXML_SPREADSHEET_DRAWING_NAMESPACES:
                    continue
                if local == "row" and item.text is not None:
                    row_value = int(item.text)
                elif local == "col" and item.text is not None:
                    column_value = int(item.text)
        if row_value is not None and column_value is not None and row_value >= 0 and column_value >= 0:
            cell = f"{_xlsx_column_name(column_value)}{row_value + 1}"
        else:
            cell = None
        cells.extend([cell] * relationship_occurrences)
    return cells


def _pptx_visible_graphic_frames(
    root: ElementTree.Element,
) -> list[ElementTree.Element]:
    root_namespace, root_name = _xml_name(root.tag)
    if (
        root_namespace not in OOXML_PRESENTATION_NAMESPACES
        or root_name != "sld"
    ):
        raise ValueError("ooxml_slide_root_invalid")
    common_slides = _direct_xml_children(
        root, OOXML_PRESENTATION_NAMESPACES, "cSld"
    )
    if len(common_slides) > 1:
        raise ValueError("ooxml_slide_content_ambiguous")
    if not common_slides:
        return []
    shape_trees = _direct_xml_children(
        common_slides[0], OOXML_PRESENTATION_NAMESPACES, "spTree"
    )
    if len(shape_trees) > 1:
        raise ValueError("ooxml_slide_shape_tree_ambiguous")
    if not shape_trees:
        return []

    frames: list[ElementTree.Element] = []

    def collect(container: ElementTree.Element) -> None:
        for child in container:
            namespace, local_name = _xml_name(child.tag)
            if namespace not in OOXML_PRESENTATION_NAMESPACES:
                continue
            if local_name == "graphicFrame":
                frames.append(child)
            elif local_name == "grpSp":
                collect(child)

    collect(shape_trees[0])
    return frames


def _graphic_frame_relationship(
    frame: ElementTree.Element,
    *,
    payload_namespaces: set[str],
    payload_local_name: str,
    relationship_attribute_name: str,
) -> str | None:
    graphics = _direct_xml_children(frame, OOXML_DRAWING_NAMESPACES, "graphic")
    if len(graphics) > 1:
        raise ValueError("ooxml_graphic_carrier_ambiguous")
    if not graphics:
        return None
    graphic_data = _direct_xml_children(
        graphics[0], OOXML_DRAWING_NAMESPACES, "graphicData"
    )
    if len(graphic_data) > 1:
        raise ValueError("ooxml_graphic_data_ambiguous")
    if not graphic_data or graphic_data[0].attrib.get("uri") not in payload_namespaces:
        return None
    payloads = _direct_xml_children(
        graphic_data[0], payload_namespaces, payload_local_name
    )
    if len(payloads) > 1:
        raise ValueError("ooxml_graphic_payload_ambiguous")
    if not payloads:
        return None
    relationship_id = _relationship_attribute(
        payloads[0], relationship_attribute_name
    )
    if not relationship_id:
        raise ValueError("ooxml_graphic_relationship_missing")
    return relationship_id


def _pptx_graphic_relationship_ids(
    archive: zipfile.ZipFile,
    slide_part: str,
    *,
    payload_namespaces: set[str],
    payload_local_name: str,
    relationship_attribute_name: str,
) -> list[str]:
    if slide_part not in archive.namelist():
        return []
    raw = archive.read(slide_part)
    validate_xml_bytes(raw)
    try:
        root = ElementTree.fromstring(raw)
    except ElementTree.ParseError as exc:
        raise ValueError("ooxml_slide_xml_invalid") from exc
    relationship_ids: list[str] = []
    for frame in _pptx_visible_graphic_frames(root):
        relationship_id = _graphic_frame_relationship(
            frame,
            payload_namespaces=payload_namespaces,
            payload_local_name=payload_local_name,
            relationship_attribute_name=relationship_attribute_name,
        )
        if relationship_id:
            relationship_ids.append(relationship_id)
    return relationship_ids


def _xlsx_anchor_cell(anchor: ElementTree.Element) -> str | None:
    origins = _direct_xml_children(
        anchor, OOXML_SPREADSHEET_DRAWING_NAMESPACES, "from"
    )
    if len(origins) > 1:
        raise ValueError("ooxml_drawing_anchor_origin_ambiguous")
    if not origins:
        return None
    row_nodes = _direct_xml_children(
        origins[0], OOXML_SPREADSHEET_DRAWING_NAMESPACES, "row"
    )
    column_nodes = _direct_xml_children(
        origins[0], OOXML_SPREADSHEET_DRAWING_NAMESPACES, "col"
    )
    if len(row_nodes) > 1 or len(column_nodes) > 1:
        raise ValueError("ooxml_drawing_anchor_origin_ambiguous")
    if not row_nodes or not column_nodes:
        return None
    try:
        row_value = int(row_nodes[0].text or "")
        column_value = int(column_nodes[0].text or "")
    except ValueError as exc:
        raise ValueError("ooxml_drawing_anchor_origin_invalid") from exc
    if row_value < 0 or column_value < 0:
        raise ValueError("ooxml_drawing_anchor_origin_invalid")
    return f"{_xlsx_column_name(column_value)}{row_value + 1}"


def _xlsx_verified_chart_placements(
    archive: zipfile.ZipFile,
    drawing_part: str,
) -> list[dict[str, str | None]]:
    if drawing_part not in archive.namelist():
        return []
    raw = archive.read(drawing_part)
    validate_xml_bytes(raw)
    try:
        root = ElementTree.fromstring(raw)
    except ElementTree.ParseError as exc:
        raise ValueError("ooxml_drawing_xml_invalid") from exc
    root_namespace, root_name = _xml_name(root.tag)
    if (
        root_namespace not in OOXML_SPREADSHEET_DRAWING_NAMESPACES
        or root_name != "wsDr"
    ):
        raise ValueError("ooxml_drawing_root_invalid")
    placements: list[dict[str, str | None]] = []
    for anchor in root:
        anchor_namespace, anchor_name = _xml_name(anchor.tag)
        if (
            anchor_namespace not in OOXML_SPREADSHEET_DRAWING_NAMESPACES
            or anchor_name
            not in {"oneCellAnchor", "twoCellAnchor", "absoluteAnchor"}
        ):
            continue
        frames = _direct_xml_children(
            anchor, OOXML_SPREADSHEET_DRAWING_NAMESPACES, "graphicFrame"
        )
        if len(frames) > 1:
            raise ValueError("ooxml_drawing_chart_carrier_ambiguous")
        if not frames:
            continue
        relationship_id = _graphic_frame_relationship(
            frames[0],
            payload_namespaces=OOXML_CHART_NAMESPACES,
            payload_local_name="chart",
            relationship_attribute_name="id",
        )
        if relationship_id:
            placements.append({
                "relationship_id": relationship_id,
                "cell": _xlsx_anchor_cell(anchor),
            })
    return placements


def referenced_ooxml_media(
    archive: zipfile.ZipFile,
    *,
    media_prefixes: tuple[str, ...],
    include_unresolved_pptx_inheritance: bool = False,
) -> list[dict[str, Any]]:
    """Find only media reachable through actually-used OOXML relationships.

    Presence under ``word/media`` or ``xl/media`` alone is not evidence that an
    image appears in the document.  This walker starts at package relationships,
    follows only relationship IDs present in each reachable source part, and
    records each concrete use separately.
    """
    relationships = _ooxml_relationships(archive)
    _, drawing_sheets = _xlsx_sheet_context(archive, relationships)
    slide_context = _pptx_slide_context(archive, relationships)
    queue = [""]
    reached = {""}
    results: list[dict[str, Any]] = []
    while queue:
        source_part = queue.pop(0)
        used = (
            {
                row["id"]: 1
                for row in relationships.get("", [])
                if _ooxml_relationship_kind(row["type"]) == "officeDocument"
            }
            if source_part == ""
            else _ooxml_used_relationship_ids(archive, source_part)
        )
        relationship_uses = (
            [] if source_part == ""
            else _ooxml_relationship_uses(archive, source_part)
        )
        for row in relationships.get(source_part, []):
            relationship_kind = _ooxml_relationship_kind(row["type"])
            implicit_presentation_part = relationship_kind in {
                "slideLayout",
                "slideMaster",
            }
            if (
                (row["id"] not in used and not implicit_presentation_part)
                or row["target_mode"].casefold() == "external"
            ):
                continue
            target = _resolve_ooxml_target(source_part, row["target"])
            if target is None:
                raise ValueError("ooxml_relationship_target_invalid")
            if (
                target.startswith(media_prefixes)
                and target in archive.namelist()
                and relationship_kind == "image"
                and row["id"] in used
            ):
                if (
                    target.startswith("word/media/")
                    and not _word_visual_source_part_allowed(
                        archive, source_part
                    )
                ):
                    continue
                if (
                    target.startswith("xl/media/")
                    and source_part not in drawing_sheets
                ):
                    continue
                if (
                    target.startswith("ppt/media/")
                    and source_part not in slide_context
                ):
                    continue
                pptx_display_scope: str | None = None
                if target.startswith("ppt/media/"):
                    if source_part.startswith("ppt/slides/"):
                        pptx_display_scope = "direct_slide"
                    elif source_part.startswith(
                        ("ppt/slideLayouts/", "ppt/slideMasters/")
                    ):
                        pptx_display_scope = "inherited_visibility_unresolved"
                    else:
                        pptx_display_scope = "non_slide_visibility_unresolved"
                    if (
                        pptx_display_scope != "direct_slide"
                        and not include_unresolved_pptx_inheritance
                    ):
                        continue
                uses = [
                    item for item in relationship_uses
                    if item["relationship_id"] == row["id"]
                ]
                cells = (
                    _xlsx_drawing_cells(archive, source_part, row["id"])
                    if source_part in drawing_sheets else []
                )
                if len(uses) != used[row["id"]]:
                    raise ValueError("ooxml_relationship_occurrence_count_invalid")
                if source_part in drawing_sheets and len(cells) != used[row["id"]]:
                    raise ValueError("ooxml_drawing_cell_occurrence_count_invalid")
                occurrences = list(zip(uses, cells or [None] * len(uses)))
                slide_numbers: tuple[int | None, ...] = (
                    slide_context.get(source_part) or (None,)
                )
                for occurrence, (relationship_use, cell) in enumerate(occurrences, 1):
                    usage_kind = relationship_use["usage_kind"]
                    for slide_number in slide_numbers:
                        locator_parts = [
                            f"part={source_part}",
                            f"relationship={row['id']}",
                            f"occurrence={occurrence}",
                            f"usage={usage_kind}",
                        ]
                        if slide_number is not None:
                            locator_parts.append(f"slide={slide_number}")
                        result: dict[str, Any] = {
                            "member": target,
                            "source_part": source_part,
                            "relationship_id": row["id"],
                            "relationship_occurrence": occurrence,
                            "usage_kind": usage_kind,
                            "locator_text": ";".join(locator_parts),
                        }
                        if slide_number is not None:
                            result["slide_number"] = slide_number
                        if pptx_display_scope is not None:
                            result["pptx_display_scope"] = pptx_display_scope
                        if source_part in drawing_sheets:
                            result["sheet_name"] = drawing_sheets[source_part]
                        if cell is not None:
                            result["cell"] = cell
                        results.append(result)
                continue
            if target in relationships and target not in reached:
                reached.add(target)
                queue.append(target)
    return results


def _ooxml_chart_payload(raw: bytes, source_member: str) -> dict[str, Any]:
    """Extract explicit chart labels, cached values, and series from chart XML."""
    validate_xml_bytes(raw)
    try:
        root = ElementTree.fromstring(raw)
    except ElementTree.ParseError as exc:
        raise ValueError("ooxml_chart_xml_invalid") from exc

    root_namespace, root_name = _xml_name(root.tag)
    if (
        root_namespace not in OOXML_CHART_NAMESPACES
        or root_name != "chartSpace"
    ):
        raise ValueError("ooxml_chart_root_invalid")

    def local(item: ElementTree.Element) -> str:
        return _xml_name(item.tag)[1]

    def chart_item(item: ElementTree.Element, name: str) -> bool:
        namespace, local_name = _xml_name(item.tag)
        return namespace in OOXML_CHART_NAMESPACES and local_name == name

    def chart_text_value(item: ElementTree.Element) -> bool:
        namespace, local_name = _xml_name(item.tag)
        return (
            namespace in OOXML_CHART_NAMESPACES and local_name == "v"
        ) or (
            namespace in OOXML_DRAWING_NAMESPACES and local_name == "t"
        )

    def point_values(container: ElementTree.Element) -> dict[str, Any]:
        indexed: list[tuple[int, Any]] = []
        values_in_document_order: list[Any] = []
        invalid_index = False
        duplicate_index = False
        seen_indexes: set[int] = set()
        multi_level_count = sum(
            1 for item in container.iter() if chart_item(item, "lvl")
        )
        for point in container.iter():
            if not chart_item(point, "pt"):
                continue
            value_node = next(
                (
                    child_item for child_item in point.iter()
                    if chart_text_value(child_item)
                    and child_item.text is not None
                ),
                None,
            )
            if value_node is None:
                continue
            value: Any = value_node.text.strip()
            values_in_document_order.append(value)
            try:
                index = int(point.attrib["idx"])
            except (KeyError, TypeError, ValueError):
                invalid_index = True
                continue
            if index < 0:
                invalid_index = True
                continue
            if index in seen_indexes:
                duplicate_index = True
            seen_indexes.add(index)
            indexed.append((index, value))
        index_valid = (
            bool(indexed)
            and not invalid_index
            and not duplicate_index
            and multi_level_count <= 1
            and len(indexed) == len(values_in_document_order)
        )
        ordered = sorted(indexed, key=lambda item: item[0])
        return {
            "values": (
                [value for _, value in ordered]
                if index_valid else values_in_document_order
            ),
            "indexed": ordered if index_valid else [],
            "index_valid": index_valid,
            "multi_level_count": multi_level_count,
        }

    formulas = [
        item.text.strip()
        for item in root.iter()
        if chart_item(item, "f") and item.text and item.text.strip()
    ]
    cached_labels = [
        item.text.strip()
        for item in root.iter()
        if chart_text_value(item) and item.text and item.text.strip()
    ]
    chart_types = list(dict.fromkeys(
        local(item)
        for item in root.iter()
        if _xml_name(item.tag)[0] in OOXML_CHART_NAMESPACES
        and local(item).casefold().endswith("chart")
        and local(item).casefold() not in {"chart", "chartspace"}
    ))
    title_values: list[str] = []
    for title in (
        item for item in root.iter() if chart_item(item, "title")
    ):
        title_values.extend(
            item.text.strip()
            for item in title.iter()
            if chart_text_value(item) and item.text and item.text.strip()
        )
    series: list[dict[str, Any]] = []
    for series_index, series_node in enumerate(
        (item for item in root.iter() if chart_item(item, "ser")),
        1,
    ):
        names: list[str] = []
        category_points: list[dict[str, Any]] = []
        value_points: list[dict[str, Any]] = []
        series_formulas: list[str] = []
        for child_item in series_node:
            if _xml_name(child_item.tag)[0] not in OOXML_CHART_NAMESPACES:
                continue
            child_kind = local(child_item)
            if child_kind == "tx":
                names.extend(
                    value.text.strip()
                    for value in child_item.iter()
                    if chart_text_value(value)
                    and value.text and value.text.strip()
                )
            elif child_kind in {"cat", "xVal"}:
                category_points.append(point_values(child_item))
            elif child_kind in {"val", "yVal", "bubbleSize"}:
                value_points.append(point_values(child_item))
            series_formulas.extend(
                value.text.strip()
                for value in child_item.iter()
                if chart_item(value, "f")
                and value.text and value.text.strip()
            )
        categories = [
            value
            for point_set in category_points
            for value in point_set["values"]
        ]
        values = [
            value
            for point_set in value_points
            for value in point_set["values"]
        ]
        category_indexes = [
            index
            for point_set in category_points
            for index, _ in point_set["indexed"]
        ]
        value_indexes = [
            index
            for point_set in value_points
            for index, _ in point_set["indexed"]
        ]
        cache_structure_valid = (
            len(category_points) <= 1
            and len(value_points) <= 1
            and all(
                point_set["index_valid"]
                for point_set in category_points + value_points
            )
        )
        if (
            values
            and categories
            and cache_structure_valid
            and category_indexes == value_indexes
        ):
            cache_status = "paired_cached_values"
            points = [
                {"category": category, "value": value}
                for category, value in zip(categories, values)
            ]
        elif values and categories and cache_structure_valid:
            cache_status = "unresolved_cache_index_mismatch"
            points = []
        elif values and categories:
            cache_status = "unresolved_cache_structure"
            points = []
        elif (
            values
            and not categories
            and len(value_points) == 1
            and all(point_set["index_valid"] for point_set in value_points)
        ):
            cache_status = "values_without_categories"
            points = []
        elif values:
            cache_status = "unresolved_cache_structure"
            points = []
        elif values or categories:
            cache_status = "unresolved_cache_length_mismatch"
            points = []
        else:
            cache_status = "cached_values_missing"
            points = []
        series.append({
            "series_index": series_index,
            "name": " / ".join(dict.fromkeys(names)) if names else None,
            "categories": categories,
            "values": values,
            "points": points,
            "cache_status": cache_status,
            "category_indexes": category_indexes,
            "value_indexes": value_indexes,
            "formulas": list(dict.fromkeys(series_formulas)),
        })
    return {
        "source_member": source_member,
        "xml_sha256": digest_bytes(raw),
        "title": " / ".join(dict.fromkeys(title_values)) if title_values else None,
        "chart_types": chart_types,
        "formulas": formulas,
        "cached_labels": cached_labels,
        "series": series,
    }


def _chart_summary_text(payload: dict[str, Any]) -> str:
    lines = ["図表（ファイル内の保存済み構造）"]
    title = payload.get("title")
    if title:
        lines.append(f"タイトル: {title}")
    source_member = payload.get("source_member")
    if source_member:
        lines.append(f"参照元: {source_member}")
    series = payload.get("series", [])
    lines.append(f"系列数: {len(series)}")
    for item in series:
        name = item.get("name") or f"系列{item.get('series_index')}"
        lines.append(f"系列: {name}; 保存値状態: {item.get('cache_status')}")
    return "\n".join(lines)


def _chart_series_text(payload: dict[str, Any], series: dict[str, Any]) -> str:
    name = series.get("name") or f"系列{series.get('series_index')}"
    lines = [
        f"図表系列: {name}",
        f"保存値状態: {series.get('cache_status')}",
    ]
    points = series.get("points", [])
    if points:
        lines.extend(
            f"{point['category']}: {point['value']}"
            for point in points
        )
    elif series.get("cache_status") == "values_without_categories":
        lines.append(
            "カテゴリなし保存値: "
            + ", ".join(str(value) for value in series.get("values", []))
        )
    formulas = series.get("formulas", [])
    if formulas:
        lines.append("参照式（未実行）: " + " | ".join(str(value) for value in formulas))
    return "\n".join(lines)


def referenced_ooxml_charts(
    archive: zipfile.ZipFile,
) -> list[dict[str, Any]]:
    """Return only chart parts reached through concrete document placements."""
    relationships = _ooxml_relationships(archive)
    _, drawing_sheets = _xlsx_sheet_context(archive, relationships)
    slide_context = _pptx_slide_context(archive, relationships)
    queue = [""]
    reached = {""}
    results: list[dict[str, Any]] = []
    while queue:
        source_part = queue.pop(0)
        used = (
            {
                row["id"]: 1
                for row in relationships.get("", [])
                if _ooxml_relationship_kind(row["type"]) == "officeDocument"
            }
            if source_part == ""
            else _ooxml_used_relationship_ids(archive, source_part)
        )
        if source_part in drawing_sheets:
            verified_chart_placements = _xlsx_verified_chart_placements(
                archive, source_part
            )
        elif source_part in slide_context:
            verified_chart_placements = [
                {"relationship_id": relationship_id, "cell": None}
                for relationship_id in _pptx_graphic_relationship_ids(
                    archive,
                    source_part,
                    payload_namespaces=OOXML_CHART_NAMESPACES,
                    payload_local_name="chart",
                    relationship_attribute_name="id",
                )
            ]
        else:
            verified_chart_placements = []
        for row in relationships.get(source_part, []):
            relationship_kind = _ooxml_relationship_kind(row["type"])
            implicit_presentation_part = relationship_kind in {
                "slideLayout",
                "slideMaster",
            }
            if (
                (row["id"] not in used and not implicit_presentation_part)
                or row["target_mode"].casefold() == "external"
            ):
                continue
            target = _resolve_ooxml_target(source_part, row["target"])
            if target is None:
                raise ValueError("ooxml_relationship_target_invalid")
            if (
                relationship_kind == "chart"
                and target in archive.namelist()
            ):
                placements = [
                    item for item in verified_chart_placements
                    if item["relationship_id"] == row["id"]
                ]
                if not placements:
                    continue
                if source_part in drawing_sheets:
                    if not target.startswith("xl/charts/"):
                        raise ValueError("ooxml_chart_target_invalid")
                elif source_part in slide_context:
                    if not target.startswith("ppt/charts/"):
                        raise ValueError("ooxml_chart_target_invalid")
                else:
                    continue
                slide_numbers: tuple[int | None, ...] = (
                    slide_context.get(source_part) or (None,)
                )
                for occurrence, placement in enumerate(placements, 1):
                    cell = placement["cell"]
                    for slide_number in slide_numbers:
                        locator_parts = [
                            f"part={source_part}",
                            f"relationship={row['id']}",
                            f"occurrence={occurrence}",
                        ]
                        result: dict[str, Any] = {
                            "member": target,
                            "source_part": source_part,
                            "relationship_id": row["id"],
                            "relationship_occurrence": occurrence,
                        }
                        if source_part in drawing_sheets:
                            result["sheet_name"] = drawing_sheets[source_part]
                            locator_parts.append(
                                "sheet=" + urllib.parse.quote(
                                    drawing_sheets[source_part], safe="-._~"
                                )
                            )
                        if cell is not None:
                            result["cell"] = cell
                            locator_parts.append(f"cell={cell}")
                        if slide_number is not None:
                            result["slide_number"] = slide_number
                            locator_parts.append(f"slide={slide_number}")
                        result["locator_text"] = ";".join(locator_parts)
                        results.append(result)
                continue
            if target in relationships and target not in reached:
                reached.add(target)
                queue.append(target)
    return results


def referenced_pptx_diagrams(
    archive: zipfile.ZipFile,
) -> list[dict[str, Any]]:
    """Return slide-bound, explicitly referenced SmartArt data records."""
    relationships = _ooxml_relationships(archive)
    slide_context = _pptx_slide_context(archive, relationships)
    records: list[dict[str, Any]] = []
    queue = [""]
    reached = {""}
    while queue:
        source_part = queue.pop(0)
        used = (
            {
                row["id"]: 1
                for row in relationships.get("", [])
                if _ooxml_relationship_kind(row["type"]) == "officeDocument"
            }
            if source_part == ""
            else _ooxml_used_relationship_ids(archive, source_part)
        )
        verified_diagram_ids = (
            _pptx_graphic_relationship_ids(
                archive,
                source_part,
                payload_namespaces=OOXML_DIAGRAM_NAMESPACES,
                payload_local_name="relIds",
                relationship_attribute_name="dm",
            )
            if source_part in slide_context else []
        )
        for row in relationships.get(source_part, []):
            relationship_kind = _ooxml_relationship_kind(row["type"])
            implicit_presentation_part = relationship_kind in {
                "slideLayout",
                "slideMaster",
            }
            if (
                (row["id"] not in used and not implicit_presentation_part)
                or row["target_mode"].casefold() == "external"
            ):
                continue
            target = _resolve_ooxml_target(source_part, row["target"])
            if target is None:
                raise ValueError("ooxml_relationship_target_invalid")
            if (
                relationship_kind == "diagramData"
                and target.startswith("ppt/diagrams/")
                and target.casefold().endswith(".xml")
                and target in archive.namelist()
            ):
                if source_part not in slide_context:
                    continue
                relationship_occurrences = sum(
                    relationship_id == row["id"]
                    for relationship_id in verified_diagram_ids
                )
                if not relationship_occurrences:
                    continue
                raw = archive.read(target)
                validate_xml_bytes(raw)
                try:
                    root = ElementTree.fromstring(raw)
                except ElementTree.ParseError as exc:
                    raise ValueError("ooxml_diagram_xml_invalid") from exc
                root_namespace, root_name = _xml_name(root.tag)
                if (
                    root_namespace not in OOXML_DIAGRAM_NAMESPACES
                    or root_name != "dataModel"
                ):
                    raise ValueError("ooxml_diagram_data_root_invalid")
                points: list[dict[str, Any]] = []
                connections: list[dict[str, str]] = []
                for item in root.iter():
                    item_namespace, kind = _xml_name(item.tag)
                    if item_namespace not in OOXML_DIAGRAM_NAMESPACES:
                        continue
                    if kind == "pt":
                        text_values = [
                            value.text.strip()
                            for value in item.iter()
                            if _xml_name(value.tag)[0] in OOXML_DRAWING_NAMESPACES
                            and _xml_name(value.tag)[1] == "t"
                            and value.text and value.text.strip()
                        ]
                        points.append({
                            "model_id": item.attrib.get("modelId"),
                            "type": item.attrib.get("type"),
                            "text": "\n".join(text_values),
                        })
                    elif kind == "cxn":
                        connections.append({
                            key: item.attrib[key]
                            for key in ("modelId", "srcId", "destId", "type")
                            if key in item.attrib
                        })
                if not any(point.get("text") for point in points):
                    continue
                for occurrence in range(1, relationship_occurrences + 1):
                    for slide_number in slide_context.get(source_part, ()):
                        records.append({
                            "slide_number": slide_number,
                            "source_member": target,
                            "source_part": source_part,
                            "relationship_id": row["id"],
                            "relationship_occurrence": occurrence,
                            "xml_sha256": digest_bytes(raw),
                            "points": points,
                            "connections": connections,
                        })
                continue
            if target in relationships and target not in reached:
                reached.add(target)
                queue.append(target)
    return records


def discover_password_candidates(root: Path) -> tuple[str, ...]:
    """Derive generic Office password candidates from path-visible aliases and dates."""
    aliases: set[str] = set()
    dates: set[str] = set()
    embedded: set[str] = set()
    for path in root.rglob("*"):
        if not path.is_file():
            continue
        visible = unicodedata.normalize("NFC", path.as_posix())
        dates.update(DATE_TOKEN_PATTERN.findall(visible))
        for alias, date_token in ALIAS_DATE_PATTERN.findall(visible):
            aliases.add(alias)
            dates.add(date_token)
            embedded.add(f"{alias}{date_token}")
    candidates = set(embedded)
    for alias in aliases:
        for date_token in dates:
            for extension in ("docx", "xlsx", "pptx"):
                candidates.add(f"DA-{alias}-{date_token}-{extension}")
                candidates.add(f"DA-{alias.upper()}-{date_token}-{extension}")
    return tuple(sorted(candidates))


def json_value(value: Any) -> Any:
    if isinstance(value, (datetime, date)):
        return value.isoformat()
    if isinstance(value, bytes):
        return {"bytes_sha256": digest_bytes(value), "length": len(value)}
    if value is None or isinstance(value, (str, int, float, bool)):
        return value
    if isinstance(value, dict):
        return {str(key): json_value(item) for key, item in value.items()}
    if isinstance(value, (list, tuple)):
        return [json_value(item) for item in value]
    return str(value)


def value_type(value: Any) -> str:
    if value is None:
        return "null"
    if isinstance(value, bool):
        return "boolean"
    if isinstance(value, int):
        return "integer"
    if isinstance(value, float):
        return "number"
    if isinstance(value, datetime):
        return "datetime"
    if isinstance(value, date):
        return "date"
    if isinstance(value, str):
        return "text"
    if isinstance(value, bytes):
        return "binary"
    return "mixed"


def content(*, raw_text: str | None = None, raw_value: Any = None,
            content_ref: str | None = None, mime_type: str | None = None) -> dict[str, Any]:
    if raw_text is not None:
        payload: dict[str, Any] = {"raw_text": raw_text, "normalized_text": normalize_text(raw_text)}
        hashed = {"raw_text": raw_text}
        kind = "text"
        original_length = len(raw_text)
    elif content_ref is not None:
        payload = {"content_ref": content_ref}
        hashed = {"content_ref": content_ref}
        kind = "binary"
        original_length = None
    else:
        clean = json_value(raw_value)
        payload = {"raw_value": clean, "normalized_value": clean}
        hashed = {"raw_value": clean}
        kind = value_type(raw_value)
        original_length = None
    payload.update({"value_type": kind, "sha256": digest_value(hashed), "is_truncated": False})
    if mime_type:
        payload["mime_type"] = mime_type
    if original_length is not None:
        payload["original_length"] = original_length
    return payload


def nfc_path(path: Path) -> str:
    return unicodedata.normalize("NFC", path.as_posix())


class DeferredVisualStoreError(RuntimeError):
    """The private per-document visual spool violated its integrity contract."""


class Probe:
    def __init__(
        self,
        root: Path,
        run_at: str,
        max_items: int | None,
        *,
        diagnostic: bool = True,
        extractor: str = EXTRACTOR,
        extractor_version: str = EXTRACTOR_VERSION,
        record_sink: Callable[[str, dict[str, Any]], None] | None = None,
        retain_records: bool = True,
        password_candidates: tuple[str, ...] = (),
        visual_observation_mode: str = "immediate",
    ) -> None:
        if visual_observation_mode not in VISUAL_OBSERVATION_MODES:
            raise ValueError(
                "visual_observation_mode must be immediate, "
                "deferred_per_document, or suppressed"
            )
        self.root = root.resolve()
        self.run_at = run_at
        self.max_items = max_items
        self.diagnostic = diagnostic
        self.extractor = extractor
        self.extractor_version = extractor_version
        self.record_sink = record_sink
        self.retain_records = retain_records
        self.password_candidates = password_candidates
        self.visual_observation_mode = visual_observation_mode
        self.documents: list[dict[str, Any]] = []
        self.evidence: list[dict[str, Any]] = []
        self.relations: list[dict[str, Any]] = []
        self._leaf_counts: dict[str, int] = {}
        self._embedded_visual_usage: dict[str, dict[str, Any]] = {}
        self._current_document: dict[str, Any] | None = None
        self._deferred_visual_tasks: list[dict[str, Any]] = []
        self._suppressed_visual_tasks: list[dict[str, Any]] = []
        self._visual_spool_root: Path | None = None
        self._visual_spool_root_identity: tuple[int, int] | None = None
        self._visual_spool_by_sha256: dict[str, Path] = {}
        self._visual_spool_bytes = 0

    def emit(self, kind: str, record: dict[str, Any]) -> None:
        if self.retain_records:
            getattr(self, kind).append(record)
        if self.record_sink is not None:
            self.record_sink(kind, record)

    def relative_path(self, path: Path) -> str:
        resolved = path.resolve()
        try:
            relative = resolved.relative_to(self.root)
        except ValueError as exc:
            raise ValueError(f"input is outside --root: {resolved}") from exc
        return nfc_path(relative)

    def add_document(self, path: Path, parser: str) -> dict[str, Any]:
        rel = self.relative_path(path)
        source_sha = digest_file(path)
        doc_id = stable_id("doc", {"relative_path": rel, "source_sha256": source_sha})
        stat = path.stat()
        if self.diagnostic:
            initial_status = "partial"
            initial_warnings = [
                f"diagnostic sample: at most {self.max_items} leaf items per document",
                "must not be used as the answer pipeline index",
            ]
        else:
            initial_status = "success"
            initial_warnings = []
        record = {
            "schema_version": SCHEMA_VERSION,
            "record_type": "document",
            "document_id": doc_id,
            "source": {
                "relative_path": rel,
                "file_name": unicodedata.normalize("NFC", path.name),
                "extension": path.suffix.lower().lstrip("."),
                "media_type": mimetypes.guess_type(path.name)[0] or "application/octet-stream",
                "size_bytes": stat.st_size,
                "sha256": source_sha,
                "modified_at": datetime.fromtimestamp(stat.st_mtime, timezone.utc).isoformat(),
            },
            "extraction": {
                "status": initial_status,
                "parser": parser,
                "parser_version": self.extractor_version,
                "extracted_at": self.run_at,
                "warnings": initial_warnings,
                "errors": [],
            },
        }
        if self._current_document is not None:
            raise RuntimeError("only one source document may be active per Probe instance")
        self._current_document = record
        if self.retain_records:
            self.documents.append(record)
        self._leaf_counts[doc_id] = 0
        return record

    def finalize_document(self) -> dict[str, Any]:
        if self._current_document is None:
            raise RuntimeError("no active source document to finalize")
        document = self._current_document
        if self.record_sink is not None:
            self.record_sink("documents", document)
        self._current_document = None
        return document

    def may_add_leaf(self, doc_id: str) -> bool:
        if self.max_items is not None and self._leaf_counts[doc_id] >= self.max_items:
            return False
        self._leaf_counts[doc_id] += 1
        return True

    def limit_reached(self, doc_id: str) -> bool:
        return self.max_items is not None and self._leaf_counts[doc_id] >= self.max_items

    @staticmethod
    def mark_partial(document: dict[str, Any], warning: str) -> None:
        document["extraction"]["status"] = "partial"
        if warning not in document["extraction"]["warnings"]:
            document["extraction"]["warnings"].append(warning)

    def record_failure(self, path: Path, error: Exception) -> None:
        existing = self._current_document or next(
            (item for item in reversed(self.documents) if item["source"]["relative_path"] == self.relative_path(path)),
            None,
        )
        document = existing or self.add_document(path, f"{path.suffix.lower().lstrip('.') or 'unknown'}-parser")
        document["extraction"]["status"] = "failed"
        document["extraction"]["errors"] = [f"{type(error).__name__}: {error}"]

    def add_evidence(
        self,
        document_id: str,
        evidence_type: str,
        location: dict[str, Any],
        item_content: dict[str, Any],
        *,
        parent_id: str | None = None,
        ordinal: int | None = None,
        style: dict[str, Any] | None = None,
        geometry: dict[str, Any] | None = None,
        native_properties: dict[str, Any] | None = None,
        method: str = "native_parser",
        confidence: float = 1.0,
        deterministic: bool = True,
        warning: str | None = None,
    ) -> dict[str, Any]:
        identity = {
            "document_id": document_id,
            "evidence_type": evidence_type,
            "location": location,
            "content_sha256": item_content["sha256"],
        }
        record: dict[str, Any] = {
            "schema_version": SCHEMA_VERSION,
            "record_type": "evidence",
            "evidence_id": stable_id("ev", identity),
            "document_id": document_id,
            "evidence_type": evidence_type,
            "location": location,
            "content": item_content,
            "provenance": {
                "extraction_method": method,
                "extractor": self.extractor,
                "extractor_version": self.extractor_version,
                "extracted_at": self.run_at,
                "deterministic": deterministic,
                "confidence": confidence,
                "warnings": [warning] if warning else [],
            },
        }
        if parent_id:
            record["parent_evidence_id"] = parent_id
        if ordinal is not None:
            record["ordinal"] = ordinal
        if style:
            record["style"] = style
        if geometry:
            record["geometry"] = geometry
        if native_properties:
            record["native_properties"] = native_properties
        self.emit("evidence", record)
        if parent_id:
            self.add_relation(
                "structural", "contains",
                {"record_type": "evidence", "record_id": parent_id},
                {"record_type": "evidence", "record_id": record["evidence_id"]},
            )
        return record

    def add_relation(
        self,
        relation_class: str,
        relation_type: str,
        from_ref: dict[str, str],
        to_ref: dict[str, str],
        *,
        properties: dict[str, Any] | None = None,
        supporting_evidence_ids: list[str] | None = None,
        rule_or_model: str = "native containment",
    ) -> None:
        relation_properties = properties or {}
        supporting_ids = supporting_evidence_ids or []
        identity = {
            "class": relation_class,
            "type": relation_type,
            "from": from_ref,
            "to": to_ref,
            "generator": self.extractor,
            "generator_version": self.extractor_version,
        }
        self.emit("relations", {
            "schema_version": SCHEMA_VERSION,
            "record_type": "relation",
            "relation_id": stable_id("rel", identity),
            "relation_class": relation_class,
            "relation_type": relation_type,
            "from_ref": from_ref,
            "to_ref": to_ref,
            "properties": relation_properties,
            "supporting_evidence_ids": supporting_ids,
            "provenance": {
                "generated_by": self.extractor,
                "generator_version": self.extractor_version,
                "generated_at": self.run_at,
                "deterministic": True,
                "confidence": 1.0,
                "rule_or_model": rule_or_model,
                "warnings": [],
            },
            "status": "verified",
        })

    def contain_document(self, doc_id: str, evidence_id: str) -> None:
        self.add_relation(
            "structural", "contains",
            {"record_type": "document", "record_id": doc_id},
            {"record_type": "evidence", "record_id": evidence_id},
        )

    @staticmethod
    def _visual_materialization_contract(
        visual_origin: dict[str, Any],
    ) -> tuple[str, int]:
        materialization = visual_origin.get("materialization")
        if not isinstance(materialization, dict):
            raise DeferredVisualStoreError(
                "visual origin has no materialization contract"
            )
        expected_sha256 = materialization.get("rendered_sha256")
        expected_size = materialization.get("rendered_size_bytes")
        if (
            not isinstance(expected_sha256, str)
            or not re.fullmatch(r"[0-9a-f]{64}", expected_sha256)
            or isinstance(expected_size, bool)
            or not isinstance(expected_size, int)
            or expected_size <= 0
        ):
            raise DeferredVisualStoreError(
                "visual materialization digest or size is invalid"
            )
        return expected_sha256, expected_size

    @staticmethod
    def _visual_observation_location(
        location_prefix: dict[str, Any],
    ) -> dict[str, Any]:
        return {
            **location_prefix,
            "locator_text": "visual_observation=whole_image",
        }

    @staticmethod
    def _release_active_paddle_worker() -> None:
        from local_image_ocr import active_paddle_session

        paddle_session = active_paddle_session()
        if paddle_session is not None:
            # The memo and build-scoped session remain live. Only Paddle's
            # native process is retired before Gemma is allowed to allocate.
            paddle_session.release_idle_worker()

    def _ensure_visual_spool_root(self) -> Path:
        if self._visual_spool_root is None:
            root = Path(tempfile.mkdtemp(prefix="aiec-visual-spool-"))
            os.chmod(root, 0o700)
            root_stat = root.lstat()
            if not stat.S_ISDIR(root_stat.st_mode) or stat.S_ISLNK(root_stat.st_mode):
                raise DeferredVisualStoreError(
                    "private visual spool root is not a real directory"
                )
            self._visual_spool_root = root
            self._visual_spool_root_identity = (
                root_stat.st_dev,
                root_stat.st_ino,
            )
        return self._visual_spool_root

    @staticmethod
    def _write_all_descriptor(descriptor: int, raw: bytes) -> None:
        view = memoryview(raw)
        written = 0
        while written < len(view):
            count = os.write(descriptor, view[written:])
            if count <= 0:
                raise OSError("short write to private visual spool")
            written += count

    def _spool_visual_bytes(self, raw: bytes, sha256: str) -> Path:
        existing = self._visual_spool_by_sha256.get(sha256)
        if existing is not None:
            return existing
        root = self._ensure_visual_spool_root()
        identity = self._visual_spool_root_identity
        if identity is None:
            raise DeferredVisualStoreError(
                "private visual spool root identity is missing"
            )
        name = f"{sha256}.image"
        destination = root / name
        root_flags = (
            os.O_RDONLY
            | getattr(os, "O_DIRECTORY", 0)
            | getattr(os, "O_CLOEXEC", 0)
            | getattr(os, "O_NOFOLLOW", 0)
        )
        file_flags = (
            os.O_WRONLY
            | os.O_CREAT
            | os.O_EXCL
            | getattr(os, "O_CLOEXEC", 0)
            | getattr(os, "O_NOFOLLOW", 0)
        )
        try:
            root_descriptor = os.open(root, root_flags)
        except OSError as exc:
            raise DeferredVisualStoreError(
                "private visual spool root cannot be opened safely for writing"
            ) from exc
        try:
            opened_root = os.fstat(root_descriptor)
            if (
                not stat.S_ISDIR(opened_root.st_mode)
                or (opened_root.st_dev, opened_root.st_ino) != identity
            ):
                raise DeferredVisualStoreError(
                    "private visual spool root identity changed before writing"
                )
            try:
                descriptor = os.open(
                    name, file_flags, 0o600, dir_fd=root_descriptor
                )
            except OSError as exc:
                raise DeferredVisualStoreError(
                    "private visual spool file cannot be created exclusively"
                ) from exc
            try:
                os.fchmod(descriptor, 0o600)
                self._write_all_descriptor(descriptor, raw)
                os.fsync(descriptor)
                written = os.fstat(descriptor)
                if not stat.S_ISREG(written.st_mode) or written.st_size != len(raw):
                    raise DeferredVisualStoreError(
                        "private visual spool file size changed while writing"
                    )
            except BaseException:
                try:
                    os.close(descriptor)
                finally:
                    try:
                        os.unlink(name, dir_fd=root_descriptor)
                    except FileNotFoundError:
                        pass
                raise
            else:
                os.close(descriptor)
        except DeferredVisualStoreError:
            raise
        except OSError as exc:
            raise DeferredVisualStoreError(
                "private visual spool write could not be completed safely"
            ) from exc
        finally:
            os.close(root_descriptor)
        self._visual_spool_by_sha256[sha256] = destination
        self._visual_spool_bytes += len(raw)
        return destination

    def _cleanup_visual_spool(self) -> None:
        root = self._visual_spool_root
        if root is None:
            if (
                self._visual_spool_root_identity is not None
                or self._visual_spool_by_sha256
                or self._visual_spool_bytes
            ):
                raise DeferredVisualStoreError(
                    "private visual spool cleanup state is inconsistent"
                )
            self._deferred_visual_tasks.clear()
            return

        identity = self._visual_spool_root_identity
        if identity is None:
            raise DeferredVisualStoreError(
                "private visual spool root identity is missing"
            )
        expected_names: set[str] = set()
        for expected_sha256, path in self._visual_spool_by_sha256.items():
            expected_name = f"{expected_sha256}.image"
            if path.parent != root or path.name != expected_name:
                raise DeferredVisualStoreError(
                    "private visual spool file map is inconsistent"
                )
            expected_names.add(expected_name)

        flags = (
            os.O_RDONLY
            | getattr(os, "O_DIRECTORY", 0)
            | getattr(os, "O_CLOEXEC", 0)
            | getattr(os, "O_NOFOLLOW", 0)
        )
        try:
            descriptor = os.open(root, flags)
        except OSError as exc:
            raise DeferredVisualStoreError(
                "private visual spool root cannot be opened safely for cleanup"
            ) from exc
        try:
            opened = os.fstat(descriptor)
            if (
                not stat.S_ISDIR(opened.st_mode)
                or (opened.st_dev, opened.st_ino) != identity
            ):
                raise DeferredVisualStoreError(
                    "private visual spool root identity changed before cleanup"
                )
            actual_names = set(os.listdir(descriptor))
            if actual_names != expected_names:
                raise DeferredVisualStoreError(
                    "private visual spool contents changed before cleanup"
                )
            for name in sorted(expected_names):
                entry = os.stat(name, dir_fd=descriptor, follow_symlinks=False)
                if stat.S_ISDIR(entry.st_mode):
                    raise DeferredVisualStoreError(
                        "private visual spool entry became a directory"
                    )
            for name in sorted(expected_names):
                os.unlink(name, dir_fd=descriptor)
            if os.listdir(descriptor):
                raise DeferredVisualStoreError(
                    "private visual spool is not empty after cleanup"
                )
        except DeferredVisualStoreError:
            raise
        except OSError as exc:
            raise DeferredVisualStoreError(
                "private visual spool cannot be cleaned safely"
            ) from exc
        finally:
            os.close(descriptor)

        try:
            current = root.lstat()
        except OSError as exc:
            raise DeferredVisualStoreError(
                "private visual spool root disappeared before removal"
            ) from exc
        if (
            not stat.S_ISDIR(current.st_mode)
            or stat.S_ISLNK(current.st_mode)
            or (current.st_dev, current.st_ino) != identity
        ):
            raise DeferredVisualStoreError(
                "private visual spool root identity changed before removal"
            )
        try:
            root.rmdir()
        except OSError as exc:
            raise DeferredVisualStoreError(
                "private visual spool root cannot be removed safely"
            ) from exc

        self._deferred_visual_tasks.clear()
        self._visual_spool_by_sha256.clear()
        self._visual_spool_bytes = 0
        self._visual_spool_root = None
        self._visual_spool_root_identity = None

    def _read_verified_spooled_visual(
        self,
        task: dict[str, Any],
    ) -> bytes:
        from local_visual_observation import read_checked_image_bytes

        path = task.get("spool_path")
        expected_sha256 = task.get("expected_sha256")
        expected_size = task.get("expected_size")
        if not isinstance(path, Path):
            raise DeferredVisualStoreError("visual task has no private spool path")
        try:
            raw = read_checked_image_bytes(path)
        except Exception as exc:
            raise DeferredVisualStoreError(
                "private visual spool file cannot be read safely"
            ) from exc
        actual_sha256 = digest_bytes(raw)
        if len(raw) != expected_size or actual_sha256 != expected_sha256:
            raise DeferredVisualStoreError(
                "private visual spool file differs from the queued materialization"
            )
        return raw

    def _schedule_local_visual_observation(
        self,
        image_path: Path,
        document: dict[str, Any],
        *,
        parent_id: str,
        location_prefix: dict[str, Any],
        visual_origin: dict[str, Any],
        ordinal: int,
        exact_location: dict[str, Any] | None = None,
        deadline_at: float | None = None,
        excluded_normalized_texts: set[str] | None = None,
    ) -> bool:
        """Run now, suppress for a child projection, or queue for root flush."""
        location = (
            dict(exact_location)
            if exact_location is not None
            else self._visual_observation_location(location_prefix)
        )
        remaining_timeout: float | None = None
        if deadline_at is not None:
            remaining_timeout = float(deadline_at) - time.monotonic()
            if remaining_timeout <= 0:
                self.mark_partial(
                    document,
                    "visual meaning skipped because its document time limit was reached",
                )
                return False
        if self.visual_observation_mode == "immediate":
            immediate_timeout = remaining_timeout
            if immediate_timeout is not None:
                from local_visual_observation import MAX_TIMEOUT_SECONDS

                if immediate_timeout >= MAX_TIMEOUT_SECONDS:
                    # Preserve the historical observe_path call contract when
                    # its own fixed 180-second cap is already stricter.
                    immediate_timeout = None
            return self._add_local_visual_observation(
                image_path,
                document,
                parent_id=parent_id,
                location_prefix=location_prefix,
                visual_origin=visual_origin,
                ordinal=ordinal,
                exact_location=location,
                excluded_normalized_texts=excluded_normalized_texts,
                timeout=immediate_timeout,
                deadline_at=deadline_at,
            )
        descriptor = {
            "image_path": Path(image_path),
            "parent_id": parent_id,
            "location": location,
            "ordinal": ordinal,
        }
        if self.visual_observation_mode == "suppressed":
            self._suppressed_visual_tasks.append(descriptor)
            return False

        expected_sha256, expected_size = self._visual_materialization_contract(
            visual_origin
        )
        from local_image_ocr import read_checked_image_bytes
        from local_visual_observation import MAX_IMAGE_BYTES

        try:
            raw = read_checked_image_bytes(Path(image_path))
        except Exception as exc:
            raise DeferredVisualStoreError(
                "visual source cannot be read safely"
            ) from exc
        actual_sha256 = digest_bytes(raw)
        if len(raw) != expected_size or actual_sha256 != expected_sha256:
            raise DeferredVisualStoreError(
                "visual source differs from its materialization contract"
            )
        if len(raw) > MAX_IMAGE_BYTES:
            kind = str(visual_origin.get("kind", "visual_image"))
            self.mark_partial(
                document,
                f"{kind} visual meaning unavailable: image exceeds the "
                "local visual observation safety limit",
            )
            return False
        is_new = expected_sha256 not in self._visual_spool_by_sha256
        would_exceed = (
            len(self._deferred_visual_tasks) >= MAX_DEFERRED_VISUAL_TASKS
            or (
                is_new
                and self._visual_spool_bytes + len(raw)
                > MAX_DEFERRED_VISUAL_SPOOL_BYTES
            )
        )
        if would_exceed and self._deferred_visual_tasks:
            self._release_active_paddle_worker()
            self._flush_deferred_visual_observations()
            self._cleanup_visual_spool()
        if len(raw) > MAX_DEFERRED_VISUAL_SPOOL_BYTES:
            self.mark_partial(
                document,
                "visual meaning skipped because one image exceeds the private spool limit",
            )
            return False
        spool_path = self._spool_visual_bytes(raw, expected_sha256)
        self._deferred_visual_tasks.append({
            **descriptor,
            "spool_path": spool_path,
            "expected_sha256": expected_sha256,
            "expected_size": expected_size,
            "visual_origin": visual_origin,
            "deadline_at": deadline_at,
            "excluded_normalized_texts": set(excluded_normalized_texts or ()),
        })
        return True

    def _flush_deferred_visual_observations(self) -> None:
        if self.visual_observation_mode != "deferred_per_document":
            raise RuntimeError("only a deferred root Probe may flush visual tasks")
        from local_visual_observation import MAX_TIMEOUT_SECONDS

        tasks = list(self._deferred_visual_tasks)
        self._deferred_visual_tasks.clear()
        for task in tasks:
            document = self._current_document
            if document is None:
                raise DeferredVisualStoreError(
                    "visual tasks cannot outlive their source document"
                )
            deadline_at = task.get("deadline_at")
            if deadline_at is not None:
                remaining = float(deadline_at) - time.monotonic()
                if remaining <= 0:
                    self.mark_partial(
                        document,
                        "visual meaning skipped because its document time limit was reached",
                    )
                    continue
            else:
                remaining = None
            raw = self._read_verified_spooled_visual(task)
            if deadline_at is not None:
                remaining = float(deadline_at) - time.monotonic()
                if remaining <= 0:
                    self.mark_partial(
                        document,
                        "visual meaning skipped because its document time limit was reached",
                    )
                    continue
            self._add_local_visual_observation(
                task["image_path"],
                document,
                parent_id=task["parent_id"],
                location_prefix={},
                visual_origin=task["visual_origin"],
                ordinal=task["ordinal"],
                exact_location=task["location"],
                image_bytes=raw,
                timeout=(
                    min(float(remaining), MAX_TIMEOUT_SECONDS)
                    if remaining is not None
                    else MAX_TIMEOUT_SECONDS
                ),
                deadline_at=deadline_at,
                release_paddle=False,
                excluded_normalized_texts=task["excluded_normalized_texts"],
            )

    def extract(self, path: Path) -> None:
        if (
            self.visual_observation_mode == "deferred_per_document"
            and (
                self._deferred_visual_tasks
                or self._visual_spool_root is not None
                or self._current_document is not None
            )
        ):
            raise RuntimeError("deferred visual document scope is already active")
        try:
            suffix = path.suffix.lower()
            if suffix in DIRECT_TEXT_SUFFIXES and path.stat().st_size > MAX_DIRECT_TEXT_BYTES:
                self.extract_large_text(path)
            elif suffix == ".docx":
                self.extract_docx(path)
            elif suffix == ".xlsx":
                self.extract_xlsx(path)
            elif suffix == ".pptx":
                self.extract_pptx(path)
            elif suffix == ".pdf":
                self.extract_pdf(path)
            elif suffix in IMAGE_SUFFIXES:
                self.extract_image(path)
            elif suffix in {".csv", ".tsv"}:
                self.extract_delimited(path)
            elif suffix == ".json":
                self.extract_json(path)
            elif suffix == ".xml":
                self.extract_xml(path)
            elif suffix == ".ipynb":
                self.extract_notebook(path)
            elif suffix in PLAIN_TEXT_SUFFIXES:
                self.extract_plain_text(path)
            else:
                self.extract_other(path)
            if (
                self.visual_observation_mode == "deferred_per_document"
                and self._deferred_visual_tasks
            ):
                self._release_active_paddle_worker()
                self._flush_deferred_visual_observations()
            self.finalize_document()
        finally:
            if self.visual_observation_mode == "deferred_per_document":
                self._cleanup_visual_spool()

    def extract_large_text(self, path: Path) -> None:
        """Retain searchable text from a large file with bounded memory.

        Structural parsers for JSON, XML, notebooks and delimited files first
        materialize the complete source. For a large source, preserving the
        text in deterministic chunks is more useful than either crashing or
        pretending that the structure was fully parsed. The explicit partial
        status lets a later specialized reader replace this fallback.
        """
        doc = self.add_document(path, "bounded-text-stream")
        doc_id = doc["document_id"]
        self.mark_partial(
            doc,
            "large text-like file read with bounded streaming; native structure remains unresolved",
        )
        detected = detect_text_file_encoding(path)
        decoding = "utf-8" if detected == "utf-8-replacement" else detected
        replacement_count = 0
        character_offset = 0
        read_block_count = 0
        chunk_count = 0
        stopped_at_item_limit = False
        with path.open("r", encoding=decoding, errors="replace", newline="") as handle:
            while True:
                if read_block_count >= MAX_STREAM_TEXT_READ_BLOCKS:
                    if handle.read(1):
                        self.mark_partial(
                            doc,
                            "streaming extraction stopped after "
                            f"{MAX_STREAM_TEXT_READ_BLOCKS} read blocks",
                        )
                    break
                block = handle.read(STREAM_TEXT_READ_CHARS)
                if not block:
                    break
                read_block_count += 1
                replacement_count += block.count("\ufffd")
                for question_chunk in exact_text_chunks(block):
                    if not self.may_add_leaf(doc_id):
                        self.mark_partial(
                            doc,
                            "streaming extraction stopped at the configured item limit",
                        )
                        stopped_at_item_limit = True
                        break
                    chunk_count += 1
                    start = character_offset + question_chunk.start
                    end = character_offset + question_chunk.end
                    evidence_type = (
                        "code_block"
                        if path.suffix.lower() in CODE_SUFFIXES else "text_block"
                    )
                    ev = self.add_evidence(
                        doc_id,
                        evidence_type,
                        {
                            "object_index": chunk_count,
                            "locator_text": f"characters={start + 1}-{end}",
                        },
                        content(raw_text=question_chunk.text),
                        ordinal=chunk_count,
                        native_properties={
                            "encoding": detected,
                            "character_start": start,
                            "character_end": end,
                            "character_offset_basis": "zero_based_half_open",
                            "source_stream_read_block": read_block_count,
                            "source_structure_status": "unresolved",
                        },
                        method="bounded_streaming_text",
                        warning=(
                            "native structure was not parsed for this large text-like source"
                        ),
                    )
                    self.contain_document(doc_id, ev["evidence_id"])
                character_offset += len(block)
                if stopped_at_item_limit:
                    break
        if not chunk_count:
            self.mark_partial(doc, "large text-like file contained no decodable text")
        if replacement_count:
            self.mark_partial(
                doc,
                f"streaming decoder inserted {replacement_count} replacement character(s)",
            )

    def office_source(self, path: Path) -> tuple[str | io.BytesIO, bool]:
        if zipfile.is_zipfile(path):
            return str(path), False
        try:
            import msoffcrypto
        except ImportError as exc:
            raise RuntimeError("msoffcrypto-tool is required for password-protected Office files") from exc
        for candidate in self.password_candidates:
            try:
                output = io.BytesIO()
                with path.open("rb") as handle:
                    office = msoffcrypto.OfficeFile(handle)
                    office.load_key(password=candidate)
                    office.decrypt(output)
                output.seek(0)
                return output, True
            except Exception:
                continue
        raise ValueError("password-protected Office file could not be decrypted with derived candidates")

    def extract_docx(self, path: Path) -> None:
        try:
            from docx import Document
            from docx.oxml.table import CT_Tbl
            from docx.oxml.text.paragraph import CT_P
            from docx.table import Table
            from docx.text.paragraph import Paragraph
        except ImportError:
            self.extract_docx_ooxml(path)
            return

        source, decrypted = self.office_source(path)
        validate_ooxml_archive(
            source,
            required_members=DOCX_REQUIRED_OOXML_MEMBERS,
        )
        parsed = Document(source)
        doc = self.add_document(path, "python-docx")
        doc_id = doc["document_id"]
        if decrypted:
            doc["extraction"]["warnings"].append("password-protected Office source decrypted in memory")
        paragraph_index = 0
        table_index = 0
        body_order = 0
        preceding_heading: dict[str, str] | None = None
        for child in parsed.element.body.iterchildren():
            if isinstance(child, CT_P):
                paragraph_index += 1
                body_order += 1
                paragraph = Paragraph(child, parsed)
                if not paragraph.text or not self.may_add_leaf(doc_id):
                    continue
                style: dict[str, Any] = {}
                if paragraph.style and paragraph.style.style_id:
                    style["source_style_id"] = paragraph.style.style_id
                evidence_type = "heading" if paragraph.style and paragraph.style.name.lower().startswith("heading") else "paragraph"
                ev = self.add_evidence(
                    doc_id, evidence_type, {"paragraph_index": paragraph_index}, content(raw_text=paragraph.text),
                    ordinal=paragraph_index, style=style or None,
                    native_properties={
                        "paragraph_style_name": paragraph.style.name if paragraph.style else None,
                        "body_order": body_order,
                    },
                )
                self.contain_document(doc_id, ev["evidence_id"])
                for run_index, run in enumerate(paragraph.runs, 1):
                    if not run.text or not self.may_add_leaf(doc_id):
                        continue
                    run_style: dict[str, Any] = {}
                    for key in ("bold", "italic", "underline"):
                        value = getattr(run, key, None)
                        if value is not None:
                            run_style[key] = bool(value)
                    if run.font.name:
                        run_style["font_name"] = run.font.name
                    if run.font.size:
                        run_style["font_size_pt"] = float(run.font.size.pt)
                    color = getattr(getattr(run.font, "color", None), "rgb", None)
                    if color:
                        run_style["font_color_argb"] = f"FF{color}"[-8:].upper()
                    highlight = getattr(run.font, "highlight_color", None)
                    if run_style or highlight is not None:
                        self.add_evidence(
                            doc_id,
                            "style_span",
                            {"paragraph_index": paragraph_index, "object_index": run_index},
                            content(raw_text=run.text),
                            parent_id=ev["evidence_id"],
                            ordinal=run_index,
                            style=run_style or None,
                            native_properties={"highlight_color": str(highlight) if highlight is not None else None},
                        )
                if evidence_type == "heading":
                    preceding_heading = {"text": paragraph.text, "evidence_id": ev["evidence_id"]}
            elif isinstance(child, CT_Tbl):
                table_index += 1
                body_order += 1
                table = Table(child, parsed)
                native_properties: dict[str, Any] = {"body_order": body_order}
                if preceding_heading is not None:
                    native_properties.update({
                        "preceding_heading_text": preceding_heading["text"],
                        "preceding_heading_evidence_id": preceding_heading["evidence_id"],
                    })
                table_ev = self.add_evidence(
                    doc_id, "table", {"table_index": table_index},
                    content(raw_value={"rows": len(table.rows), "columns": len(table.columns)}),
                    ordinal=table_index,
                    native_properties=native_properties,
                )
                self.contain_document(doc_id, table_ev["evidence_id"])
                if preceding_heading is not None:
                    self.add_relation(
                        "structural", "section_contains",
                        {"record_type": "evidence", "record_id": preceding_heading["evidence_id"]},
                        {"record_type": "evidence", "record_id": table_ev["evidence_id"]},
                    )
                for row_index, row in enumerate(table.rows, 1):
                    for column_index, cell in enumerate(row.cells, 1):
                        if not self.may_add_leaf(doc_id):
                            break
                        self.add_evidence(
                            doc_id, "table_cell",
                            {"table_index": table_index, "row_index": row_index, "column_index": column_index},
                            content(raw_text=cell.text), parent_id=table_ev["evidence_id"],
                            ordinal=column_index,
                        )
                    if self.limit_reached(doc_id):
                        break

        seen_parts: set[str] = set()
        for section_index, section in enumerate(parsed.sections, 1):
            for evidence_type, part in (("header", section.header), ("footer", section.footer)):
                part_name = str(getattr(getattr(part, "part", None), "partname", f"{evidence_type}-{section_index}"))
                if part_name in seen_parts:
                    continue
                seen_parts.add(part_name)
                values = [paragraph.text for paragraph in part.paragraphs if paragraph.text]
                for table in part.tables:
                    values.extend(
                        " | ".join(cell.text for cell in row.cells)
                        for row in table.rows
                        if any(cell.text for cell in row.cells)
                    )
                text_value = "\n".join(values).strip()
                if text_value and self.may_add_leaf(doc_id):
                    ev = self.add_evidence(
                        doc_id,
                        evidence_type,
                        {"section": f"section-{section_index}", "source_member": part_name},
                        content(raw_text=text_value),
                        ordinal=section_index,
                    )
                    self.contain_document(doc_id, ev["evidence_id"])

        comments = getattr(parsed, "comments", None)
        if comments is not None:
            for comment_index, comment in enumerate(comments, 1):
                text_value = getattr(comment, "text", "")
                if text_value and self.may_add_leaf(doc_id):
                    ev = self.add_evidence(
                        doc_id,
                        "comment",
                        {"object_index": comment_index, "locator_text": f"comment={comment_index}"},
                        content(raw_text=text_value),
                        ordinal=comment_index,
                        native_properties={
                            "author": getattr(comment, "author", None),
                            "initials": getattr(comment, "initials", None),
                        },
                    )
                    self.contain_document(doc_id, ev["evidence_id"])
        if isinstance(source, io.BytesIO):
            source.seek(0)
        with zipfile.ZipFile(source) as archive:
            self._project_ooxml_referenced_media(
                archive,
                doc,
                media_prefixes=("word/media/",),
            )

    def extract_docx_ooxml(self, path: Path) -> None:
        """Read core DOCX structures without optional Python packages."""
        source, decrypted = self.office_source(path)
        validate_ooxml_archive(
            source,
            required_members=DOCX_REQUIRED_OOXML_MEMBERS,
        )
        doc = self.add_document(path, "ooxml-stdlib-docx-fallback")
        doc_id = doc["document_id"]
        self.mark_partial(
            doc,
            "DOCX standard-library fallback preserves body text/table order but rich run formatting, fields, revisions, footnotes, and section layout remain unresolved",
        )
        if decrypted:
            doc["extraction"]["warnings"].append(
                "password-protected Office source decrypted in memory"
            )
        if isinstance(source, io.BytesIO):
            source.seek(0)
        with zipfile.ZipFile(source) as archive:
            relationships = _ooxml_relationships(archive)
            _require_ooxml_office_document_binding(
                archive, relationships, "word/document.xml"
            )
            root = _ooxml_xml_root(
                archive,
                "word/document.xml",
                namespaces=OOXML_WORDPROCESSING_NAMESPACES,
                local_names={"document"},
            )
            bodies = _direct_xml_children(
                root, OOXML_WORDPROCESSING_NAMESPACES, "body"
            )
            if len(bodies) != 1:
                raise ValueError("ooxml_word_body_invalid")
            styles = _word_style_catalog(archive, relationships)
            paragraph_index = 0
            table_index = 0
            body_order = 0
            preceding_heading: dict[str, str] | None = None

            for block in _word_block_children(bodies[0]):
                _, block_name = _xml_name(block.tag)
                body_order += 1
                if block_name == "p":
                    paragraph_index += 1
                    text_value = _word_paragraph_text(block)
                    if not text_value or not self.may_add_leaf(doc_id):
                        continue
                    style_id = _word_paragraph_style_id(block)
                    style_row = styles.get(style_id or "", {})
                    style_name = style_row.get("name")
                    evidence_type = (
                        "heading"
                        if _word_style_is_heading(style_id, styles)
                        or _word_paragraph_has_outline_level(block)
                        else "paragraph"
                    )
                    ev = self.add_evidence(
                        doc_id,
                        evidence_type,
                        {"paragraph_index": paragraph_index},
                        content(raw_text=text_value),
                        ordinal=paragraph_index,
                        style=(
                            {"source_style_id": style_id}
                            if style_id else None
                        ),
                        native_properties={
                            "paragraph_style_name": style_name,
                            "body_order": body_order,
                            "fallback_detail_status": "core_text_and_style_reference_only",
                        },
                        method="ooxml_stdlib_docx_fallback",
                    )
                    self.contain_document(doc_id, ev["evidence_id"])
                    if evidence_type == "heading":
                        preceding_heading = {
                            "text": text_value,
                            "evidence_id": ev["evidence_id"],
                        }
                    continue

                if block_name != "tbl":
                    continue
                table_index += 1
                rows = _direct_xml_children(
                    block, OOXML_WORDPROCESSING_NAMESPACES, "tr"
                )
                parsed_rows: list[list[str]] = []
                for row in rows:
                    parsed_cells: list[str] = []
                    for cell in _direct_xml_children(
                        row, OOXML_WORDPROCESSING_NAMESPACES, "tc"
                    ):
                        values = [
                            _word_paragraph_text(item)
                            for item in _word_block_children(cell)
                            if _xml_name(item.tag)[1] == "p"
                        ]
                        parsed_cells.append("\n".join(values))
                    parsed_rows.append(parsed_cells)
                native_properties: dict[str, Any] = {
                    "body_order": body_order,
                    "fallback_detail_status": "merged_cell_and_nested_table_semantics_unresolved",
                }
                if preceding_heading is not None:
                    native_properties.update({
                        "preceding_heading_text": preceding_heading["text"],
                        "preceding_heading_evidence_id": preceding_heading[
                            "evidence_id"
                        ],
                    })
                table_ev = self.add_evidence(
                    doc_id,
                    "table",
                    {"table_index": table_index},
                    content(raw_value={
                        "rows": len(parsed_rows),
                        "columns": max(
                            (len(row) for row in parsed_rows), default=0
                        ),
                    }),
                    ordinal=table_index,
                    native_properties=native_properties,
                    method="ooxml_stdlib_docx_fallback",
                )
                self.contain_document(doc_id, table_ev["evidence_id"])
                if preceding_heading is not None:
                    self.add_relation(
                        "structural",
                        "section_contains",
                        {
                            "record_type": "evidence",
                            "record_id": preceding_heading["evidence_id"],
                        },
                        {
                            "record_type": "evidence",
                            "record_id": table_ev["evidence_id"],
                        },
                    )
                for row_index, row in enumerate(parsed_rows, 1):
                    for column_index, cell_text in enumerate(row, 1):
                        if not self.may_add_leaf(doc_id):
                            break
                        self.add_evidence(
                            doc_id,
                            "table_cell",
                            {
                                "table_index": table_index,
                                "row_index": row_index,
                                "column_index": column_index,
                            },
                            content(raw_text=cell_text),
                            parent_id=table_ev["evidence_id"],
                            ordinal=column_index,
                            native_properties={
                                "fallback_detail_status": "plain_cell_text_only"
                            },
                            method="ooxml_stdlib_docx_fallback",
                        )
                    if self.limit_reached(doc_id):
                        break

            used_relationships = _ooxml_used_relationship_ids(
                archive, "word/document.xml"
            )
            seen_parts: set[str] = set()
            section_index = 0
            for row in relationships.get("word/document.xml", []):
                kind = _ooxml_relationship_kind(row["type"])
                if (
                    kind not in {"header", "footer"}
                    or row["target_mode"].casefold() == "external"
                    or row["id"] not in used_relationships
                ):
                    continue
                member = _resolve_ooxml_target(
                    "word/document.xml", row["target"]
                )
                if (
                    member is None
                    or member in seen_parts
                    or not _word_visual_source_part_allowed(archive, member)
                ):
                    if member is None:
                        raise ValueError("ooxml_word_section_binding_invalid")
                    continue
                seen_parts.add(member)
                section_index += 1
                part_root = _ooxml_xml_root(
                    archive,
                    member,
                    namespaces=OOXML_WORDPROCESSING_NAMESPACES,
                    local_names={"hdr" if kind == "header" else "ftr"},
                )
                values: list[str] = []
                for part_block in _word_block_children(part_root):
                    _, part_block_name = _xml_name(part_block.tag)
                    if part_block_name == "p":
                        value = _word_paragraph_text(part_block)
                        if value:
                            values.append(value)
                    elif part_block_name == "tbl":
                        for table_row in _direct_xml_children(
                            part_block,
                            OOXML_WORDPROCESSING_NAMESPACES,
                            "tr",
                        ):
                            row_values: list[str] = []
                            for cell in _direct_xml_children(
                                table_row,
                                OOXML_WORDPROCESSING_NAMESPACES,
                                "tc",
                            ):
                                row_values.append("\n".join(
                                    _word_paragraph_text(item)
                                    for item in _word_block_children(cell)
                                    if _xml_name(item.tag)[1] == "p"
                                ))
                            if any(row_values):
                                values.append(" | ".join(row_values))
                text_value = "\n".join(values).strip()
                if text_value and self.may_add_leaf(doc_id):
                    ev = self.add_evidence(
                        doc_id,
                        kind,
                        {
                            "section": f"relationship-{section_index}",
                            "source_member": member,
                        },
                        content(raw_text=text_value),
                        ordinal=section_index,
                        native_properties={
                            "relationship_id": row["id"],
                            "reference_occurrences": used_relationships[
                                row["id"]
                            ],
                        },
                        method="ooxml_stdlib_docx_fallback",
                    )
                    self.contain_document(doc_id, ev["evidence_id"])

            comment_rows = [
                row for row in relationships.get("word/document.xml", [])
                if _ooxml_relationship_kind(row["type"]) == "comments"
                and row["target_mode"].casefold() != "external"
            ]
            if len(comment_rows) > 1:
                raise ValueError("ooxml_word_comments_binding_ambiguous")
            if comment_rows:
                comments_member = _resolve_ooxml_target(
                    "word/document.xml", comment_rows[0]["target"]
                )
                if comments_member is None:
                    raise ValueError("ooxml_word_comments_binding_invalid")
                comments_root = _ooxml_xml_root(
                    archive,
                    comments_member,
                    namespaces=OOXML_WORDPROCESSING_NAMESPACES,
                    local_names={"comments"},
                )
                for comment_index, comment in enumerate(
                    _direct_xml_children(
                        comments_root,
                        OOXML_WORDPROCESSING_NAMESPACES,
                        "comment",
                    ),
                    1,
                ):
                    text_value = "\n".join(
                        _word_paragraph_text(item)
                        for item in _word_block_children(comment)
                        if _xml_name(item.tag)[1] == "p"
                    ).strip()
                    if not text_value or not self.may_add_leaf(doc_id):
                        continue
                    ev = self.add_evidence(
                        doc_id,
                        "comment",
                        {
                            "object_index": comment_index,
                            "source_member": comments_member,
                            "locator_text": f"comment={comment_index}",
                        },
                        content(raw_text=text_value),
                        ordinal=comment_index,
                        native_properties={
                            "author": _xml_attribute(
                                comment,
                                OOXML_WORDPROCESSING_NAMESPACES,
                                "author",
                            ),
                            "initials": _xml_attribute(
                                comment,
                                OOXML_WORDPROCESSING_NAMESPACES,
                                "initials",
                            ),
                        },
                        method="ooxml_stdlib_docx_fallback",
                    )
                    self.contain_document(doc_id, ev["evidence_id"])

            self._project_ooxml_referenced_media(
                archive,
                doc,
                media_prefixes=("word/media/",),
            )

    def extract_xlsx(self, path: Path) -> None:
        try:
            from openpyxl import load_workbook
        except ImportError:
            self.extract_xlsx_ooxml(path)
            return

        source, decrypted = self.office_source(path)
        validate_ooxml_archive(
            source,
            required_members=frozenset({
                "[Content_Types].xml",
                "xl/workbook.xml",
                "xl/_rels/workbook.xml.rels",
            }),
        )
        workbook = load_workbook(source, data_only=False, read_only=False)
        if isinstance(source, io.BytesIO):
            source.seek(0)
        # ``data_only=True`` exposes the result saved in the XLSX package.  It
        # does not calculate the formula, so every downstream representation
        # labels it as a stored, not-recalculated value.
        cached_workbook = load_workbook(source, data_only=True, read_only=True)
        doc = self.add_document(path, "openpyxl+ooxml")
        doc_id = doc["document_id"]
        if decrypted:
            doc["extraction"]["warnings"].append("password-protected Office source decrypted in memory")
        sheet_ids: dict[str, str] = {}
        for sheet_index, sheet in enumerate(workbook.worksheets, 1):
            cached_sheet = cached_workbook[sheet.title]
            cached_rows = iter(cached_sheet.iter_rows())
            sheet_ev = self.add_evidence(
                doc_id, "worksheet", {"sheet_name": sheet.title},
                content(raw_value={"title": sheet.title, "max_row": sheet.max_row, "max_column": sheet.max_column}),
                ordinal=sheet_index,
            )
            sheet_ids[sheet.title] = sheet_ev["evidence_id"]
            self.contain_document(doc_id, sheet_ev["evidence_id"])
            for row in sheet.iter_rows():
                cached_row = next(cached_rows, ())
                cached_by_coordinate = {
                    item.coordinate: item.value
                    for item in cached_row
                    if getattr(item, "coordinate", None)
                }
                for cell_obj in row:
                    if cell_obj.value is None:
                        continue
                    if not self.may_add_leaf(doc_id):
                        break
                    style: dict[str, Any] = {
                        "source_style_id": str(cell_obj.style_id),
                        "number_format": cell_obj.number_format,
                    }
                    if cell_obj.font:
                        if cell_obj.font.bold is not None:
                            style["bold"] = bool(cell_obj.font.bold)
                        if cell_obj.font.italic is not None:
                            style["italic"] = bool(cell_obj.font.italic)
                        if cell_obj.font.underline:
                            style["underline"] = True
                        if cell_obj.font.name:
                            style["font_name"] = cell_obj.font.name
                        if cell_obj.font.sz:
                            style["font_size_pt"] = float(cell_obj.font.sz)
                    if cell_obj.fill and cell_obj.fill.fill_type:
                        style["fill_type"] = cell_obj.fill.fill_type
                    is_formula = cell_obj.data_type == "f"
                    cached_value = (
                        cached_by_coordinate.get(cell_obj.coordinate)
                        if is_formula else None
                    )
                    cell_value = cached_value if is_formula else cell_obj.value
                    native_properties: dict[str, Any] = {
                        "data_type": cell_obj.data_type,
                    }
                    if is_formula:
                        native_properties.update({
                            "cached_value": json_value(cached_value),
                            "cached_value_available": cached_value is not None,
                            "cached_value_status": FORMULA_CACHED_VALUE_STATUS,
                        })
                    cell_ev = self.add_evidence(
                        doc_id, "table_cell", {"sheet_name": sheet.title, "cell": cell_obj.coordinate},
                        content(raw_value=cell_value), parent_id=sheet_ev["evidence_id"],
                        style=style,
                        native_properties=native_properties,
                    )
                    if is_formula:
                        self.add_evidence(
                            doc_id, "formula", {"sheet_name": sheet.title, "cell": cell_obj.coordinate},
                            content(raw_text=str(cell_obj.value)), parent_id=cell_ev["evidence_id"],
                            native_properties={
                                "cached_value": json_value(cached_value),
                                "cached_value_available": cached_value is not None,
                                "cached_value_status": FORMULA_CACHED_VALUE_STATUS,
                            },
                        )
                    if cell_obj.comment is not None:
                        self.add_evidence(
                            doc_id,
                            "comment",
                            {"sheet_name": sheet.title, "cell": cell_obj.coordinate},
                            content(raw_text=cell_obj.comment.text),
                            parent_id=cell_ev["evidence_id"],
                            native_properties={"author": cell_obj.comment.author},
                        )
                if self.limit_reached(doc_id):
                    break
            for merged_index, merged in enumerate(sheet.merged_cells.ranges, 1):
                merged_ev = self.add_evidence(
                    doc_id, "merged_range", {"sheet_name": sheet.title, "range": str(merged)},
                    content(raw_text=str(merged)), parent_id=sheet_ev["evidence_id"], ordinal=merged_index,
                )
                del merged_ev
            for pivot_index, pivot in enumerate(getattr(sheet, "_pivots", []), 1):
                name = getattr(pivot, "name", None) or f"pivot_{pivot_index}"
                self.add_evidence(
                    doc_id, "pivot_table",
                    {"sheet_name": sheet.title, "object_index": pivot_index, "object_id": str(name)},
                    content(raw_value={"name": str(name)}), parent_id=sheet_ev["evidence_id"], ordinal=pivot_index,
                )
            auto_filter = getattr(sheet, "auto_filter", None)
            if auto_filter is not None and getattr(auto_filter, "ref", None):
                ev = self.add_evidence(
                    doc_id,
                    "filter",
                    {"sheet_name": sheet.title, "range": str(auto_filter.ref)},
                    content(raw_value={
                        "ref": str(auto_filter.ref),
                        "filter_columns": [json_value(item) for item in getattr(auto_filter, "filterColumn", [])],
                    }),
                    parent_id=sheet_ev["evidence_id"],
                )
            validations = getattr(getattr(sheet, "data_validations", None), "dataValidation", [])
            for validation_index, validation in enumerate(validations, 1):
                ev = self.add_evidence(
                    doc_id,
                    "data_validation",
                    {"sheet_name": sheet.title, "object_index": validation_index,
                     "range": str(getattr(validation, "sqref", "")) or "unknown"},
                    content(raw_value={
                        "type": getattr(validation, "type", None),
                        "formula1": getattr(validation, "formula1", None),
                        "formula2": getattr(validation, "formula2", None),
                        "operator": getattr(validation, "operator", None),
                    }),
                    parent_id=sheet_ev["evidence_id"],
                    ordinal=validation_index,
                )

        defined_names = getattr(workbook, "defined_names", None)
        if defined_names is not None:
            values = defined_names.values() if hasattr(defined_names, "values") else defined_names.definedName
            for name_index, defined_name in enumerate(values, 1):
                ev = self.add_evidence(
                    doc_id,
                    "defined_name",
                    {"object_index": name_index, "object_id": str(getattr(defined_name, "name", name_index))},
                    content(raw_value={
                        "name": getattr(defined_name, "name", None),
                        "attr_text": getattr(defined_name, "attr_text", None),
                        "local_sheet_id": getattr(defined_name, "localSheetId", None),
                    }),
                    ordinal=name_index,
                )
                self.contain_document(doc_id, ev["evidence_id"])

        if isinstance(source, io.BytesIO):
            source.seek(0)
        with zipfile.ZipFile(source) as archive:
            self._project_ooxml_charts(archive, doc)
            self._project_ooxml_referenced_media(
                archive,
                doc,
                media_prefixes=("xl/media/",),
            )
        cached_workbook.close()
        workbook.close()

    def extract_xlsx_ooxml(self, path: Path) -> None:
        """Read core XLSX structure using only OOXML and the standard library.

        This fallback deliberately reports ``partial`` because it does not
        resolve the complete Excel formatting/comment/pivot object model.  It
        still preserves the native sheet name, cell coordinate, formula,
        merge, filter, validation, chart, and embedded-image source binding so
        downstream SearchUnits do not collapse the workbook into flat text.
        """
        source, decrypted = self.office_source(path)
        validate_ooxml_archive(
            source,
            required_members=frozenset({
                "[Content_Types].xml",
                "xl/workbook.xml",
                "xl/_rels/workbook.xml.rels",
            }),
        )
        doc = self.add_document(path, "ooxml-stdlib-xlsx-fallback")
        doc_id = doc["document_id"]
        self.mark_partial(
            doc,
            "openpyxl unavailable; OOXML fallback preserves core workbook structure but limits comments, pivots, and resolved formatting",
        )
        if decrypted:
            doc["extraction"]["warnings"].append("password-protected Office source decrypted in memory")

        def child(element: ElementTree.Element, local_name: str) -> ElementTree.Element | None:
            return next((item for item in element if item.tag.endswith("}" + local_name) or item.tag == local_name), None)

        def attr_by_local_name(element: ElementTree.Element, local_name: str) -> str | None:
            return next((value for key, value in element.attrib.items() if key == local_name or key.endswith("}" + local_name)), None)

        def numeric(raw: str) -> Any:
            # Python floats would round the OOXML decimal lexeme. Preserve
            # decimal/scientific values verbatim; arbitrary-size integers are
            # exact, and every numeric cell also records ``raw_lexeme`` below.
            if not re.fullmatch(r"[+-]?\d+", raw):
                return raw
            try:
                return int(raw)
            except ValueError:
                return raw

        def cell_value(
            cell: ElementTree.Element,
            shared: list[str],
        ) -> tuple[Any, str, str | None, Any, str | None]:
            kind = cell.attrib.get("t", "n")
            formula_node = child(cell, "f")
            value_node = child(cell, "v")
            formula = formula_node.text if formula_node is not None and formula_node.text is not None else None
            raw_value = value_node.text if value_node is not None and value_node.text is not None else ""
            cached: Any = None
            if formula is not None:
                if not raw_value:
                    cached = None
                elif kind == "b":
                    cached = raw_value == "1"
                elif kind in {"str", "e", "d"}:
                    cached = raw_value
                else:
                    cached = numeric(raw_value)
                return cached, "f", formula, cached, raw_value or None
            if kind == "inlineStr":
                value = "".join(item.text or "" for item in cell.iter() if item.tag.endswith("}t") or item.tag == "t")
                return value, "inlineStr", None, None, None
            if kind == "s":
                try:
                    return shared[int(raw_value)], "s", None, None, None
                except (ValueError, IndexError):
                    return raw_value, "s", None, None, None
            if kind == "b":
                return raw_value == "1", "b", None, None, None
            if kind in {"str", "e", "d"}:
                return raw_value, kind, None, None, None
            return numeric(raw_value), kind, None, None, raw_value or None

        if isinstance(source, io.BytesIO):
            source.seek(0)
        with zipfile.ZipFile(source) as archive:
            names = set(archive.namelist())
            if "xl/workbook.xml" not in names:
                raise ValueError("XLSX package has no xl/workbook.xml")

            shared: list[str] = []
            if "xl/sharedStrings.xml" in names:
                shared_root = ElementTree.fromstring(archive.read("xl/sharedStrings.xml"))
                for item in shared_root.iter():
                    if item.tag.endswith("}si") or item.tag == "si":
                        shared.append("".join(
                            node.text or "" for node in item.iter()
                            if node.tag.endswith("}t") or node.tag == "t"
                        ))

            custom_formats: dict[str, str] = {}
            cell_formats: list[dict[str, str]] = []
            if "xl/styles.xml" in names:
                styles_root = ElementTree.fromstring(archive.read("xl/styles.xml"))
                for item in styles_root.iter():
                    if item.tag.endswith("}numFmt") or item.tag == "numFmt":
                        if item.attrib.get("numFmtId") and item.attrib.get("formatCode"):
                            custom_formats[item.attrib["numFmtId"]] = item.attrib["formatCode"]
                    if item.tag.endswith("}cellXfs") or item.tag == "cellXfs":
                        cell_formats = [dict(node.attrib) for node in item if node.tag.endswith("}xf") or node.tag == "xf"]

            relationships: dict[str, str] = {}
            rels_name = "xl/_rels/workbook.xml.rels"
            if rels_name in names:
                rels_root = ElementTree.fromstring(archive.read(rels_name))
                for relation in rels_root:
                    identifier = relation.attrib.get("Id")
                    target = relation.attrib.get("Target")
                    if not identifier or not target or relation.attrib.get("TargetMode") == "External":
                        continue
                    member = posixpath.normpath(
                        target.lstrip("/") if target.startswith("/")
                        else posixpath.join("xl", target)
                    )
                    if member.startswith("xl/") and ".." not in PurePosixPath(member).parts:
                        relationships[identifier] = member

            workbook_root = ElementTree.fromstring(archive.read("xl/workbook.xml"))
            sheets = [
                item for item in workbook_root.iter()
                if item.tag.endswith("}sheet") or item.tag == "sheet"
            ]
            for sheet_index, sheet in enumerate(sheets, 1):
                title = sheet.attrib.get("name") or f"Sheet{sheet_index}"
                relation_id = attr_by_local_name(sheet, "id")
                member = relationships.get(relation_id or "", f"xl/worksheets/sheet{sheet_index}.xml")
                if member not in names:
                    self.mark_partial(doc, f"worksheet OOXML part unavailable for sheet index {sheet_index}")
                    continue
                sheet_root = ElementTree.fromstring(archive.read(member))
                cells = [
                    item for item in sheet_root.iter()
                    if item.tag.endswith("}c") or item.tag == "c"
                ]
                row_numbers: list[int] = []
                column_numbers: list[int] = []
                for cell in cells:
                    coordinate = cell.attrib.get("r", "")
                    match = re.fullmatch(r"([A-Z]{1,4})([1-9][0-9]*)", coordinate)
                    if not match:
                        continue
                    letters, row_text = match.groups()
                    column = 0
                    for character in letters:
                        column = column * 26 + ord(character) - ord("A") + 1
                    column_numbers.append(column)
                    row_numbers.append(int(row_text))
                sheet_ev = self.add_evidence(
                    doc_id, "worksheet", {"sheet_name": title},
                    content(raw_value={
                        "title": title,
                        "max_row": max(row_numbers, default=0),
                        "max_column": max(column_numbers, default=0),
                    }),
                    ordinal=sheet_index,
                    native_properties={"source_member": member, "state": sheet.attrib.get("state", "visible")},
                )
                self.contain_document(doc_id, sheet_ev["evidence_id"])
                for cell in cells:
                    coordinate = cell.attrib.get("r")
                    if not coordinate or not self.may_add_leaf(doc_id):
                        continue
                    value, data_type, formula, cached, raw_lexeme = cell_value(cell, shared)
                    if value in {None, ""} and formula is None:
                        continue
                    style_id = cell.attrib.get("s", "0")
                    style: dict[str, Any] = {"source_style_id": style_id}
                    try:
                        style_record = cell_formats[int(style_id)]
                    except (ValueError, IndexError):
                        style_record = {}
                    number_format_id = style_record.get("numFmtId")
                    if number_format_id is not None:
                        style["number_format"] = custom_formats.get(number_format_id, f"builtin:{number_format_id}")
                    native_properties: dict[str, Any] = {"data_type": data_type, "source_member": member}
                    if raw_lexeme is not None:
                        native_properties[
                            "cached_raw_lexeme" if formula is not None else "raw_lexeme"
                        ] = raw_lexeme
                    if cached is not None:
                        native_properties["cached_value"] = cached
                    if formula is not None:
                        native_properties.update({
                            "cached_value_available": cached is not None,
                            "cached_value_status": FORMULA_CACHED_VALUE_STATUS,
                        })
                    cell_ev = self.add_evidence(
                        doc_id, "table_cell", {"sheet_name": title, "cell": coordinate},
                        content(raw_value=value), parent_id=sheet_ev["evidence_id"],
                        style=style, native_properties=native_properties,
                    )
                    if formula is not None:
                        self.add_evidence(
                            doc_id, "formula", {"sheet_name": title, "cell": coordinate},
                            content(raw_text="=" + formula), parent_id=cell_ev["evidence_id"],
                            native_properties={
                                "cached_value": cached,
                                "cached_value_available": cached is not None,
                                "cached_value_status": FORMULA_CACHED_VALUE_STATUS,
                                **({"cached_raw_lexeme": raw_lexeme} if raw_lexeme is not None else {}),
                            },
                        )
                for merged_index, merged in enumerate(
                    (item for item in sheet_root.iter() if item.tag.endswith("}mergeCell") or item.tag == "mergeCell"),
                    1,
                ):
                    reference = merged.attrib.get("ref")
                    if reference:
                        self.add_evidence(
                            doc_id, "merged_range", {"sheet_name": title, "range": reference},
                            content(raw_text=reference), parent_id=sheet_ev["evidence_id"], ordinal=merged_index,
                        )
                auto_filter = next(
                    (item for item in sheet_root.iter() if item.tag.endswith("}autoFilter") or item.tag == "autoFilter"),
                    None,
                )
                if auto_filter is not None and auto_filter.attrib.get("ref"):
                    self.add_evidence(
                        doc_id, "filter", {"sheet_name": title, "range": auto_filter.attrib["ref"]},
                        content(raw_value={
                            "ref": auto_filter.attrib["ref"],
                            "filter_columns": [dict(item.attrib) for item in auto_filter],
                        }),
                        parent_id=sheet_ev["evidence_id"],
                    )
                for validation_index, validation in enumerate(
                    (item for item in sheet_root.iter() if item.tag.endswith("}dataValidation") or item.tag == "dataValidation"),
                    1,
                ):
                    formula1 = child(validation, "formula1")
                    formula2 = child(validation, "formula2")
                    reference = validation.attrib.get("sqref", "unknown")
                    self.add_evidence(
                        doc_id, "data_validation",
                        {"sheet_name": title, "object_index": validation_index, "range": reference},
                        content(raw_value={
                            "type": validation.attrib.get("type"),
                            "formula1": formula1.text if formula1 is not None else None,
                            "formula2": formula2.text if formula2 is not None else None,
                            "operator": validation.attrib.get("operator"),
                        }),
                        parent_id=sheet_ev["evidence_id"], ordinal=validation_index,
                    )

            defined_names = [
                item for item in workbook_root.iter()
                if item.tag.endswith("}definedName") or item.tag == "definedName"
            ]
            for name_index, defined_name in enumerate(defined_names, 1):
                name = defined_name.attrib.get("name") or str(name_index)
                ev = self.add_evidence(
                    doc_id, "defined_name", {"object_index": name_index, "object_id": name},
                    content(raw_value={
                        "name": name,
                        "attr_text": defined_name.text,
                        "local_sheet_id": defined_name.attrib.get("localSheetId"),
                    }), ordinal=name_index,
                )
                self.contain_document(doc_id, ev["evidence_id"])

            self._project_ooxml_charts(archive, doc)
            self._project_ooxml_referenced_media(
                archive,
                doc,
                media_prefixes=("xl/media/",),
            )

    def extract_pptx(self, path: Path) -> None:
        try:
            from pptx import Presentation
        except ImportError:
            self.extract_pptx_ooxml(path)
            return

        source, decrypted = self.office_source(path)
        validate_ooxml_archive(
            source,
            required_members=PPTX_REQUIRED_OOXML_MEMBERS,
        )
        presentation = Presentation(source)
        doc = self.add_document(path, "python-pptx")
        doc_id = doc["document_id"]
        if decrypted:
            doc["extraction"]["warnings"].append("password-protected Office source decrypted in memory")

        def walk_shapes(shapes: Any, prefix: str = "") -> Any:
            for local_index, candidate in enumerate(shapes, 1):
                key = f"{prefix}.{local_index}" if prefix else str(local_index)
                yield key, candidate
                nested = getattr(candidate, "shapes", None)
                if nested is not None:
                    yield from walk_shapes(nested, key)

        for slide_number, slide in enumerate(presentation.slides, 1):
            slide_ev = self.add_evidence(
                doc_id, "slide", {"slide_number": slide_number},
                content(raw_value={"slide_number": slide_number, "shape_count": len(slide.shapes)}),
                ordinal=slide_number,
            )
            self.contain_document(doc_id, slide_ev["evidence_id"])
            for shape_index, (shape_path, shape) in enumerate(walk_shapes(slide.shapes), 1):
                if not self.may_add_leaf(doc_id):
                    break
                shape_locator_id = str(shape.shape_id) if "." not in shape_path else f"{shape.shape_id}:{shape_path}"
                text_value = getattr(shape, "text", "")
                shape_content = content(raw_text=text_value) if text_value else content(
                    raw_value={"shape_type": str(shape.shape_type), "name": shape.name}
                )
                geometry = {
                    "coordinate_space": "slide", "unit": "emu",
                    "x": shape.left, "y": shape.top, "width": shape.width, "height": shape.height,
                }
                shape_ev = self.add_evidence(
                    doc_id, "shape",
                    {"slide_number": slide_number, "shape_id": shape_locator_id, "object_index": shape_index},
                    shape_content, parent_id=slide_ev["evidence_id"], ordinal=shape_index,
                    geometry=geometry,
                    native_properties={"name": shape.name, "shape_type": str(shape.shape_type), "shape_path": shape_path},
                )
                if getattr(shape, "has_table", False):
                    table_ev = self.add_evidence(
                        doc_id, "table",
                        {"slide_number": slide_number, "shape_id": shape_locator_id},
                        content(raw_value={"rows": len(shape.table.rows), "columns": len(shape.table.columns)}),
                        parent_id=shape_ev["evidence_id"],
                    )
                    for row_index, row in enumerate(shape.table.rows, 1):
                        for column_index, cell_obj in enumerate(row.cells, 1):
                            if not self.may_add_leaf(doc_id):
                                break
                            self.add_evidence(
                                doc_id, "table_cell",
                                {"slide_number": slide_number, "shape_id": shape_locator_id,
                                 "row_index": row_index, "column_index": column_index},
                                content(raw_text=cell_obj.text), parent_id=table_ev["evidence_id"],
                                ordinal=column_index,
                            )
                        if self.limit_reached(doc_id):
                            break
                if getattr(shape, "has_chart", False):
                    chart = shape.chart
                    chart_part = getattr(chart, "part", None)
                    chart_blob = getattr(chart_part, "blob", None)
                    chart_member = str(
                        getattr(chart_part, "partname", "")
                    ).lstrip("/")
                    if not isinstance(chart_blob, bytes) or not chart_member:
                        self.mark_partial(
                            doc,
                            f"slide {slide_number} chart {shape_locator_id} lacks a source-bound OOXML part",
                        )
                    else:
                        try:
                            chart_payload = _ooxml_chart_payload(
                                chart_blob, chart_member
                            )
                        except ValueError as exc:
                            self.mark_partial(
                                doc,
                                f"slide {slide_number} chart {shape_locator_id} could not be parsed: {exc}",
                            )
                        else:
                            chart_location = {
                                "slide_number": slide_number,
                                "shape_id": shape_locator_id,
                                "source_member": chart_member,
                                "locator_text": (
                                    f"slide={slide_number};shape={shape_locator_id};"
                                    f"chart={chart_member}"
                                ),
                            }
                            chart_ev = self.add_evidence(
                                doc_id,
                                "chart",
                                chart_location,
                                content(
                                    raw_text=_chart_summary_text(chart_payload)
                                ),
                                parent_id=shape_ev["evidence_id"],
                                native_properties={
                                    "ooxml_part": chart_member,
                                    "xml_sha256": chart_payload["xml_sha256"],
                                    "chart_payload": chart_payload,
                                    "cached_value_status": FORMULA_CACHED_VALUE_STATUS,
                                },
                                method="verified_ooxml_chart_cache",
                            )
                            for series_index, series in enumerate(
                                chart_payload["series"], 1
                            ):
                                if (
                                    series["cache_status"].startswith("unresolved_")
                                    or series["cache_status"] == "cached_values_missing"
                                ):
                                    self.mark_partial(
                                        doc,
                                        f"slide {slide_number} chart {shape_locator_id} series {series_index} has {series['cache_status']}",
                                    )
                                self.add_evidence(
                                    doc_id,
                                    "chart_series",
                                    {
                                        **chart_location,
                                        "series_index": series_index,
                                        "locator_text": (
                                            f"{chart_location['locator_text']};"
                                            f"series={series_index}"
                                        ),
                                    },
                                    content(
                                        raw_text=_chart_series_text(
                                            chart_payload, series
                                        )
                                    ),
                                    parent_id=chart_ev["evidence_id"],
                                    ordinal=series_index,
                                    native_properties={
                                        "ooxml_part": chart_member,
                                        "xml_sha256": chart_payload["xml_sha256"],
                                        "series": series,
                                        "cached_value_status": FORMULA_CACHED_VALUE_STATUS,
                                    },
                                    method="verified_ooxml_chart_cache",
                                )
                image = getattr(shape, "image", None)
                if image is not None:
                    blob = image.blob
                    image_location = {
                        "slide_number": slide_number, "shape_id": shape_locator_id,
                    }
                    image_ref = (
                        f"{doc['source']['relative_path']}#slide={slide_number};"
                        f"shape={shape.shape_id}"
                    )
                    image_source_name = (
                        getattr(image, "filename", None)
                        or f"image.{getattr(image, 'ext', 'png')}"
                    )
                    visual_origin = self._embedded_visual_origin(
                        blob,
                        doc,
                        location_prefix=image_location,
                        source_name=image_source_name,
                        visual_origin_kind="office_embedded_image",
                    )
                    image_ev = self.add_evidence(
                        doc_id,
                        "image",
                        image_location,
                        content(
                            content_ref=image_ref,
                            mime_type=getattr(image, "content_type", None),
                        ),
                        parent_id=shape_ev["evidence_id"],
                        native_properties={
                            "embedded_sha256": digest_bytes(blob),
                            "size_bytes": len(blob),
                            "file_name": getattr(image, "filename", None),
                            "visual_origin": {
                                **visual_origin,
                            },
                        },
                    )
                    self._project_embedded_image_bytes(
                        blob, doc, parent_id=image_ev["evidence_id"],
                        location_prefix=image_location, content_ref=image_ref,
                        source_name=image_source_name,
                        visual_origin=visual_origin,
                    )
            if getattr(slide, "has_notes_slide", False):
                notes_frame = getattr(slide.notes_slide, "notes_text_frame", None)
                notes_text = getattr(notes_frame, "text", "") if notes_frame is not None else ""
                if notes_text.strip() and self.may_add_leaf(doc_id):
                    self.add_evidence(
                        doc_id,
                        "speaker_note",
                        {"slide_number": slide_number, "locator_text": "speaker-notes"},
                        content(raw_text=notes_text),
                        parent_id=slide_ev["evidence_id"],
                    )
        if isinstance(source, io.BytesIO):
            source.seek(0)
        with zipfile.ZipFile(source) as archive:
            # ``python-pptx`` exposes ordinary picture shapes through
            # ``shape.image`` above, but not slide backgrounds, image-backed
            # shape fills, or inherited layout/master artwork.  Walk the
            # package relationships for those displayed visual sources while
            # excluding only direct-slide pictures already projected above.
            self._project_ooxml_referenced_media(
                archive,
                doc,
                media_prefixes=("ppt/media/",),
                skip_direct_pptx_pictures=True,
            )
            for diagram_index, diagram in enumerate(
                referenced_pptx_diagrams(archive),
                1,
            ):
                diagram_location = {
                    "slide_number": diagram["slide_number"],
                    "source_member": diagram["source_member"],
                    "object_index": diagram_index,
                    "locator_text": (
                        f"slide={diagram['slide_number']};"
                        f"smartart={diagram['source_member']};"
                        f"relationship={diagram['relationship_id']};"
                        f"occurrence={diagram['relationship_occurrence']}"
                    ),
                }
                diagram_ev = self.add_evidence(
                    doc_id,
                    "shape",
                    diagram_location,
                    content(raw_text="SmartArt（ファイル内の明示構造）"),
                    ordinal=diagram_index,
                    native_properties={
                        "ooxml_part": diagram["source_member"],
                        "xml_sha256": diagram["xml_sha256"],
                        "point_count": len(diagram["points"]),
                        "connection_count": len(diagram["connections"]),
                        "ooxml_relationship": {
                            "source_part": diagram["source_part"],
                            "relationship_id": diagram["relationship_id"],
                            "relationship_occurrence": diagram[
                                "relationship_occurrence"
                            ],
                        },
                    },
                )
                self.contain_document(doc_id, diagram_ev["evidence_id"])
                model_ids = [point.get("model_id") for point in diagram["points"]]
                model_id_contract_valid = (
                    all(isinstance(value, str) and value for value in model_ids)
                    and len(model_ids) == len(set(model_ids))
                )
                if not model_id_contract_valid:
                    self.mark_partial(
                        doc,
                        f"SmartArt {diagram['source_member']} has missing or duplicate model ids",
                    )
                point_evidence: dict[str, dict[str, Any]] = {}
                for point_index, point in enumerate(diagram["points"], 1):
                    point_text = str(point.get("text") or "").strip()
                    if not point_text:
                        continue
                    model_id = point.get("model_id")
                    point_ev = self.add_evidence(
                        doc_id,
                        "text_block",
                        {
                            "slide_number": diagram["slide_number"],
                            "source_member": diagram["source_member"],
                            "object_index": point_index,
                            "object_id": (
                                str(model_id) if model_id else f"point-{point_index}"
                            ),
                            "locator_text": (
                                f"{diagram_location['locator_text']};"
                                f"point={urllib.parse.quote(str(model_id or point_index), safe='-._~')}"
                            ),
                        },
                        content(raw_text=point_text),
                        parent_id=diagram_ev["evidence_id"],
                        ordinal=point_index,
                        native_properties={
                            "smartart_model_id": model_id,
                            "smartart_point_type": point.get("type"),
                            "ooxml_part": diagram["source_member"],
                            "xml_sha256": diagram["xml_sha256"],
                        },
                    )
                    if (
                        model_id_contract_valid
                        and isinstance(model_id, str)
                    ):
                        point_evidence[model_id] = point_ev
                connection_groups: dict[
                    tuple[str, str],
                    list[tuple[dict[str, Any], dict[str, Any]]],
                ] = {}
                for connection_index, connection in enumerate(
                    diagram["connections"], 1
                ):
                    source_id = connection.get("srcId")
                    target_id = connection.get("destId")
                    source_ev = point_evidence.get(str(source_id))
                    target_ev = point_evidence.get(str(target_id))
                    if source_ev is None or target_ev is None:
                        self.mark_partial(
                            doc,
                            f"SmartArt {diagram['source_member']} connection {connection_index} has an unresolved endpoint",
                        )
                        continue
                    connection_type = str(connection.get("type") or "unspecified")
                    connection_ev = self.add_evidence(
                        doc_id,
                        "text_block",
                        {
                            "slide_number": diagram["slide_number"],
                            "source_member": diagram["source_member"],
                            "object_index": connection_index,
                            "locator_text": (
                                f"{diagram_location['locator_text']};"
                                f"connection={connection_index}"
                            ),
                        },
                        content(raw_text=(
                            "SmartArtの明示接続: "
                            f"{source_ev['content']['raw_text']} -> "
                            f"{target_ev['content']['raw_text']} "
                            f"(原形式type={connection_type})"
                        )),
                        parent_id=diagram_ev["evidence_id"],
                        ordinal=connection_index,
                        native_properties={
                            "smartart_connection": connection,
                            "ooxml_part": diagram["source_member"],
                            "xml_sha256": diagram["xml_sha256"],
                            "semantic_interpretation_performed": False,
                        },
                    )
                    connection_groups.setdefault(
                        (str(source_id), str(target_id)), []
                    ).append((connection, connection_ev))
                for (source_id, target_id), grouped_connections in sorted(
                    connection_groups.items()
                ):
                    source_ev = point_evidence[source_id]
                    target_ev = point_evidence[target_id]
                    self.add_relation(
                        "structural",
                        "diagram_connection",
                        {
                            "record_type": "evidence",
                            "record_id": source_ev["evidence_id"],
                        },
                        {
                            "record_type": "evidence",
                            "record_id": target_ev["evidence_id"],
                        },
                        properties={
                            "raw_connections": [
                                connection
                                for connection, _ in grouped_connections
                            ],
                            "source_member": diagram["source_member"],
                            "slide_number": diagram["slide_number"],
                            "semantic_interpretation_performed": False,
                        },
                        supporting_evidence_ids=[
                            diagram_ev["evidence_id"],
                            *[
                                connection_ev["evidence_id"]
                                for _, connection_ev in grouped_connections
                            ],
                        ],
                        rule_or_model="native SmartArt srcId/destId connection",
                    )

    def extract_pptx_ooxml(self, path: Path) -> None:
        """Read slide text, tables and package visuals without python-pptx."""
        source, decrypted = self.office_source(path)
        validate_ooxml_archive(
            source,
            required_members=PPTX_REQUIRED_OOXML_MEMBERS,
        )
        doc = self.add_document(path, "ooxml-stdlib-pptx-fallback")
        doc_id = doc["document_id"]
        self.mark_partial(
            doc,
            "PPTX standard-library fallback preserves presentation-order slide text, tables, local geometry, charts, referenced media, and explicit SmartArt; themes, animations, inherited layout text, and group transforms remain unresolved",
        )
        if decrypted:
            doc["extraction"]["warnings"].append(
                "password-protected Office source decrypted in memory"
            )
        if isinstance(source, io.BytesIO):
            source.seek(0)
        with zipfile.ZipFile(source) as archive:
            relationships = _ooxml_relationships(archive)
            _require_ooxml_office_document_binding(
                archive, relationships, "ppt/presentation.xml"
            )
            slide_context = _pptx_slide_context(archive, relationships)
            ordered_slides = sorted(
                slide_context.items(), key=lambda item: item[1]
            )
            for slide_part, slide_numbers in ordered_slides:
                if len(slide_numbers) != 1:
                    raise ValueError("ooxml_slide_binding_ambiguous")
                slide_number = slide_numbers[0]
                root = _ooxml_xml_root(
                    archive,
                    slide_part,
                    namespaces=OOXML_PRESENTATION_NAMESPACES,
                    local_names={"sld"},
                )
                tree = _pptx_shape_tree(root)
                walked: list[
                    tuple[str, ElementTree.Element, bool]
                ] = []

                def walk_shapes(
                    container: ElementTree.Element,
                    prefix: str = "",
                    nested_in_group: bool = False,
                ) -> None:
                    local_index = 0
                    for candidate in container:
                        namespace, local_name = _xml_name(candidate.tag)
                        if (
                            namespace not in OOXML_PRESENTATION_NAMESPACES
                            or local_name not in {
                                "sp", "pic", "graphicFrame", "cxnSp", "grpSp",
                            }
                        ):
                            continue
                        local_index += 1
                        shape_path = (
                            f"{prefix}.{local_index}"
                            if prefix else str(local_index)
                        )
                        walked.append(
                            (shape_path, candidate, nested_in_group)
                        )
                        if local_name == "grpSp":
                            walk_shapes(candidate, shape_path, True)

                if tree is not None:
                    walk_shapes(tree)
                slide_ev = self.add_evidence(
                    doc_id,
                    "slide",
                    {"slide_number": slide_number},
                    content(raw_value={
                        "slide_number": slide_number,
                        "shape_count": len(walked),
                    }),
                    ordinal=slide_number,
                    native_properties={
                        "source_member": slide_part,
                        "shape_count_includes_group_descendants": True,
                        "fallback_detail_status": "direct_slide_content_only",
                    },
                    method="ooxml_stdlib_pptx_fallback",
                )
                self.contain_document(doc_id, slide_ev["evidence_id"])
                seen_shape_ids: set[str] = set()
                group_geometry_warning_added = False
                for shape_index, (
                    shape_path, shape, nested_in_group
                ) in enumerate(walked, 1):
                    if not self.may_add_leaf(doc_id):
                        break
                    _, shape_type = _xml_name(shape.tag)
                    shape_id, shape_name = _pptx_shape_identity(shape)
                    if shape_id:
                        if shape_id in seen_shape_ids:
                            raise ValueError("ooxml_slide_shape_id_ambiguous")
                        seen_shape_ids.add(shape_id)
                        shape_locator_id = (
                            shape_id if "." not in shape_path
                            else f"{shape_id}:{shape_path}"
                        )
                    else:
                        shape_locator_id = f"unresolved:{shape_path}"
                        self.mark_partial(
                            doc,
                            f"slide {slide_number} shape {shape_path} has no canonical non-visual id",
                        )
                    geometry, local_geometry_only = _pptx_shape_geometry(
                        shape, nested_in_group=nested_in_group
                    )
                    if local_geometry_only and not group_geometry_warning_added:
                        self.mark_partial(
                            doc,
                            f"slide {slide_number} contains grouped shapes whose slide-space transform is unresolved",
                        )
                        group_geometry_warning_added = True
                    table = (
                        _pptx_table(shape)
                        if shape_type == "graphicFrame" else None
                    )
                    text_value = (
                        "" if table is not None else _pptx_shape_text(shape)
                    )
                    shape_content = (
                        content(raw_text=text_value)
                        if text_value else content(raw_value={
                            "shape_type": shape_type,
                            "name": shape_name,
                        })
                    )
                    shape_ev = self.add_evidence(
                        doc_id,
                        "shape",
                        {
                            "slide_number": slide_number,
                            "shape_id": shape_locator_id,
                            "object_index": shape_index,
                        },
                        shape_content,
                        parent_id=slide_ev["evidence_id"],
                        ordinal=shape_index,
                        geometry=geometry,
                        native_properties={
                            "name": shape_name,
                            "shape_type": shape_type,
                            "shape_path": shape_path,
                            "source_member": slide_part,
                            "geometry_status": (
                                "group_local_transform_unresolved"
                                if local_geometry_only else
                                "direct_slide_geometry_or_absent"
                            ),
                            "fallback_detail_status": "core_shape_content_only",
                        },
                        method="ooxml_stdlib_pptx_fallback",
                    )
                    if table is None:
                        continue
                    table_rows = _pptx_table_rows(table)
                    table_ev = self.add_evidence(
                        doc_id,
                        "table",
                        {
                            "slide_number": slide_number,
                            "shape_id": shape_locator_id,
                        },
                        content(raw_value={
                            "rows": len(table_rows),
                            "columns": max(
                                (len(row) for row in table_rows), default=0
                            ),
                        }),
                        parent_id=shape_ev["evidence_id"],
                        native_properties={
                            "source_member": slide_part,
                            "fallback_detail_status": "plain_cell_text_only",
                        },
                        method="ooxml_stdlib_pptx_fallback",
                    )
                    for row_index, row in enumerate(table_rows, 1):
                        for column_index, cell_text in enumerate(row, 1):
                            if not self.may_add_leaf(doc_id):
                                break
                            self.add_evidence(
                                doc_id,
                                "table_cell",
                                {
                                    "slide_number": slide_number,
                                    "shape_id": shape_locator_id,
                                    "row_index": row_index,
                                    "column_index": column_index,
                                },
                                content(raw_text=cell_text),
                                parent_id=table_ev["evidence_id"],
                                ordinal=column_index,
                                native_properties={
                                    "fallback_detail_status": "plain_cell_text_only"
                                },
                                method="ooxml_stdlib_pptx_fallback",
                            )
                        if self.limit_reached(doc_id):
                            break

                notes_rows = [
                    row for row in relationships.get(slide_part, [])
                    if _ooxml_relationship_kind(row["type"]) == "notesSlide"
                    and row["target_mode"].casefold() != "external"
                ]
                if len(notes_rows) > 1:
                    raise ValueError("ooxml_notes_slide_binding_ambiguous")
                if notes_rows:
                    notes_member = _resolve_ooxml_target(
                        slide_part, notes_rows[0]["target"]
                    )
                    if notes_member is None:
                        raise ValueError("ooxml_notes_slide_binding_invalid")
                    notes_root = _ooxml_xml_root(
                        archive,
                        notes_member,
                        namespaces=OOXML_PRESENTATION_NAMESPACES,
                        local_names={"notes"},
                    )
                    common = _direct_xml_children(
                        notes_root, OOXML_PRESENTATION_NAMESPACES, "cSld"
                    )
                    if len(common) > 1:
                        raise ValueError("ooxml_notes_slide_content_ambiguous")
                    notes_values: list[str] = []
                    if common:
                        trees = _direct_xml_children(
                            common[0], OOXML_PRESENTATION_NAMESPACES, "spTree"
                        )
                        if len(trees) > 1:
                            raise ValueError("ooxml_notes_slide_content_ambiguous")
                        if trees:
                            for note_shape in _direct_xml_children(
                                trees[0], OOXML_PRESENTATION_NAMESPACES, "sp"
                            ):
                                value = _pptx_shape_text(note_shape)
                                if value.strip():
                                    notes_values.append(value)
                    notes_text = "\n".join(notes_values).strip()
                    if notes_text and self.may_add_leaf(doc_id):
                        self.add_evidence(
                            doc_id,
                            "speaker_note",
                            {
                                "slide_number": slide_number,
                                "source_member": notes_member,
                                "locator_text": "speaker-notes",
                            },
                            content(raw_text=notes_text),
                            parent_id=slide_ev["evidence_id"],
                            native_properties={
                                "fallback_detail_status": "plain_notes_text_only"
                            },
                            method="ooxml_stdlib_pptx_fallback",
                        )

            self._project_ooxml_charts(archive, doc)
            self._project_ooxml_referenced_media(
                archive,
                doc,
                media_prefixes=("ppt/media/",),
                skip_direct_pptx_pictures=False,
            )
            self._project_pptx_diagrams(archive, doc)

    def _project_pptx_diagrams(
        self,
        archive: zipfile.ZipFile,
        doc: dict[str, Any],
    ) -> None:
        """Project literal SmartArt points and raw source connections."""
        doc_id = doc["document_id"]
        for diagram_index, diagram in enumerate(
            referenced_pptx_diagrams(archive),
            1,
        ):
            diagram_location = {
                "slide_number": diagram["slide_number"],
                "source_member": diagram["source_member"],
                "object_index": diagram_index,
                "locator_text": (
                    f"slide={diagram['slide_number']};"
                    f"smartart={diagram['source_member']};"
                    f"relationship={diagram['relationship_id']};"
                    f"occurrence={diagram['relationship_occurrence']}"
                ),
            }
            diagram_ev = self.add_evidence(
                doc_id,
                "shape",
                diagram_location,
                content(raw_text="SmartArt（ファイル内の明示構造）"),
                ordinal=diagram_index,
                native_properties={
                    "ooxml_part": diagram["source_member"],
                    "xml_sha256": diagram["xml_sha256"],
                    "point_count": len(diagram["points"]),
                    "connection_count": len(diagram["connections"]),
                    "ooxml_relationship": {
                        "source_part": diagram["source_part"],
                        "relationship_id": diagram["relationship_id"],
                        "relationship_occurrence": diagram[
                            "relationship_occurrence"
                        ],
                    },
                },
            )
            self.contain_document(doc_id, diagram_ev["evidence_id"])
            model_ids = [point.get("model_id") for point in diagram["points"]]
            model_id_contract_valid = (
                all(isinstance(value, str) and value for value in model_ids)
                and len(model_ids) == len(set(model_ids))
            )
            if not model_id_contract_valid:
                self.mark_partial(
                    doc,
                    f"SmartArt {diagram['source_member']} has missing or duplicate model ids",
                )
            point_evidence: dict[str, dict[str, Any]] = {}
            for point_index, point in enumerate(diagram["points"], 1):
                point_text = str(point.get("text") or "").strip()
                if not point_text:
                    continue
                model_id = point.get("model_id")
                point_ev = self.add_evidence(
                    doc_id,
                    "text_block",
                    {
                        "slide_number": diagram["slide_number"],
                        "source_member": diagram["source_member"],
                        "object_index": point_index,
                        "object_id": (
                            str(model_id) if model_id else f"point-{point_index}"
                        ),
                        "locator_text": (
                            f"{diagram_location['locator_text']};"
                            f"point={urllib.parse.quote(str(model_id or point_index), safe='-._~')}"
                        ),
                    },
                    content(raw_text=point_text),
                    parent_id=diagram_ev["evidence_id"],
                    ordinal=point_index,
                    native_properties={
                        "smartart_model_id": model_id,
                        "smartart_point_type": point.get("type"),
                        "ooxml_part": diagram["source_member"],
                        "xml_sha256": diagram["xml_sha256"],
                    },
                )
                if model_id_contract_valid and isinstance(model_id, str):
                    point_evidence[model_id] = point_ev
            connection_groups: dict[
                tuple[str, str],
                list[tuple[dict[str, Any], dict[str, Any]]],
            ] = {}
            for connection_index, connection in enumerate(
                diagram["connections"], 1
            ):
                source_id = connection.get("srcId")
                target_id = connection.get("destId")
                source_ev = point_evidence.get(str(source_id))
                target_ev = point_evidence.get(str(target_id))
                if source_ev is None or target_ev is None:
                    self.mark_partial(
                        doc,
                        f"SmartArt {diagram['source_member']} connection {connection_index} has an unresolved endpoint",
                    )
                    continue
                connection_type = str(
                    connection.get("type") or "unspecified"
                )
                connection_ev = self.add_evidence(
                    doc_id,
                    "text_block",
                    {
                        "slide_number": diagram["slide_number"],
                        "source_member": diagram["source_member"],
                        "object_index": connection_index,
                        "locator_text": (
                            f"{diagram_location['locator_text']};"
                            f"connection={connection_index}"
                        ),
                    },
                    content(raw_text=(
                        "SmartArtの明示接続: "
                        f"{source_ev['content']['raw_text']} -> "
                        f"{target_ev['content']['raw_text']} "
                        f"(原形式type={connection_type})"
                    )),
                    parent_id=diagram_ev["evidence_id"],
                    ordinal=connection_index,
                    native_properties={
                        "smartart_connection": connection,
                        "ooxml_part": diagram["source_member"],
                        "xml_sha256": diagram["xml_sha256"],
                        "semantic_interpretation_performed": False,
                    },
                )
                connection_groups.setdefault(
                    (str(source_id), str(target_id)), []
                ).append((connection, connection_ev))
            for (source_id, target_id), grouped_connections in sorted(
                connection_groups.items()
            ):
                source_ev = point_evidence[source_id]
                target_ev = point_evidence[target_id]
                self.add_relation(
                    "structural",
                    "diagram_connection",
                    {
                        "record_type": "evidence",
                        "record_id": source_ev["evidence_id"],
                    },
                    {
                        "record_type": "evidence",
                        "record_id": target_ev["evidence_id"],
                    },
                    properties={
                        "raw_connections": [
                            connection for connection, _ in grouped_connections
                        ],
                        "source_member": diagram["source_member"],
                        "slide_number": diagram["slide_number"],
                        "semantic_interpretation_performed": False,
                    },
                    supporting_evidence_ids=[
                        diagram_ev["evidence_id"],
                        *[
                            connection_ev["evidence_id"]
                            for _, connection_ev in grouped_connections
                        ],
                    ],
                    rule_or_model="native SmartArt srcId/destId connection",
                )

    def extract_pdf(self, path: Path) -> None:
        from local_pdf_page_renderer import (
            DEFAULT_DPI,
            MAX_PDF_DOCUMENT_NATIVE_TEXT_CHARS,
            MAX_PDF_DOCUMENT_NATIVE_SECONDS,
            MAX_PDF_DOCUMENT_RENDERED_BYTES,
            MAX_PDF_DOCUMENT_RENDERED_PIXELS,
            MAX_PDF_DOCUMENT_SECONDS,
            inspect_pdf_snapshot,
            read_pdf_snapshot_page,
            render_pdf_snapshot_page,
            snapshot_pdf,
        )

        parser = "pdfkit-jxa+local-page-render+adaptive-local-image-reader"
        doc = self.add_document(path, parser)
        doc_id = doc["document_id"]
        if not self.diagnostic:
            doc["extraction"]["warnings"].append(
                "PDF native text is preserved per page; locally rendered pages are also read as visual sources"
            )
        pages_without_text = 0
        visually_read_pages = 0
        total_rendered_bytes = 0
        total_rendered_pixels = 0
        total_native_text_chars = 0
        native_seconds = 0.0
        visual_seconds = 0.0
        visual_deadline_at = time.monotonic() + MAX_PDF_DOCUMENT_SECONDS
        processed_pages = 0
        visual_budget_exhausted = False
        with snapshot_pdf(path) as snapshot, tempfile.TemporaryDirectory(
            prefix="aiec-pdf-pages-"
        ) as temporary:
            if (
                snapshot.source_sha256 != doc["source"]["sha256"]
                or snapshot.source_size_bytes != doc["source"]["size_bytes"]
            ):
                raise RuntimeError("PDF source changed before its private snapshot")
            inspection = inspect_pdf_snapshot(snapshot)
            page_count = inspection["page_count"]
            pages = inspection["pages"]
            visual_root = Path(temporary)
            for page_number in range(1, page_count + 1):
                if native_seconds >= MAX_PDF_DOCUMENT_NATIVE_SECONDS:
                    self.mark_partial(
                        doc,
                        f"PDF native text stopped before page {page_number} because its independent time limit was reached",
                    )
                    break
                rendered_path = visual_root / f"page-{page_number:06d}.png"
                planned = pages[page_number - 1]
                native_started = time.monotonic()
                try:
                    page_info = read_pdf_snapshot_page(
                        snapshot,
                        page_number,
                        dpi=DEFAULT_DPI,
                    )
                except Exception as exc:
                    pages_without_text += 1
                    page_ev = self.add_evidence(
                        doc_id,
                        "page",
                        {"page_number": page_number},
                        content(
                            content_ref=(
                                f"{doc['source']['relative_path']}#page={page_number}"
                            ),
                            mime_type="application/pdf",
                        ),
                        ordinal=page_number,
                        geometry={
                            "coordinate_space": "page",
                            "unit": "pt",
                            "coordinate_origin": "top_left",
                            "x": 0,
                            "y": 0,
                            "width": planned["page_width_pt"],
                            "height": planned["page_height_pt"],
                        },
                        native_properties={"text_layer_present": False},
                        warning="PDFKit page text read failed; no native or visual text was asserted",
                    )
                    self.contain_document(doc_id, page_ev["evidence_id"])
                    self.mark_partial(
                        doc,
                        f"page {page_number} local text and visual reading unavailable: "
                        f"{type(exc).__name__}: {str(exc)[:300]}",
                    )
                    processed_pages = page_number
                    continue
                finally:
                    native_seconds += time.monotonic() - native_started
                if page_info["source_sha256"] != doc["source"]["sha256"]:
                    raise RuntimeError("PDF page text is not bound to the source document")
                page_width = float(page_info["page_width_pt"])
                page_height = float(page_info["page_height_pt"])
                if (
                    abs(page_width - float(planned["page_width_pt"])) > 1e-6
                    or abs(page_height - float(planned["page_height_pt"])) > 1e-6
                ):
                    raise RuntimeError("PDF page geometry changed after inspection")
                text_value = str(page_info["native_text"])
                if (
                    total_native_text_chars + len(text_value)
                    > MAX_PDF_DOCUMENT_NATIVE_TEXT_CHARS
                ):
                    self.mark_partial(
                        doc,
                        f"PDF native text stopped at page {page_number} because the document character limit was reached",
                    )
                    break
                total_native_text_chars += len(text_value)
                geometry = {
                    "coordinate_space": "page", "unit": "pt",
                    "coordinate_origin": "top_left",
                    "x": 0, "y": 0, "width": page_width, "height": page_height,
                }
                if text_value.strip():
                    item_content = content(raw_text=text_value)
                    warning = None
                else:
                    pages_without_text += 1
                    item_content = content(
                        content_ref=f"{doc['source']['relative_path']}#page={page_number}",
                        mime_type="application/pdf",
                    )
                    warning = "no native text layer; local visual reading requested"
                page_ev = self.add_evidence(
                    doc_id, "page", {"page_number": page_number}, item_content,
                    ordinal=page_number, geometry=geometry,
                    native_properties={"text_layer_present": bool(text_value.strip())},
                    warning=warning,
                )
                self.contain_document(doc_id, page_ev["evidence_id"])
                processed_pages = page_number

                planned_pixels = (
                    int(planned["render_width_px"])
                    * int(planned["render_height_px"])
                )
                if (
                    visual_budget_exhausted
                    or visual_seconds >= MAX_PDF_DOCUMENT_SECONDS
                    or total_rendered_pixels + planned_pixels
                    > MAX_PDF_DOCUMENT_RENDERED_PIXELS
                ):
                    if not visual_budget_exhausted:
                        self.mark_partial(
                            doc,
                            f"PDF visual reading stopped before page {page_number} because its time or pixel limit was reached; native text remains available",
                        )
                    visual_budget_exhausted = True
                    continue
                visual_started = time.monotonic()
                try:
                    materialization = render_pdf_snapshot_page(
                        snapshot,
                        page_number,
                        rendered_path,
                        dpi=DEFAULT_DPI,
                        require_native_text=True,
                    )
                    if not isinstance(materialization, dict):
                        raise DeferredVisualStoreError(
                            "PDF page materialization contract is invalid"
                        )
                    try:
                        materialized_source_sha256 = materialization["source_sha256"]
                        materialized_page_width = float(materialization["page_width_pt"])
                        materialized_page_height = float(materialization["page_height_pt"])
                        materialized_native_text = str(materialization["native_text"])
                        _, materialized_size = self._visual_materialization_contract({
                            "materialization": materialization,
                        })
                    except DeferredVisualStoreError:
                        raise
                    except (KeyError, TypeError, ValueError) as exc:
                        raise DeferredVisualStoreError(
                            "PDF page materialization contract is invalid"
                        ) from exc
                    if materialized_source_sha256 != doc["source"]["sha256"]:
                        raise DeferredVisualStoreError(
                            "rendered PDF page is not bound to the source document"
                        )
                    if (
                        abs(materialized_page_width - page_width) > 1e-6
                        or abs(materialized_page_height - page_height) > 1e-6
                        or materialized_native_text != text_value
                    ):
                        raise DeferredVisualStoreError(
                            "PDF page changed between text read and visual render"
                        )
                    try:
                        from local_image_ocr import read_checked_image_bytes

                        rendered_raw = read_checked_image_bytes(rendered_path)
                    except Exception as exc:
                        raise DeferredVisualStoreError(
                            "rendered PDF page cannot be read safely"
                        ) from exc
                    if (
                        len(rendered_raw) != materialized_size
                        or digest_bytes(rendered_raw)
                        != materialization["rendered_sha256"]
                    ):
                        raise DeferredVisualStoreError(
                            "rendered PDF page differs from its materialization contract"
                        )
                    del rendered_raw
                    total_rendered_bytes += materialized_size
                    if total_rendered_bytes > MAX_PDF_DOCUMENT_RENDERED_BYTES:
                        rendered_path.unlink(missing_ok=True)
                        self.mark_partial(
                            doc,
                            f"PDF visual reading stopped at page {page_number} because the document byte limit was reached",
                        )
                        visual_budget_exhausted = True
                        continue
                    total_rendered_pixels += planned_pixels
                    projected = self._project_local_image_evidence(
                        rendered_path,
                        doc,
                        parent_id=page_ev["evidence_id"],
                        location_prefix={"page_number": page_number},
                        content_ref=(
                            f"{doc['source']['relative_path']}#page={page_number};"
                            f"render=full_page;dpi={DEFAULT_DPI}"
                        ),
                        visual_origin_kind="pdf_page_image",
                        materialization=materialization,
                        native_text=text_value,
                        visual_deadline_at=visual_deadline_at,
                    )
                    if projected:
                        visually_read_pages += 1
                except DeferredVisualStoreError:
                    raise
                except Exception as exc:
                    self.mark_partial(
                        doc,
                        f"page {page_number} local visual reading unavailable: "
                        f"{type(exc).__name__}: {str(exc)[:300]}",
                    )
                finally:
                    visual_seconds += time.monotonic() - visual_started
            if processed_pages < page_count:
                self.mark_partial(
                    doc,
                    f"{page_count - processed_pages} PDF page(s) were not processed because a native-text safety limit was reached",
                )
        if pages_without_text and visually_read_pages < pages_without_text:
            self.mark_partial(
                doc,
                f"{pages_without_text - visually_read_pages} page(s) without native text "
                "did not yield local visual text Evidence",
            )

    @staticmethod
    def _normalized_visible_text(value: object) -> str:
        if not isinstance(value, str):
            return ""
        return re.sub(r"\s+", "", unicodedata.normalize("NFKC", value)).casefold()

    @classmethod
    def _normalized_visible_lines(cls, value: object) -> set[str]:
        if not isinstance(value, str):
            return set()
        return {
            normalized
            for line in value.splitlines()
            if (normalized := cls._normalized_visible_text(line))
        }

    @staticmethod
    def _merge_visual_location(
        source_location: dict[str, Any], child_location: dict[str, Any]
    ) -> dict[str, Any]:
        """Keep both the containing image identity and its child locator."""
        merged = dict(source_location)
        child = dict(child_location)
        if "object_index" in merged and "object_index" in child:
            merged["image_object_index"] = merged.pop("object_index")
        source_locator = merged.pop("locator_text", None)
        child_locator = child.pop("locator_text", None)
        for key, value in child.items():
            if key in merged and merged[key] != value:
                merged[f"visual_{key}"] = value
            else:
                merged[key] = value
        locators = [
            str(value).strip()
            for value in (source_locator, child_locator)
            if isinstance(value, str) and value.strip()
        ]
        if locators:
            merged["locator_text"] = ";".join(locators)
        return merged

    def _project_local_image_evidence(
        self,
        image_path: Path,
        document: dict[str, Any],
        *,
        parent_id: str,
        location_prefix: dict[str, Any],
        content_ref: str,
        visual_origin_kind: str,
        materialization: dict[str, Any],
        native_text: str = "",
        reuse_parent_image: bool = False,
        visual_deadline_at: float | None = None,
    ) -> int:
        """Project the existing audited image reader into one container document.

        The temporary raster never becomes an independent source document. Every
        projected record remains bound to the original PDF/Office source and to
        the materialized-image digest recorded in ``visual_origin``.
        """
        child = Probe(
            image_path.parent,
            self.run_at,
            None,
            diagnostic=False,
            extractor=self.extractor,
            extractor_version=self.extractor_version,
            visual_observation_mode="suppressed",
        )
        child.extract(image_path)
        if len(child.documents) != 1:
            raise DeferredVisualStoreError(
                "local image reader did not produce one source document"
            )
        child_document = child.documents[0]
        expected_sha256, expected_size = self._visual_materialization_contract({
            "materialization": materialization,
        })
        child_source = child_document.get("source")
        child_images = [
            item for item in child.evidence
            if item.get("evidence_type") == "image"
        ]
        if not isinstance(child_source, dict) or len(child_images) != 1:
            raise DeferredVisualStoreError(
                "local image reader output cannot be bound to the materialization"
            )
        child_image_source_sha256 = (
            child_images[0].get("native_properties", {}).get("source_sha256")
            if isinstance(child_images[0].get("native_properties"), dict)
            else None
        )
        if (
            child_source.get("sha256") != expected_sha256
            or child_source.get("size_bytes") != expected_size
            or child_image_source_sha256 != expected_sha256
        ):
            raise DeferredVisualStoreError(
                "local image reader output differs from the materialization contract"
            )
        id_map: dict[str, str] = {}
        projected_text = 0
        native_lines = self._normalized_visible_lines(native_text)
        visual_origin = {
            "kind": visual_origin_kind,
            "source_relative_path": document["source"]["relative_path"],
            "source_sha256": document["source"]["sha256"],
            "source_location": dict(location_prefix),
            "materialization": materialization,
        }
        if visual_origin_kind in {
            "office_embedded_image",
            "notebook_embedded_image",
        }:
            self.mark_partial(
                document,
                "embedded-image OCR is provisional because display selection, "
                "crop, transparency, or transforms are unresolved",
            )
        container_location_locked = reuse_parent_image
        for item in child.evidence:
            evidence_type = item["evidence_type"]
            if evidence_type == "image" and reuse_parent_image:
                id_map[item["evidence_id"]] = parent_id
                continue
            raw_text = item.get("content", {}).get("raw_text")
            if evidence_type == "text_block" and native_lines:
                candidate_key = self._normalized_visible_text(raw_text)
                if candidate_key and candidate_key in native_lines:
                    continue
            location = self._merge_visual_location(
                location_prefix, item.get("location", {})
            )
            if evidence_type == "image":
                if container_location_locked:
                    raise RuntimeError("local image reader produced multiple image containers")
                visual_origin["source_location"] = dict(location)
                container_location_locked = True
            elif not container_location_locked:
                raise RuntimeError("local image reader emitted child Evidence before its image container")
            item_parent = item.get("parent_evidence_id")
            mapped_parent = id_map.get(item_parent, parent_id)
            item_content = (
                content(content_ref=content_ref, mime_type="image/png")
                if evidence_type == "image"
                else dict(item["content"])
            )
            native_properties = dict(item.get("native_properties", {}))
            native_properties["visual_origin"] = visual_origin
            provenance = item.get("provenance", {})
            method = (
                "verified_local_visual_materialization"
                if evidence_type == "image"
                else provenance.get("extraction_method", "unknown")
            )
            display_transform_unresolved = (
                evidence_type == "ocr_line"
                and visual_origin_kind in {
                    "office_embedded_image",
                    "notebook_embedded_image",
                }
            )
            if display_transform_unresolved:
                native_properties["display_transform_resolved"] = False
            downgrade_office_high = (
                display_transform_unresolved
                and native_properties.get("agreement_type")
                == "independent_agreement"
                and native_properties.get("quality_tier") == "high"
            )
            if downgrade_office_high:
                native_properties["embedded_source_agreement_type"] = (
                    native_properties.get("agreement_type")
                )
                native_properties["agreement_type"] = (
                    "display_transform_unresolved"
                )
                native_properties["quality_tier"] = "provisional"
                native_properties["provisional_marker"] = (
                    PROVISIONAL_OCR_MARKER
                )
                method = "adaptive_local_ocr_provisional"
            warnings = list(provenance.get("warnings", []))
            if display_transform_unresolved:
                warnings.append(
                    "container display selection/crop/transparency/transform is unresolved; "
                    "raw embedded-image OCR is provisional only"
                )
            projected = self.add_evidence(
                document["document_id"],
                evidence_type,
                location,
                item_content,
                parent_id=mapped_parent,
                ordinal=item.get("ordinal"),
                style=item.get("style"),
                geometry=item.get("geometry"),
                native_properties=native_properties,
                method=method,
                confidence=(
                    0.0
                    if downgrade_office_high
                    else float(provenance.get("confidence", 0.0))
                ),
                deterministic=bool(provenance.get("deterministic", False)),
                warning=("; ".join(str(value) for value in warnings[:4]) or None),
            )
            id_map[item["evidence_id"]] = projected["evidence_id"]
            if evidence_type in {"ocr_line", "text_block"} and raw_text:
                projected_text += 1
        for visual_task in child._suppressed_visual_tasks:
            child_parent = visual_task.get("parent_id")
            mapped_visual_parent = id_map.get(child_parent)
            if mapped_visual_parent is None:
                raise RuntimeError(
                    "suppressed child visual task has no projected parent Evidence"
                )
            child_location = visual_task.get("location")
            if not isinstance(child_location, dict):
                raise RuntimeError("suppressed child visual task has no location")
            child_ordinal = visual_task.get("ordinal")
            if isinstance(child_ordinal, bool) or not isinstance(child_ordinal, int):
                raise RuntimeError("suppressed child visual task has no ordinal")
            visual_retained_or_queued = self._schedule_local_visual_observation(
                image_path,
                document,
                parent_id=mapped_visual_parent,
                location_prefix={},
                exact_location=self._merge_visual_location(
                    location_prefix, child_location
                ),
                visual_origin=visual_origin,
                ordinal=child_ordinal,
                deadline_at=visual_deadline_at,
                excluded_normalized_texts=native_lines,
            )
            if visual_retained_or_queued:
                projected_text += 1
        child_status = child_document.get("extraction", {}).get("status")
        if child_status != "success":
            for child_warning in child_document.get("extraction", {}).get("warnings", []):
                self.mark_partial(
                    document,
                    f"{visual_origin_kind} {location_prefix}: {str(child_warning)[:300]}",
                )
        return projected_text

    def _add_local_visual_observation(
        self,
        image_path: Path,
        document: dict[str, Any],
        *,
        parent_id: str,
        location_prefix: dict[str, Any],
        visual_origin: dict[str, Any],
        ordinal: int,
        exact_location: dict[str, Any] | None = None,
        image_bytes: bytes | None = None,
        timeout: float | None = None,
        release_paddle: bool = True,
        excluded_normalized_texts: set[str] | None = None,
        deadline_at: float | None = None,
    ) -> bool:
        """Retain local VLM meaning as visibly provisional discovery Evidence."""
        kind = str(visual_origin.get("kind", "visual_image"))
        try:
            if release_paddle:
                # This build is deliberately sequential on memory-constrained
                # Macs. Keep memoized OCR results, but retire Paddle's native
                # process before loading Gemma for visual interpretation.
                Probe._release_active_paddle_worker()
            materialization = visual_origin.get("materialization", {})
            expected_sha256 = (
                materialization.get("rendered_sha256")
                if isinstance(materialization, dict) else None
            )
            observe_kwargs: dict[str, Any] = {
                "expected_input_sha256": expected_sha256,
            }
            if timeout is not None:
                observe_kwargs["timeout"] = timeout
            if image_bytes is None:
                from local_visual_observation import observe_path

                visual = observe_path(image_path, **observe_kwargs)
            else:
                from local_visual_observation import observe_image

                visual = observe_image(image_bytes, **observe_kwargs)
            if deadline_at is not None and time.monotonic() > float(deadline_at):
                self.mark_partial(
                    document,
                    "visual meaning skipped because its document time limit was reached",
                )
                return False
            if (
                excluded_normalized_texts
                and self._normalized_visible_text(visual.get("text"))
                in excluded_normalized_texts
            ):
                return False
            self.add_evidence(
                document["document_id"],
                "text_block",
                (
                    dict(exact_location)
                    if exact_location is not None
                    else self._visual_observation_location(location_prefix)
                ),
                content(raw_text=visual["text"]),
                parent_id=parent_id,
                ordinal=ordinal,
                native_properties={
                    "quality_tier": "provisional",
                    "provisional_marker": PROVISIONAL_OCR_MARKER,
                    "question_independent": True,
                    "observation_type": visual["observation_type"],
                    "structured_observation": visual["observation"],
                    "model": visual["model"],
                    "model_digest": visual["model_digest"],
                    "prompt_sha256": visual["prompt_sha256"],
                    "input_image_sha256": visual["input_image_sha256"],
                    "model_output_sha256": visual["model_output_sha256"],
                    "runner": visual["runner"],
                    "runner_version": visual["runner_version"],
                    "host": visual["host"],
                    "temperature": visual["temperature"],
                    "strict_json": visual["strict_json"],
                    "external_network_used": False,
                    "downloads_performed": False,
                    "visual_origin": visual_origin,
                },
                method="local_vlm_visual_observation_provisional",
                confidence=0.0,
                deterministic=False,
                warning="visual meaning is provisional discovery Evidence only",
            )
            self.mark_partial(
                document,
                f"{kind} provisional visual observation retained for discovery only",
            )
            return True
        except Exception as exc:
            self.mark_partial(
                document,
                f"{kind} visual meaning unavailable: "
                f"{type(exc).__name__}: {str(exc)[:300]}",
            )
            return False

    def _project_ooxml_charts(
        self,
        archive: zipfile.ZipFile,
        document: dict[str, Any],
    ) -> int:
        """Project relationship-bound native chart caches as searchable text."""
        projected = 0
        for chart_index, placement in enumerate(
            referenced_ooxml_charts(archive),
            1,
        ):
            member = str(placement["member"])
            raw = archive.read(member)
            try:
                payload = _ooxml_chart_payload(raw, member)
            except ValueError as exc:
                self.mark_partial(
                    document,
                    f"referenced chart {member} could not be parsed: {exc}",
                )
                continue
            location: dict[str, Any] = {
                "source_member": member,
                "object_index": chart_index,
                "locator_text": str(placement["locator_text"]),
            }
            for key in ("sheet_name", "cell", "slide_number"):
                if key in placement:
                    location[key] = placement[key]
            relationship_binding = {
                "source_part": placement["source_part"],
                "relationship_id": placement["relationship_id"],
                "relationship_occurrence": placement["relationship_occurrence"],
            }
            chart_ev = self.add_evidence(
                document["document_id"],
                "chart",
                location,
                content(raw_text=_chart_summary_text(payload)),
                ordinal=chart_index,
                native_properties={
                    "ooxml_part": member,
                    "xml_sha256": payload["xml_sha256"],
                    "chart_payload": payload,
                    "ooxml_relationship": relationship_binding,
                    "cached_value_status": FORMULA_CACHED_VALUE_STATUS,
                },
                method="verified_ooxml_chart_cache",
            )
            self.contain_document(
                document["document_id"], chart_ev["evidence_id"]
            )
            projected += 1
            for series_index, series in enumerate(payload["series"], 1):
                cache_status = series["cache_status"]
                if (
                    cache_status.startswith("unresolved_")
                    or cache_status == "cached_values_missing"
                ):
                    self.mark_partial(
                        document,
                        f"chart {member} series {series_index} has {cache_status}",
                    )
                self.add_evidence(
                    document["document_id"],
                    "chart_series",
                    {
                        **location,
                        "series_index": series_index,
                        "locator_text": (
                            f"{location['locator_text']};series={series_index}"
                        ),
                    },
                    content(raw_text=_chart_series_text(payload, series)),
                    parent_id=chart_ev["evidence_id"],
                    ordinal=series_index,
                    native_properties={
                        "ooxml_part": member,
                        "xml_sha256": payload["xml_sha256"],
                        "series": series,
                        "ooxml_relationship": relationship_binding,
                        "cached_value_status": FORMULA_CACHED_VALUE_STATUS,
                    },
                    method="verified_ooxml_chart_cache",
                )
        return projected

    def _project_ooxml_referenced_media(
        self,
        archive: zipfile.ZipFile,
        document: dict[str, Any],
        *,
        media_prefixes: tuple[str, ...],
        skip_direct_pptx_pictures: bool = False,
    ) -> int:
        """Project only image placements reachable from the OOXML part graph."""
        projected = 0
        placements = referenced_ooxml_media(
            archive,
            media_prefixes=media_prefixes,
        )
        if any(prefix.startswith("ppt/media/") for prefix in media_prefixes):
            unresolved_inherited_count = (
                _count_unresolved_pptx_inherited_media(archive)
            )
            if unresolved_inherited_count:
                self.mark_partial(
                    document,
                    "PowerPoint layout/master images were not projected because "
                    "their effective slide visibility is unresolved "
                    f"(candidate occurrences={unresolved_inherited_count})",
                )
        if skip_direct_pptx_pictures:
            placements = [
                placement for placement in placements
                if not (
                    placement.get("usage_kind") == "picture"
                    and str(placement.get("source_part", "")).startswith(
                        "ppt/slides/"
                    )
                )
            ]
        for image_index, placement in enumerate(placements, 1):
            member = str(placement["member"])
            raw = archive.read(member)
            image_location: dict[str, Any] = {
                "source_member": member,
                "object_index": image_index,
                "locator_text": str(placement["locator_text"]),
            }
            for key in ("sheet_name", "cell", "slide_number"):
                if key in placement:
                    image_location[key] = placement[key]
            image_ref = (
                f"{document['source']['relative_path']}::{member};"
                f"part={placement['source_part']};"
                f"relationship={placement['relationship_id']};"
                f"occurrence={placement['relationship_occurrence']}"
            )
            relationship_binding = {
                "source_part": placement["source_part"],
                "relationship_id": placement["relationship_id"],
                "relationship_occurrence": placement["relationship_occurrence"],
                "usage_kind": placement.get("usage_kind", "other"),
            }
            if "slide_number" in placement:
                relationship_binding["slide_number"] = placement["slide_number"]
            visual_origin = self._embedded_visual_origin(
                raw,
                document,
                location_prefix=image_location,
                source_name=member,
                visual_origin_kind="office_embedded_image",
            )
            image_ev = self.add_evidence(
                document["document_id"],
                "image",
                image_location,
                content(
                    content_ref=image_ref,
                    mime_type=mimetypes.guess_type(member)[0],
                ),
                ordinal=image_index,
                native_properties={
                    "embedded_sha256": digest_bytes(raw),
                    "size_bytes": len(raw),
                    "ooxml_relationship": relationship_binding,
                    "visual_origin": visual_origin,
                },
            )
            self.contain_document(document["document_id"], image_ev["evidence_id"])
            projected += self._project_embedded_image_bytes(
                raw,
                document,
                parent_id=image_ev["evidence_id"],
                location_prefix=image_location,
                content_ref=image_ref,
                source_name=member,
                visual_origin=visual_origin,
            )
        return projected

    def _embedded_visual_origin(
        self,
        raw: bytes,
        document: dict[str, Any],
        *,
        location_prefix: dict[str, Any],
        source_name: str,
        visual_origin_kind: str,
    ) -> dict[str, Any]:
        image_sha256 = digest_bytes(raw)
        return {
            "kind": visual_origin_kind,
            "source_relative_path": document["source"]["relative_path"],
            "source_sha256": document["source"]["sha256"],
            "source_location": dict(location_prefix),
            "materialization": {
                "runner": "verified_embedded_image_copy",
                "runner_version": self.extractor_version,
                "external_network_used": False,
                "source_sha256": document["source"]["sha256"],
                "embedded_sha256": image_sha256,
                "rendered_sha256": image_sha256,
                "rendered_size_bytes": len(raw),
                "source_name": source_name,
                "display_transform_resolved": visual_origin_kind not in {
                    "office_embedded_image",
                    "notebook_embedded_image",
                },
                "display_transform_status": (
                    "unresolved"
                    if visual_origin_kind in {
                        "office_embedded_image",
                        "notebook_embedded_image",
                    }
                    else "identity"
                ),
            },
        }

    def _project_embedded_image_bytes(
        self,
        raw: bytes,
        document: dict[str, Any],
        *,
        parent_id: str,
        location_prefix: dict[str, Any],
        content_ref: str,
        source_name: str,
        visual_origin_kind: str = "office_embedded_image",
        visual_origin: dict[str, Any] | None = None,
    ) -> int:
        suffix = Path(source_name).suffix.casefold()
        if suffix not in IMAGE_SUFFIXES:
            self.mark_partial(
                document,
                f"embedded visual {source_name} uses an unsupported local OCR format",
            )
            return 0
        document_id = str(document["document_id"])
        usage = self._embedded_visual_usage.setdefault(
            document_id,
            {
                "count": 0,
                "bytes": 0,
                "started_at": time.monotonic(),
                "exhausted": False,
            },
        )
        if usage["exhausted"]:
            return 0
        elapsed = time.monotonic() - float(usage["started_at"])
        limit_reason: str | None = None
        if int(usage["count"]) + 1 > MAX_EMBEDDED_VISUALS_PER_DOCUMENT:
            limit_reason = "image count"
        elif int(usage["bytes"]) + len(raw) > MAX_EMBEDDED_VISUAL_BYTES_PER_DOCUMENT:
            limit_reason = "total image bytes"
        elif elapsed > MAX_EMBEDDED_VISUAL_SECONDS_PER_DOCUMENT:
            limit_reason = "processing time"
        if limit_reason is not None:
            usage["exhausted"] = True
            self.mark_partial(
                document,
                f"embedded visual reading stopped before {source_name} because the document {limit_reason} limit was reached",
            )
            return 0
        usage["count"] = int(usage["count"]) + 1
        usage["bytes"] = int(usage["bytes"]) + len(raw)
        with tempfile.TemporaryDirectory(prefix="aiec-embedded-image-") as temporary:
            image_path = Path(temporary) / f"source{suffix}"
            image_path.write_bytes(raw)
            expected_origin = self._embedded_visual_origin(
                raw,
                document,
                location_prefix=location_prefix,
                source_name=source_name,
                visual_origin_kind=visual_origin_kind,
            )
            if visual_origin is not None and visual_origin != expected_origin:
                raise ValueError("embedded visual origin differs from its source bytes")
            bound_origin = expected_origin if visual_origin is None else visual_origin
            try:
                return self._project_local_image_evidence(
                    image_path,
                    document,
                    parent_id=parent_id,
                    location_prefix=location_prefix,
                    content_ref=content_ref,
                    visual_origin_kind=visual_origin_kind,
                    materialization=bound_origin["materialization"],
                    reuse_parent_image=True,
                    visual_deadline_at=(
                        float(usage["started_at"])
                        + MAX_EMBEDDED_VISUAL_SECONDS_PER_DOCUMENT
                    ),
                )
            except DeferredVisualStoreError:
                raise
            except Exception as exc:
                self.mark_partial(
                    document,
                    f"embedded visual {source_name} local reading unavailable: "
                    f"{type(exc).__name__}: {str(exc)[:300]}",
                )
                return 0

    def extract_image(self, path: Path) -> None:
        """Preserve every located reading and distinguish its support tier."""
        doc = self.add_document(path, "adaptive-local-image-reader-v0.7.0")
        doc_id = doc["document_id"]
        try:
            from local_image_ocr import (
                MAX_UNLOCATED_TRANSCRIPT_TOKENS,
                UNLOCATED_TRANSCRIPT_MODEL,
                UNLOCATED_TRANSCRIPT_PROMPT_SHA256,
                extract,
            )
            observation = extract(path)
        except Exception as exc:
            self.mark_partial(doc, f"local image OCR unavailable: {type(exc).__name__}: {exc}")
            visual_origin = {
                "kind": "standalone_image",
                "source_relative_path": doc["source"]["relative_path"],
                "source_sha256": doc["source"]["sha256"],
                "source_location": {"object_index": 1},
                "materialization": {
                    "runner": "verified_source_image_bytes",
                    "runner_version": self.extractor_version,
                    "external_network_used": False,
                    "source_sha256": doc["source"]["sha256"],
                    "rendered_sha256": doc["source"]["sha256"],
                    "rendered_size_bytes": doc["source"]["size_bytes"],
                },
            }
            image_ev = self.add_evidence(
                doc_id,
                "image",
                {"object_index": 1},
                content(content_ref=doc["source"]["relative_path"], mime_type=doc["source"]["media_type"]),
                ordinal=1,
                native_properties={
                    "source_sha256": doc["source"]["sha256"],
                    "visual_origin": visual_origin,
                },
                method="verified_image_bytes",
            )
            self.contain_document(doc_id, image_ev["evidence_id"])
            self._schedule_local_visual_observation(
                path, doc, parent_id=image_ev["evidence_id"],
                location_prefix={"object_index": 1},
                visual_origin=visual_origin, ordinal=2,
            )
            return

        width = observation["dimensions"]["width_px"]
        height = observation["dimensions"]["height_px"]
        visual_origin = {
            "kind": "standalone_image",
            "source_relative_path": doc["source"]["relative_path"],
            "source_sha256": doc["source"]["sha256"],
            "source_location": {"object_index": 1},
            "materialization": {
                "runner": "verified_source_image_bytes",
                "runner_version": self.extractor_version,
                "external_network_used": False,
                "source_sha256": doc["source"]["sha256"],
                "rendered_sha256": observation["input_sha256"],
                "rendered_size_bytes": path.stat().st_size,
            },
        }
        image_ev = self.add_evidence(
            doc_id,
            "image",
            {"object_index": 1},
            content(content_ref=doc["source"]["relative_path"], mime_type=doc["source"]["media_type"]),
            ordinal=1,
            geometry={
                "coordinate_space": "image",
                "coordinate_origin": "top_left",
                "unit": "px",
                "x": 0,
                "y": 0,
                "width": width,
                "height": height,
            },
            native_properties={
                "source_sha256": observation["input_sha256"],
                "image_format": observation["image_format"],
                "orientation": observation["orientation"],
                "orientation_source": observation["orientation_source"],
                "source_dimensions": observation["source_dimensions"],
                "ocr_input_sha256": observation["ocr_input_sha256"],
                "ocr_input_dimensions": observation["ocr_input_dimensions"],
                "ocr_input_orientation": observation["ocr_input_orientation"],
                "coordinate_frame_policy": observation["coordinate_frame_policy"],
                "canonicalization": observation["canonicalization"],
                "ocr_engines": observation["engines"],
                "independent_ocr_engines": observation["independent_engines"],
                "unresolved_ocr_line_count": observation["unresolved_count"],
                "visual_origin": visual_origin,
            },
            method="verified_image_bytes",
        )
        self.contain_document(doc_id, image_ev["evidence_id"])
        read_lines = observation.get("read_lines", observation["consensus_lines"])
        for line_index, line in enumerate(read_lines, 1):
            bbox = line["bbox"]
            confidence_values = [
                value for value in (line["primary_confidence"], line["audit_confidence"])
                if value is not None
            ]
            agreement_type = line.get("agreement_type", "independent_agreement")
            expected_quality_tier = OCR_QUALITY_BY_AGREEMENT.get(agreement_type)
            if expected_quality_tier is None:
                raise ValueError(f"unsupported OCR agreement type: {agreement_type!r}")
            quality_tier = line.get("quality_tier", expected_quality_tier)
            if quality_tier != expected_quality_tier:
                raise ValueError(
                    "OCR quality tier disagrees with engine independence: "
                    f"{agreement_type!r} cannot be {quality_tier!r}"
                )
            upstream_marker = line.get("provisional_marker")
            if quality_tier == "provisional":
                if upstream_marker != PROVISIONAL_OCR_MARKER:
                    raise ValueError("provisional OCR marker is missing or not canonical")
            elif upstream_marker is not None:
                raise ValueError("high OCR evidence must not carry a provisional marker")
            provenance = line.get("provenance")
            if not isinstance(provenance, dict):
                raise ValueError("OCR observation provenance is missing")
            primary_pass = provenance.get("primary_pass")
            primary_engine = ocr_engine(primary_pass)
            primary_group = ocr_independence_group(primary_pass)
            if provenance.get("primary_engine") != primary_engine:
                raise ValueError("OCR primary engine is invalid")
            if provenance.get("primary_independence_group") != primary_group:
                raise ValueError("OCR primary independence group is invalid")
            audit_pass = provenance.get("audit_pass")
            audit_engine = None
            audit_group = None
            if audit_pass is not None:
                audit_engine = ocr_engine(audit_pass)
                audit_group = ocr_independence_group(audit_pass)
                if provenance.get("audit_engine") != audit_engine:
                    raise ValueError("OCR audit engine is invalid")
                if provenance.get("audit_independence_group") != audit_group:
                    raise ValueError("OCR audit independence group is invalid")
            if agreement_type == "independent_agreement" and (
                audit_group is None or primary_group == audit_group
            ):
                raise ValueError(
                    "high OCR evidence requires distinct line-level engine groups"
                )
            if agreement_type == "same_engine_agreement" and (
                audit_group is None or primary_group != audit_group
            ):
                raise ValueError(
                    "same-engine OCR evidence requires one line-level engine group"
                )
            if agreement_type == "provisional_single_pass" and audit_group is not None:
                raise ValueError("single-pass OCR evidence must not claim an audit pass")
            validate_ocr_supporters(
                line,
                provenance,
                primary_pass=primary_pass,
                primary_engine=primary_engine,
                primary_group=primary_group,
                audit_pass=audit_pass,
                audit_engine=audit_engine,
                audit_group=audit_group,
                agreement_type=agreement_type,
            )
            if quality_tier == "high" and observation["independent_engines"] is not True:
                raise ValueError("high OCR evidence requires independent engine groups")
            bbox_coordinate_system = line.get("bbox_coordinate_system")
            if bbox_coordinate_system not in OCR_BBOX_COORDINATE_SYSTEMS:
                raise ValueError("OCR bbox coordinate system is missing or unsupported")
            if (
                quality_tier == "high"
                and bbox_coordinate_system
                != "source_orientation_1_top_left_normalized_1000"
            ):
                raise ValueError(
                    "high OCR evidence requires a shared source-orientation-1 frame"
                )
            self.add_evidence(
                doc_id,
                "ocr_line",
                {"object_index": line_index},
                content(raw_text=line["text"]),
                parent_id=image_ev["evidence_id"],
                ordinal=line_index,
                geometry={
                    "coordinate_space": "image",
                    "coordinate_origin": "top_left",
                    "unit": "normalized_1000",
                    "x": bbox[0],
                    "y": bbox[1],
                    "width": bbox[2],
                    "height": bbox[3],
                },
                native_properties={
                    "consensus_method": "spatial-nfc-exact-or-provisional",
                    "agreement_type": agreement_type,
                    "quality_tier": quality_tier,
                    "bbox_coordinate_system": bbox_coordinate_system,
                    **(
                        {"provisional_marker": PROVISIONAL_OCR_MARKER}
                        if quality_tier == "provisional" else {}
                    ),
                    "spatial_overlap": line["overlap"],
                    "primary_confidence": line["primary_confidence"],
                    "audit_confidence": line["audit_confidence"],
                    "independent_engines": observation["independent_engines"],
                    "observation_provenance": provenance,
                    "visual_origin": visual_origin,
                },
                method=(
                    "dual_local_ocr_consensus"
                    if quality_tier == "high"
                    else "adaptive_local_ocr_provisional"
                ),
                confidence=min(confidence_values) if confidence_values else 0.0,
                deterministic=False,
            )
        unlocated = observation.get("unlocated_transcript")
        if unlocated is not None:
            if not isinstance(unlocated, dict):
                raise ValueError("unlocated transcript must be an object")
            transcript = unlocated.get("text")
            model_digest = unlocated.get("model_digest")
            prompt_sha256 = unlocated.get("prompt_sha256")
            normalized_model_digest = (
                model_digest.removeprefix("sha256:")
                if isinstance(model_digest, str) else ""
            )
            if not isinstance(transcript, str) or not transcript.strip():
                raise ValueError("unlocated transcript text is missing")
            if (
                unlocated.get("location_status") != "unlocated"
                or unlocated.get("quality_tier") != "provisional"
                or unlocated.get("provisional_marker") != PROVISIONAL_OCR_MARKER
                or unlocated.get("question_independent") is not True
                or unlocated.get("transcript_type")
                != "whole_image_faithful_transcript"
                or unlocated.get("model") != UNLOCATED_TRANSCRIPT_MODEL
                or unlocated.get("runner") != "ollama_loopback_chat"
                or unlocated.get("host") != "127.0.0.1"
                or not isinstance(model_digest, str)
                or not re.fullmatch(r"[0-9a-f]{64}", normalized_model_digest)
                or prompt_sha256 != UNLOCATED_TRANSCRIPT_PROMPT_SHA256
                or unlocated.get("temperature") != 0
                or unlocated.get("num_predict")
                != MAX_UNLOCATED_TRANSCRIPT_TOKENS
            ):
                raise ValueError("unlocated transcript provenance is invalid")
            transcript_value = transcript.strip()
            marker_prefix = f"{PROVISIONAL_OCR_MARKER}\n"
            transcript_chunks = exact_text_chunks(
                transcript_value,
                max_chars=MAX_QUESTION_EVIDENCE_CHARS - len(marker_prefix),
            )
            transcript_sha256 = digest_bytes(transcript_value.encode("utf-8"))
            base_ordinal = len(read_lines)
            for chunk_index, question_chunk in enumerate(transcript_chunks, 1):
                self.add_evidence(
                    doc_id,
                    "text_block",
                    {
                        "object_index": base_ordinal + chunk_index,
                        "locator_text": (
                            "location_status=unlocated;source=image;"
                            f"chunk={chunk_index}/{len(transcript_chunks)};"
                            f"characters={question_chunk.start + 1}-{question_chunk.end}"
                        ),
                    },
                    content(raw_text=marker_prefix + question_chunk.text),
                    parent_id=image_ev["evidence_id"],
                    ordinal=base_ordinal + chunk_index,
                    native_properties={
                        "location_status": "unlocated",
                        "quality_tier": "provisional",
                        "provisional_marker": PROVISIONAL_OCR_MARKER,
                        "transcript_type": unlocated.get("transcript_type"),
                        "question_independent": True,
                        "model": unlocated.get("model"),
                        "model_digest": model_digest,
                        "prompt_sha256": prompt_sha256,
                        "runner": unlocated.get("runner"),
                        "host": "127.0.0.1",
                        "temperature": unlocated.get("temperature"),
                        "num_predict": unlocated.get("num_predict"),
                        "transcript_sha256": transcript_sha256,
                        "transcript_chunk_index": chunk_index,
                        "transcript_chunk_count": len(transcript_chunks),
                        "character_start": question_chunk.start,
                        "character_end": question_chunk.end,
                        "character_offset_basis": "zero_based_half_open",
                        "visual_origin": visual_origin,
                    },
                    method="local_vlm_unlocated_transcript_provisional",
                    confidence=0.0,
                    deterministic=False,
                    warning=(
                        "whole-image transcript has no coordinates and is provisional only"
                    ),
                )
            self.mark_partial(
                doc,
                "unlocated whole-image transcript retained as provisional Evidence",
            )
        if not read_lines:
            self.mark_partial(doc, "adaptive local OCR produced no located text observations")
        if observation["unresolved_count"]:
            self.mark_partial(
                doc,
                f"{observation['unresolved_count']} OCR reading(s) remain provisional but are retained",
            )
        self._schedule_local_visual_observation(
            path, doc, parent_id=image_ev["evidence_id"],
            location_prefix={"object_index": 1},
            visual_origin=visual_origin,
            ordinal=len(read_lines) + 1,
        )

    def extract_delimited(self, path: Path) -> None:
        text_value, encoding = read_text(path)
        delimiter = "\t" if path.suffix.lower() == ".tsv" else ","
        reader = csv.reader(io.StringIO(text_value, newline=""), delimiter=delimiter)
        doc = self.add_document(path, "python-csv")
        doc_id = doc["document_id"]
        try:
            raw_headers = next(reader)
        except StopIteration:
            self.mark_partial(doc, "delimited file is empty")
            return
        headers = [value.strip() or f"column_{index}" for index, value in enumerate(raw_headers, 1)]
        table_ev = self.add_evidence(
            doc_id,
            "table",
            {"table_index": 1, "locator_text": "delimited-table"},
            content(raw_value={"delimiter": delimiter, "headers": raw_headers}),
            ordinal=1,
            native_properties={"encoding": encoding, "column_count": len(headers)},
        )
        self.contain_document(doc_id, table_ev["evidence_id"])
        for source_row_number, row in enumerate(reader, 2):
            if not any(value != "" for value in row):
                continue
            if not self.may_add_leaf(doc_id):
                break
            values = list(row)
            if len(values) < len(headers):
                values.extend([""] * (len(headers) - len(values)))
            rendered = " / ".join(
                f"{headers[index] if index < len(headers) else f'column_{index + 1}'}:{value}"
                for index, value in enumerate(values)
            )
            self.add_evidence(
                doc_id,
                "table_row",
                {"table_index": 1, "row_index": source_row_number, "locator_text": f"row={source_row_number}"},
                content(raw_text=rendered),
                parent_id=table_ev["evidence_id"],
                ordinal=source_row_number,
                native_properties={"headers": headers, "values": values, "source_row_number": source_row_number},
            )
        if encoding == "utf-8-replacement":
            self.mark_partial(doc, "input required replacement characters during decoding")

    def extract_json(self, path: Path) -> None:
        text_value, encoding = read_text(path)
        parsed = json.loads(text_value)
        doc = self.add_document(path, "python-json")
        doc_id = doc["document_id"]

        def walk(value: Any, pointer: str, parent_id: str | None, ordinal: int) -> None:
            if isinstance(value, (dict, list)):
                container = self.add_evidence(
                    doc_id,
                    "metadata",
                    {"locator_text": pointer or "/"},
                    content(raw_value={
                        "container_type": "object" if isinstance(value, dict) else "array",
                        "item_count": len(value),
                    }),
                    parent_id=parent_id,
                    ordinal=ordinal,
                    native_properties={"json_pointer": pointer or "/"},
                )
                if parent_id is None:
                    self.contain_document(doc_id, container["evidence_id"])
                items = value.items() if isinstance(value, dict) else enumerate(value)
                for child_ordinal, (key, child) in enumerate(items, 1):
                    token = str(key).replace("~", "~0").replace("/", "~1")
                    walk(child, f"{pointer}/{token}", container["evidence_id"], child_ordinal)
                return
            if not self.may_add_leaf(doc_id):
                return
            item_content = content(raw_text=value) if isinstance(value, str) else content(raw_value=value)
            ev = self.add_evidence(
                doc_id,
                "field",
                {"locator_text": pointer or "/"},
                item_content,
                parent_id=parent_id,
                ordinal=ordinal,
                native_properties={"json_pointer": pointer or "/", "encoding": encoding},
            )
            if parent_id is None:
                self.contain_document(doc_id, ev["evidence_id"])

        walk(parsed, "", None, 1)

    def extract_xml(self, path: Path) -> None:
        text_value, encoding = read_text(path)
        root = ElementTree.fromstring(text_value)
        doc = self.add_document(path, "xml.etree.ElementTree")
        doc_id = doc["document_id"]

        def walk(element: ElementTree.Element, parent_path: str, parent_id: str | None, ordinal: int) -> None:
            element_path = f"{parent_path}/{element.tag}[{ordinal}]"
            container = self.add_evidence(
                doc_id,
                "metadata",
                {"locator_text": element_path},
                content(raw_value={"tag": element.tag, "attribute_count": len(element.attrib), "child_count": len(element)}),
                parent_id=parent_id,
                ordinal=ordinal,
                native_properties={"xml_path": element_path, "encoding": encoding},
            )
            if parent_id is None:
                self.contain_document(doc_id, container["evidence_id"])
            for attr_index, (name, value) in enumerate(sorted(element.attrib.items()), 1):
                if self.may_add_leaf(doc_id):
                    self.add_evidence(
                        doc_id,
                        "field",
                        {"locator_text": f"{element_path}/@{name}"},
                        content(raw_text=value),
                        parent_id=container["evidence_id"],
                        ordinal=attr_index,
                        native_properties={"xml_path": f"{element_path}/@{name}"},
                    )
            if element.text and element.text.strip() and self.may_add_leaf(doc_id):
                self.add_evidence(
                    doc_id,
                    "field",
                    {"locator_text": f"{element_path}/text()"},
                    content(raw_text=element.text),
                    parent_id=container["evidence_id"],
                    native_properties={"xml_path": f"{element_path}/text()"},
                )
            for child_index, child in enumerate(list(element), 1):
                walk(child, element_path, container["evidence_id"], child_index)

        walk(root, "", None, 1)

    def extract_plain_text(self, path: Path) -> None:
        text_value, encoding = read_text(path)
        parser = "plain-code" if path.suffix.lower() in CODE_SUFFIXES else "plain-text"
        doc = self.add_document(path, parser)
        doc_id = doc["document_id"]
        if not text_value:
            self.mark_partial(doc, "text file is empty")
            return
        if path.suffix.lower() in CODE_SUFFIXES:
            lines = text_value.splitlines(keepends=True)
            for block_index, start in enumerate(range(0, len(lines), 80), 1):
                if not self.may_add_leaf(doc_id):
                    break
                block = "".join(lines[start:start + 80])
                ev = self.add_evidence(
                    doc_id,
                    "code_block",
                    {"code_line_start": start + 1, "code_line_end": min(start + 80, len(lines)),
                     "locator_text": f"lines={start + 1}-{min(start + 80, len(lines))}"},
                    content(raw_text=block),
                    ordinal=block_index,
                    native_properties={"encoding": encoding, "language": path.suffix.lower().lstrip(".")},
                )
                self.contain_document(doc_id, ev["evidence_id"])
            return

        blocks = re.split(r"\n[\t ]*\n+", text_value)
        paragraph_index = 0
        for block in blocks:
            if not block:
                continue
            if not self.may_add_leaf(doc_id):
                break
            paragraph_index += 1
            stripped = block.strip()
            evidence_type = "heading" if path.suffix.lower() == ".md" and re.match(r"^#{1,6}\s+", stripped) else "paragraph"
            ev = self.add_evidence(
                doc_id,
                evidence_type,
                {"paragraph_index": paragraph_index},
                content(raw_text=block),
                ordinal=paragraph_index,
                native_properties={"encoding": encoding},
            )
            self.contain_document(doc_id, ev["evidence_id"])
        if encoding == "utf-8-replacement":
            self.mark_partial(doc, "input required replacement characters during decoding")

    def extract_notebook(self, path: Path) -> None:
        text_value, encoding = read_text(path)
        notebook = json.loads(text_value)
        doc = self.add_document(path, "nbformat-json")
        doc_id = doc["document_id"]
        image_count = 0

        suffix_by_mime = {
            "image/png": ".png",
            "image/jpeg": ".jpg",
            "image/tiff": ".tiff",
            "image/bmp": ".bmp",
        }

        def add_notebook_image(
            payload: Any,
            *,
            cell_index: int,
            mime_type: str,
            locator_kind: str,
            attachment_name: str | None = None,
        ) -> str:
            nonlocal image_count
            image_count += 1
            encoded = payload if isinstance(payload, str) else (
                "".join(payload) if isinstance(payload, list) else ""
            )
            blob: bytes | None = None
            try:
                compact = re.sub(r"\s", "", encoded)
                blob = base64.b64decode(compact, validate=True)
                if not blob:
                    raise ValueError("empty notebook image")
                blob_sha = digest_bytes(blob)
                size_bytes = len(blob)
            except Exception:
                blob_sha = digest_bytes(encoded.encode("ascii", errors="ignore"))
                size_bytes = 0
                self.mark_partial(
                    doc,
                    f"notebook cell {cell_index} contains malformed {mime_type} image data",
                )
            locator = f"cell={cell_index};{locator_kind}={image_count}"
            if attachment_name is not None:
                locator += ";attachment=" + urllib.parse.quote(
                    attachment_name,
                    safe="-._~",
                )
            image_location = {
                "notebook_cell_index": cell_index,
                "object_index": image_count,
                "locator_text": locator,
            }
            image_ref = f"{doc['source']['relative_path']}#{locator}"
            suffix = suffix_by_mime.get(mime_type.casefold())
            image_source_name = f"notebook-image-{image_count}{suffix or '.bin'}"
            visual_origin = (
                self._embedded_visual_origin(
                    blob,
                    doc,
                    location_prefix=image_location,
                    source_name=image_source_name,
                    visual_origin_kind="notebook_embedded_image",
                )
                if blob is not None else {"kind": "notebook_embedded_image"}
            )
            native_properties: dict[str, Any] = {
                "embedded_sha256": blob_sha,
                "size_bytes": size_bytes,
                "visual_origin": visual_origin,
            }
            if attachment_name is not None:
                native_properties["attachment_name"] = attachment_name
            image_ev = self.add_evidence(
                doc_id,
                "image",
                image_location,
                content(content_ref=image_ref, mime_type=mime_type),
                ordinal=image_count,
                native_properties=native_properties,
            )
            self.contain_document(doc_id, image_ev["evidence_id"])
            if blob is not None and suffix is not None:
                self._project_embedded_image_bytes(
                    blob,
                    doc,
                    parent_id=image_ev["evidence_id"],
                    location_prefix=image_location,
                    content_ref=image_ref,
                    source_name=image_source_name,
                    visual_origin_kind="notebook_embedded_image",
                    visual_origin=visual_origin,
                )
            elif blob is not None:
                self.mark_partial(
                    doc,
                    f"notebook cell {cell_index} image type {mime_type} is not supported by the local reader",
                )
            return f"[embedded image sha256={blob_sha}]"

        def strip_data_uri(source_text: str, cell_index: int) -> str:
            def replace(match: re.Match[str]) -> str:
                return add_notebook_image(
                    match.group(2),
                    cell_index=cell_index,
                    mime_type=match.group(1),
                    locator_kind="source-image",
                )

            return DATA_URI_PATTERN.sub(replace, source_text)

        for cell_index, cell in enumerate(notebook.get("cells", []), 1):
            source_value = cell.get("source", [])
            source_text = (
                "".join(source_value)
                if isinstance(source_value, list)
                else str(source_value)
            )
            cell_type = cell.get("cell_type")
            attachment_references = (
                {
                    unicodedata.normalize("NFC", urllib.parse.unquote(match.group(1)))
                    for match in NOTEBOOK_ATTACHMENT_PATTERN.finditer(source_text)
                }
                if cell_type == "markdown" else set()
            )
            attachments = cell.get("attachments", {})
            if isinstance(attachments, dict) and attachment_references:
                normalized_attachments: dict[str, list[tuple[str, Any]]] = {}
                for raw_name, payloads in attachments.items():
                    if not isinstance(raw_name, str):
                        continue
                    normalized_attachments.setdefault(
                        unicodedata.normalize("NFC", raw_name),
                        [],
                    ).append((raw_name, payloads))
                for attachment_name in sorted(attachment_references):
                    matches = normalized_attachments.get(attachment_name, [])
                    if len(matches) != 1:
                        self.mark_partial(
                            doc,
                            f"notebook cell {cell_index} attachment {attachment_name!r} is missing or Unicode-ambiguous",
                        )
                        continue
                    raw_name, payloads = matches[0]
                    if not isinstance(payloads, dict):
                        self.mark_partial(
                            doc,
                            f"notebook cell {cell_index} attachment {attachment_name!r} has invalid payloads",
                        )
                        continue
                    for mime_type in sorted(payloads):
                        if isinstance(mime_type, str) and mime_type.startswith("image/"):
                            add_notebook_image(
                                payloads[mime_type],
                                cell_index=cell_index,
                                mime_type=mime_type,
                                locator_kind="attachment-image",
                                attachment_name=raw_name,
                            )
            if cell_type == "markdown":
                source_text = strip_data_uri(source_text, cell_index)
            if source_text and self.may_add_leaf(doc_id):
                ev = self.add_evidence(
                    doc_id,
                    "notebook_cell",
                    {"notebook_cell_index": cell_index, "locator_text": f"cell={cell_index}"},
                    content(raw_text=source_text),
                    ordinal=cell_index,
                    native_properties={"cell_type": cell.get("cell_type"), "encoding": encoding},
                )
                self.contain_document(doc_id, ev["evidence_id"])
            output_index = 0
            for output in cell.get("outputs", []):
                output_index += 1
                data = output.get("data", {})
                text_output = output.get("text", "")
                if isinstance(text_output, list):
                    text_output = "".join(text_output)
                if not text_output:
                    plain = data.get("text/plain", "")
                    text_output = "".join(plain) if isinstance(plain, list) else plain
                if text_output and self.may_add_leaf(doc_id):
                    ev = self.add_evidence(
                        doc_id,
                        "text_block",
                        {"notebook_cell_index": cell_index, "object_index": output_index,
                         "locator_text": f"cell={cell_index};output={output_index}"},
                        content(raw_text=str(text_output)),
                        ordinal=output_index,
                        native_properties={"output_type": output.get("output_type")},
                    )
                    self.contain_document(doc_id, ev["evidence_id"])
                for mime_type, payload in data.items():
                    if not mime_type.startswith("image/"):
                        continue
                    add_notebook_image(
                        payload,
                        cell_index=cell_index,
                        mime_type=mime_type,
                        locator_kind=f"output={output_index};output-image",
                    )
        if image_count:
            doc["extraction"]["warnings"].append(
                f"{image_count} embedded image(s) routed through local visual reading"
            )

    def extract_other(self, path: Path) -> None:
        doc = self.add_document(path, "metadata-only")
        doc["extraction"]["status"] = "deferred"
        warning = "unsupported format; content extraction deferred"
        if warning not in doc["extraction"]["warnings"]:
            doc["extraction"]["warnings"].append(warning)
        ev = self.add_evidence(
            doc["document_id"], "metadata", {"locator_text": "file"},
            content(content_ref=doc["source"]["relative_path"], mime_type=doc["source"]["media_type"]),
            warning="unsupported format; content extraction deferred",
        )
        self.contain_document(doc["document_id"], ev["evidence_id"])

    def write(self, output: Path) -> None:
        if not self.retain_records:
            raise RuntimeError("write() is unavailable when retain_records is false")
        output.mkdir(parents=True, exist_ok=True)
        for name, records in (
            ("documents.jsonl", self.documents),
            ("evidence.jsonl", self.evidence),
            ("relations.jsonl", self.relations),
        ):
            with (output / name).open("w", encoding="utf-8", newline="\n") as handle:
                for record in sorted(records, key=lambda item: next(
                    item[key] for key in ("document_id", "evidence_id", "relation_id") if key in item
                )):
                    handle.write(canonical_json(record) + "\n")


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--root", required=True, type=Path, help="root used for normalized relative paths")
    parser.add_argument("--input", required=True, type=Path, nargs="+", help="files to probe")
    parser.add_argument("--out", required=True, type=Path, help="isolated diagnostic output directory")
    parser.add_argument("--max-items", type=int, default=40, help="leaf sample limit per document")
    parser.add_argument(
        "--run-at", default=datetime.now(timezone.utc).isoformat(),
        help="ISO-8601 extraction time; pass the same value to compare byte-for-byte reruns",
    )
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    if args.max_items < 1:
        raise SystemExit("--max-items must be at least 1")
    probe = Probe(
        args.root,
        args.run_at,
        args.max_items,
        password_candidates=discover_password_candidates(args.root.resolve()),
    )
    for path in args.input:
        if not path.is_file():
            raise SystemExit(f"input is not a file: {path}")
        probe.extract(path)
    probe.write(args.out)
    print(canonical_json({
        "documents": len(probe.documents),
        "evidence": len(probe.evidence),
        "relations": len(probe.relations),
        "output": str(args.out.resolve()),
    }))


if __name__ == "__main__":
    main()
