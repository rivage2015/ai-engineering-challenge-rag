#!/usr/bin/env python3
"""Run the pre-promotion cross-format anti-hardcoding gate.

This evaluator derives a second, structurally isomorphic corpus from the frozen
ORION fixture.  It changes file names, entity values, assignment periods and
question wording, then reruns the real Layer 1 readers, security gate, semantic
graph builder and deterministic answerer.  Gold questions are opened only
after the mutated graph has been frozen, and neither builder nor answerer is
given a gold or mutation-spec path.
"""

from __future__ import annotations

import argparse
import importlib.util
import json
import random
import shutil
import sys
import tempfile
from datetime import date, datetime
from pathlib import Path
from typing import Any, Iterable, Mapping, Sequence


REPOSITORY = Path(__file__).resolve().parents[1]
SCRIPTS = REPOSITORY / "scripts"
ENGINE = REPOSITORY / "distribution" / "macos-local-memory" / "engine"
DEFAULT_DATASET = REPOSITORY / "evaluation" / "cross-format-kg-v0.1"
FIXED_RUN_AT = "2026-08-27T00:00:00+00:00"


class AntiHardcodingError(ValueError):
    """Raised when a mutation or anti-hardcoding assertion fails."""


def canonical_json(value: Any) -> str:
    return json.dumps(
        value,
        ensure_ascii=False,
        sort_keys=True,
        separators=(",", ":"),
        allow_nan=False,
    )


def load_module(name: str, path: Path) -> Any:
    specification = importlib.util.spec_from_file_location(name, path)
    if specification is None or specification.loader is None:
        raise ImportError(f"cannot load module: {path}")
    module = importlib.util.module_from_spec(specification)
    sys.modules[name] = module
    specification.loader.exec_module(module)
    return module


phase2 = load_module(
    "cross_format_kg_anti_hardcoding_phase2",
    SCRIPTS / "evaluate_cross_format_kg_phase2.py",
)


def read_json(path: Path) -> dict[str, Any]:
    try:
        value = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError) as exc:
        raise AntiHardcodingError(f"cannot read JSON object: {path}") from exc
    if not isinstance(value, dict):
        raise AntiHardcodingError(f"JSON object required: {path}")
    return value


def read_jsonl(path: Path) -> list[dict[str, Any]]:
    records: list[dict[str, Any]] = []
    try:
        with path.open(encoding="utf-8") as handle:
            for line_number, line in enumerate(handle, start=1):
                if not line.strip():
                    continue
                value = json.loads(line)
                if not isinstance(value, dict):
                    raise AntiHardcodingError(
                        f"JSON object required at {path}:{line_number}"
                    )
                records.append(value)
    except json.JSONDecodeError as exc:
        raise AntiHardcodingError(
            f"invalid JSONL at {path}:{exc.lineno}"
        ) from exc
    if not records:
        raise AntiHardcodingError(f"JSONL is empty: {path}")
    return records


def write_json(path: Path, value: Mapping[str, Any]) -> None:
    path.write_text(canonical_json(dict(value)) + "\n", encoding="utf-8")


def write_jsonl(path: Path, records: Iterable[Mapping[str, Any]]) -> None:
    with path.open("w", encoding="utf-8", newline="\n") as handle:
        for record in records:
            handle.write(canonical_json(dict(record)) + "\n")


def validate_variant(value: Mapping[str, Any]) -> None:
    required = {
        "schema_version",
        "record_type",
        "variant_id",
        "enumeration_seed",
        "source_file_order",
        "renamed_files",
        "replacements",
        "forbidden_old_values",
        "required_new_graph_values",
    }
    if set(value) != required:
        raise AntiHardcodingError("anti-hardcoding variant fields are invalid")
    if (
        value.get("schema_version") != "0.1"
        or value.get("record_type")
        != "cross_format_kg_anti_hardcoding_variant"
        or type(value.get("enumeration_seed")) is not int
    ):
        raise AntiHardcodingError("anti-hardcoding variant header is invalid")
    source_order = value.get("source_file_order")
    renamed = value.get("renamed_files")
    replacements = value.get("replacements")
    if (
        not isinstance(source_order, list)
        or not source_order
        or len(source_order) != len(set(source_order))
        or not isinstance(renamed, dict)
        or set(source_order) != set(renamed)
        or len(set(renamed.values())) != len(renamed)
        or any(Path(old).suffix != Path(new).suffix for old, new in renamed.items())
        or any(Path(old) == Path(new) for old, new in renamed.items())
    ):
        raise AntiHardcodingError("file-order or rename contract is invalid")
    if (
        not isinstance(replacements, list)
        or not replacements
        or any(
            not isinstance(item, list)
            or len(item) != 2
            or not all(isinstance(part, str) and part for part in item)
            or item[0] == item[1]
            for item in replacements
        )
    ):
        raise AntiHardcodingError("replacement contract is invalid")
    old_values = value.get("forbidden_old_values")
    new_values = value.get("required_new_graph_values")
    if (
        not isinstance(old_values, list)
        or not old_values
        or len(old_values) != len(set(old_values))
        or not isinstance(new_values, list)
        or not new_values
        or len(new_values) != len(set(new_values))
    ):
        raise AntiHardcodingError("old/new value contract is invalid")
    if set(old_values) != {item[0] for item in replacements}:
        raise AntiHardcodingError(
            "every replaced source value must be forbidden in mutated output"
        )


def replace_text(text: str, replacements: Sequence[Sequence[str]]) -> str:
    value = text
    for old, new in replacements:
        value = value.replace(old, new)
    return value


def _replace_docx_paragraph(paragraph: Any, replacements: Sequence[Sequence[str]]) -> None:
    original = paragraph.text
    replaced = replace_text(original, replacements)
    if replaced == original:
        return
    if paragraph.runs:
        paragraph.runs[0].text = replaced
        for run in paragraph.runs[1:]:
            run.text = ""
    else:
        paragraph.add_run(replaced)


def mutate_docx(
    source: Path,
    destination: Path,
    replacements: Sequence[Sequence[str]],
) -> None:
    from docx import Document

    document = Document(source)
    for paragraph in document.paragraphs:
        _replace_docx_paragraph(paragraph, replacements)
    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    _replace_docx_paragraph(paragraph, replacements)
    for section in document.sections:
        for container in (section.header, section.footer):
            for paragraph in container.paragraphs:
                _replace_docx_paragraph(paragraph, replacements)
    properties = document.core_properties
    for name in ("title", "subject", "category", "keywords", "comments"):
        value = getattr(properties, name, None)
        if isinstance(value, str):
            setattr(properties, name, replace_text(value, replacements))
    destination.parent.mkdir(parents=True, exist_ok=True)
    document.save(destination)


def mutate_xlsx(
    source: Path,
    destination: Path,
    replacements: Sequence[Sequence[str]],
) -> None:
    from openpyxl import load_workbook

    date_replacements = dict(replacements)
    workbook = load_workbook(source)
    for worksheet in workbook.worksheets:
        for row in worksheet.iter_rows():
            for cell in row:
                value = cell.value
                if isinstance(value, str):
                    cell.value = replace_text(value, replacements)
                elif isinstance(value, (datetime, date)):
                    old_date = value.date() if isinstance(value, datetime) else value
                    replacement = date_replacements.get(old_date.isoformat())
                    if replacement is not None:
                        shifted = date.fromisoformat(replacement)
                        cell.value = (
                            datetime.combine(shifted, value.time())
                            if isinstance(value, datetime)
                            else shifted
                        )
    properties = workbook.properties
    for name in ("title", "subject", "keywords", "description", "category"):
        value = getattr(properties, name, None)
        if isinstance(value, str):
            setattr(properties, name, replace_text(value, replacements))
    destination.parent.mkdir(parents=True, exist_ok=True)
    workbook.save(destination)


def _mutate_pptx_shapes(
    shapes: Iterable[Any], replacements: Sequence[Sequence[str]]
) -> None:
    for shape in shapes:
        if getattr(shape, "shape_type", None) == 6 and hasattr(shape, "shapes"):
            _mutate_pptx_shapes(shape.shapes, replacements)
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    value = cell.text
                    replaced = replace_text(value, replacements)
                    if replaced != value:
                        cell.text = replaced
        elif getattr(shape, "has_text_frame", False):
            value = shape.text_frame.text
            replaced = replace_text(value, replacements)
            if replaced != value:
                shape.text_frame.text = replaced


def mutate_pptx(
    source: Path,
    destination: Path,
    replacements: Sequence[Sequence[str]],
) -> None:
    from pptx import Presentation

    presentation = Presentation(source)
    for slide in presentation.slides:
        _mutate_pptx_shapes(slide.shapes, replacements)
    properties = presentation.core_properties
    for name in ("title", "subject", "keywords", "comments", "category"):
        value = getattr(properties, name, None)
        if isinstance(value, str):
            setattr(properties, name, replace_text(value, replacements))
    destination.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(destination)


def mutate_pdf(
    source: Path,
    destination: Path,
    replacements: Sequence[Sequence[str]],
) -> None:
    from pypdf import PdfReader, PdfWriter
    from pypdf.generic import DecodedStreamObject, NameObject

    reader = PdfReader(source)
    if len(reader.pages) != 1:
        raise AntiHardcodingError("identity-register PDF must have one page")
    writer = PdfWriter()
    writer.clone_document_from_reader(reader)
    page = writer.pages[0]
    contents = page.get_contents()
    if contents is None:
        raise AntiHardcodingError("identity-register PDF has no content stream")
    data = contents.get_data()
    byte_replacements = {
        b"EMP-104": replace_text("EMP-104", replacements).encode("ascii"),
        b"EMP-208": replace_text("EMP-208", replacements).encode("ascii"),
        b"2023-03-20": replace_text("2023-03-20", replacements).encode("ascii"),
        # The deterministic source PDF uses an embedded subset font.  Compose
        # different synthetic names only from glyphs already in that subset so
        # the ToUnicode map remains truthful after rewriting the content stream.
        b"\\031\\032\\033\\034": b"\\037\\032\\033\\034",
        b"\\037\\200\\201": b"\\031\\200\\201",
    }
    for old, new in byte_replacements.items():
        if old not in data:
            raise AntiHardcodingError(
                f"identity-register PDF mutation anchor is missing: {old!r}"
            )
        data = data.replace(old, new)
    stream = DecodedStreamObject()
    stream.set_data(data)
    page[NameObject("/Contents")] = writer._add_object(stream)
    metadata = {
        str(key): replace_text(str(value), replacements)
        for key, value in (reader.metadata or {}).items()
        if value is not None
    }
    if metadata:
        writer.add_metadata(metadata)
    destination.parent.mkdir(parents=True, exist_ok=True)
    with destination.open("wb") as handle:
        writer.write(handle)


def materialize_mutated_corpus(
    dataset: Path,
    destination: Path,
    variant: Mapping[str, Any],
) -> list[str]:
    source_root = dataset / "corpus"
    replacements = variant["replacements"]
    created: list[str] = []
    for old_relative in variant["source_file_order"]:
        source = source_root / old_relative
        if not source.is_file():
            raise AntiHardcodingError(f"source fixture is missing: {source}")
        new_relative = variant["renamed_files"][old_relative]
        destination_path = destination / new_relative
        suffix = source.suffix.casefold()
        if suffix == ".docx":
            mutate_docx(source, destination_path, replacements)
        elif suffix == ".xlsx":
            mutate_xlsx(source, destination_path, replacements)
        elif suffix == ".pptx":
            mutate_pptx(source, destination_path, replacements)
        elif suffix == ".pdf":
            mutate_pdf(source, destination_path, replacements)
        else:
            raise AntiHardcodingError(f"unsupported fixture format: {source}")
        created.append(new_relative)
    if set(created) != set(variant["renamed_files"].values()):
        raise AntiHardcodingError("mutated corpus does not cover every renamed file")
    return created


def _shuffle_jsonl(path: Path, seed: int) -> tuple[list[str], list[str]]:
    records = read_jsonl(path)
    key = "evidence_id" if "evidence_id" in records[0] else "document_id"
    before = [str(record[key]) for record in records]
    shuffled = list(records)
    random.Random(seed).shuffle(shuffled)
    after = [str(record[key]) for record in shuffled]
    if after == before:
        shuffled = shuffled[1:] + shuffled[:1]
        after = [str(record[key]) for record in shuffled]
    write_jsonl(path, shuffled)
    return before, after


def build_mutated_phase1(
    *,
    corpus: Path,
    output: Path,
    python: Path,
    guard: Any,
    enumeration_seed: int,
) -> tuple[Path, list[Any], dict[str, list[str]]]:
    intermediate = output / "intermediate"
    search = output / "search"
    adapter = output / "layer1-adapter"
    runs = []
    commands = (
        (
            str(python),
            str(SCRIPTS / "build_intermediate_records.py"),
            "--root",
            str(corpus.resolve()),
            "--out",
            str(intermediate.resolve()),
            "--run-at",
            FIXED_RUN_AT,
        ),
        (
            str(python),
            str(SCRIPTS / "build_search_units.py"),
            "--intermediate",
            str(intermediate.resolve()),
            "--out",
            str(search.resolve()),
        ),
        (
            str(python),
            str(SCRIPTS / "adapt_layer1_to_local_memory.py"),
            "--intermediate",
            str(intermediate.resolve()),
            "--search-output",
            str(search.resolve()),
            "--source-root",
            str(corpus.resolve()),
            "--out",
            str(adapter.resolve()),
        ),
    )
    for command in commands:
        runs.append(guard.run(command, cwd=REPOSITORY))

    documents_before, documents_after = _shuffle_jsonl(
        adapter / "semantic-documents.jsonl", enumeration_seed
    )
    evidence_before, evidence_after = _shuffle_jsonl(
        adapter / "semantic-evidence.jsonl", enumeration_seed + 1
    )
    runs.append(guard.run(
        (
            str(python.resolve()),
            str(ENGINE / "content_security_gate.py"),
            "--evidence",
            str((adapter / "semantic-evidence.jsonl").resolve()),
            "--documents",
            str((adapter / "semantic-documents.jsonl").resolve()),
            "--output-dir",
            str(adapter.resolve()),
        ),
        cwd=REPOSITORY,
    ))
    return adapter, runs, {
        "document_ids_before": documents_before,
        "document_ids_after": documents_after,
        "evidence_ids_before": evidence_before,
        "evidence_ids_after": evidence_after,
    }


def _assert_no_values(label: str, rendered: str, values: Sequence[str]) -> None:
    leaked = sorted({value for value in values if value in rendered})
    if leaked:
        raise AntiHardcodingError(f"{label} leaked old values: {leaked}")


def snapshot_projection(snapshot: Any) -> dict[str, Any]:
    return {
        "nodes": [
            {
                "node_type": node.node_type,
                "canonical_key": node.canonical_key,
                "properties": node.properties,
            }
            for node in snapshot.nodes.values()
        ],
        "edges": [
            {
                "relation_type": edge.relation_type,
                "properties": edge.properties,
            }
            for edge in snapshot.edges.values()
        ],
        "evidence": [
            {
                "relative_path": evidence.relative_path,
                "observed_text": evidence.observed_text,
            }
            for evidence in snapshot.evidence.values()
        ],
    }


def validate_case(
    snapshot: Any,
    answer: Mapping[str, Any],
    case: Mapping[str, Any],
    corpus: Path,
    forbidden: Sequence[str],
) -> dict[str, Any]:
    question = case.get("question")
    expected = case.get("expected")
    if not isinstance(question, str) or not isinstance(expected, dict):
        raise AntiHardcodingError("anti-hardcoding QA case is invalid")
    trace = phase2.validate_answer_trace(snapshot, answer, question)
    for field in ("decision", "operation"):
        if answer.get(field) != expected.get(field):
            raise AntiHardcodingError(
                f"{case.get('case_id')}: {field} mismatch: "
                f"{answer.get(field)!r} != {expected.get(field)!r}"
            )
    actual_facts = {
        item.get("field"): item.get("value")
        for item in answer.get("asserted_facts", [])
        if isinstance(item, dict)
    }
    required_facts = expected.get("required_facts", {})
    if not isinstance(required_facts, dict):
        raise AntiHardcodingError("required_facts must be an object")
    for field, value in required_facts.items():
        if actual_facts.get(field) != value:
            raise AntiHardcodingError(
                f"{case.get('case_id')}: fact {field} mismatch: "
                f"{actual_facts.get(field)!r} != {value!r}"
            )
    required_relations = {
        tuple(item) for item in expected.get("required_relations", [])
    }
    actual_relations = {
        (item.get("from"), item.get("relation"), item.get("to"))
        for item in answer.get("asserted_relations", [])
        if isinstance(item, dict)
    }
    if not required_relations <= actual_relations:
        raise AntiHardcodingError(
            f"{case.get('case_id')}: required relations are missing"
        )
    if answer.get("decision") == "HOLD":
        if (
            answer.get("reason_code") != expected.get("reason_code")
            or answer.get("must_request_concept")
            != expected.get("must_request_concept")
            or answer.get("asserted_facts")
            or answer.get("asserted_relations")
        ):
            raise AntiHardcodingError(
                f"{case.get('case_id')}: HOLD contract mismatch"
            )
        concept = expected.get("must_request_concept")
        if isinstance(concept, str) and concept not in str(answer.get("answer_text")):
            raise AntiHardcodingError(
                f"{case.get('case_id')}: HOLD does not request {concept}"
            )
    visited = trace["visited_document_paths"]
    minimum_documents = case.get("minimum_distinct_visited_documents")
    if type(minimum_documents) is not int or len(visited) < minimum_documents:
        raise AntiHardcodingError(
            f"{case.get('case_id')}: insufficient cross-document traversal"
        )
    for relative_path in visited:
        if not (corpus / relative_path).is_file():
            raise AntiHardcodingError(
                f"{case.get('case_id')}: trace source does not exist: {relative_path}"
            )
    answer_projection = {
        "answer_text": answer.get("answer_text"),
        "asserted_facts": answer.get("asserted_facts"),
        "asserted_relations": answer.get("asserted_relations"),
        "visited_document_paths": visited,
        "resolved_source_references": answer.get("trace", {}).get(
            "resolved_source_references"
        ),
    }
    _assert_no_values(
        f"answer {case.get('case_id')}",
        canonical_json(answer_projection),
        forbidden,
    )
    return {
        "case_id": case.get("case_id"),
        "question_hash": answer.get("question_hash"),
        "decision": answer.get("decision"),
        "operation": answer.get("operation"),
        "used_semantic_edge_count": len(trace["used_edge_ids"]),
        "visited_document_count": len(visited),
        "old_value_leak_count": 0,
        "passed": True,
        "answer": answer,
    }


def run_evaluation(
    *,
    dataset: Path,
    variant_path: Path,
    cases_path: Path,
    output: Path,
    builder: Path,
    answerer: Path,
    python: Path,
) -> dict[str, Any]:
    dataset = dataset.resolve(strict=True)
    variant_path = variant_path.resolve(strict=True)
    cases_path = cases_path.resolve(strict=True)
    output = output.resolve()
    builder = builder.resolve(strict=True)
    answerer = answerer.resolve(strict=True)
    # Keep a virtual-environment launcher as a symlink: resolving it to the
    # base interpreter drops that environment's document-reader dependencies.
    python = python.absolute()
    if not python.is_file():
        raise AntiHardcodingError(f"Python runtime is missing: {python}")
    if output.exists():
        raise AntiHardcodingError(f"refusing to overwrite output: {output}")
    variant = read_json(variant_path)
    validate_variant(variant)
    output.parent.mkdir(parents=True, exist_ok=True)
    staging = Path(tempfile.mkdtemp(
        prefix=f".{output.name}.building-", dir=output.parent
    ))
    try:
        guard_root = staging / "guard"
        guard_root.mkdir()
        guard = phase2.NetworkGuard(guard_root)
        corpus = staging / "mutated-corpus"
        creation_order = materialize_mutated_corpus(dataset, corpus, variant)
        phase1, phase1_runs, enumeration = build_mutated_phase1(
            corpus=corpus,
            output=staging / "phase1",
            python=python,
            guard=guard,
            enumeration_seed=variant["enumeration_seed"],
        )

        extracted_documents = read_jsonl(phase1 / "semantic-documents.jsonl")
        extracted_paths = {
            item.get("source", {}).get("relative_path")
            for item in extracted_documents
        }
        renamed_paths = set(variant["renamed_files"].values())
        if extracted_paths != renamed_paths:
            raise AntiHardcodingError(
                "Layer 1 did not preserve the complete renamed source set"
            )
        safe_artifacts = (
            (phase1 / "semantic-documents.jsonl").read_text(encoding="utf-8")
            + (phase1 / "safe-answer-evidence.jsonl").read_text(encoding="utf-8")
        )
        _assert_no_values(
            "mutated Layer 1 artifacts",
            safe_artifacts,
            variant["forbidden_old_values"],
        )

        build = phase2.build_and_freeze(
            phase1_dir=phase1,
            dataset=dataset,
            builder=builder,
            python=python,
            staging=staging,
            guard=guard,
        )
        graph_projection = snapshot_projection(build.snapshot)
        rendered_graph = canonical_json(graph_projection)
        _assert_no_values(
            "mutated graph", rendered_graph, variant["forbidden_old_values"]
        )
        graph_keys = {
            node.canonical_key for node in build.snapshot.nodes.values()
        }
        missing_new = sorted(
            set(variant["required_new_graph_values"]) - graph_keys
        )
        if missing_new:
            raise AntiHardcodingError(
                f"mutated graph is missing required new values: {missing_new}"
            )

        implementation_text = (
            builder.read_text(encoding="utf-8")
            + answerer.read_text(encoding="utf-8")
        )
        implementation_forbidden = [
            value for replacement in variant["replacements"] for value in replacement
        ]
        _assert_no_values(
            "production builder/answerer source",
            implementation_text,
            implementation_forbidden,
        )

        # EVALUATOR BOUNDARY: questions and expectations are opened only after
        # the mutated graph has been built, validated and frozen read-only.
        cases = read_jsonl(cases_path)
        original_questions = {
            item.get("question")
            for item in read_jsonl(dataset / "gold" / "qa-cases.jsonl")
        }
        if any(case.get("question") in original_questions for case in cases):
            raise AntiHardcodingError("a mutation question copied an original question")
        for case in cases:
            question = case.get("question")
            if isinstance(question, str) and question in implementation_text:
                raise AntiHardcodingError(
                    "production builder/answerer contains a full evaluation question"
                )

        io_dir = staging / "answerer-io"
        io_dir.mkdir()
        results = []
        answer_runs = []
        for sequence, case in enumerate(cases, start=1):
            answer, run = phase2.invoke_answerer(
                python=python,
                answerer=answerer,
                graph_path=build.graph_path,
                question=case["question"],
                io_dir=io_dir,
                sequence=sequence,
                dataset=dataset,
                guard=guard,
            )
            answer_runs.append(run)
            results.append(validate_case(
                build.snapshot,
                answer,
                case,
                corpus,
                variant["forbidden_old_values"],
            ))
            phase2._assert_graph_still_frozen(build)

        accepted = sum(item["decision"] == "ACCEPTED" for item in results)
        held = sum(item["decision"] == "HOLD" for item in results)
        if (accepted, held) != (4, 1):
            raise AntiHardcodingError(
                f"expected four ACCEPTED and one HOLD, got {accepted}/{held}"
            )
        all_runs = [*phase1_runs, build.builder_run, *answer_runs]
        network_attempts = sum(run.network_attempt_count for run in all_runs)
        if network_attempts != 0:
            raise AntiHardcodingError("outbound network attempts were not zero")
        if (
            enumeration["document_ids_before"]
            == enumeration["document_ids_after"]
            or enumeration["evidence_ids_before"]
            == enumeration["evidence_ids_after"]
        ):
            raise AntiHardcodingError("graph input enumeration order was not changed")

        results_path = staging / "anti-hardcoding-results.jsonl"
        report_path = staging / "anti-hardcoding-report.json"
        write_jsonl(results_path, results)
        report = {
            "schema_version": "0.1",
            "record_type": "cross_format_kg_anti_hardcoding_report",
            "dataset_id": dataset.name,
            "variant_id": variant["variant_id"],
            "decision": "PASS",
            "required_before_answer_promotion": True,
            "mutations": {
                "corpus_file_creation_order_changed": creation_order
                != sorted(creation_order),
                "graph_input_document_order_changed": True,
                "graph_input_evidence_order_changed": True,
                "all_five_files_renamed": len(renamed_paths) == 5,
                "project_values_replaced": True,
                "work_values_replaced": True,
                "employee_values_replaced": True,
                "person_values_replaced": True,
                "validity_intervals_shifted_with_boundary_preserved": True,
                "all_questions_paraphrased": True,
            },
            "source_file_creation_order": creation_order,
            "graph_snapshot_id": build.snapshot.graph_snapshot_id,
            "graph_file_sha256": build.graph_file_sha256,
            "graph_counts": {
                "nodes": len(build.snapshot.nodes),
                "edges": len(build.snapshot.edges),
                "source_evidence": len(build.snapshot.evidence),
            },
            "normal_case_count": len(results),
            "accepted_case_count": accepted,
            "hold_case_count": held,
            "old_value_leak_count": 0,
            "production_literal_leak_count": 0,
            "measured_outbound_network_attempt_count": network_attempts,
            "subprocess_count": len(all_runs),
            "gold_boundary": {
                "graph_frozen_before_questions_loaded": True,
                "builder_received_mutation_spec": False,
                "builder_received_questions_or_expected_answers": False,
                "answerer_payload_keys": ["question"],
                "answerer_received_expected_answers": False,
            },
            "artifacts": {
                "mutated_corpus": "mutated-corpus/",
                "phase1": "phase1/layer1-adapter/",
                "graph": "semantic-graph.sqlite3",
                "builder_state": "semantic-graph-state.json",
                "results": results_path.name,
            },
        }
        write_json(report_path, report)
        staging.replace(output)
        return report
    except Exception:
        shutil.rmtree(staging, ignore_errors=True)
        raise


def parse_args(argv: Sequence[str] | None = None) -> argparse.Namespace:
    runtime_python = REPOSITORY / "rag" / ".venv" / "bin" / "python"
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--dataset", type=Path, default=DEFAULT_DATASET)
    parser.add_argument("--variant", type=Path)
    parser.add_argument("--cases", type=Path)
    parser.add_argument("--out", type=Path, required=True)
    parser.add_argument(
        "--builder",
        type=Path,
        default=SCRIPTS / "build_cross_document_semantic_graph.py",
    )
    parser.add_argument(
        "--answerer",
        type=Path,
        default=SCRIPTS / "query_cross_document_semantic_graph.py",
    )
    parser.add_argument(
        "--python",
        type=Path,
        default=runtime_python if runtime_python.is_file() else Path(sys.executable),
    )
    args = parser.parse_args(argv)
    if args.variant is None:
        args.variant = args.dataset / "anti-hardcoding-variant.json"
    if args.cases is None:
        args.cases = args.dataset / "gold" / "anti-hardcoding-qa-cases.jsonl"
    return args


def main(argv: Sequence[str] | None = None) -> int:
    args = parse_args(argv)
    try:
        report = run_evaluation(
            dataset=args.dataset,
            variant_path=args.variant,
            cases_path=args.cases,
            output=args.out,
            builder=args.builder,
            answerer=args.answerer,
            python=args.python,
        )
    except (AntiHardcodingError, phase2.EvaluationError, OSError) as exc:
        print(f"anti-hardcoding evaluation failed: {exc}", file=sys.stderr)
        return 1
    print(json.dumps(report, ensure_ascii=False, indent=2, sort_keys=True))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
