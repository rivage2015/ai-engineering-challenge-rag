"""Additional fail-closed deterministic rules for score-candidate rollout.

Every rule is a full-question grammar with opaque entity/field/value slots.
Rules read only named source artifacts and never read validation answers,
predictions, or question IDs.  Any non-unique source, table, or interpretation
returns ``None`` so the existing RAG answer remains untouched.
"""

from __future__ import annotations

import ast
import colorsys
import csv
import hashlib
import io
import json
import math
import re
import statistics
import unicodedata
import xml.etree.ElementTree as ET
import zipfile
from collections import Counter
from datetime import date, datetime
from decimal import Decimal, ROUND_HALF_UP
from fractions import Fraction
from pathlib import Path, PurePosixPath
from typing import Any, Iterable, Mapping, Sequence

from structured_candidate import (
    StructuredCandidateAnswer,
    StructuredCandidateDecision,
    _candidate_values,
    _cell_text,
    _legal_core,
    _location_matches,
    _normalized,
    _read_source_tables,
)


DATE_RANGE = re.compile(
    r"^(?P<location>.+?)の(?P<container>[^、]+?\.(?:xlsx|csv|tsv))において、"
    r"(?P<start>\d{4}-\d{2}-\d{2})から(?P<end>\d{4}-\d{2}-\d{2})の間に"
    r"(?P<left>開始日)または(?P<right>終了日)が設定されている"
    r"(?P<target>タスクID)をすべて挙げてください。?$"
)
ASSIGNEE_COUNT = re.compile(
    r"^(?P<location>.+?)の(?P<container>.+?)において、"
    r"(?P<person>.+?)さんが担当者に含まれるタスクIDはいくつありますか。?$"
)
PROJECT_PERSON_ASSIGNMENT_ROLE = re.compile(
    r"^(?P<location>.+?)案件において、(?P<person>.+?)さんは"
    r"どの役割としてアサインされていますか。?$"
)
PHASE_LATEST = re.compile(
    r"^(?P<location>.+?)の(?P<container>[^、]+?\.xlsx)において、"
    r"フェーズNo(?P<phase>\d+)にて最後に開始するタスク名は何ですか。?$"
)
BUFFER_SUM = re.compile(
    r"^(?P<location>.+?)の(?P<container>[^、]+?\.xlsx)において、"
    r"バッファとして使用した工数の合計は何時間ですか。?$"
)
CHECKPOINT_TASKS = re.compile(
    r"^(?P<location>.+?)のチェックポイント(?P<number>\d+)として設定されている"
    r"内容に関連するタスクIDを教えてください。?$"
)
MISSING_ROWS_MAX = re.compile(
    r"^分析データの中で、1つでも欠損値がある行数が最も多い案件を、"
    r"主略称で答えてください。?$"
)
STANDARDIZED_SHARE = re.compile(
    r"^(?P<location>.+?)の分析対象データにおいて、標準化された"
    r"(?P<measure>[A-Za-z_][A-Za-z0-9_]*)が(?P<z>-?\d+(?:\.\d+)?)未満の行のうち、"
    r"(?P<category>[A-Za-z_][A-Za-z0-9_]*)=(?P<value>[A-Za-z0-9_\-]+)に該当し、かつ"
    r"(?P=measure)が(?P=category)=(?P=value)全体の平均を上回る行の割合は何%ですか。"
    r"小数第(?P<digits>\d+)位まで答えてください。?$"
)
INTERACTION_COLUMNS = re.compile(
    r"^(?P<location>.+?)の分析出力 metrics\.json の "
    r"feature_selection\.selected_columns に含まれている列のうち、"
    r"分析コードで生成された数値交互作用特徴量の列名をすべて答えてください。?$"
)
GB_PARAMS = re.compile(
    r"^(?P<location>.+?)の分析コードにおいて、今回の学習で勾配ブースティング法の"
    r"モデルに実際に渡される n_estimators、learning_rate、random_state はそれぞれ"
    r"いくつですか。設定ファイルに明示されていない値がある場合も、実行時にコード上で"
    r"適用される値を含めて答えてください。?$"
)
XLSX_CHART_SERIES_COLUMN = re.compile(
    r"^(?P<location>.+?)の(?P<container>[^,、。]+?\.xlsx)の"
    r"(?P<sheet>[^,、。]+?)にあるグラフ"
    r"(?P<chart>[0-9０-９]+)はどのカラムを可視化したものですか。?$"
)
REGRESSION_PREDICTION = re.compile(
    r"^(?P<location>.+?)の(?P<container>[^、]+?\.xlsx)にて算出された回帰係数を使って"
    r"(?P<id_field>id)=(?P<id_value>-?\d+)を予測した場合の予測値はいくらになりますか。"
    r"小数第(?P<digits>\d+)位まで求めてください。?$"
)
NEGATIVE_CORRELATION = re.compile(
    r"^(?P<location>.+?)の顧客データにおいて、目的変数と最も強い負の相関を持つ"
    r"カラムは何ですか。?$"
)
HIGHLIGHT_ROWS = re.compile(
    r"^(?P<location>.+?)の(?P<container>[^、]+?\.xlsx)において、"
    r"(?:(?P<sheet>[^、]+?)シートで)?オレンジ(?:色)?にハイライトされている行の"
    r"(?P<target>タスク名|タスクID)をすべて(?:答えて|教えて)ください。?$"
)
BLUE_SUM = re.compile(
    r"^(?P<location>.+?)の(?P<container>[^、]+?\.xlsx)において、"
    r"青色ハイライト部分の合計値を求めてください。四捨五入して整数で答えてください。?$"
)
YELLOW_INTERSECTION = re.compile(
    r"^(?P<location>.+?)の(?P<container>[^、]+?\.xlsx)において、"
    r"黄色ハイライトが交差している2つのセルの値の差の絶対値を計算してください。?$"
)
COHORT_GROUP_MEAN_ARGMAX = re.compile(
    r"^(?P<location>.+?)のプロジェクトデータ（(?P<container>[^（）]+?\.(?:csv|xlsx))）において、"
    r"(?P<filter_field>[A-Za-z_][A-Za-z0-9_]*)=(?P<filter_value>[^、]+?)の"
    r"(?P<cohort>[^、]+?)の中で、(?P<measure>[A-Za-z_][A-Za-z0-9_]*)の"
    r"平均値が最も高い(?P<group_field>[^、は]+?)は何(?P<unit>歳|年|日|時間)ですか。?$"
)
MULTI_FILTER_MEAN_HALF_UP = re.compile(
    r"^(?P<location>.+?)の分析対象データにおいて、"
    r"(?P<filters>(?:[A-Za-z_][A-Za-z0-9_]*=[^、]+、)+"
    r"[A-Za-z_][A-Za-z0-9_]*=[^、]+)に該当する"
    r"(?P<measure>[A-Za-z_][A-Za-z0-9_]*)の平均を算出してください。"
    r"四捨五入して整数値で出してください。?$"
)
NOTEBOOK_HEATMAP_MIN_ABS_CORRELATION = re.compile(
    r"^(?P<location>.+?)の\s*(?P<container>[^、]+?\.ipynb)\s*にある"
    r"(?P<visual>[^。]*?ヒートマップ)の図で可視化されている特徴量のうち、"
    r"(?P<target>[A-Za-z_][A-Za-z0-9_]*)との相関係数の絶対値が最も小さい"
    r"特徴量名を答えてください。?$"
)
REPORT_METRIC_DELTA = re.compile(
    r"^(?P<location>.+?)案件において、(?P<report_kind>.+?報告)資料に記載された"
    r"(?P<metric>[^、]+?)スコアの詳細値と、最終分析出力"
    r"(?P<container>[^、]+?\.json)に記録されている(?P=metric)スコアの詳細値を用いて、"
    r"改善幅を小数第(?P<digits>\d+)位まで答えてください。?$"
)
CONTRACT_HOURS_RATIO_TAX_DELTA = re.compile(
    r"^(?P<location>.+?)の契約条件において、仮に実績工数が見込工数の"
    r"(?P<denominator>[0-9０-９]+)分の(?P<numerator>[0-9０-９]+)だった場合、"
    r"最終請求金額（税込）は見込金額（税込）よりいくら少なくなりますか。?$"
)
EXCEL_AUTOFILTER_CONDITIONS = re.compile(
    r"^(?P<location>.+?)の(?P<container>[^、]+?\.xlsx)において、"
    r"(?P<sheet>[^、]+?)シートでフィルター?で抽出されている条件を"
    r"(?:教えて|答えて)ください。?$"
)
PIVOT_AVERAGE_ARGMAX_CONDITIONS_AND_AGGREGATE = re.compile(
    r"^(?P<location>.+?)の(?P<container>[^、。]+?\.xlsx)内の\s*"
    r"(?P<pivot_marker>PivotTable|ピボットテーブル)\s*"
    r"で集計されている表から、(?P<metric_surface>[^、。]+?)の平均が"
    r"最も高いものの抽出条件と集計内容を"
    r"(?:答えて|教えて)ください。?$"
)
PIVOT_AVERAGE_ARGMAX_CONDITIONS = re.compile(
    r"^(?P<location>.+?)の(?P<container>[^、。]+?\.xlsx)の"
    r"(?P<sheet>[^、。]+?)シートにおいて、平均(?P<metric_surface>[^、。]+?)が"
    r"最も高い層の抽出条件を"
    r"(?:答えて|教えて)ください。?$"
)
PPTX_OLD_LATEST_VISIBLE_TEXT_DIFF = re.compile(
    r"^(?P<location>.+)の(?P<document_key>[^の、。]+?)について、"
    r"oldフォルダ内の旧版と提案フォルダ直下の最新版を比較し、"
    r"変更された箇所を変更前と変更後で(?:答えて|教えて)"
    r"ください。?$"
)
PYTHON_CATEGORICAL_DTYPE_UNIQUE_RULE = re.compile(
    r"^(?P<location>.+?)の分析コードにおいて、"
    r"(?P<label>[A-Za-z_][A-Za-z0-9_]*)は\s*dtype\s*と"
    r"ユニーク数の条件でどのように判定していますか。?$",
    flags=re.IGNORECASE,
)
DOCX_HIGHLIGHTED_TEXT = re.compile(
    r"^(?P<location>.+)の(?P<document_key>[^\s、。()（）]+?)資料"
    r"[（(](?P<format>docx)[）)]において、"
    r"(?P<color>[^、。]+?)でハイライトされている部分を"
    r"すべて抜き出してください。?$",
    flags=re.IGNORECASE,
)
PPTX_SHAPE_FILL_TEXT = re.compile(
    r"^(?P<location>.+)の(?P<container>[^、。]+?)[PpＰｐ]"
    r"(?P<slide>[0-9０-９]+)において、"
    r"(?P<color>[^、。]+?)で強調されている箇所の文字列を"
    r"抜き出してください。?$"
)
ALL_PROJECT_MILESTONE_CUTOFF = re.compile(
    r"^(?P<milestone_left>[^,、。]+?)\s*"
    r"(?:、?\s*(?:または|もしくは|又は|あるいは)\s*)"
    r"(?P<milestone_right>[^,、。]+?)\s*"
    r"(?:が|の実施日が)\s*"
    r"(?P<cutoff>(?:20\d{2}年\d{1,2}月\d{1,2}日|"
    r"20\d{2}[-/.]\d{1,2}[-/.]\d{1,2}))\s*"
    r"(?P<comparator>以前|まで)に(?:実施|開催)された案件を(?:、)?"
    r"(?:主略称(?:で|にて)(?:すべて|全て)|"
    r"(?:すべて|全て)主略称(?:で|にて))"
    r"(?:挙げて|答えて|列挙して)ください。?$"
)
ALL_PROJECT_PAID_GROSS_TAX_SUM = re.compile(
    r"^(?:全案件|すべての案件|全ての案件)(?:で|において)"
    r"(?:支払った|支払済みの)(?P<gross>税込金額)をもとに、"
    r"(?P<tax>消費税額)の総額を(?:計算|算出)してください。?$"
)

ORANGE_RGB = frozenset({"FFF2E0D0", "FFFFF0E6", "FFFFC000", "FFFFA500", "FFFF9900"})
BLUE_RGB = frozenset({"FF00B0F0", "FF0000FF", "FF4472C4", "FF5B9BD5"})
YELLOW_RGB = frozenset({"FFFFFF00", "FFFFFF99", "FFFFEB3B"})

_WORD_HIGHLIGHT_ALIASES = {
    "黄": "yellow",
    "黄色": "yellow",
    "yellow": "yellow",
    "赤": "red",
    "赤色": "red",
    "red": "red",
    "青": "blue",
    "青色": "blue",
    "blue": "blue",
    "緑": "green",
    "緑色": "green",
    "green": "green",
    "シアン": "cyan",
    "cyan": "cyan",
    "マゼンタ": "magenta",
    "magenta": "magenta",
    "黒": "black",
    "黒色": "black",
    "black": "black",
    "白": "white",
    "白色": "white",
    "white": "white",
}
_PPTX_COLOR_ALIASES = {
    "赤": "red",
    "赤色": "red",
    "red": "red",
    "オレンジ": "orange",
    "橙": "orange",
    "orange": "orange",
    "黄": "yellow",
    "黄色": "yellow",
    "yellow": "yellow",
    "緑": "green",
    "緑色": "green",
    "green": "green",
    "シアン": "cyan",
    "水色": "cyan",
    "cyan": "cyan",
    "青": "blue",
    "青色": "blue",
    "blue": "blue",
    "マゼンタ": "magenta",
    "magenta": "magenta",
}
_PPTX_HUE_CENTERS = {
    "red": 0.0,
    "orange": 30.0,
    "yellow": 60.0,
    "green": 120.0,
    "cyan": 180.0,
    "blue": 240.0,
    "magenta": 300.0,
}

GRAPH_RULE_VERSION = "1.4"
_GRAPH_RULES = (
    (DATE_RANGE, "date_range_identifier_list", ("retrieve", "filter", "boolean_test", "project", "deduplicate", "list"), ("list", "identifier", "all", None)),
    (ASSIGNEE_COUNT, "assignee_task_count", ("retrieve", "filter", "project", "deduplicate", "count"), ("scalar", "integer", "single", None)),
    (
        PROJECT_PERSON_ASSIGNMENT_ROLE,
        "project_person_assignment_role",
        (
            "retrieve",
            "select_authoritative",
            "verify_complete",
            "filter",
            "boolean_test",
            "project",
        ),
        ("scalar", "string", "single", None),
    ),
    (PHASE_LATEST, "phase_latest_start_task", ("retrieve", "filter", "sort", "project"), ("scalar", "string", "single", None)),
    (BUFFER_SUM, "buffer_effort_sum", ("retrieve", "filter", "sum"), ("scalar", "number", "single", "時間")),
    (CHECKPOINT_TASKS, "checkpoint_milestone_task_join", ("retrieve", "filter", "project", "retrieve", "filter", "project", "deduplicate", "list"), ("list", "identifier", "all", None)),
    (MISSING_ROWS_MAX, "missing_rows_argmax_project", ("retrieve", "calculate", "argmax_all", "project"), ("scalar", "identifier", "single", None)),
    (STANDARDIZED_SHARE, "standardized_conditional_share", ("retrieve", "mean", "calculate", "filter", "mean", "filter", "count", "calculate"), ("scalar", "number", "single", "%")),
    (INTERACTION_COLUMNS, "metrics_code_interaction_list", ("retrieve", "retrieve", "filter", "verify", "project", "list"), ("list", "string", "all", None)),
    (GB_PARAMS, "runtime_model_parameter_resolution", ("retrieve", "retrieve", "verify", "project"), ("key_value", "number", "multiple", None)),
    (
        XLSX_CHART_SERIES_COLUMN,
        "xlsx_chart_series_column",
        ("retrieve", "select", "select", "resolve", "verify", "project"),
        ("scalar", "identifier", "single", None),
    ),
    (REGRESSION_PREDICTION, "regression_standardize_predict", ("retrieve", "mean", "calculate", "sum"), ("scalar", "number", "single", None)),
    (NEGATIVE_CORRELATION, "strongest_negative_correlation", ("retrieve", "calculate", "filter", "argmin_all", "project"), ("scalar", "identifier", "single", None)),
    (HIGHLIGHT_ROWS, "highlighted_row_projection", ("retrieve", "filter", "project", "deduplicate", "list"), ("list", "identifier", "all", None)),
    (BLUE_SUM, "highlighted_numeric_sum", ("retrieve", "filter", "sum", "calculate"), ("scalar", "number", "single", None)),
    (YELLOW_INTERSECTION, "highlight_intersection_difference", ("retrieve", "filter", "calculate", "absolute_distance"), ("scalar", "number", "single", None)),
    (COHORT_GROUP_MEAN_ARGMAX, "cohort_group_mean_argmax", ("retrieve", "filter", "filter", "group", "mean", "argmax_all", "project"), ("scalar", "number", "single", None)),
    (MULTI_FILTER_MEAN_HALF_UP, "multi_filter_mean_half_up", ("retrieve", "filter", "mean", "calculate"), ("scalar", "integer", "single", None)),
    (NOTEBOOK_HEATMAP_MIN_ABS_CORRELATION, "notebook_heatmap_min_abs_correlation", ("retrieve", "retrieve", "calculate", "sort", "filter", "argmin_all", "project"), ("scalar", "identifier", "single", None)),
    (REPORT_METRIC_DELTA, "report_metrics_decimal_delta", ("retrieve", "retrieve", "project", "calculate"), ("scalar", "number", "single", None)),
    (CONTRACT_HOURS_RATIO_TAX_DELTA, "contract_hours_ratio_tax_delta", ("retrieve", "verify", "calculate", "calculate"), ("scalar", "integer", "single", "円")),
    (EXCEL_AUTOFILTER_CONDITIONS, "excel_autofilter_conditions", ("retrieve", "verify", "project", "list"), ("key_value", "string", "multiple", None)),
    (
        PIVOT_AVERAGE_ARGMAX_CONDITIONS_AND_AGGREGATE,
        "pivot_average_argmax_conditions_and_aggregate",
        (
            "retrieve",
            "verify",
            "project",
            "group",
            "mean",
            "verify_complete",
            "argmax_all",
            "project",
        ),
        ("key_value", "string", "multiple", None),
    ),
    (
        PIVOT_AVERAGE_ARGMAX_CONDITIONS,
        "pivot_average_argmax_conditions",
        (
            "retrieve",
            "verify",
            "project",
            "group",
            "mean",
            "verify_complete",
            "argmax_all",
            "project",
        ),
        ("key_value", "string", "multiple", None),
    ),
    (
        PPTX_OLD_LATEST_VISIBLE_TEXT_DIFF,
        "pptx_old_latest_visible_text_diff",
        (
            "retrieve",
            "retrieve",
            "verify",
            "project",
            "compare",
            "filter",
            "list",
        ),
        ("list", "change_record", "all", None),
    ),
    (
        PYTHON_CATEGORICAL_DTYPE_UNIQUE_RULE,
        "python_categorical_dtype_unique_rule",
        (
            "retrieve",
            "resolve_entrypoint",
            "parse",
            "trace",
            "filter",
            "calculate",
            "classify",
            "project",
        ),
        ("scalar", "string", "single", None),
    ),
    (DOCX_HIGHLIGHTED_TEXT, "docx_highlighted_text_projection", ("retrieve", "match", "verify", "filter", "group", "project", "list"), ("list", "string", "all", None)),
    (PPTX_SHAPE_FILL_TEXT, "pptx_shape_fill_text_projection", ("retrieve", "verify", "select", "filter", "project", "list"), ("list", "string", "all", None)),
    (
        ALL_PROJECT_MILESTONE_CUTOFF,
        "all_project_milestone_cutoff_primary_alias",
        (
            "enumerate_projects",
            "retrieve",
            "match",
            "extract_date",
            "verify_complete",
            "filter",
            "resolve_primary_alias",
            "sort",
            "list",
        ),
        ("list", "identifier", "all", None),
    ),
    (
        ALL_PROJECT_PAID_GROSS_TAX_SUM,
        "all_project_paid_gross_tax_sum",
        (
            "enumerate_projects",
            "retrieve",
            "decrypt",
            "select_authoritative",
            "match",
            "verify_complete",
            "deduplicate",
            "retrieve_actual",
            "verify_arithmetic",
            "override_estimate",
            "calculate",
            "sum",
        ),
        ("scalar", "integer", "single", "円"),
    ),
)


def _canonical_json(value: Any) -> str:
    return json.dumps(
        value,
        ensure_ascii=False,
        sort_keys=True,
        separators=(",", ":"),
        allow_nan=False,
    )


def _parse_equality_filters(value: str) -> tuple[tuple[str, str], ...] | None:
    filters: list[tuple[str, str]] = []
    seen: set[str] = set()
    for item in value.split("、"):
        field, separator, expected = item.partition("=")
        field = field.strip()
        expected = expected.strip()
        if (
            separator != "="
            or re.fullmatch(r"[A-Za-z_][A-Za-z0-9_]*", field) is None
            or not expected
            or field in seen
        ):
            return None
        seen.add(field)
        filters.append((field, expected))
    return tuple(filters) if len(filters) >= 2 else None


def _parse_question_date(value: str) -> date | None:
    normalized = unicodedata.normalize("NFKC", value).strip()
    match = re.fullmatch(
        r"(20\d{2})(?:年|[-/.])(\d{1,2})(?:月|[-/.])(\d{1,2})日?",
        normalized,
    )
    if match is None:
        return None
    try:
        return date(*(int(part) for part in match.groups()))
    except ValueError:
        return None


def _declared_color(
    value: str,
    aliases: Mapping[str, str],
) -> str | None:
    return aliases.get(_normalized(value))


def _contract_operators(
    rule_id: str,
    bindings: Mapping[str, str],
    default: Sequence[str],
) -> tuple[str, ...] | None:
    if rule_id == "all_project_milestone_cutoff_primary_alias":
        if (
            _parse_question_date(bindings.get("cutoff", "")) is None
            or _normalized(bindings.get("milestone_left", ""))
            == _normalized(bindings.get("milestone_right", ""))
        ):
            return None
        return tuple(default)
    if rule_id == "contract_hours_ratio_tax_delta":
        try:
            denominator = int(
                unicodedata.normalize("NFKC", bindings.get("denominator", ""))
            )
            numerator = int(
                unicodedata.normalize("NFKC", bindings.get("numerator", ""))
            )
        except ValueError:
            return None
        if not (0 <= numerator < denominator <= 1_000_000):
            return None
        return tuple(default)
    if rule_id == "docx_highlighted_text_projection":
        document_key = unicodedata.normalize(
            "NFKC", bindings.get("document_key", "")
        ).strip()
        if (
            re.fullmatch(r"[A-Za-z][A-Za-z0-9._-]*", document_key) is None
            or _normalized(bindings.get("format", "")) != "docx"
            or _declared_color(
                bindings.get("color", ""), _WORD_HIGHLIGHT_ALIASES
            )
            is None
        ):
            return None
        return tuple(default)
    if rule_id == "pptx_shape_fill_text_projection":
        try:
            slide = int(unicodedata.normalize("NFKC", bindings.get("slide", "")))
        except ValueError:
            return None
        if (
            not 1 <= slide <= 10_000
            or not bindings.get("container", "").strip()
            or _declared_color(bindings.get("color", ""), _PPTX_COLOR_ALIASES)
            is None
        ):
            return None
        return tuple(default)
    if rule_id != "multi_filter_mean_half_up":
        return tuple(default)
    filters = _parse_equality_filters(bindings.get("filters", ""))
    if filters is None:
        return None
    return ("retrieve", *("filter" for _ in filters), "mean", "calculate")


def graph_contract_for_question(question: str) -> dict[str, Any] | None:
    """Compile one complete extended grammar into a typed operation graph."""

    for pattern, rule_id, operators, output_shape in _GRAPH_RULES:
        match = pattern.fullmatch(question)
        if match is None:
            continue
        bindings = {
            key: value
            for key, value in sorted(match.groupdict().items())
            if value is not None
        }
        resolved_operators = _contract_operators(rule_id, bindings, operators)
        if resolved_operators is None:
            return None
        nodes: list[dict[str, Any]] = []
        previous = "input_question"
        for index, operator in enumerate(resolved_operators, 1):
            output_ref = f"value_{index:03d}"
            nodes.append(
                {
                    "operation_id": f"op_{index:03d}_{operator}",
                    "operator": operator,
                    "input_refs": [previous],
                    "output_ref": output_ref,
                }
            )
            previous = output_ref
        container, value_type, cardinality, unit = output_shape
        if rule_id == "highlighted_row_projection" and bindings.get("target") == "タスク名":
            value_type = "string"
        if rule_id == "highlighted_numeric_sum":
            value_type = "integer"
        if rule_id == "cohort_group_mean_argmax":
            unit = bindings["unit"]
        display_precision = None
        if rule_id in {
            "standardized_conditional_share",
            "regression_standardize_predict",
            "report_metrics_decimal_delta",
        }:
            display_precision = {
                "mode": "decimal_places",
                "digits": int(bindings["digits"]),
            }
        required_keys = None
        if rule_id == "runtime_model_parameter_resolution":
            required_keys = ["n_estimators", "learning_rate", "random_state"]
        scope = {
            "location": bindings.get("location"),
            "container": bindings.get("container"),
        }
        if rule_id == "contract_hours_ratio_tax_delta":
            scope["container"] = "契約書.docx"
        if rule_id == "docx_highlighted_text_projection":
            scope.update(
                {
                    "container": "*.docx",
                    "document_key": bindings["document_key"],
                    "style_channel": "word_run_highlight",
                    "color": _declared_color(
                        bindings["color"], _WORD_HIGHLIGHT_ALIASES
                    ),
                }
            )
        if rule_id == "pptx_shape_fill_text_projection":
            scope.update(
                {
                    "slide": int(
                        unicodedata.normalize("NFKC", bindings["slide"])
                    ),
                    "style_channel": "shape_solid_fill",
                    "color": _declared_color(
                        bindings["color"], _PPTX_COLOR_ALIASES
                    ),
                }
            )
        if rule_id == "pptx_old_latest_visible_text_diff":
            scope.update(
                {
                    "container": "*.pptx",
                    "document_key": bindings["document_key"],
                    "old_location": "proposal/old",
                    "latest_location": "proposal/direct",
                    "comparison_channel": "visible_text",
                }
            )
        if rule_id == "python_categorical_dtype_unique_rule":
            scope.update(
                {
                    "container": "analysis_project",
                    "entrypoint_source": "README",
                    "classification_label": bindings["label"],
                }
            )
        if rule_id == "project_person_assignment_role":
            scope.update(
                {
                    "container": "02.計画/*.xlsx",
                    "source_kind": "authoritative_complete_roster",
                    "required_fields": ["役割", "氏名"],
                    "membership_subject": bindings["person"],
                }
            )
        if rule_id == "all_project_paid_gross_tax_sum":
            scope.update(
                {
                    "location": "プロジェクト/*",
                    "container": "02.計画/*.xlsx",
                    "payment_state": "completed",
                    "amount_basis": bindings["gross"],
                    "derived_measure": bindings["tax"],
                    "completion_source": "authoritative_project_plan",
                    "actual_override_source": "unique_current_final_report",
                    "actual_override_model": "time_and_materials",
                }
            )
        if "sheet" in bindings:
            scope["sheet"] = bindings["sheet"]
        if rule_id == "xlsx_chart_series_column":
            scope.update(
                {
                    "chart_index": int(
                        unicodedata.normalize("NFKC", bindings["chart"])
                    ),
                    "source_channel": "xlsx_ooxml_chart_series",
                }
            )
        core = {
            "graph_rule_version": GRAPH_RULE_VERSION,
            "rule_id": rule_id,
            "question_sha256": hashlib.sha256(question.encode("utf-8")).hexdigest(),
            "bindings": bindings,
            "scope": scope,
            "operation_graph": {
                "external_inputs": [
                    {
                        "input_ref": "input_question",
                        "input_type": "source_records",
                        "source": "question_scope",
                    }
                ],
                "nodes": nodes,
                "edges": [
                    {"from": nodes[index - 1]["output_ref"], "to": nodes[index]["operation_id"]}
                    for index in range(1, len(nodes))
                ],
            },
            "requested_output": {
                "source_operation_ref": nodes[-1]["operation_id"],
                "cardinality": cardinality,
                "answer_shape": {
                    "container": container,
                    "value_type": value_type,
                    "unit": unit,
                },
                "display_precision": display_precision,
                "required_keys": required_keys,
            },
        }
        return {
            "graph_contract_id": "xgraph_" + hashlib.sha256(
                _canonical_json(core).encode("utf-8")
            ).hexdigest()[:32],
            **core,
        }

    # Analysis-artifact joins are isolated from row/table rules because they
    # certify one selected run across leaderboard, JSON configuration, and
    # Python constructor flow before exposing a scalar parameter.
    from analysis_artifact_rules import (
        graph_contract_for_question as analysis_artifact_contract,
    )

    contract = analysis_artifact_contract(question)
    if contract is not None:
        return contract

    from excel_native_rules import (
        graph_contract_for_question as excel_native_contract,
    )

    contract = excel_native_contract(question)
    if contract is not None:
        return contract

    from xlsx_highlight_projection_rules import (
        graph_contract_for_question as xlsx_highlight_contract,
    )

    contract = xlsx_highlight_contract(question)
    if contract is not None:
        return contract

    from xlsx_pivot_highlight_rules import (
        graph_contract_for_question as xlsx_pivot_highlight_contract,
    )

    contract = xlsx_pivot_highlight_contract(question)
    if contract is not None:
        return contract

    from xlsx_histogram_rules import (
        graph_contract_for_question as xlsx_histogram_contract,
    )

    contract = xlsx_histogram_contract(question)
    if contract is not None:
        return contract

    from xlsx_formula_ml_rules import (
        graph_contract_for_question as xlsx_formula_ml_contract,
    )

    contract = xlsx_formula_ml_contract(question)
    if contract is not None:
        return contract

    from xlsx_version_diff_rules import (
        graph_contract_for_question as xlsx_version_diff_contract,
    )

    contract = xlsx_version_diff_contract(question)
    if contract is not None:
        return contract

    from notebook_version_diff_rules import (
        graph_contract_for_question as notebook_version_diff_contract,
    )

    contract = notebook_version_diff_contract(question)
    if contract is not None:
        return contract

    from report_metric_delta_graph_rules import (
        graph_contract_for_question as report_metric_delta_contract,
    )

    contract = report_metric_delta_contract(question)
    if contract is not None:
        return contract

    from contract_contact_graph_rules import (
        graph_contract_for_question as contract_contact_contract,
    )

    contract = contract_contact_contract(question)
    if contract is not None:
        return contract

    from cross_document_finance_rules import (
        graph_contract_for_question as cross_document_finance_contract,
    )

    contract = cross_document_finance_contract(question)
    if contract is not None:
        return contract

    from cross_project_portfolio_rules import (
        graph_contract_for_question as cross_project_portfolio_contract,
    )

    contract = cross_project_portfolio_contract(question)
    if contract is not None:
        return contract

    from cross_project_personnel_graph_rules import (
        graph_contract_for_question as cross_project_personnel_contract,
    )

    contract = cross_project_personnel_contract(question)
    if contract is not None:
        return contract

    from model_comparison_graph_rules import (
        graph_contract_for_question as model_comparison_contract,
    )

    contract = model_comparison_contract(question)
    if contract is not None:
        return contract

    from reported_feature_correlation_graph_rules import (
        graph_contract_for_question as reported_feature_correlation_contract,
    )

    contract = reported_feature_correlation_contract(question)
    if contract is not None:
        return contract

    from priority_task_owner_graph_rules import (
        graph_contract_for_question as priority_task_owner_contract,
    )

    contract = priority_task_owner_contract(question)
    if contract is not None:
        return contract

    from document_answerability_rules import (
        graph_contract_for_question as answerability_contract,
    )

    contract = answerability_contract(question)
    if contract is not None:
        return contract

    from glossary_evidence_rules import (
        graph_contract_for_question as glossary_evidence_contract,
    )

    contract = glossary_evidence_contract(question)
    if contract is not None:
        return contract

    from notebook_correlation_rules import (
        graph_contract_for_question as notebook_correlation_contract,
    )

    contract = notebook_correlation_contract(question)
    if contract is not None:
        return contract

    from notebook_axis_tick_rules import (
        graph_contract_for_question as notebook_axis_tick_contract,
    )

    contract = notebook_axis_tick_contract(question)
    if contract is not None:
        return contract

    from xlsx_role_task_graph_rules import (
        graph_contract_for_question as xlsx_role_task_contract,
    )

    contract = xlsx_role_task_contract(question)
    if contract is not None:
        return contract

    from pptx_schedule_rules import (
        graph_contract_for_question as pptx_schedule_contract,
    )

    contract = pptx_schedule_contract(question)
    if contract is not None:
        return contract

    from pptx_scope_exclusion_rules import (
        graph_contract_for_question as pptx_scope_exclusion_contract,
    )

    contract = pptx_scope_exclusion_contract(question)
    if contract is not None:
        return contract

    from pptx_feature_legend_rules import (
        graph_contract_for_question as pptx_feature_legend_contract,
    )

    contract = pptx_feature_legend_contract(question)
    if contract is not None:
        return contract

    from pptx_revision_summary_rules import (
        graph_contract_for_question as pptx_revision_summary_contract,
    )

    contract = pptx_revision_summary_contract(question)
    if contract is not None:
        return contract

    from pdf_operational_role_rules import (
        graph_contract_for_question as pdf_operational_role_contract,
    )

    contract = pdf_operational_role_contract(question)
    if contract is not None:
        return contract

    from pdf_native_style_rules import (
        graph_contract_for_question as pdf_native_style_contract,
    )

    contract = pdf_native_style_contract(question)
    if contract is not None:
        return contract

    from pdf_investment_coefficient_rules import (
        graph_contract_for_question as pdf_investment_coefficient_contract,
    )

    contract = pdf_investment_coefficient_contract(question)
    if contract is not None:
        return contract

    from pdf_highlight_trend_rules import (
        graph_contract_for_question as pdf_highlight_trend_contract,
    )

    contract = pdf_highlight_trend_contract(question)
    if contract is not None:
        return contract

    from pdf_action_transition_rules import (
        graph_contract_for_question as pdf_action_transition_contract,
    )

    contract = pdf_action_transition_contract(question)
    if contract is not None:
        return contract

    from pdf_action_content_graph_rules import (
        graph_contract_for_question as pdf_action_content_contract,
    )

    contract = pdf_action_content_contract(question)
    if contract is not None:
        return contract

    from docx_page_structure_rules import (
        graph_contract_for_question as docx_page_structure_contract,
    )

    contract = docx_page_structure_contract(question)
    if contract is not None:
        return contract

    from docx_native_style_rules import (
        graph_contract_for_question as docx_native_style_contract,
    )

    contract = docx_native_style_contract(question)
    if contract is not None:
        return contract

    from docx_mixed_content_rules import (
        graph_contract_for_question as docx_mixed_contract,
    )

    contract = docx_mixed_contract(question)
    if contract is not None:
        return contract

    from pptx_mixed_content_rules import (
        graph_contract_for_question as pptx_mixed_contract,
    )

    contract = pptx_mixed_contract(question)
    if contract is not None:
        return contract

    from pptx_version_diff_rules import (
        graph_contract_for_question as pptx_version_diff_contract,
    )

    contract = pptx_version_diff_contract(question)
    if contract is not None:
        return contract

    from pptx_spatial_rules import (
        graph_contract_for_question as pptx_spatial_contract,
    )

    contract = pptx_spatial_contract(question)
    if contract is not None:
        return contract

    # The proposal metric grammar is isolated because it audits authored PPTX
    # runs rather than tabular records.  It still participates in the same
    # mandatory graph gate before any retrieval or answer generation.
    from proposal_metric_rules import graph_contract_for_question as proposal_contract

    contract = proposal_contract(question)
    if contract is not None:
        return contract

    # PDF visual grammars live in an independent, fail-closed module because
    # their evidence pipeline is raster/OCR based.  Keep the import lazy to
    # preserve this module's use by lightweight tabular-only callers.
    from pdf_visual_rules import graph_contract_for_pdf_question

    return graph_contract_for_pdf_question(question)


def validate_graph_contract(question: str, contract: Mapping[str, Any]) -> bool:
    """Independently rebuild the contract; no self-declared graph is trusted."""

    expected = graph_contract_for_question(question)
    return expected is not None and _canonical_json(expected) == _canonical_json(contract)


def _decision(
    answer: str,
    paths: Sequence[Path],
    root: Path,
    operations: int,
) -> StructuredCandidateDecision:
    relative = tuple(
        sorted(
            {
                unicodedata.normalize("NFC", path.relative_to(root).as_posix())
                for path in paths
            }
        )
    )
    digest = hashlib.sha256()
    for path in sorted(set(paths), key=lambda item: item.as_posix()):
        digest.update(path.read_bytes())
    return StructuredCandidateDecision(
        "resolved",
        "certified_extended",
        StructuredCandidateAnswer(
            answer=answer,
            source_paths=relative,
            source_sha256=digest.hexdigest(),
            operation_count=operations,
            output_count=1,
        ),
    )


def _all_files(engine: Any, suffixes: set[str] | None = None) -> list[Path]:
    return [
        path
        for path in engine.source_root.rglob("*")
        if path.is_file()
        and not path.is_symlink()
        and not path.name.startswith("~$")
        and (suffixes is None or path.suffix.casefold() in suffixes)
    ]


def _project_matches(engine: Any, path: Path, location: str) -> bool:
    candidates = _candidate_values(location, engine.glossary)
    relative = path.relative_to(engine.source_root)
    return _location_matches(relative.parts[:-1], candidates)


def _named_paths(
    engine: Any,
    location: str,
    container: str,
    suffixes: set[str],
) -> list[Path]:
    candidates = _candidate_values(container, engine.glossary)
    names = {_normalized(value) for value in candidates}
    matches = []
    for path in _all_files(engine, suffixes):
        if not _project_matches(engine, path, location):
            continue
        if _normalized(path.name) in names or _normalized(path.stem) in names:
            matches.append(path)
    return matches


def _plan_paths(engine: Any, location: str) -> list[Path]:
    matches = []
    for path in _all_files(engine, {".xlsx"}):
        if not _project_matches(engine, path, location):
            continue
        parts = [_normalized(part) for part in path.relative_to(engine.source_root).parts]
        if any("計画" in part for part in parts[:-1]):
            matches.append(path)
    return matches


def _unique_table(
    paths: Sequence[Path],
    required_fields: Sequence[str],
) -> tuple[Any, list[Path]] | None:
    tables = []
    for path in paths:
        tables.extend(_read_source_tables(path, required_fields))
    by_digest: dict[str, list[Any]] = {}
    for table in tables:
        by_digest.setdefault(table.table_sha256, []).append(table)
    if len(by_digest) != 1:
        return None
    duplicates = next(iter(by_digest.values()))
    table = min(duplicates, key=lambda item: (item.path.as_posix(), item.table_name))
    return table, [item.path for item in duplicates]


def _records(table: Any) -> list[dict[str, str]]:
    return [dict(zip(table.headers, row)) for row in table.rows]


def _active_sheet_records(
    path: Path,
    required_fields: Sequence[str],
) -> list[dict[str, str]] | None:
    """Read the workbook's authoritative active sheet in one forward pass.

    Planning sheets legitimately use blank cells for inherited phase labels and
    optional effort.  The stricter generic table decoder rejects those blanks,
    so these two planning rules use a purpose-specific sequential reader while
    still requiring one exact header and one exact source workbook.
    """

    from openpyxl import load_workbook

    workbook = load_workbook(path, read_only=True, data_only=True)
    try:
        worksheet = workbook.active
        iterator = worksheet.iter_rows(values_only=True)
        headers: tuple[str, ...] | None = None
        for row_number, raw_row in enumerate(iterator, 1):
            cells = tuple(_cell_text(value) for value in raw_row)
            if all(field in cells for field in required_fields):
                if len([value for value in cells if value]) != len(
                    set(value for value in cells if value)
                ):
                    return None
                last = max(index for index, value in enumerate(cells) if value)
                headers = cells[: last + 1]
                break
            if row_number >= 100:
                return None
        if headers is None:
            return None
        records: list[dict[str, str]] = []
        for raw_row in iterator:
            values = tuple(
                _cell_text(raw_row[index]) if index < len(raw_row) else ""
                for index in range(len(headers))
            )
            if any(values):
                records.append(dict(zip(headers, values)))
        return records
    finally:
        workbook.close()


def _date_range(engine: Any, match: re.Match[str]) -> StructuredCandidateDecision | None:
    paths = _named_paths(engine, match["location"], match["container"], {".xlsx", ".csv", ".tsv"})
    resolved = _unique_table(paths, [match["left"], match["right"], match["target"]])
    if resolved is None:
        return None
    table, sources = resolved
    start = datetime.fromisoformat(match["start"])
    end = datetime.fromisoformat(match["end"])
    values = []
    for row in _records(table):
        dates = []
        for field in (match["left"], match["right"]):
            try:
                dates.append(datetime.fromisoformat(row[field]))
            except ValueError:
                return None
        if any(start <= value <= end for value in dates):
            values.append(row[match["target"]])
    if not values or len(values) != len(set(values)):
        return None
    return _decision("、".join(values), sources, engine.source_root, 2)


def _assignee_count(engine: Any, match: re.Match[str]) -> StructuredCandidateDecision | None:
    named = _named_paths(engine, match["location"], match["container"], {".xlsx"})
    paths = named or _plan_paths(engine, match["location"])
    resolved = _unique_table(paths, ["タスクID", "担当者"])
    if resolved is None:
        return None
    table, sources = resolved
    person = _normalized(match["person"])
    task_ids = [
        row["タスクID"]
        for row in _records(table)
        if person in _normalized(row["担当者"])
    ]
    if not task_ids or len(task_ids) != len(set(task_ids)):
        return None
    return _decision(str(len(task_ids)), sources, engine.source_root, 2)


_COMPLETE_ROSTER_SHEET_MARKERS = (
    "リソース",
    "体制",
    "roster",
    "team",
)
_COMPLETE_ROSTER_SUBJECT_MARKERS = (
    "役割",
    "体制",
    "roster",
    "team",
    "member",
)
_COMPLETE_ROSTER_BOUNDARY_MARKERS = (
    "分担",
    "配分",
    "主担当",
    "一覧",
    "roster",
    "team",
)
_INCOMPLETE_ROSTER_VALUES = frozenset(
    {"-", "--", "—", "todo", "tbd", "na", "n/a", "未定", "未確定"}
)


def _person_identity(value: object) -> str:
    return re.sub(r"\s+", "", _normalized(value))


def _is_complete_roster_heading(sheet_name: str, title: str) -> bool:
    sheet = _normalized(sheet_name)
    heading = _normalized(title)
    return (
        any(marker in sheet for marker in _COMPLETE_ROSTER_SHEET_MARKERS)
        and any(marker in heading for marker in _COMPLETE_ROSTER_SUBJECT_MARKERS)
        and any(marker in heading for marker in _COMPLETE_ROSTER_BOUNDARY_MARKERS)
    )


def _complete_assignment_roster(
    path: Path,
) -> tuple[tuple[str, str], ...] | None:
    """Read one source-declared complete person-to-role roster.

    A role/name pair found in an arbitrary task or meeting table is not enough
    to prove absence.  The table must live on a resource/team sheet, carry an
    explicit complete-roster heading, contain atomic role/name rows, and have
    one unambiguous boundary.  More than one qualifying table is deliberately
    treated as an unresolved source conflict.
    """

    from openpyxl import load_workbook

    try:
        workbook = load_workbook(path, read_only=True, data_only=True)
    except Exception:
        return None
    candidates: list[tuple[tuple[str, str], ...]] = []
    malformed = False
    try:
        for worksheet in workbook.worksheets:
            rows = [
                tuple(_cell_text(value) for value in row)
                for row in worksheet.iter_rows(values_only=True)
            ]
            for header_index, row in enumerate(rows):
                normalized = tuple(_normalized(value) for value in row)
                if normalized.count("役割") != 1 or normalized.count("氏名") != 1:
                    continue
                role_index = normalized.index("役割")
                person_index = normalized.index("氏名")
                preceding = rows[max(0, header_index - 3) : header_index]
                title = " ".join(
                    value
                    for prior_row in preceding
                    for value in prior_row
                    if value.strip()
                )
                if not _is_complete_roster_heading(worksheet.title, title):
                    continue
                roster: list[tuple[str, str]] = []
                for body_row in rows[header_index + 1 :]:
                    if not any(value.strip() for value in body_row):
                        break
                    role = body_row[role_index].strip() if role_index < len(body_row) else ""
                    person = body_row[person_index].strip() if person_index < len(body_row) else ""
                    if (
                        not role
                        or not person
                        or _normalized(role) in _INCOMPLETE_ROSTER_VALUES
                        or _normalized(person) in _INCOMPLETE_ROSTER_VALUES
                        or re.search(r"[/／、,;；\n]", person)
                    ):
                        malformed = True
                        break
                    roster.append((role, person))
                if malformed:
                    break
                identities = [_person_identity(person) for _, person in roster]
                if (
                    not roster
                    or any(not identity for identity in identities)
                    or len(identities) != len(set(identities))
                ):
                    malformed = True
                    break
                candidates.append(tuple(roster))
            if malformed:
                break
    finally:
        workbook.close()
    if malformed or len(candidates) != 1:
        return None
    return candidates[0]


def _project_person_assignment_role(
    engine: Any,
    match: re.Match[str],
) -> StructuredCandidateDecision | None:
    projects = _all_project_directories(engine)
    if projects is None:
        return None
    location_candidates = _candidate_values(match["location"], engine.glossary)
    scoped = [
        project
        for project in projects
        if _location_matches((project.name,), location_candidates)
    ]
    if len(scoped) != 1:
        return None
    project = scoped[0]
    plan = _authoritative_project_plan_workbook(project)
    if plan is None:
        return None
    roster = _complete_assignment_roster(plan)
    if roster is None:
        return None
    target = _person_identity(match["person"])
    roles = [role for role, person in roster if _person_identity(person) == target]
    if len(roles) > 1:
        return None
    answer = roles[0] if roles else "アサインされていない"
    return _decision(answer, [plan], engine.source_root, 6)


def _phase_latest(engine: Any, match: re.Match[str]) -> StructuredCandidateDecision | None:
    paths = _named_paths(engine, match["location"], match["container"], {".xlsx"})
    if len(paths) != 1:
        return None
    records = _active_sheet_records(paths[0], ["フェーズNo.", "タスク名", "開始日"])
    if records is None:
        return None
    phase = None
    candidates: list[tuple[datetime, str]] = []
    for row in records:
        if row["フェーズNo."]:
            try:
                phase = int(Decimal(row["フェーズNo."]))
            except Exception:
                return None
        if phase != int(match["phase"]):
            continue
        try:
            candidates.append((datetime.fromisoformat(row["開始日"]), row["タスク名"]))
        except ValueError:
            return None
    if not candidates:
        return None
    latest = max(value[0] for value in candidates)
    answers = [name for date, name in candidates if date == latest]
    if len(answers) != 1:
        return None
    return _decision(answers[0], paths, engine.source_root, 3)


def _buffer_sum(engine: Any, match: re.Match[str]) -> StructuredCandidateDecision | None:
    paths = _named_paths(engine, match["location"], match["container"], {".xlsx"})
    if len(paths) != 1:
        return None
    records = _active_sheet_records(paths[0], ["種別", "工数(h)"])
    if records is None:
        return None
    total = sum(
        Decimal(row["工数(h)"])
        for row in records
        if "バッファ" in row["種別"] and row["工数(h)"]
    )
    rendered = format(total, "f").rstrip("0").rstrip(".") + "時間"
    return _decision(rendered, paths, engine.source_root, 2)


def _unique_table_from_matrix(path: Path, name: str, matrix: Sequence[Sequence[Any]], fields: Sequence[str]) -> Any | None:
    from structured_candidate import _source_sha256, _table_from_matrix

    return _table_from_matrix(path, _source_sha256(path), name, matrix, fields)


def _expand_task_tokens(value: str) -> list[str]:
    tokens: list[str] = []
    for start, end in re.findall(r"(T\d+)\s*[～~\-]\s*(T\d+)", value, flags=re.I):
        left = int(re.search(r"\d+", start).group())
        right = int(re.search(r"\d+", end).group())
        width = len(re.search(r"\d+", start).group())
        tokens.extend(f"T{number:0{width}d}" for number in range(left, right + 1))
    masked = re.sub(r"T\d+\s*[～~\-]\s*T\d+", " ", value, flags=re.I)
    tokens.extend(re.findall(r"T\d+", masked, flags=re.I))
    return list(dict.fromkeys(token.upper() for token in tokens))


def _checkpoint_tasks(engine: Any, match: re.Match[str]) -> StructuredCandidateDecision | None:
    paths = _plan_paths(engine, match["location"])
    if len(paths) != 1:
        return None
    from openpyxl import load_workbook

    workbook = load_workbook(paths[0], read_only=True, data_only=True)
    try:
        matrices = {
            ws.title: [tuple(row) for row in ws.iter_rows(values_only=True)]
            for ws in workbook.worksheets
        }
    finally:
        workbook.close()
    cp = f"CP{int(match['number'])}"
    related_ms: list[str] = []
    direct_tasks: list[str] = []
    for matrix in matrices.values():
        for row in matrix:
            cells = ["" if value is None else str(value) for value in row]
            if cp not in cells:
                continue
            joined = " ".join(cells)
            related_ms.extend(re.findall(r"MS\d+", joined, flags=re.I))
            direct_tasks.extend(_expand_task_tokens(joined))
    tasks = list(dict.fromkeys(direct_tasks))
    for milestone in dict.fromkeys(value.upper() for value in related_ms):
        for matrix in matrices.values():
            for row in matrix:
                cells = ["" if value is None else str(value) for value in row]
                if milestone in cells:
                    tasks.extend(_expand_task_tokens(" ".join(cells)))
    tasks = list(dict.fromkeys(tasks))
    if not tasks:
        return None
    return _decision("、".join(tasks), paths, engine.source_root, 3)


def _primary_alias(glossary: Any, canonical: str) -> str | None:
    for alias, canons in getattr(glossary, "entries", {}).items():
        if canonical in canons:
            return str(alias)
    return None


def _unique_primary_alias_for_project(
    glossary: Any,
    project: str,
) -> tuple[str, int] | None:
    """Resolve one source-declared primary alias without guessing.

    A primary declaration is accepted only when the alias maps to this one
    canonical project in both the general and primary glossary maps.  This
    prevents a shared alias or two competing primary aliases from silently
    choosing a project label.
    """

    project_key = _normalized(project)
    entries = getattr(glossary, "entries", {})
    primary_entries = getattr(glossary, "primary_entries", {})
    if not isinstance(entries, Mapping) or not isinstance(primary_entries, Mapping):
        return None
    candidates: list[tuple[str, int]] = []
    for ordinal, (raw_alias, raw_primary_canons) in enumerate(primary_entries.items()):
        alias = str(raw_alias).strip()
        declared_canons = entries.get(raw_alias)
        if (
            not alias
            or not isinstance(raw_primary_canons, Sequence)
            or isinstance(raw_primary_canons, (str, bytes))
            or not isinstance(declared_canons, Sequence)
            or isinstance(declared_canons, (str, bytes))
        ):
            continue
        primary_keys = {_normalized(value) for value in raw_primary_canons}
        declared_keys = {_normalized(value) for value in declared_canons}
        if primary_keys == {project_key} and declared_keys == {project_key}:
            candidates.append((alias, ordinal))
    unique = list(dict.fromkeys(candidates))
    return unique[0] if len(unique) == 1 else None


def _all_project_directories(engine: Any) -> list[Path] | None:
    roots = []
    if _normalized(engine.source_root.name) == _normalized("プロジェクト"):
        roots.append(engine.source_root)
    roots.extend(
        path
        for path in engine.source_root.rglob("*")
        if path.is_dir()
        and not path.is_symlink()
        and _normalized(path.name) == _normalized("プロジェクト")
    )
    roots = list(dict.fromkeys(roots))
    if len(roots) != 1:
        return None
    projects = [
        path
        for path in roots[0].iterdir()
        if path.is_dir() and not path.is_symlink() and not path.name.startswith(".")
    ]
    if not projects:
        return None
    return sorted(projects, key=lambda path: (_normalized(path.name), path.name))


def _project_plan_workbooks(project: Path) -> list[Path]:
    paths: list[Path] = []
    for path in project.rglob("*"):
        if (
            not path.is_file()
            or path.is_symlink()
            or path.name.startswith(("~$", "."))
            or path.suffix.casefold() not in {".xlsx", ".xlsm"}
        ):
            continue
        parent_parts = [_normalized(part) for part in path.relative_to(project).parts[:-1]]
        if any(
            "計画" in part or part in {"plan", "planning"}
            for part in parent_parts
        ):
            paths.append(path)
    return sorted(paths, key=lambda path: unicodedata.normalize("NFC", path.as_posix()))


def _password_rule_sources(engine: Any) -> list[Path]:
    matches: list[Path] = []
    for path in _all_files(engine, {".docx"}):
        if "パスワード導出規則" not in _normalized(path.stem):
            continue
        paragraphs = _docx_paragraphs(path)
        if not paragraphs:
            continue
        text = _normalized("\n".join(paragraphs))
        if all(
            token in text
            for token in (
                "da-[案件略号]-[開始年月日8桁]-[拡張子コード]",
                "主略称",
                "yyyymmdd",
            )
        ):
            matches.append(path)
    return matches


def _encrypted_workbook_bytes(
    engine: Any,
    project: Path,
    path: Path,
    primary_alias: str,
) -> tuple[bytes, list[Path]] | None:
    """Decrypt a plan only through the source-backed generic password rule."""

    rule_sources = _password_rule_sources(engine)
    if len(rule_sources) != 1:
        return None
    dated_sources: list[Path] = []
    dates: list[str] = []
    for source in project.rglob("*"):
        if not source.is_file() or source.is_symlink() or source.name.startswith("~$"):
            continue
        found = re.findall(r"20\d{6}", unicodedata.normalize("NFKC", source.name))
        if found:
            dated_sources.append(source)
            dates.extend(found)
    dates = sorted(set(dates))
    if not dates:
        return None
    try:
        from extract import password_candidates, try_decrypt
    except ImportError:
        return None
    decrypted = try_decrypt(
        path,
        password_candidates(path, [primary_alias], dates),
    )
    if decrypted is None:
        return None
    evidence = [path, *rule_sources, *dated_sources]
    return decrypted, list(dict.fromkeys(evidence))


def _source_date(value: Any) -> date | None:
    if isinstance(value, datetime):
        return value.date()
    if isinstance(value, date):
        return value
    if not isinstance(value, str):
        return None
    normalized = unicodedata.normalize("NFKC", value).strip()
    match = re.fullmatch(
        r"(20\d{2})[-/.](\d{1,2})[-/.](\d{1,2})"
        r"(?:[ T]\d{1,2}:\d{2}(?::\d{2}(?:\.\d+)?)?)?",
        normalized,
    )
    if match is None:
        match = re.fullmatch(r"(20\d{2})年(\d{1,2})月(\d{1,2})日", normalized)
    if match is None:
        return None
    try:
        return date(*(int(part) for part in match.groups()))
    except ValueError:
        return None


_NON_EVENT_MILESTONE_MARKERS = (
    "資料",
    "作成",
    "準備",
    "論点",
    "指摘",
    "反映",
    "状況",
    "予定",
    "事前",
    "事後",
)
_EVENT_MILESTONE_MARKERS = ("実施", "開催", "完了", "承認")


def _milestone_cell_matches(value: Any, milestone: str) -> bool:
    cell = _normalized(value)
    target = _normalized(milestone)
    if not cell or not target:
        return False
    if cell == target:
        return True
    position = cell.find(target)
    if position < 0:
        return False
    remainder = cell[:position] + cell[position + len(target):]
    if any(marker in remainder for marker in _NON_EVENT_MILESTONE_MARKERS):
        return False
    return (
        any(marker in remainder for marker in _EVENT_MILESTONE_MARKERS)
        or remainder.strip(" -_:/、。()（）") in {"会", "会議"}
    )


def _milestone_dates_from_workbook(
    source: Path | io.BytesIO,
    milestones: Sequence[str],
) -> tuple[tuple[set[date], ...], bool] | None:
    from openpyxl import load_workbook

    try:
        workbook = load_workbook(source, read_only=True, data_only=True)
    except Exception:
        return None
    dates_by_branch: list[set[date]] = [set() for _ in milestones]
    ambiguous = False
    try:
        for worksheet in workbook.worksheets:
            for row in worksheet.iter_rows(values_only=True):
                matched_branches = {
                    index
                    for index, milestone in enumerate(milestones)
                    if any(_milestone_cell_matches(value, milestone) for value in row)
                }
                if not matched_branches:
                    continue
                row_dates = {
                    parsed
                    for value in row
                    if (parsed := _source_date(value)) is not None
                }
                if len(row_dates) == 1:
                    for index in matched_branches:
                        dates_by_branch[index].update(row_dates)
                elif len(row_dates) > 1:
                    ambiguous = True
    finally:
        workbook.close()
    return tuple(dates_by_branch), ambiguous


def _all_project_milestone_cutoff(
    engine: Any,
    match: re.Match[str],
) -> StructuredCandidateDecision | None:
    cutoff = _parse_question_date(match["cutoff"])
    projects = _all_project_directories(engine)
    if cutoff is None or projects is None:
        return None
    milestones = (match["milestone_left"].strip(), match["milestone_right"].strip())
    resolved: list[tuple[int, int, date, str]] = []
    evidence: list[Path] = []
    for project in projects:
        alias_resolution = _unique_primary_alias_for_project(
            engine.glossary,
            project.name,
        )
        paths = _project_plan_workbooks(project)
        if alias_resolution is None or not paths:
            return None
        alias, alias_ordinal = alias_resolution
        project_dates_by_branch: list[set[date]] = [set() for _ in milestones]
        for path in paths:
            if zipfile.is_zipfile(path):
                source: Path | io.BytesIO = path
                sources = [path]
            else:
                decrypted = _encrypted_workbook_bytes(
                    engine,
                    project,
                    path,
                    alias,
                )
                if decrypted is None:
                    return None
                data, sources = decrypted
                source = io.BytesIO(data)
            extracted = _milestone_dates_from_workbook(source, milestones)
            if extracted is None:
                return None
            dates_by_branch, ambiguous = extracted
            if ambiguous:
                return None
            for index, dates in enumerate(dates_by_branch):
                project_dates_by_branch[index].update(dates)
            evidence.extend(sources)
        if any(len(dates) > 1 for dates in project_dates_by_branch):
            return None
        matched = [
            (index, next(iter(dates)))
            for index, dates in enumerate(project_dates_by_branch)
            if dates
        ]
        if not matched:
            return None
        branch_index, milestone_date = matched[0]
        if milestone_date <= cutoff:
            resolved.append((branch_index, alias_ordinal, milestone_date, alias))
    if not resolved:
        return None
    aliases = [
        alias
        for _, _, _, alias in sorted(
            resolved,
            key=lambda item: (
                item[0],
                item[1],
                item[2],
                _normalized(item[3]),
                item[3],
            ),
        )
    ]
    if len(aliases) != len(set(aliases)):
        return None
    return _decision(
        "、".join(aliases),
        evidence,
        engine.source_root,
        9,
    )


_PLAN_REVISION_SUFFIX = re.compile(
    r"^(?P<base>.*?)(?:[_\-\s]*(?:r|rev|v)(?P<revision>\d+))$",
    flags=re.IGNORECASE,
)
_NON_AUTHORITATIVE_CONTRACT = re.compile(
    r"(?:draft|ドラフト|old|旧版|旧)",
    flags=re.IGNORECASE,
)
_NON_AUTHORITATIVE_REPORT = re.compile(
    r"(?:draft|ドラフト|old|旧版|旧|backup|archive|previous|prev)",
    flags=re.IGNORECASE,
)
_PAYMENT_MARKER = re.compile(r"支払|入金|精算")
_PAYMENT_INCOMPLETE = re.compile(r"^(?:未着手|未完了|open|保留|中止)$", re.I)
_YEN_AMOUNT = re.compile(
    r"(?:[\u00a5￥]\s*(?P<prefix>[0-9]+(?:,[0-9]{3})*)|"
    r"(?P<suffix>[0-9]+(?:,[0-9]{3})*)\s*円)"
)
_PERCENT_RATE = re.compile(
    r"消費税率[^0-9]{0,12}(?P<rate>[0-9]+(?:\.[0-9]+)?)\s*%"
)
_CONTRACT_HOURLY_RATE = re.compile(
    r"(?:時間単価|1\s*時間(?:当たり|あたり))[^0-9]{0,24}"
    r"(?P<rate>[0-9]+(?:,[0-9]{3})*)\s*円"
)
_REPORT_MONEY = re.compile(
    r"(?:[\u00a5￥]\s*(?P<prefix>[0-9]+(?:,[0-9]{3})*)|"
    r"JPY\s*(?P<jpy_prefix>[0-9]+(?:,[0-9]{3})*)|"
    r"(?P<suffix>[0-9]+(?:,[0-9]{3})*)\s*(?:JPY|円))",
    flags=re.IGNORECASE,
)
_REPORT_HOURS = re.compile(
    r"(?<![0-9.])(?P<hours>[0-9]+(?:\.[0-9]+)?)\s*(?:時間|h(?![A-Za-z]))",
    flags=re.IGNORECASE,
)
_REPORT_NON_ACTUAL = re.compile(
    r"(?:精算想定値|例示|実績工数の最終確定値[^\n]{0,40}含まれない|"
    r"見込工数[^\n]{0,40}(?:用いた|基づく)|実請求時|実際の請求書)",
    flags=re.IGNORECASE | re.DOTALL,
)


def _authoritative_project_plan_workbook(project: Path) -> Path | None:
    """Select only a source-declared latest plan revision."""

    paths = _project_plan_workbooks(project)
    if not paths:
        return None
    grouped: dict[tuple[str, str], list[tuple[int | None, Path]]] = {}
    for path in paths:
        stem = unicodedata.normalize("NFKC", path.stem).strip()
        match = _PLAN_REVISION_SUFFIX.fullmatch(stem)
        if match is None:
            base = _normalized(stem)
            revision = None
        else:
            base = _normalized(match["base"])
            revision = int(match["revision"])
        parent = _normalized(path.parent.relative_to(project).as_posix())
        grouped.setdefault((parent, base), []).append((revision, path))
    if len(grouped) != 1:
        return None
    candidates = next(iter(grouped.values()))
    numbered = [(revision, path) for revision, path in candidates if revision is not None]
    if numbered:
        latest = max(revision for revision, _ in numbered)
        winners = [path for revision, path in numbered if revision == latest]
        return winners[0] if len(winners) == 1 else None
    return candidates[0][1] if len(candidates) == 1 else None


def _authoritative_project_contract(project: Path) -> Path | None:
    paths = [
        path
        for path in project.rglob("*")
        if path.is_file()
        and not path.is_symlink()
        and not path.name.startswith(("~$", "."))
        and path.suffix.casefold() == ".docx"
        and any(
            "契約" in _normalized(part)
            for part in path.relative_to(project).parts[:-1]
        )
    ]
    finals = [path for path in paths if _NON_AUTHORITATIVE_CONTRACT.search(path.stem) is None]
    return finals[0] if len(finals) == 1 else None


def _authoritative_project_final_report(project: Path) -> Path | None:
    """Bind one current final report without relying on project identity."""

    paths: list[Path] = []
    for path in project.rglob("*"):
        if (
            not path.is_file()
            or path.is_symlink()
            or path.name.startswith(("~$", "."))
            or path.suffix.casefold() not in {".pdf", ".pptx"}
        ):
            continue
        relative = path.relative_to(project)
        if not any("報告" in _normalized(part) for part in relative.parts[:-1]):
            continue
        stem = _normalized(path.stem)
        if "最終" not in stem or "報告" not in stem:
            continue
        if any(
            _NON_AUTHORITATIVE_REPORT.search(part) is not None
            for part in relative.parts
        ):
            continue
        paths.append(path)
    return paths[0] if len(paths) == 1 else None


def _docx_source_units(source: Path | io.BytesIO) -> list[str] | None:
    """Return paragraph and table-row units without flattening row structure."""

    try:
        if isinstance(source, Path) and source.stat().st_size > _DIRECT_SOURCE_MAX_BYTES:
            return None
        if isinstance(source, io.BytesIO):
            source.seek(0)
        with zipfile.ZipFile(source) as archive:
            root = ET.fromstring(archive.read("word/document.xml"))
    except Exception:
        return None
    namespace = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    body = root.find("w:body", namespace)
    if body is None:
        return None
    word = "{" + namespace["w"] + "}"

    def text_of(element: ET.Element) -> str:
        return unicodedata.normalize(
            "NFKC",
            "".join(node.text or "" for node in element.iter(word + "t")),
        ).strip()

    units: list[str] = []
    for child in body:
        if child.tag == word + "p":
            value = text_of(child)
            if value:
                units.append(value)
        elif child.tag == word + "tbl":
            for row in child.findall("w:tr", namespace):
                cells = [text_of(cell) for cell in row.findall("w:tc", namespace)]
                if any(cells):
                    units.append(" | ".join(cells))
    return units or None


def _yen_values(value: str) -> tuple[int, ...]:
    values: list[int] = []
    for match in _YEN_AMOUNT.finditer(unicodedata.normalize("NFKC", value)):
        token = match["prefix"] or match["suffix"]
        parsed = int(token.replace(",", ""))
        if parsed > 0:
            values.append(parsed)
    return tuple(values)


def _tax_triples(values: Sequence[int]) -> set[tuple[int, int, int]]:
    triples: set[tuple[int, int, int]] = set()
    for net in values:
        for tax in values:
            if tax >= net:
                continue
            gross = net + tax
            if gross in values:
                triples.add((net, tax, gross))
    return triples


def _contract_tax_facts(
    source: Path | io.BytesIO,
) -> tuple[Fraction, dict[date, set[int]]] | None:
    units = _docx_source_units(source)
    if not units:
        return None
    rates: set[Fraction] = set()
    gross_by_date: dict[date, set[int]] = {}
    for unit in units:
        for capture in _PERCENT_RATE.findall(unit):
            rate = _source_decimal(capture)
            if rate is None or not 0 < rate <= 100:
                return None
            rates.add(Fraction(rate) / 100)
        triples = _tax_triples(_yen_values(unit))
        for net, tax, gross in triples:
            rates.add(Fraction(tax, net))
            dates = {
                parsed
                for token in re.findall(
                    r"20\d{2}(?:[-/.]\d{1,2}[-/.]\d{1,2}|年\d{1,2}月\d{1,2}日)",
                    unit,
                )
                if (parsed := _source_date(token)) is not None
            }
            for payment_date in dates:
                gross_by_date.setdefault(payment_date, set()).add(gross)
    if len(rates) != 1:
        return None
    rate = next(iter(rates))
    if not 0 < rate <= 1:
        return None
    return rate, gross_by_date


def _contract_payment_model(
    source: Path | io.BytesIO,
) -> tuple[str, int | None] | None:
    """Return an explicit fixed/T&M model and the unique T&M hourly rate."""

    units = _docx_source_units(source)
    if not units:
        return None
    joined = "\n".join(units)
    normalized = _normalized(joined)
    time_materials = (
        re.search(r"time\s*[_&]\s*materials", normalized, re.I) is not None
        or (
            "実績工数" in normalized
            and re.search(r"精算|算定|事後", normalized) is not None
        )
    )
    fixed = "固定価格" in normalized
    if time_materials == fixed:
        return None
    if fixed:
        return "fixed", None
    rates: set[int] = set()
    for unit in units:
        normalized_unit = unicodedata.normalize("NFKC", unit)
        rates.update(
            int(capture.replace(",", ""))
            for capture in _CONTRACT_HOURLY_RATE.findall(normalized_unit)
        )
    if len(rates) != 1:
        return None
    return "time_and_materials", next(iter(rates))


def _pptx_report_units(path: Path) -> tuple[str, ...] | None:
    """Read authored visible slide text, preserving one unit per slide."""

    try:
        if path.stat().st_size > _DIRECT_SOURCE_MAX_BYTES:
            return None
        from pptx import Presentation

        presentation = Presentation(str(path))
    except Exception:
        return None

    def hidden(shape: Any) -> bool | None:
        element = getattr(shape, "_element", None)
        if element is None:
            return None
        properties = next(
            (
                node
                for node in element.iter()
                if node.tag == "{" + _PRESENTATION_NAMESPACE + "}cNvPr"
            ),
            None,
        )
        if properties is None:
            return None
        value = _normalized(properties.get("hidden", "0"))
        if value not in {"", "0", "false", "1", "true"}:
            return None
        return value in {"1", "true"}

    units: list[str] = []
    try:
        for slide in presentation.slides:
            parts: list[str] = []

            def visit(shapes: Iterable[Any]) -> bool:
                for shape in shapes:
                    is_hidden = hidden(shape)
                    if is_hidden is None:
                        return False
                    if is_hidden:
                        continue
                    if getattr(shape, "has_text_frame", False):
                        value = unicodedata.normalize(
                            "NFKC", shape.text_frame.text
                        ).strip()
                        if value:
                            parts.append(value)
                    if getattr(shape, "has_table", False):
                        for row in shape.table.rows:
                            cells = [
                                unicodedata.normalize("NFKC", cell.text).strip()
                                for cell in row.cells
                            ]
                            if any(cells):
                                parts.append(" ".join(cells))
                    nested = getattr(shape, "shapes", None)
                    if nested is not None and not visit(nested):
                        return False
                return True

            if not visit(slide.shapes):
                return None
            if parts:
                units.append("\n".join(parts))
    except Exception:
        return None
    return tuple(units) if units else None


def _pdf_report_units(engine: Any, path: Path) -> tuple[str, ...] | None:
    """Read every PDF page, using hash-bound OCR only where text is absent."""

    try:
        from pypdf import PdfReader

        reader = PdfReader(str(path), strict=False)
        if reader.is_encrypted or not reader.pages:
            return None
        native: list[str] = []
        needs_ocr: set[int] = set()
        for page_number, page in enumerate(reader.pages, 1):
            try:
                text = unicodedata.normalize(
                    "NFKC", page.extract_text() or ""
                ).strip()
            except Exception:
                text = ""
            if len(re.sub(r"\s+", "", text)) >= 80:
                native.append(text)
            else:
                needs_ocr.add(page_number)
    except Exception:
        return None
    if not needs_ocr:
        return tuple(native) if native else None

    try:
        from pdf_visual_rules import _all_pdf_pages, _page_runs

        source_sha256 = hashlib.sha256(path.read_bytes()).hexdigest()
        pages = _all_pdf_pages(engine, path, source_sha256)
    except Exception:
        return None
    if not pages or len(pages) != len(reader.pages):
        return None
    runs: list[str] = list(native)
    for page in pages:
        if page.page_number not in needs_ocr:
            continue
        page_runs = _page_runs(page)
        if not page_runs:
            return None
        for run in page_runs:
            value = "\n".join(
                line.text
                for line in sorted(run, key=lambda item: item.sequence)
                if line.text.strip()
            )
            if value:
                runs.append(value)
    return tuple(runs) if runs else None


def _report_money_values(value: str) -> set[int]:
    values: set[int] = set()
    for match in _REPORT_MONEY.finditer(unicodedata.normalize("NFKC", value)):
        token = match["prefix"] or match["jpy_prefix"] or match["suffix"]
        parsed = int(token.replace(",", ""))
        if parsed > 0:
            values.add(parsed)
    return values


def _report_actual_candidates(
    value: str,
    hourly_rate: int,
    tax_rate: Fraction,
) -> tuple[set[tuple[Fraction, int, int, int, int]], bool, bool]:
    """Return arithmetic candidates, whether a claim was seen, and disclaimer."""

    normalized = unicodedata.normalize("NFKC", value)
    compact = _normalized(normalized)
    saw_claim = (
        re.search(r"(?:実績\s*)?工数", compact) is not None
        and "時間単価" in compact
        and re.search(r"税抜|税別", compact) is not None
        and "税込" in compact
        and re.search(r"最終請求|請求情報|費用・請求|精算", compact)
        is not None
    )
    if not saw_claim:
        return set(), False, False
    if _REPORT_NON_ACTUAL.search(normalized) is not None:
        return set(), True, True
    money = _report_money_values(normalized)
    if hourly_rate not in money:
        return set(), True, False
    hours: set[Fraction] = set()
    for capture in _REPORT_HOURS.findall(normalized):
        parsed = _source_decimal(capture)
        if parsed is not None and parsed > 0:
            hours.add(Fraction(parsed))
    requires_explicit_tax = "消費税額" in compact
    candidates: set[tuple[Fraction, int, int, int, int]] = set()
    for actual_hours in hours:
        net_value = actual_hours * hourly_rate
        tax_value = net_value * tax_rate
        gross_value = net_value + tax_value
        if any(value.denominator != 1 for value in (net_value, tax_value, gross_value)):
            continue
        net = net_value.numerator
        tax = tax_value.numerator
        gross = gross_value.numerator
        if net not in money or gross not in money:
            continue
        if requires_explicit_tax and tax not in money:
            continue
        candidates.add((actual_hours, hourly_rate, net, tax, gross))
    return candidates, True, False


def _certified_actual_gross_values(
    engine: Any,
    report: Path,
    hourly_rate: int,
    tax_rate: Fraction,
) -> tuple[int, ...] | None:
    if report.suffix.casefold() == ".pptx":
        units = _pptx_report_units(report)
    elif report.suffix.casefold() == ".pdf":
        units = _pdf_report_units(engine, report)
    else:
        return None
    if not units:
        return None
    candidates: set[tuple[Fraction, int, int, int, int]] = set()
    for unit in units:
        found, saw_claim, disclaimed = _report_actual_candidates(
            unit,
            hourly_rate,
            tax_rate,
        )
        if disclaimed:
            continue
        if saw_claim and not found:
            return None
        candidates.update(found)
    if len(candidates) > 1:
        return None
    if not candidates:
        return ()
    return (next(iter(candidates))[4],)


def _completed_payment_events_from_workbook(
    source: Path | io.BytesIO,
    contract_gross_by_date: Mapping[date, set[int]],
) -> tuple[tuple[tuple[date, int], ...], bool] | None:
    from openpyxl import load_workbook

    try:
        if isinstance(source, io.BytesIO):
            source.seek(0)
        workbook = load_workbook(source, read_only=True, data_only=True)
    except Exception:
        return None
    saw_payment = False
    amounts_by_date: dict[date, set[int]] = {}
    try:
        for worksheet in workbook.worksheets:
            if worksheet.sheet_state != "visible":
                continue
            for row in worksheet.iter_rows(values_only=True):
                cells = [_cell_text(value) for value in row if value is not None]
                normalized_cells = [_normalized(value) for value in cells]
                joined = " | ".join(cells)
                if _PAYMENT_MARKER.search(joined) is None:
                    continue
                row_dates = {
                    parsed
                    for value in row
                    if (parsed := _source_date(value)) is not None
                }
                if not row_dates:
                    continue
                saw_payment = True
                has_complete = "完了" in normalized_cells
                has_incomplete = any(
                    _PAYMENT_INCOMPLETE.fullmatch(value) is not None
                    for value in normalized_cells
                )
                if has_incomplete and has_complete:
                    return None
                if not has_complete:
                    continue
                row_amounts = set(_yen_values(joined))
                if len(row_amounts) > 1:
                    return None
                due_date = max(row_dates)
                if row_amounts:
                    amount = next(iter(row_amounts))
                else:
                    contract_amounts = {
                        amount
                        for row_date in row_dates
                        for amount in contract_gross_by_date.get(row_date, set())
                    }
                    # A bare completed phase or checkpoint may mention
                    # settlement without itself being a payment record.  The
                    # source must either share a contract date or declare its
                    # row type as the exact payment category before a missing
                    # amount can be joined to the contract table.
                    strong_without_amount = "支払" in normalized_cells
                    if not contract_amounts and strong_without_amount:
                        contract_amounts = {
                            amount
                            for amounts in contract_gross_by_date.values()
                            for amount in amounts
                        }
                    if not contract_amounts and not strong_without_amount:
                        continue
                    if len(contract_amounts) != 1:
                        return None
                    amount = next(iter(contract_amounts))
                amounts_by_date.setdefault(due_date, set()).add(amount)
    finally:
        workbook.close()
    if any(len(amounts) != 1 for amounts in amounts_by_date.values()):
        return None
    events = tuple(
        (payment_date, next(iter(amounts)))
        for payment_date, amounts in sorted(amounts_by_date.items())
    )
    return events, saw_payment


def _source_or_decrypted(
    engine: Any,
    project: Path,
    path: Path,
    primary_alias: str,
) -> tuple[Path | io.BytesIO, list[Path]] | None:
    if zipfile.is_zipfile(path):
        return path, [path]
    decrypted = _encrypted_workbook_bytes(engine, project, path, primary_alias)
    if decrypted is None:
        return None
    data, evidence = decrypted
    return io.BytesIO(data), evidence


def _all_project_paid_gross_tax_sum(
    engine: Any,
) -> StructuredCandidateDecision | None:
    projects = _all_project_directories(engine)
    if projects is None:
        return None
    total_tax = 0
    completed_count = 0
    evidence: list[Path] = []
    for project in projects:
        alias_resolution = _unique_primary_alias_for_project(
            engine.glossary,
            project.name,
        )
        plan_path = _authoritative_project_plan_workbook(project)
        contract_path = _authoritative_project_contract(project)
        if alias_resolution is None or plan_path is None or contract_path is None:
            return None
        alias, _ = alias_resolution
        contract_source = _source_or_decrypted(
            engine,
            project,
            contract_path,
            alias,
        )
        plan_source = _source_or_decrypted(
            engine,
            project,
            plan_path,
            alias,
        )
        if contract_source is None or plan_source is None:
            return None
        contract_stream, contract_evidence = contract_source
        plan_stream, plan_evidence = plan_source
        contract_facts = _contract_tax_facts(contract_stream)
        if contract_facts is None:
            return None
        tax_rate, gross_by_date = contract_facts
        payment_model = _contract_payment_model(contract_stream)
        if payment_model is None:
            return None
        paid = _completed_payment_events_from_workbook(plan_stream, gross_by_date)
        if paid is None:
            return None
        events, saw_payment = paid
        if not saw_payment or not events:
            return None
        project_gross = [gross for _, gross in events]
        pricing_model, hourly_rate = payment_model
        if pricing_model == "time_and_materials":
            report_path = _authoritative_project_final_report(project)
            if report_path is None or hourly_rate is None:
                return None
            actual_values = _certified_actual_gross_values(
                engine,
                report_path,
                hourly_rate,
                tax_rate,
            )
            if actual_values is None:
                return None
            if actual_values:
                project_gross = list(actual_values)
            evidence.append(report_path)
        for gross in project_gross:
            tax = Fraction(gross) * tax_rate / (1 + tax_rate)
            if tax.denominator != 1 or tax.numerator <= 0:
                return None
            total_tax += tax.numerator
            completed_count += 1
        evidence.extend(contract_evidence)
        evidence.extend(plan_evidence)
    if completed_count == 0 or total_tax <= 0:
        return None
    return _decision(f"{total_tax:,}円", evidence, engine.source_root, 12)


def _missing_rows(engine: Any) -> StructuredCandidateDecision | None:
    counts: list[tuple[int, str, Path]] = []
    for path in _all_files(engine, {".csv"}):
        relative = unicodedata.normalize("NFC", path.relative_to(engine.source_root).as_posix())
        if not re.search(r"^プロジェクト/[^/]+/03\.データ/train\.csv$", relative):
            continue
        with path.open(encoding="utf-8-sig", newline="") as handle:
            rows = csv.reader(handle)
            header = next(rows, None)
            if not header:
                continue
            count = 0
            for row in rows:
                padded = row + [""] * max(0, len(header) - len(row))
                if any(
                    not value.strip() or value.strip().casefold() in {"na", "nan", "null", "none"}
                    for value in padded[: len(header)]
                ):
                    count += 1
        project = relative.split("/")[1]
        counts.append((count, project, path))
    if not counts:
        return None
    maximum = max(value[0] for value in counts)
    winners = [value for value in counts if value[0] == maximum]
    if len(winners) != 1:
        return None
    alias = _primary_alias(engine.glossary, winners[0][1])
    if alias is None:
        return None
    return _decision(alias, [winners[0][2]], engine.source_root, 2)


def _project_train_csv(engine: Any, location: str) -> Path | None:
    matches = []
    for path in _all_files(engine, {".csv"}):
        if not _project_matches(engine, path, location):
            continue
        relative = unicodedata.normalize("NFC", path.relative_to(engine.source_root).as_posix())
        if relative.endswith("/03.データ/train.csv"):
            matches.append(path)
    return matches[0] if len(matches) == 1 else None


def _csv_records(path: Path) -> tuple[list[str], list[dict[str, str]]] | None:
    with path.open(encoding="utf-8-sig", newline="") as handle:
        reader = csv.DictReader(handle)
        if not reader.fieldnames:
            return None
        return list(reader.fieldnames), [dict(row) for row in reader]


def _standardized_share(engine: Any, match: re.Match[str]) -> StructuredCandidateDecision | None:
    path = _project_train_csv(engine, match["location"])
    loaded = _csv_records(path) if path else None
    if loaded is None:
        return None
    headers, rows = loaded
    measure = match["measure"]
    category = match["category"]
    if measure not in headers or category not in headers:
        return None
    try:
        values = [float(row[measure]) for row in rows]
    except ValueError:
        return None
    mean = statistics.fmean(values)
    deviation = statistics.pstdev(values)
    if deviation == 0:
        return None
    category_values = [
        float(row[measure]) for row in rows if row[category] == match["value"]
    ]
    if not category_values:
        return None
    category_mean = statistics.fmean(category_values)
    base = [
        row
        for row, value in zip(rows, values)
        if (value - mean) / deviation < float(match["z"])
    ]
    if not base:
        return None
    numerator = sum(
        row[category] == match["value"] and float(row[measure]) > category_mean
        for row in base
    )
    digits = int(match["digits"])
    value = Decimal(numerator * 100) / Decimal(len(base))
    quantizer = Decimal(1).scaleb(-digits)
    rendered = format(value.quantize(quantizer, rounding=ROUND_HALF_UP), f".{digits}f") + "%"
    return _decision(rendered, [path], engine.source_root, 6)


def _analysis_files(engine: Any, location: str, name: str) -> list[Path]:
    return [
        path
        for path in _all_files(engine)
        if _project_matches(engine, path, location)
        and path.name == name
        and "04.分析" in unicodedata.normalize("NFC", path.as_posix())
    ]


def _interaction_columns(engine: Any, match: re.Match[str]) -> StructuredCandidateDecision | None:
    metrics = [
        path for path in _analysis_files(engine, match["location"], "metrics.json")
        if "analysis_outputs" in path.parts and "analysis_project" not in path.parts
    ]
    features = _analysis_files(engine, match["location"], "features.py")
    if len(metrics) != 1 or len(features) != 1:
        return None
    data = json.loads(metrics[0].read_text(encoding="utf-8"))
    selected = data.get("feature_selection", {}).get("selected_columns")
    if not isinstance(selected, list) or not all(isinstance(value, str) for value in selected):
        return None
    source = features[0].read_text(encoding="utf-8")
    if 'feature_name = f"{left}__x__{right}"' not in source:
        return None
    answers = [value for value in selected if "__x__" in value]
    if not answers:
        return None
    return _decision("、".join(answers), [metrics[0], features[0]], engine.source_root, 3)


def _literal_from_call(node: ast.AST, function: str, fallback_index: int) -> Any | None:
    for child in ast.walk(node):
        if not isinstance(child, ast.Call):
            continue
        name = child.func.id if isinstance(child.func, ast.Name) else (
            child.func.attr if isinstance(child.func, ast.Attribute) else None
        )
        if name != function or len(child.args) <= fallback_index:
            continue
        try:
            return ast.literal_eval(child.args[fallback_index])
        except Exception:
            continue
    return None


def _gb_params(engine: Any, match: re.Match[str]) -> StructuredCandidateDecision | None:
    configs = _analysis_files(engine, match["location"], "project_config.json")
    modeling = _analysis_files(engine, match["location"], "modeling.py")
    data_path = _project_train_csv(engine, match["location"])
    if len(configs) != 1 or len(modeling) != 1 or data_path is None:
        return None
    config = json.loads(configs[0].read_text(encoding="utf-8"))
    if config.get("model_type") != "gradient_boosting":
        return None
    loaded = _csv_records(data_path)
    if loaded is None or config.get("target_column") not in loaded[0]:
        return None
    target = str(config["target_column"])
    unique = {row[target] for row in loaded[1]}
    classification = 1 < len(unique) <= 20
    if not classification:
        return None
    tree = ast.parse(modeling[0].read_text(encoding="utf-8"))
    n_estimators = _literal_from_call(tree, "to_int", 1)
    learning_rate = _literal_from_call(tree, "to_float", 1)
    random_state = config.get("random_state", 42)
    params = config.get("model_params") or {}
    n_estimators = params.get("n_estimators", n_estimators)
    learning_rate = params.get("learning_rate", learning_rate)
    if n_estimators is None or learning_rate is None:
        return None
    answer = (
        f"n_estimators: {n_estimators}、learning_rate: {learning_rate}、"
        f"random_state: {random_state}"
    )
    return _decision(answer, [configs[0], modeling[0], data_path], engine.source_root, 4)


def _ooxml_relationship_target(
    archive: zipfile.ZipFile,
    owner_part: str,
    relationship_id: str,
    relationship_suffix: str,
) -> str | None:
    """Resolve one internal OOXML relationship without accepting path escape."""

    owner = PurePosixPath(owner_part)
    rels_part = str(owner.parent / "_rels" / f"{owner.name}.rels")
    try:
        relationships = ET.fromstring(archive.read(rels_part))
    except (KeyError, ET.ParseError):
        return None
    matches = [
        relation
        for relation in relationships
        if relation.tag.endswith("}Relationship")
        and relation.get("Id") == relationship_id
        and relation.get("TargetMode") != "External"
        and str(relation.get("Type") or "").endswith(relationship_suffix)
    ]
    if len(matches) != 1:
        return None
    target = _package_target_part(owner_part, str(matches[0].get("Target") or ""))
    if target is None or target not in archive.namelist():
        return None
    return target


def _chart_series_label(chart_root: ET.Element) -> str | None:
    series = [
        element
        for element in chart_root.iter()
        if element.tag.rsplit("}", 1)[-1] == "series"
    ]
    if len(series) != 1:
        return None
    tx_nodes = [
        child
        for child in list(series[0])
        if child.tag.rsplit("}", 1)[-1] == "tx"
    ]
    if len(tx_nodes) != 1:
        return None
    formulas = [
        value.text.strip()
        for value in tx_nodes[0].iter()
        if value.tag.rsplit("}", 1)[-1] == "f"
        and isinstance(value.text, str)
        and value.text.strip()
    ]
    # ChartEx stores its internal data-schema field under the _xlchart
    # namespace.  A conventional series title may be arbitrary display text,
    # so it is not accepted as proof of a worksheet column here.
    if len(formulas) != 1 or not formulas[0].startswith("_xlchart."):
        return None
    labels = []
    for value in tx_nodes[0].iter():
        if value.tag.rsplit("}", 1)[-1] != "v" or value.text is None:
            continue
        text = unicodedata.normalize("NFC", value.text).strip()
        if text:
            labels.append(text)
    labels = list(dict.fromkeys(labels))
    return labels[0] if len(labels) == 1 else None


def _xlsx_chart_series_column(
    engine: Any,
    match: re.Match[str],
) -> StructuredCandidateDecision | None:
    paths = _named_paths(engine, match["location"], match["container"], {".xlsx"})
    if len(paths) != 1 or paths[0].stat().st_size > _DIRECT_SOURCE_MAX_BYTES:
        return None
    workbook_path = paths[0]
    office_rel = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    expected_sheet = _normalized(match["sheet"])
    chart_number = int(unicodedata.normalize("NFKC", match["chart"]))
    expected_chart = _normalized(f"グラフ {chart_number}").replace(" ", "")
    try:
        with zipfile.ZipFile(workbook_path) as archive:
            infos = archive.infolist()
            if len(infos) > 20_000 or sum(item.file_size for item in infos) > 1_024**3:
                return None
            workbook = ET.fromstring(archive.read("xl/workbook.xml"))
            sheet_nodes = [
                sheet
                for sheet in workbook.iter()
                if sheet.tag.rsplit("}", 1)[-1] == "sheet"
                and _normalized(sheet.get("name")) == expected_sheet
            ]
            if len(sheet_nodes) != 1:
                return None
            sheet_relation = sheet_nodes[0].get("{" + office_rel + "}id")
            if not sheet_relation:
                return None
            sheet_part = _ooxml_relationship_target(
                archive,
                "xl/workbook.xml",
                sheet_relation,
                "/worksheet",
            )
            if sheet_part is None:
                return None
            sheet_root = ET.fromstring(archive.read(sheet_part))
            drawing_nodes = [
                node
                for node in sheet_root.iter()
                if node.tag.rsplit("}", 1)[-1] == "drawing"
            ]
            if len(drawing_nodes) != 1:
                return None
            drawing_relation = drawing_nodes[0].get("{" + office_rel + "}id")
            if not drawing_relation:
                return None
            drawing_part = _ooxml_relationship_target(
                archive,
                sheet_part,
                drawing_relation,
                "/drawing",
            )
            if drawing_part is None:
                return None
            drawing_root = ET.fromstring(archive.read(drawing_part))
            anchors = []
            for anchor in list(drawing_root):
                names = [
                    _normalized(node.get("name")).replace(" ", "")
                    for node in anchor.iter()
                    if node.tag.rsplit("}", 1)[-1] == "cNvPr"
                    and node.get("name")
                ]
                if names == [expected_chart]:
                    anchors.append(anchor)
            if len(anchors) != 1:
                return None
            chart_nodes = [
                node
                for node in anchors[0].iter()
                if node.tag.rsplit("}", 1)[-1] == "chart"
                and node.get("{" + office_rel + "}id")
            ]
            if len(chart_nodes) != 1:
                return None
            chart_relation = chart_nodes[0].get("{" + office_rel + "}id")
            chart_parts = [
                _ooxml_relationship_target(
                    archive, drawing_part, chart_relation, suffix
                )
                for suffix in ("/chartEx", "/chart")
            ]
            chart_parts = list(dict.fromkeys(part for part in chart_parts if part))
            if len(chart_parts) != 1:
                return None
            chart_root = ET.fromstring(archive.read(chart_parts[0]))
            label = _chart_series_label(chart_root)
            if label is None:
                return None
    except (KeyError, OSError, ET.ParseError, zipfile.BadZipFile, ValueError):
        return None
    return _decision(label, [workbook_path], engine.source_root, 6)


def _regression_prediction(engine: Any, match: re.Match[str]) -> StructuredCandidateDecision | None:
    paths = _named_paths(engine, match["location"], match["container"], {".xlsx"})
    if len(paths) != 1:
        return None
    from openpyxl import load_workbook

    workbook = load_workbook(paths[0], read_only=True, data_only=True)
    try:
        if "回帰分析" not in workbook.sheetnames or "train" not in workbook.sheetnames:
            return None
        regression = workbook["回帰分析"]
        coefficient_column = None
        header_row = None
        for row in regression.iter_rows():
            for cell in row:
                if cell.value == "係数":
                    coefficient_column = cell.column
                    header_row = cell.row
                    break
            if coefficient_column:
                break
        if coefficient_column is None or header_row is None:
            return None
        coefficients: dict[str, Decimal] = {}
        for row in range(header_row + 1, regression.max_row + 1):
            label = regression.cell(row, 1).value
            value = regression.cell(row, coefficient_column).value
            if label is not None and value is not None:
                coefficients[str(label)] = Decimal(str(value))
        intercept = coefficients.pop("切片", None)
        if intercept is None or not coefficients:
            return None
        train = workbook["train"]
        row_iterator = train.iter_rows(values_only=True)
        first_row = next(row_iterator, None)
        if first_row is None:
            return None
        headers = [str(value) if value is not None else "" for value in first_row]
        if match["id_field"] not in headers or any(field not in headers for field in coefficients):
            return None
        columns = {field: headers.index(field) for field in coefficients}
        id_column = headers.index(match["id_field"])
        target_found = False
        numeric: dict[str, list[Decimal]] = {field: [] for field in coefficients}
        target_values: dict[str, Decimal] = {}
        for row in row_iterator:
            if not any(value is not None for value in row):
                continue
            if str(row[id_column]) == match["id_value"]:
                target_found = True
            for field, column in columns.items():
                value = row[column]
                if value is None:
                    return None
                decimal_value = Decimal(str(value))
                numeric[field].append(decimal_value)
                if str(row[id_column]) == match["id_value"]:
                    target_values[field] = decimal_value
        if not target_found:
            return None
    finally:
        workbook.close()
    prediction = intercept
    for field, coefficient in coefficients.items():
        values = numeric[field]
        mean = sum(values) / Decimal(len(values))
        variance = sum((value - mean) ** 2 for value in values) / Decimal(len(values))
        deviation = variance.sqrt()
        if deviation == 0:
            return None
        prediction += coefficient * ((target_values[field] - mean) / deviation)
    digits = int(match["digits"])
    quantizer = Decimal(1).scaleb(-digits)
    answer = format(prediction.quantize(quantizer, rounding=ROUND_HALF_UP), f".{digits}f")
    return _decision(answer, paths, engine.source_root, 5)


def _negative_correlation(engine: Any, match: re.Match[str]) -> StructuredCandidateDecision | None:
    path = _project_train_csv(engine, match["location"])
    configs = _analysis_files(engine, match["location"], "project_config.json")
    loaded = _csv_records(path) if path else None
    if loaded is None or len(configs) != 1:
        return None
    headers, rows = loaded
    target = json.loads(configs[0].read_text(encoding="utf-8")).get("target_column")
    if target not in headers:
        return None
    numeric: dict[str, list[float]] = {}
    for field in headers:
        try:
            numeric[field] = [float(row[field]) for row in rows]
        except (TypeError, ValueError):
            continue
    if target not in numeric:
        return None
    target_values = numeric[target]
    correlations: list[tuple[float, str]] = []
    for field, values in numeric.items():
        if field == target or len(values) != len(target_values):
            continue
        x_mean = statistics.fmean(values)
        y_mean = statistics.fmean(target_values)
        numerator = sum((x - x_mean) * (y - y_mean) for x, y in zip(values, target_values))
        left = math.sqrt(sum((x - x_mean) ** 2 for x in values))
        right = math.sqrt(sum((y - y_mean) ** 2 for y in target_values))
        if left and right:
            correlations.append((numerator / (left * right), field))
    negative = [item for item in correlations if item[0] < 0]
    if not negative:
        return None
    negative.sort()
    if len(negative) > 1 and math.isclose(negative[0][0], negative[1][0], abs_tol=1e-15):
        return None
    return _decision(negative[0][1], [path, configs[0]], engine.source_root, 3)


def _rgb(cell: Any) -> str | None:
    color = cell.fill.fgColor
    return color.rgb if cell.fill.fill_type == "solid" and color.type == "rgb" else None


def _unique_workbook(engine: Any, location: str, container: str) -> Path | None:
    paths = _named_paths(engine, location, container, {".xlsx"})
    return paths[0] if len(paths) == 1 else None


def _header_position(worksheet: Any, field: str) -> tuple[int, int] | None:
    matches = [
        (cell.row, cell.column)
        for row in worksheet.iter_rows(min_row=1, max_row=min(20, worksheet.max_row))
        for cell in row
        if str(cell.value).strip() == field
    ]
    return matches[0] if len(matches) == 1 else None


def _highlight_rows(engine: Any, match: re.Match[str]) -> StructuredCandidateDecision | None:
    path = _unique_workbook(engine, match["location"], match["container"])
    if path is None:
        return None
    from openpyxl import load_workbook

    workbook = load_workbook(path, read_only=False, data_only=True)
    try:
        if match["sheet"]:
            sheets = [ws for ws in workbook.worksheets if _normalized(ws.title) == _normalized(match["sheet"])]
        else:
            sheets = list(workbook.worksheets)
        answers: list[str] = []
        for worksheet in sheets:
            position = _header_position(worksheet, match["target"])
            if position is None:
                continue
            header_row, target_column = position
            highlighted_rows = sorted(
                {
                    cell.row
                    for row in worksheet.iter_rows(min_row=header_row + 1)
                    for cell in row
                    if _rgb(cell) in ORANGE_RGB
                }
            )
            for row in highlighted_rows:
                value = worksheet.cell(row, target_column).value
                if value is not None:
                    answers.append(str(value))
    finally:
        workbook.close()
    if not answers or len(answers) != len(set(answers)):
        return None
    return _decision("、".join(answers), [path], engine.source_root, 2)


def _blue_sum(engine: Any, match: re.Match[str]) -> StructuredCandidateDecision | None:
    path = _unique_workbook(engine, match["location"], match["container"])
    if path is None:
        return None
    from openpyxl import load_workbook

    workbook = load_workbook(path, read_only=False, data_only=True)
    values: list[Decimal] = []
    try:
        for worksheet in workbook.worksheets:
            for row in worksheet.iter_rows():
                for cell in row:
                    if _rgb(cell) not in BLUE_RGB or cell.value is None:
                        continue
                    if isinstance(cell.value, bool) or not isinstance(cell.value, (int, float, Decimal)):
                        return None
                    values.append(Decimal(str(cell.value)))
    finally:
        workbook.close()
    if not values:
        return None
    total = sum(values)
    answer = str(total.quantize(Decimal("1"), rounding=ROUND_HALF_UP))
    return _decision(answer, [path], engine.source_root, 2)


def _yellow_intersection(engine: Any, match: re.Match[str]) -> StructuredCandidateDecision | None:
    path = _unique_workbook(engine, match["location"], match["container"])
    if path is None:
        return None
    from openpyxl import load_workbook

    workbook = load_workbook(path, read_only=False, data_only=True)
    intersections: list[Decimal] = []
    try:
        for worksheet in workbook.worksheets:
            row_yellow: Counter[int] = Counter()
            row_values: Counter[int] = Counter()
            column_yellow: Counter[int] = Counter()
            column_values: Counter[int] = Counter()
            for row in worksheet.iter_rows():
                for cell in row:
                    if cell.value is not None:
                        row_values[cell.row] += 1
                        column_values[cell.column] += 1
                    if _rgb(cell) in YELLOW_RGB:
                        row_yellow[cell.row] += 1
                        column_yellow[cell.column] += 1
            full_rows = {
                row for row, count in row_yellow.items()
                if count > 1 and count >= 0.9 * row_values[row]
            }
            full_columns = {
                column for column, count in column_yellow.items()
                if count > 1 and count >= 0.9 * column_values[column]
            }
            for row in full_rows:
                for column in full_columns:
                    cell = worksheet.cell(row, column)
                    if _rgb(cell) in YELLOW_RGB and isinstance(cell.value, (int, float, Decimal)) and not isinstance(cell.value, bool):
                        intersections.append(Decimal(str(cell.value)))
    finally:
        workbook.close()
    if len(intersections) != 2:
        return None
    difference = abs(intersections[0] - intersections[1])
    answer = format(difference, "f").rstrip("0").rstrip(".")
    return _decision(answer, [path], engine.source_root, 2)


_GENERIC_FIELD_ALIASES: Mapping[str, tuple[str, ...]] = {
    "年齢": ("age",),
    "性別": ("gender", "sex"),
}
_GENERIC_COHORT_ALIASES: Mapping[str, tuple[str, ...]] = {
    "女性": ("female", "woman", "women"),
    "男性": ("male", "man", "men"),
}


def _lexical_candidates(value: str, glossary: Any, aliases: Mapping[str, tuple[str, ...]]) -> set[str]:
    candidates = set(_candidate_values(value, glossary))
    candidates.update(aliases.get(value, ()))
    return {_normalized(candidate) for candidate in candidates}


def _resolve_header(headers: Sequence[str], surface: str, glossary: Any) -> str | None:
    candidates = _lexical_candidates(surface, glossary, _GENERIC_FIELD_ALIASES)
    matches = [header for header in headers if _normalized(header) in candidates]
    return matches[0] if len(matches) == 1 else None


def _resolve_categorical_value(
    records: Sequence[Mapping[str, str]],
    headers: Sequence[str],
    surface: str,
    glossary: Any,
    excluded_fields: set[str],
) -> tuple[str, str] | None:
    candidates = _lexical_candidates(surface, glossary, _GENERIC_COHORT_ALIASES)
    matches: list[tuple[str, str]] = []
    for field in headers:
        if field in excluded_fields:
            continue
        values = {str(row[field]) for row in records if str(row[field]).strip()}
        for value in values:
            if _normalized(value) in candidates:
                matches.append((field, value))
    unique = list(dict.fromkeys(matches))
    return unique[0] if len(unique) == 1 else None


def _render_decimal(value: Decimal) -> str:
    rendered = format(value, "f")
    return rendered.rstrip("0").rstrip(".") if "." in rendered else rendered


def _cohort_group_mean_argmax(
    engine: Any,
    match: re.Match[str],
) -> StructuredCandidateDecision | None:
    paths = _named_paths(
        engine,
        match["location"],
        match["container"],
        {".csv", ".xlsx"},
    )
    resolved = _unique_table(paths, [match["filter_field"], match["measure"]])
    if resolved is None:
        return None
    table, sources = resolved
    records = _records(table)
    group_field = _resolve_header(table.headers, match["group_field"], engine.glossary)
    if group_field is None:
        return None
    cohort = _resolve_categorical_value(
        records,
        table.headers,
        match["cohort"],
        engine.glossary,
        {match["filter_field"], match["measure"], group_field},
    )
    if cohort is None:
        return None
    cohort_field, cohort_value = cohort
    filtered = [
        row
        for row in records
        if _normalized(row[match["filter_field"]]) == _normalized(match["filter_value"])
        and _normalized(row[cohort_field]) == _normalized(cohort_value)
    ]
    if not filtered:
        return None
    grouped: dict[Decimal, list[Decimal]] = {}
    try:
        for row in filtered:
            group = Decimal(row[group_field])
            grouped.setdefault(group, []).append(Decimal(row[match["measure"]]))
    except Exception:
        return None
    means = {group: sum(values) / Decimal(len(values)) for group, values in grouped.items()}
    maximum = max(means.values())
    winners = [group for group, mean in means.items() if mean == maximum]
    if len(winners) != 1:
        return None
    return _decision(
        _render_decimal(winners[0]) + match["unit"],
        sources,
        engine.source_root,
        7,
    )


def _project_training_data(engine: Any, location: str) -> list[Path]:
    paths: list[Path] = []
    for path in _all_files(engine, {".csv", ".xlsx"}):
        if not _project_matches(engine, path, location):
            continue
        relative = unicodedata.normalize("NFC", path.relative_to(engine.source_root).as_posix())
        if relative.endswith("/03.データ/train.csv") or relative.endswith("/03.データ/train.xlsx"):
            paths.append(path)
    return paths


def _multi_filter_mean_half_up(
    engine: Any,
    match: re.Match[str],
) -> StructuredCandidateDecision | None:
    filters = _parse_equality_filters(match["filters"])
    if filters is None:
        return None
    required = [field for field, _ in filters] + [match["measure"]]
    resolved = _unique_table(_project_training_data(engine, match["location"]), required)
    if resolved is None:
        return None
    table, sources = resolved
    rows = [
        row
        for row in _records(table)
        if all(_normalized(row[field]) == _normalized(value) for field, value in filters)
    ]
    if not rows:
        return None
    try:
        values = [Decimal(row[match["measure"]]) for row in rows]
    except Exception:
        return None
    answer = str(
        (sum(values) / Decimal(len(values))).quantize(Decimal("1"), rounding=ROUND_HALF_UP)
    )
    return _decision(answer, sources, engine.source_root, len(filters) + 3)


def _call_attributes(node: ast.AST) -> set[str]:
    return {
        child.func.attr
        for child in ast.walk(node)
        if isinstance(child, ast.Call) and isinstance(child.func, ast.Attribute)
    }


def _descending_sort_is_explicit(node: ast.AST) -> bool:
    for child in ast.walk(node):
        if not isinstance(child, ast.Call) or not isinstance(child.func, ast.Attribute):
            continue
        if child.func.attr != "sort_values":
            continue
        for keyword in child.keywords:
            if keyword.arg == "ascending" and isinstance(keyword.value, ast.Constant):
                return keyword.value.value is False
    return False


def _head_index_list(node: ast.AST) -> tuple[str, int] | None:
    if not isinstance(node, ast.Call) or not isinstance(node.func, ast.Attribute):
        return None
    if node.func.attr != "tolist" or not isinstance(node.func.value, ast.Attribute):
        return None
    index = node.func.value
    if index.attr != "index" or not isinstance(index.value, ast.Call):
        return None
    head = index.value
    if not isinstance(head.func, ast.Attribute) or head.func.attr != "head":
        return None
    if not isinstance(head.func.value, ast.Name) or len(head.args) != 1:
        return None
    if not isinstance(head.args[0], ast.Constant) or not isinstance(head.args[0].value, int):
        return None
    return head.func.value.id, head.args[0].value


def _subscript_name(node: ast.AST) -> str | None:
    if not isinstance(node, ast.Subscript):
        return None
    slice_node = node.slice
    return slice_node.id if isinstance(slice_node, ast.Name) else None


def _notebook_heatmap_top_n(notebook: Path) -> int | None:
    try:
        document = json.loads(notebook.read_text(encoding="utf-8"))
    except Exception:
        return None
    candidates: set[int] = set()
    for cell in document.get("cells", []):
        if cell.get("cell_type") != "code":
            continue
        source = "".join(cell.get("source", []))
        try:
            tree = ast.parse(source)
        except SyntaxError:
            continue
        assignments: dict[str, list[ast.AST]] = {}
        for child in ast.walk(tree):
            if not isinstance(child, (ast.Assign, ast.AnnAssign)):
                continue
            targets = child.targets if isinstance(child, ast.Assign) else [child.target]
            value = child.value
            for target in targets:
                if isinstance(target, ast.Name) and value is not None:
                    assignments.setdefault(target.id, []).append(value)
        for top_name, values in assignments.items():
            for value in values:
                head_spec = _head_index_list(value)
                if head_spec is None:
                    continue
                rank_name, top_n = head_spec
                rank_values = assignments.get(rank_name, [])
                if not any(
                    {"corrwith", "abs", "sort_values"} <= _call_attributes(rank_value)
                    and _descending_sort_is_explicit(rank_value)
                    for rank_value in rank_values
                ):
                    continue
                corr_names = {
                    name
                    for name, corr_values in assignments.items()
                    if any(
                        isinstance(corr_value, ast.Call)
                        and isinstance(corr_value.func, ast.Attribute)
                        and corr_value.func.attr == "corr"
                        and _subscript_name(corr_value.func.value) == top_name
                        for corr_value in corr_values
                    )
                }
                heatmap_inputs = {
                    call.args[0].id
                    for call in ast.walk(tree)
                    if isinstance(call, ast.Call)
                    and isinstance(call.func, ast.Attribute)
                    and call.func.attr == "heatmap"
                    and call.args
                    and isinstance(call.args[0], ast.Name)
                }
                if top_n > 0 and corr_names & heatmap_inputs:
                    candidates.add(top_n)
    return next(iter(candidates)) if len(candidates) == 1 else None


def _notebook_declared_correlation_target(notebook: Path) -> str | None:
    try:
        document = json.loads(notebook.read_text(encoding="utf-8"))
    except Exception:
        return None
    literal_assignments: dict[str, set[str]] = {}
    references: list[tuple[str, str]] = []
    for cell in document.get("cells", []):
        if cell.get("cell_type") != "code":
            continue
        try:
            tree = ast.parse("".join(cell.get("source", [])))
        except SyntaxError:
            continue
        for node in ast.walk(tree):
            if isinstance(node, ast.Assign) and isinstance(node.value, ast.Constant):
                if isinstance(node.value.value, str):
                    for target in node.targets:
                        if isinstance(target, ast.Name):
                            literal_assignments.setdefault(target.id, set()).add(node.value.value)
            if not isinstance(node, ast.Call) or not isinstance(node.func, ast.Attribute):
                continue
            if node.func.attr != "corrwith" or len(node.args) != 1:
                continue
            argument = node.args[0]
            if not isinstance(argument, ast.Subscript):
                continue
            if isinstance(argument.slice, ast.Name):
                references.append(("variable", argument.slice.id))
            elif isinstance(argument.slice, ast.Constant) and isinstance(argument.slice.value, str):
                references.append(("literal", argument.slice.value))
    targets: set[str] = set()
    for kind, value in references:
        if kind == "literal":
            targets.add(value)
        else:
            targets.update(literal_assignments.get(value, ()))
    return next(iter(targets)) if len(targets) == 1 else None


def _notebook_data_paths(engine: Any, notebook: Path) -> list[Path]:
    try:
        document = json.loads(notebook.read_text(encoding="utf-8"))
    except Exception:
        return []
    literals: set[str] = set()
    for cell in document.get("cells", []):
        if cell.get("cell_type") != "code":
            continue
        try:
            tree = ast.parse("".join(cell.get("source", [])))
        except SyntaxError:
            continue
        for node in ast.walk(tree):
            if isinstance(node, ast.Constant) and isinstance(node.value, str):
                if Path(node.value).suffix.casefold() in {".csv", ".xlsx"}:
                    literals.add(node.value)
    matches: set[Path] = set()
    root = engine.source_root.resolve()
    for literal in literals:
        relative = Path(literal)
        if relative.is_absolute() or ".." in relative.parts:
            continue
        ancestor = notebook.parent
        while True:
            candidate = (ancestor / relative).resolve()
            if candidate.is_relative_to(root) and candidate.is_file() and not candidate.is_symlink():
                matches.add(candidate)
            if ancestor.resolve() == root or root not in ancestor.resolve().parents:
                break
            ancestor = ancestor.parent
    return sorted(matches, key=lambda path: path.as_posix())


def _numeric_columns(
    headers: Sequence[str],
    rows: Sequence[Mapping[str, str]],
) -> dict[str, list[float | None]]:
    columns: dict[str, list[float | None]] = {}
    for field in headers:
        values: list[float | None] = []
        valid = True
        for row in rows:
            raw = str(row[field]).strip()
            if not raw:
                values.append(None)
                continue
            try:
                value = float(raw)
            except ValueError:
                valid = False
                break
            if not math.isfinite(value):
                values.append(None)
            else:
                values.append(value)
        if valid and sum(value is not None for value in values) >= 2:
            columns[field] = values
    return columns


def _pearson(left: Sequence[float | None], right: Sequence[float | None]) -> float | None:
    pairs = [(x, y) for x, y in zip(left, right) if x is not None and y is not None]
    if len(pairs) < 2:
        return None
    xs = [pair[0] for pair in pairs]
    ys = [pair[1] for pair in pairs]
    x_mean = statistics.fmean(xs)
    y_mean = statistics.fmean(ys)
    numerator = sum((x - x_mean) * (y - y_mean) for x, y in pairs)
    denominator = math.sqrt(
        sum((x - x_mean) ** 2 for x in xs) * sum((y - y_mean) ** 2 for y in ys)
    )
    return numerator / denominator if denominator else None


def _notebook_heatmap_min_abs_correlation(
    engine: Any,
    match: re.Match[str],
) -> StructuredCandidateDecision | None:
    notebooks = _named_paths(engine, match["location"], match["container"], {".ipynb"})
    if len(notebooks) != 1:
        return None
    notebook = notebooks[0]
    top_n = _notebook_heatmap_top_n(notebook)
    declared_target = _notebook_declared_correlation_target(notebook)
    data_paths = _notebook_data_paths(engine, notebook)
    if (
        top_n is None
        or declared_target != match["target"]
        or len(data_paths) != 1
    ):
        return None
    data_path = data_paths[0]
    if data_path.suffix.casefold() == ".csv":
        loaded = _csv_records(data_path)
        if loaded is None:
            return None
        headers, rows = loaded
    else:
        resolved = _unique_table([data_path], [match["target"]])
        if resolved is None:
            return None
        table, _ = resolved
        headers, rows = list(table.headers), _records(table)
    target = match["target"]
    numeric = _numeric_columns(headers, rows)
    if target not in numeric:
        return None
    correlations: list[tuple[str, float]] = []
    for field in headers:
        if field == target or field not in numeric:
            continue
        value = _pearson(numeric[field], numeric[target])
        if value is not None:
            correlations.append((field, abs(value)))
    correlations.sort(key=lambda item: item[1], reverse=True)
    visible = correlations[:top_n]
    if not visible:
        return None
    minimum = min(value for _, value in visible)
    winners = [field for field, value in visible if math.isclose(value, minimum, abs_tol=1e-15)]
    if len(winners) != 1:
        return None
    return _decision(winners[0], [notebook, data_path], engine.source_root, 7)


def _docx_paragraphs(path: Path) -> list[str] | None:
    try:
        with zipfile.ZipFile(path) as archive:
            root = ET.fromstring(archive.read("word/document.xml"))
    except Exception:
        return None
    namespace = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    body = root.find("w:body", namespace)
    if body is None:
        return None
    paragraphs: list[str] = []
    paragraph_tag = "{" + namespace["w"] + "}p"
    for child in body:
        if child.tag != paragraph_tag:
            continue
        text = "".join(node.text or "" for node in child.findall(".//w:t", namespace)).strip()
        if text:
            paragraphs.append(text)
    return paragraphs


_DIRECT_SOURCE_MAX_BYTES = 256 * 1024 * 1024
_SOURCE_NUMBER = r"([0-9]+(?:,[0-9]{3})*(?:\.[0-9]+)?)"
_CONTRACT_UNIT = re.compile(
    r"時間単価(?:は|:)?\s*" + _SOURCE_NUMBER
    + r"\s*円\s*(?:\((?:消費)?税別\)|(?:消費)?税別)"
)
_CONTRACT_HOURS = re.compile(
    r"(?:想定総工数|見込工数)(?:は|:)?\s*" + _SOURCE_NUMBER + r"\s*時間"
)
_CONTRACT_AMOUNTS = re.compile(
    r"見込金額(?:は|:)?[、,\s]*税抜\s*"
    + _SOURCE_NUMBER
    + r"\s*円[、,\s]+消費税\s*"
    + _SOURCE_NUMBER
    + r"\s*円[、,\s]+税込\s*"
    + _SOURCE_NUMBER
    + r"\s*円"
)


def _docx_text_blocks(path: Path) -> list[str] | None:
    """Read source text in body order, including paragraphs nested in tables."""

    try:
        if path.stat().st_size > _DIRECT_SOURCE_MAX_BYTES:
            return None
        with zipfile.ZipFile(path) as archive:
            root = ET.fromstring(archive.read("word/document.xml"))
    except Exception:
        return None
    namespace = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    body = root.find("w:body", namespace)
    if body is None:
        return None
    blocks = []
    for paragraph in body.findall(".//w:p", namespace):
        text = "".join(
            node.text or "" for node in paragraph.findall(".//w:t", namespace)
        ).strip()
        if text:
            blocks.append(unicodedata.normalize("NFKC", text))
    return blocks or None


def _source_decimal(value: str) -> Decimal | None:
    try:
        parsed = Decimal(
            unicodedata.normalize("NFKC", value).replace(",", "").strip()
        )
    except Exception:
        return None
    return parsed if parsed.is_finite() else None


def _contract_billing_facts(
    path: Path,
) -> tuple[Decimal, Decimal, Decimal, Decimal, Decimal] | None:
    blocks = _docx_text_blocks(path)
    if not blocks:
        return None
    unit_values: set[Decimal] = set()
    hour_values: set[Decimal] = set()
    amount_values: set[tuple[Decimal, Decimal, Decimal]] = set()
    for block in blocks:
        for capture in _CONTRACT_UNIT.findall(block):
            parsed = _source_decimal(capture)
            if parsed is None:
                return None
            unit_values.add(parsed)
        for capture in _CONTRACT_HOURS.findall(block):
            parsed = _source_decimal(capture)
            if parsed is None:
                return None
            hour_values.add(parsed)
        for captures in _CONTRACT_AMOUNTS.findall(block):
            parsed = tuple(_source_decimal(capture) for capture in captures)
            if any(value is None for value in parsed):
                return None
            amount_values.add(parsed)  # type: ignore[arg-type]
    has_actual_settlement = any(
        re.search(r"実績工数.*?(?:事後精算|月次精算)", block) is not None
        for block in blocks
    )
    has_linear_formula = any(
        re.search(
            r"最終請求(?:額|金額).*?実績工数.*?時間単価.*?消費税.*?(?:加算|含)",
            block,
        )
        is not None
        for block in blocks
    )
    if (
        len(unit_values) != 1
        or len(hour_values) != 1
        or len(amount_values) != 1
        or not has_actual_settlement
        or not has_linear_formula
    ):
        return None
    unit = next(iter(unit_values))
    hours = next(iter(hour_values))
    expected_ex_tax, tax, expected_with_tax = next(iter(amount_values))
    if (
        unit <= 0
        or hours <= 0
        or expected_ex_tax <= 0
        or tax < 0
        or expected_with_tax <= 0
        or unit != unit.to_integral_value()
        or expected_ex_tax != expected_ex_tax.to_integral_value()
        or tax != tax.to_integral_value()
        or expected_with_tax != expected_with_tax.to_integral_value()
        or unit * hours != expected_ex_tax
        or expected_ex_tax + tax != expected_with_tax
    ):
        return None
    return unit, hours, expected_ex_tax, tax, expected_with_tax


def _contract_hours_ratio_tax_delta(
    engine: Any,
    match: re.Match[str],
) -> StructuredCandidateDecision | None:
    paths = _named_paths(engine, match["location"], "契約書", {".docx"})
    if len(paths) != 1:
        return None
    facts = _contract_billing_facts(paths[0])
    if facts is None:
        return None
    try:
        denominator = int(unicodedata.normalize("NFKC", match["denominator"]))
        numerator = int(unicodedata.normalize("NFKC", match["numerator"]))
    except ValueError:
        return None
    if not (0 <= numerator < denominator <= 1_000_000):
        return None
    expected_with_tax = facts[-1]
    difference = Fraction(expected_with_tax) * Fraction(
        denominator - numerator,
        denominator,
    )
    if difference.denominator != 1 or difference.numerator < 0:
        # The contract contains no source-backed yen rounding rule.
        return None
    answer = f"{difference.numerator:,}円"
    return _decision(answer, paths, engine.source_root, 4)


def _literal_autofilter_conditions(
    path: Path,
    sheet_name: str,
) -> tuple[tuple[str, tuple[str, ...]], ...] | None:
    try:
        source_size = path.stat().st_size
    except OSError:
        return None
    if source_size > _DIRECT_SOURCE_MAX_BYTES:
        return None
    from openpyxl import load_workbook
    from openpyxl.utils.cell import range_boundaries

    try:
        workbook = load_workbook(
            path,
            read_only=False,
            data_only=False,
            keep_links=False,
        )
    except Exception:
        return None
    try:
        worksheets = [
            worksheet
            for worksheet in workbook.worksheets
            if _normalized(worksheet.title) == _normalized(sheet_name)
        ]
        if len(worksheets) != 1:
            return None
        worksheet = worksheets[0]
        reference = getattr(worksheet.auto_filter, "ref", None)
        if not isinstance(reference, str) or not reference.strip():
            return None
        try:
            min_col, min_row, max_col, max_row = range_boundaries(
                reference.replace("$", "")
            )
        except (TypeError, ValueError):
            return None
        if min_col > max_col or min_row >= max_row:
            return None
        headers = [
            _cell_text(worksheet.cell(min_row, column).value)
            for column in range(min_col, max_col + 1)
        ]
        normalized_headers = [_normalized(header) for header in headers]
        if (
            any(not header for header in normalized_headers)
            or len(normalized_headers) != len(set(normalized_headers))
        ):
            return None
        columns = list(getattr(worksheet.auto_filter, "filterColumn", ()) or ())
        if not columns:
            return None
        conditions: list[tuple[int, str, tuple[str, ...]]] = []
        seen_column_ids: set[int] = set()
        for column in columns:
            column_id = getattr(column, "colId", None)
            if isinstance(column_id, bool) or not isinstance(column_id, int):
                return None
            if column_id in seen_column_ids or not 0 <= column_id < len(headers):
                return None
            seen_column_ids.add(column_id)
            if any(
                getattr(column, attribute, None) is not None
                for attribute in (
                    "top10",
                    "customFilters",
                    "dynamicFilter",
                    "colorFilter",
                    "iconFilter",
                )
            ):
                return None
            filters = getattr(column, "filters", None)
            if (
                filters is None
                or getattr(filters, "blank", None) is True
                or getattr(filters, "dateGroupItem", ())
            ):
                return None
            raw_values = list(getattr(filters, "filter", ()) or ())
            values = tuple(
                dict.fromkeys(unicodedata.normalize("NFKC", str(value)).strip() for value in raw_values)
            )
            header = headers[column_id]
            if (
                not values
                or any(not value for value in values)
                or any(token in header for token in ("=", "、"))
                or any(
                    token in value
                    for value in values
                    for token in ("=", "、", "または")
                )
            ):
                return None
            conditions.append((column_id, header, values))
        conditions.sort(key=lambda item: item[0])
        return tuple((header, values) for _, header, values in conditions)
    except Exception:
        return None
    finally:
        workbook.close()


def _excel_autofilter_conditions(
    engine: Any,
    match: re.Match[str],
) -> StructuredCandidateDecision | None:
    paths = _named_paths(
        engine,
        match["location"],
        match["container"],
        {".xlsx"},
    )
    if len(paths) != 1:
        return None
    conditions = _literal_autofilter_conditions(paths[0], match["sheet"])
    if not conditions:
        return None
    rendered = []
    for header, values in conditions:
        rendered.append(f"{header}=" + "または".join(values))
    return _decision("、".join(rendered), paths, engine.source_root, 4)


def _pivot_metric_identity(value: object) -> str:
    normalized = _normalized(value)
    return re.sub(
        r"(?:平均|average|avg)|[\s/／:：]+",
        "",
        normalized,
        flags=re.IGNORECASE,
    )


def _pivot_text(value: object) -> str | None:
    try:
        rendered = unicodedata.normalize("NFKC", _cell_text(value)).strip()
    except (TypeError, ValueError):
        return None
    if not rendered or any(token in rendered for token in ("=", "、", "\n", "\r")):
        return None
    return rendered


def _pivot_decimal(value: object) -> Decimal | None:
    if isinstance(value, bool) or not isinstance(value, (Decimal, int, float)):
        return None
    if isinstance(value, float) and not math.isfinite(value):
        return None
    try:
        result = Decimal(str(value))
    except Exception:
        return None
    return result if result.is_finite() else None


def _pivot_index(value: object) -> int | None:
    return value if isinstance(value, int) and not isinstance(value, bool) else None


def _pivot_average_data_fields(pivot: Any) -> tuple[tuple[Any, int, str], ...]:
    cache = getattr(pivot, "cache", None)
    cache_fields = list(getattr(cache, "cacheFields", ()) or ())
    answers: list[tuple[Any, int, str]] = []
    for data_field in list(getattr(pivot, "dataFields", ()) or ()):
        field_index = _pivot_index(getattr(data_field, "fld", None))
        if (
            field_index is None
            or not 0 <= field_index < len(cache_fields)
            or _normalized(getattr(data_field, "subtotal", "")) != "average"
            or _normalized(getattr(data_field, "showDataAs", "normal")) != "normal"
        ):
            continue
        field_name = _pivot_text(getattr(cache_fields[field_index], "name", None))
        caption = _pivot_text(getattr(data_field, "name", None))
        if field_name is None or caption is None:
            continue
        answers.append((data_field, field_index, field_name))
    return tuple(answers)


def _pivot_resolve_average_metric(
    pivot: Any,
    metric_surface: str,
) -> tuple[Any, int, str] | None:
    averages = _pivot_average_data_fields(pivot)
    if not averages:
        return None
    requested = _pivot_metric_identity(metric_surface)
    if not requested:
        return None
    exact = [
        item
        for item in averages
        if requested
        in {
            _pivot_metric_identity(item[2]),
            _pivot_metric_identity(getattr(item[0], "name", "")),
        }
    ]
    if len(exact) == 1:
        return exact[0]
    if exact:
        return None
    # A language-independent fallback is safe only when the Pivot defines one
    # and only one average measure.  This lets a question use a human-facing
    # label while refusing to guess among sibling measures.
    return averages[0] if len(averages) == 1 else None


def _pivot_source_definition(
    workbook: Any,
    pivot: Any,
) -> tuple[Any, tuple[int, int, int, int], list[Any]] | None:
    from openpyxl.utils.cell import range_boundaries

    cache = getattr(pivot, "cache", None)
    cache_source = getattr(cache, "cacheSource", None)
    worksheet_source = getattr(cache_source, "worksheetSource", None)
    if (
        _normalized(getattr(cache_source, "type", "")) != "worksheet"
        or getattr(cache_source, "connectionId", None) is not None
        or getattr(cache_source, "consolidation", None) is not None
        or worksheet_source is None
        or getattr(worksheet_source, "name", None) is not None
    ):
        return None
    source_sheet = _pivot_text(getattr(worksheet_source, "sheet", None))
    source_ref = getattr(worksheet_source, "ref", None)
    if source_sheet is None or not isinstance(source_ref, str):
        return None
    worksheets = [
        worksheet
        for worksheet in workbook.worksheets
        if _normalized(worksheet.title) == _normalized(source_sheet)
    ]
    if len(worksheets) != 1:
        return None
    try:
        boundaries = range_boundaries(source_ref.replace("$", ""))
    except (TypeError, ValueError):
        return None
    min_col, min_row, max_col, max_row = boundaries
    if (
        min_col > max_col
        or min_row >= max_row
        or (max_col - min_col + 1) * (max_row - min_row + 1) > 1_000_000
    ):
        return None
    return worksheets[0], boundaries, list(getattr(cache, "cacheFields", ()) or ())


def _pivot_row_field_definition(
    pivot: Any,
    cache_fields: Sequence[Any],
) -> tuple[tuple[int, str], ...] | None:
    row_fields: list[tuple[int, str]] = []
    seen: set[int] = set()
    for row_field in list(getattr(pivot, "rowFields", ()) or ()):
        field_index = _pivot_index(getattr(row_field, "x", None))
        if (
            field_index is None
            or field_index in seen
            or not 0 <= field_index < len(cache_fields)
        ):
            return None
        field_name = _pivot_text(getattr(cache_fields[field_index], "name", None))
        if field_name is None:
            return None
        seen.add(field_index)
        row_fields.append((field_index, field_name))
    if not 1 <= len(row_fields) <= 32:
        return None
    if len({_normalized(name) for _, name in row_fields}) != len(row_fields):
        return None
    return tuple(row_fields)


def _pivot_filters_are_literal_complete(pivot: Any) -> bool:
    if (
        list(getattr(pivot, "pageFields", ()) or ())
        or list(getattr(pivot, "filters", ()) or ())
        or getattr(pivot, "dataOnRows", False) is True
        or getattr(pivot, "multipleFieldFilters", False) is True
    ):
        return False
    for column_field in list(getattr(pivot, "colFields", ()) or ()):
        # -2 is Excel's data-field sentinel.  Any categorical column field
        # would require a two-dimensional hierarchy resolver and is refused.
        if _pivot_index(getattr(column_field, "x", None)) != -2:
            return False
    return True


def _pivot_item_indexes(item: Any) -> tuple[int, ...] | None:
    indexes: list[int] = []
    for value in list(getattr(item, "x", ()) or ()):
        index = _pivot_index(getattr(value, "v", None))
        if index is None or index < 0:
            return None
        indexes.append(index)
    return tuple(indexes)


def _pivot_display_leaves(
    worksheet: Any,
    pivot: Any,
    row_fields: Sequence[tuple[int, str]],
    data_field: Any,
) -> dict[tuple[str, ...], tuple[tuple[str, ...], Decimal]] | None:
    from openpyxl.utils.cell import range_boundaries

    location = getattr(pivot, "location", None)
    reference = getattr(location, "ref", None)
    if not isinstance(reference, str):
        return None
    try:
        min_col, min_row, max_col, max_row = range_boundaries(
            reference.replace("$", "")
        )
    except (TypeError, ValueError):
        return None
    depth = len(row_fields)
    if min_col > max_col or min_row >= max_row:
        return None
    compact = getattr(pivot, "compact", None) is True
    compact_data = getattr(pivot, "compactData", None) is True
    outline = getattr(pivot, "outline", None) is True
    if compact:
        if not compact_data or not outline or getattr(location, "firstDataCol", None) != 1:
            return None
    elif compact_data or outline or getattr(location, "firstDataCol", None) != depth:
        return None

    caption = _normalized(getattr(data_field, "name", ""))
    metric_columns = []
    for column in range(min_col, max_col + 1):
        cell = worksheet.cell(min_row, column)
        if cell.data_type == "f":
            return None
        if _normalized(_cell_text(cell.value)) == caption:
            metric_columns.append(column)
    if len(metric_columns) != 1:
        return None
    metric_column = metric_columns[0]
    if metric_column < min_col + (1 if compact else depth):
        return None

    row_items = list(getattr(pivot, "rowItems", ()) or ())
    if len(row_items) != max_row - min_row:
        return None
    state: list[str | None] = [None] * depth
    leaves: dict[tuple[str, ...], tuple[tuple[str, ...], Decimal]] = {}
    grand_count = 0
    for offset, item in enumerate(row_items, 1):
        row_number = min_row + offset
        item_type = _normalized(getattr(item, "t", "data") or "data")
        if item_type not in {"data", "default", "grand"}:
            return None
        if item_type == "grand":
            grand_count += 1
            if offset != len(row_items):
                return None
            continue
        if item_type == "default":
            continue
        item_indexes = _pivot_item_indexes(item)
        item_depth = _pivot_index(getattr(item, "r", None))
        if item_indexes is None or item_depth is None:
            return None

        if compact:
            label_cell = worksheet.cell(row_number, min_col)
            if label_cell.data_type == "f":
                return None
            label = _pivot_text(label_cell.value)
            indent = label_cell.alignment.indent
            if (
                label is None
                or isinstance(indent, bool)
                or not isinstance(indent, (int, float))
                or float(indent) != int(indent)
            ):
                return None
            label_depth = int(indent)
            if (
                label_depth != item_depth
                or not 0 <= label_depth < depth
                or len(item_indexes) != 1
                or any(state[index] is None for index in range(label_depth))
            ):
                return None
            state[label_depth] = label
            for index in range(label_depth + 1, depth):
                state[index] = None
            is_leaf = label_depth == depth - 1
        else:
            if not 0 <= item_depth < depth or len(item_indexes) != depth - item_depth:
                return None
            labels: list[str | None] = []
            for index in range(depth):
                cell = worksheet.cell(row_number, min_col + index)
                if cell.data_type == "f":
                    return None
                labels.append(_pivot_text(cell.value) if cell.value is not None else None)
            if (
                any(labels[index] is not None for index in range(item_depth))
                or any(labels[index] is None for index in range(item_depth, depth))
                or any(state[index] is None for index in range(item_depth))
            ):
                return None
            for index in range(item_depth, depth):
                state[index] = labels[index]
            is_leaf = True

        if not is_leaf or any(value is None for value in state):
            continue
        metric_cell = worksheet.cell(row_number, metric_column)
        if metric_cell.data_type == "f":
            return None
        metric_value = _pivot_decimal(metric_cell.value)
        if metric_value is None:
            return None
        rendered_labels = tuple(str(value) for value in state)
        key = tuple(_normalized(value) for value in rendered_labels)
        if key in leaves:
            return None
        leaves[key] = rendered_labels, metric_value
    if grand_count != 1 or not leaves:
        return None
    return leaves


def _pivot_raw_group_means(
    source_sheet: Any,
    source_bounds: tuple[int, int, int, int],
    row_fields: Sequence[tuple[int, str]],
    metric_field_index: int,
    metric_field_name: str,
) -> dict[tuple[str, ...], tuple[tuple[str, ...], Decimal]] | None:
    min_col, min_row, max_col, max_row = source_bounds
    headers: list[str] = []
    for column in range(min_col, max_col + 1):
        cell = source_sheet.cell(min_row, column)
        if cell.data_type == "f":
            return None
        header = _pivot_text(cell.value)
        if header is None:
            return None
        headers.append(header)
    normalized_headers = [_normalized(header) for header in headers]
    if len(normalized_headers) != len(set(normalized_headers)):
        return None
    header_columns = {
        normalized: min_col + index
        for index, normalized in enumerate(normalized_headers)
    }
    required_names = [name for _, name in row_fields] + [metric_field_name]
    if any(_normalized(name) not in header_columns for name in required_names):
        return None
    group_columns = [header_columns[_normalized(name)] for _, name in row_fields]
    metric_column = header_columns[_normalized(metric_field_name)]
    # The cache's field ordinal must identify the same source-range header.
    if metric_field_index != normalized_headers.index(_normalized(metric_field_name)):
        return None

    aggregates: dict[tuple[str, ...], tuple[tuple[str, ...], Decimal, int]] = {}
    for row_number in range(min_row + 1, max_row + 1):
        rendered: list[str] = []
        for column in group_columns:
            cell = source_sheet.cell(row_number, column)
            if cell.data_type == "f":
                return None
            value = _pivot_text(cell.value)
            if value is None:
                return None
            rendered.append(value)
        metric_cell = source_sheet.cell(row_number, metric_column)
        if metric_cell.data_type == "f":
            return None
        metric = _pivot_decimal(metric_cell.value)
        if metric is None:
            return None
        labels = tuple(rendered)
        key = tuple(_normalized(value) for value in labels)
        previous = aggregates.get(key)
        if previous is None:
            aggregates[key] = labels, metric, 1
        else:
            if previous[0] != labels:
                return None
            aggregates[key] = labels, previous[1] + metric, previous[2] + 1
    if not aggregates:
        return None
    return {
        key: (labels, total / count)
        for key, (labels, total, count) in aggregates.items()
    }


def _pivot_answer_from_workbook(
    workbook: Any,
    metric_surface: str,
    explicit_sheet: str | None,
    include_aggregate: bool,
) -> str | None:
    if explicit_sheet is None:
        worksheets = list(workbook.worksheets)
    else:
        worksheets = [
            worksheet
            for worksheet in workbook.worksheets
            if _normalized(worksheet.title) == _normalized(explicit_sheet)
        ]
        if len(worksheets) != 1:
            return None

    candidates: list[tuple[Any, Any, tuple[Any, int, str]]] = []
    for worksheet in worksheets:
        for pivot in list(getattr(worksheet, "_pivots", ()) or ()):
            metric = _pivot_resolve_average_metric(pivot, metric_surface)
            if metric is not None:
                candidates.append((worksheet, pivot, metric))
    if len(candidates) != 1:
        return None
    pivot_sheet, pivot, metric = candidates[0]
    data_field, metric_field_index, metric_field_name = metric
    if not _pivot_filters_are_literal_complete(pivot):
        return None
    source = _pivot_source_definition(workbook, pivot)
    if source is None:
        return None
    source_sheet, source_bounds, cache_fields = source
    row_fields = _pivot_row_field_definition(pivot, cache_fields)
    if row_fields is None or metric_field_index in {index for index, _ in row_fields}:
        return None
    displayed = _pivot_display_leaves(pivot_sheet, pivot, row_fields, data_field)
    raw = _pivot_raw_group_means(
        source_sheet,
        source_bounds,
        row_fields,
        metric_field_index,
        metric_field_name,
    )
    if displayed is None or raw is None or set(displayed) != set(raw):
        return None
    for key, (_, raw_mean) in raw.items():
        display_mean = displayed[key][1]
        tolerance = max(Decimal("1e-9"), abs(raw_mean) * Decimal("1e-12"))
        if abs(display_mean - raw_mean) > tolerance:
            return None

    max_raw = max(mean for _, mean in raw.values())
    raw_winners = [key for key, (_, mean) in raw.items() if mean == max_raw]
    max_display = max(mean for _, mean in displayed.values())
    display_winners = [
        key for key, (_, mean) in displayed.items() if mean == max_display
    ]
    if len(raw_winners) != 1 or display_winners != raw_winners:
        return None
    winner = raw_winners[0]
    labels = displayed[winner][0]
    field_names = [name for _, name in row_fields]
    if len(field_names) != len(labels):
        return None
    rendered = [
        f"{field_name}={label}"
        for field_name, label in zip(field_names, labels)
    ]
    if include_aggregate:
        caption = _pivot_text(getattr(data_field, "name", None))
        if caption is None:
            return None
        # The maximum aggregate is used to choose the winning group, but the
        # question asks for its extraction conditions and aggregation content.
        # Returning the numeric maximum would expose an intermediate value as
        # an additional requested output.
        return f"{'、'.join(rendered)}で抽出されたデータに対する{caption}"
    return "、".join(rendered)


def _load_pivot_workbook(path: Path) -> Any | None:
    try:
        if path.stat().st_size > _DIRECT_SOURCE_MAX_BYTES:
            return None
        from openpyxl import load_workbook

        return load_workbook(
            path,
            read_only=False,
            data_only=False,
            keep_links=False,
        )
    except Exception:
        return None


def _pivot_average_argmax(
    engine: Any,
    match: re.Match[str],
    *,
    include_aggregate: bool,
) -> StructuredCandidateDecision | None:
    paths = _named_paths(
        engine,
        match["location"],
        match["container"],
        {".xlsx"},
    )
    if len(paths) != 1:
        return None
    workbook = _load_pivot_workbook(paths[0])
    if workbook is None:
        return None
    try:
        answer = _pivot_answer_from_workbook(
            workbook,
            match["metric_surface"],
            match.groupdict().get("sheet"),
            include_aggregate,
        )
    except Exception:
        answer = None
    finally:
        workbook.close()
    if not answer:
        return None
    return _decision(answer, paths, engine.source_root, 8)


_DOCX_IDENTITY_LABEL = re.compile(
    r"(?:チェックポイント|会議ID|資料ID|文書ID|報告ID)"
    r"\s*[:：]\s*([A-Za-z][A-Za-z0-9._-]*)",
    flags=re.IGNORECASE,
)
_WORD_NAMESPACE = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
_WORD_HIGHLIGHT_VALUES = frozenset(
    {
        "black",
        "blue",
        "cyan",
        "darkblue",
        "darkcyan",
        "darkgray",
        "darkgreen",
        "darkmagenta",
        "darkred",
        "darkyellow",
        "green",
        "lightgray",
        "magenta",
        "none",
        "red",
        "white",
        "yellow",
    }
)


def _docx_identity_keys(path: Path) -> set[str] | None:
    blocks = _docx_text_blocks(path)
    if not blocks:
        return None
    return {
        _normalized(capture)
        for block in blocks
        for capture in _DOCX_IDENTITY_LABEL.findall(block)
    }


def _word_run_text(run: ET.Element) -> str:
    namespace = "{" + _WORD_NAMESPACE + "}"
    tokens: list[str] = []
    for node in run.iter():
        if node.tag == namespace + "t":
            tokens.append(node.text or "")
        elif node.tag == namespace + "tab":
            tokens.append("\t")
        elif node.tag in {namespace + "br", namespace + "cr"}:
            tokens.append("\n")
    return "".join(tokens)


def _docx_highlight_segments(path: Path, color: str) -> tuple[str, ...] | None:
    try:
        if path.stat().st_size > _DIRECT_SOURCE_MAX_BYTES:
            return None
        with zipfile.ZipFile(path) as archive:
            names = set(archive.namelist())
            if "word/document.xml" not in names:
                return None
            # Style-inherited highlighting requires the full Word cascade.
            # Refuse such documents instead of silently treating it as direct
            # run formatting.
            if "word/styles.xml" in names:
                styles = ET.fromstring(archive.read("word/styles.xml"))
                if any(
                    node.tag == "{" + _WORD_NAMESPACE + "}highlight"
                    for node in styles.iter()
                ):
                    return None
            story_parts = ["word/document.xml"]
            story_parts.extend(
                sorted(
                    name
                    for name in names
                    if re.fullmatch(r"word/(?:header|footer)\d+\.xml", name)
                )
            )
            story_parts.extend(
                name
                for name in ("word/footnotes.xml", "word/endnotes.xml")
                if name in names
            )
            roots = [ET.fromstring(archive.read(name)) for name in story_parts]
    except Exception:
        return None

    namespace = {"w": _WORD_NAMESPACE}
    highlight_attribute = "{" + _WORD_NAMESPACE + "}val"
    segments: list[str] = []
    for root in roots:
        for paragraph in root.iter("{" + _WORD_NAMESPACE + "}p"):
            current: list[str] = []
            for run in paragraph.iter("{" + _WORD_NAMESPACE + "}r"):
                text = _word_run_text(run)
                if not text:
                    continue
                highlight = run.find("./w:rPr/w:highlight", namespace)
                value = None
                if highlight is not None:
                    raw_value = highlight.get(highlight_attribute)
                    if raw_value is None:
                        return None
                    value = raw_value.casefold()
                    if value not in _WORD_HIGHLIGHT_VALUES:
                        return None
                if value == color:
                    current.append(text)
                    continue
                if current:
                    rendered = unicodedata.normalize("NFC", "".join(current)).strip()
                    if rendered:
                        segments.append(rendered)
                    current = []
            if current:
                rendered = unicodedata.normalize("NFC", "".join(current)).strip()
                if rendered:
                    segments.append(rendered)
    return tuple(segments) if segments else None


def _docx_highlighted_text(
    engine: Any,
    match: re.Match[str],
) -> StructuredCandidateDecision | None:
    document_key = _normalized(match["document_key"])
    color = _declared_color(match["color"], _WORD_HIGHLIGHT_ALIASES)
    if color is None:
        return None
    paths: list[Path] = []
    for path in _all_files(engine, {".docx"}):
        if not _project_matches(engine, path, match["location"]):
            continue
        relative_parts = [
            _normalized(part)
            for part in path.relative_to(engine.source_root).parts
        ]
        if not any("資料" in part for part in relative_parts):
            continue
        if _docx_identity_keys(path) == {document_key}:
            paths.append(path)
    # Source binding is complete before any highlight property is inspected.
    if len(paths) != 1:
        return None
    segments = _docx_highlight_segments(paths[0], color)
    if not segments:
        return None
    return _decision("、".join(segments), paths, engine.source_root, 7)


_PRESENTATION_NAMESPACE = (
    "http://schemas.openxmlformats.org/presentationml/2006/main"
)
_DRAWING_NAMESPACE = "http://schemas.openxmlformats.org/drawingml/2006/main"
_OFFICE_REL_NAMESPACE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
)
_PACKAGE_REL_NAMESPACE = "http://schemas.openxmlformats.org/package/2006/relationships"
_MARKUP_COMPAT_NAMESPACE = (
    "http://schemas.openxmlformats.org/markup-compatibility/2006"
)


def _package_target_part(base_part: str, target: str) -> str | None:
    normalized_target = target.replace("\\", "/")
    if normalized_target.startswith("/"):
        candidate = PurePosixPath(normalized_target.lstrip("/"))
    else:
        candidate = PurePosixPath(base_part).parent / normalized_target
    parts: list[str] = []
    for part in candidate.parts:
        if part in {"", "."}:
            continue
        if part == "..":
            if not parts:
                return None
            parts.pop()
            continue
        parts.append(part)
    return "/".join(parts) if parts else None


def _pptx_slide_root(path: Path, page: int) -> ET.Element | None:
    try:
        if path.stat().st_size > _DIRECT_SOURCE_MAX_BYTES:
            return None
        with zipfile.ZipFile(path) as archive:
            presentation = ET.fromstring(archive.read("ppt/presentation.xml"))
            relationships = ET.fromstring(
                archive.read("ppt/_rels/presentation.xml.rels")
            )
            slide_ids = list(
                presentation.iter("{" + _PRESENTATION_NAMESPACE + "}sldId")
            )
            if not 1 <= page <= len(slide_ids):
                return None
            relation_id = slide_ids[page - 1].get(
                "{" + _OFFICE_REL_NAMESPACE + "}id"
            )
            matches = [
                relation
                for relation in relationships.iter(
                    "{" + _PACKAGE_REL_NAMESPACE + "}Relationship"
                )
                if relation.get("Id") == relation_id
            ]
            if len(matches) != 1:
                return None
            relation = matches[0]
            if (
                relation.get("TargetMode") == "External"
                or not str(relation.get("Type", "")).endswith("/slide")
            ):
                return None
            target = relation.get("Target")
            if not isinstance(target, str):
                return None
            part = _package_target_part("ppt/presentation.xml", target)
            if part is None or part not in archive.namelist():
                return None
            return ET.fromstring(archive.read(part))
    except Exception:
        return None


def _drawingml_direct_rgb(solid_fill: ET.Element) -> str | None:
    children = list(solid_fill)
    if len(children) != 1 or children[0].tag != "{" + _DRAWING_NAMESPACE + "}srgbClr":
        return None
    color = children[0]
    if list(color):
        # Tint/shade/luminance transforms need a complete DrawingML resolver.
        return None
    value = str(color.get("val", "")).upper()
    return value if re.fullmatch(r"[0-9A-F]{6}", value) else None


def _drawingml_color_matches(rgb: str, declared_color: str) -> bool:
    center = _PPTX_HUE_CENTERS.get(declared_color)
    if center is None:
        return False
    red, green, blue = (
        int(rgb[index : index + 2], 16) / 255.0 for index in (0, 2, 4)
    )
    hue, lightness, saturation = colorsys.rgb_to_hls(red, green, blue)
    chroma = max(red, green, blue) - min(red, green, blue)
    hue_degrees = hue * 360.0
    distance = abs(hue_degrees - center)
    distance = min(distance, 360.0 - distance)
    return (
        distance <= 12.0
        and saturation >= 0.45
        and chroma >= 0.20
        and 0.12 <= lightness <= 0.88
    )


def _pptx_shape_text(shape: ET.Element) -> str:
    namespace = {
        "p": _PRESENTATION_NAMESPACE,
        "a": _DRAWING_NAMESPACE,
    }
    paragraphs: list[str] = []
    for paragraph in shape.findall("./p:txBody/a:p", namespace):
        text = "".join(
            node.text or "" for node in paragraph.findall(".//a:t", namespace)
        ).strip()
        if text:
            paragraphs.append(text)
    return "\n".join(paragraphs)


def _pptx_shape_fill_texts(
    root: ET.Element,
    declared_color: str,
) -> tuple[str, ...] | None:
    namespace = {
        "p": _PRESENTATION_NAMESPACE,
        "a": _DRAWING_NAMESPACE,
    }
    # Raster images, tables/charts, grouped inheritance and alternate content
    # can carry visually red text without a direct shape fill.  Refuse those
    # mixed channels instead of pretending this source-only rule saw them.
    unsupported = (
        root.find(".//p:pic", namespace) is not None
        or root.find(".//p:graphicFrame", namespace) is not None
        or root.find(".//p:grpSp", namespace) is not None
        or root.find(
            ".//{" + _MARKUP_COMPAT_NAMESPACE + "}AlternateContent"
        )
        is not None
    )
    if unsupported:
        return None
    answers: list[str] = []
    for shape in root.findall(".//p:sp", namespace):
        properties = shape.find("./p:nvSpPr/p:cNvPr", namespace)
        if properties is not None and _normalized(properties.get("hidden", "")) in {
            "1",
            "true",
        }:
            continue
        fills = shape.findall("./p:spPr/a:solidFill", namespace)
        if len(fills) > 1:
            return None
        if not fills:
            continue
        rgb = _drawingml_direct_rgb(fills[0])
        if rgb is None:
            return None
        if not _drawingml_color_matches(rgb, declared_color):
            continue
        text = _pptx_shape_text(shape)
        if not text:
            # The intended label may live in a separate overlaid text box.
            return None
        answers.append(text)
    return tuple(answers) if answers else None


def _pptx_shape_fill_text(
    engine: Any,
    match: re.Match[str],
) -> StructuredCandidateDecision | None:
    paths = _named_paths(
        engine,
        match["location"],
        match["container"],
        {".pptx"},
    )
    if len(paths) != 1:
        return None
    try:
        page = int(unicodedata.normalize("NFKC", match["slide"]))
    except ValueError:
        return None
    declared_color = _declared_color(match["color"], _PPTX_COLOR_ALIASES)
    if declared_color is None:
        return None
    root = _pptx_slide_root(paths[0], page)
    if root is None:
        return None
    answers = _pptx_shape_fill_texts(root, declared_color)
    if not answers:
        return None
    return _decision("、".join(answers), paths, engine.source_root, 6)


def _pptx_canonical_text(text_body: ET.Element | None) -> str | None:
    if text_body is None:
        return ""
    namespace = {"a": _DRAWING_NAMESPACE}
    if text_body.find(".//a:fld", namespace) is not None:
        # Fields can render dates, page numbers, or document properties not
        # represented by their cached a:t value.
        return None
    paragraphs: list[str] = []
    for paragraph in text_body.findall("./a:p", namespace):
        pieces: list[str] = []
        for node in paragraph.iter():
            if node.tag == "{" + _DRAWING_NAMESPACE + "}t":
                pieces.append(node.text or "")
            elif node.tag in {
                "{" + _DRAWING_NAMESPACE + "}br",
                "{" + _DRAWING_NAMESPACE + "}tab",
            }:
                pieces.append(" ")
        rendered = " ".join(
            unicodedata.normalize("NFKC", "".join(pieces)).split()
        )
        if rendered:
            paragraphs.append(rendered)
    return " / ".join(paragraphs)


def _pptx_shape_visibility(shape: ET.Element, kind: str) -> bool | None:
    namespace = {"p": _PRESENTATION_NAMESPACE}
    paths = {
        "shape": "./p:nvSpPr/p:cNvPr",
        "connector": "./p:nvCxnSpPr/p:cNvPr",
        "table": "./p:nvGraphicFramePr/p:cNvPr",
    }
    properties = shape.find(paths[kind], namespace)
    if properties is None:
        return None
    hidden = _normalized(properties.get("hidden", "0"))
    if hidden not in {"", "0", "false", "1", "true"}:
        return None
    return hidden not in {"1", "true"}


def _pptx_placeholder_topology(shape: ET.Element) -> tuple[str, str] | None:
    namespace = {"p": _PRESENTATION_NAMESPACE}
    placeholders = shape.findall(".//p:nvPr/p:ph", namespace)
    if len(placeholders) > 1:
        return None
    if not placeholders:
        return "", ""
    placeholder = placeholders[0]
    placeholder_type = str(placeholder.get("type") or "obj")
    placeholder_index = str(placeholder.get("idx") or "0")
    return placeholder_type, placeholder_index


def _pptx_slide_model(
    root: ET.Element,
    slide_number: int,
    *,
    inherited: bool = False,
) -> tuple[tuple[Any, ...], tuple[dict[str, Any], ...]] | None:
    namespace = {
        "p": _PRESENTATION_NAMESPACE,
        "a": _DRAWING_NAMESPACE,
        "mc": _MARKUP_COMPAT_NAMESPACE,
    }
    if (
        root.find(".//mc:AlternateContent", namespace) is not None
        or root.find(".//p:pic", namespace) is not None
        or root.find(".//p:grpSp", namespace) is not None
        or root.find(".//p:contentPart", namespace) is not None
        or root.find(".//p:oleObj", namespace) is not None
        or root.find(".//a:blip", namespace) is not None
    ):
        return None
    shape_tree = root.find("./p:cSld/p:spTree", namespace)
    if shape_tree is None:
        return None
    allowed_scaffolding = {
        "{" + _PRESENTATION_NAMESPACE + "}nvGrpSpPr",
        "{" + _PRESENTATION_NAMESPACE + "}grpSpPr",
        "{" + _PRESENTATION_NAMESPACE + "}extLst",
    }
    topology: list[Any] = []
    units: list[dict[str, Any]] = []
    visible_shape_index = 0
    for child in shape_tree:
        if child.tag in allowed_scaffolding:
            continue
        if child.tag == "{" + _PRESENTATION_NAMESPACE + "}sp":
            kind = "shape"
        elif child.tag == "{" + _PRESENTATION_NAMESPACE + "}cxnSp":
            kind = "connector"
        elif child.tag == "{" + _PRESENTATION_NAMESPACE + "}graphicFrame":
            kind = "table"
        else:
            return None
        visibility = _pptx_shape_visibility(child, kind)
        placeholder = _pptx_placeholder_topology(child)
        if visibility is None or placeholder is None:
            return None
        if placeholder[0] in {"dt", "ftr", "sldNum"}:
            # Cached date/footer/slide-number values in standard layouts and
            # masters are dynamic presentation metadata, not authored visible
            # text.  They routinely differ between otherwise identical files.
            if inherited:
                continue
            return None
        visible_shape_index += 1

        if kind in {"shape", "connector"}:
            text_body = child.find("./p:txBody", namespace)
            text = _pptx_canonical_text(text_body)
            if text is None:
                return None
            geometry = ""
            if kind == "shape":
                properties = child.find("./p:spPr", namespace)
                if properties is None:
                    return None
                preset = properties.find("./a:prstGeom", namespace)
                custom = properties.find("./a:custGeom", namespace)
                if preset is not None and custom is not None:
                    return None
                if preset is not None:
                    geometry = "preset:" + str(preset.get("prst") or "")
                elif custom is not None:
                    geometry = "custom"
            topology.append(
                (kind, visibility, placeholder, geometry, text_body is not None)
            )
            if visibility and text:
                units.append(
                    {
                        "key": (slide_number, visible_shape_index, "shape"),
                        "slide": slide_number,
                        "shape": visible_shape_index,
                        "kind": "shape",
                        "text": text,
                    }
                )
            continue

        graphic_data = child.find("./a:graphic/a:graphicData", namespace)
        if graphic_data is None:
            return None
        tables = graphic_data.findall("./a:tbl", namespace)
        if len(tables) != 1 or len(list(graphic_data)) != 1:
            return None
        table = tables[0]
        grid = table.find("./a:tblGrid", namespace)
        rows = table.findall("./a:tr", namespace)
        if grid is None or not rows:
            return None
        column_count = len(grid.findall("./a:gridCol", namespace))
        if column_count < 1:
            return None
        table_matrix: list[tuple[str, ...]] = []
        cell_topology: list[tuple[Any, ...]] = []
        for row_index, row in enumerate(rows):
            cells = row.findall("./a:tc", namespace)
            if len(cells) != column_count:
                return None
            rendered_row: list[str] = []
            for column_index, cell in enumerate(cells):
                text = _pptx_canonical_text(cell.find("./a:txBody", namespace))
                if text is None:
                    return None
                rendered_row.append(text)
                cell_topology.append(
                    (
                        str(cell.get("gridSpan") or "1"),
                        str(cell.get("rowSpan") or "1"),
                        _normalized(cell.get("hMerge", "0")),
                        _normalized(cell.get("vMerge", "0")),
                    )
                )
                if visibility:
                    units.append(
                        {
                            "key": (
                                slide_number,
                                visible_shape_index,
                                "table",
                                row_index,
                                column_index,
                            ),
                            "slide": slide_number,
                            "shape": visible_shape_index,
                            "kind": "table",
                            "row": row_index,
                            "column": column_index,
                            "text": text,
                        }
                    )
            table_matrix.append(tuple(rendered_row))
        topology.append(
            (
                kind,
                visibility,
                placeholder,
                len(rows),
                column_count,
                tuple(cell_topology),
            )
        )
        for unit in units:
            if (
                unit["kind"] == "table"
                and unit["slide"] == slide_number
                and unit["shape"] == visible_shape_index
            ):
                unit["table_matrix"] = tuple(table_matrix)
    return tuple(topology), tuple(units)


def _pptx_inherited_signature(root: ET.Element) -> tuple[Any, ...] | None:
    """Return style-insensitive inherited visible text and shape topology."""

    model = _pptx_slide_model(root, 0, inherited=True)
    if model is None:
        return None
    topology, units = model
    visible_text = tuple(
        sorted((tuple(unit["key"]), str(unit["text"])) for unit in units)
    )
    return topology, visible_text


def _pptx_related_part(
    archive: zipfile.ZipFile,
    names: set[str],
    source_part: str,
    relationship_suffix: str,
    expected_prefix: str,
) -> str | None:
    source = PurePosixPath(source_part)
    relationship_part = (
        source.parent / "_rels" / f"{source.name}.rels"
    ).as_posix()
    if relationship_part not in names:
        return None
    try:
        relationships = ET.fromstring(archive.read(relationship_part))
    except (KeyError, ET.ParseError):
        return None
    relation_ids: set[str] = set()
    matches: list[ET.Element] = []
    for relation in relationships.iter(
        "{" + _PACKAGE_REL_NAMESPACE + "}Relationship"
    ):
        relation_id = relation.get("Id")
        if not isinstance(relation_id, str) or relation_id in relation_ids:
            return None
        relation_ids.add(relation_id)
        if str(relation.get("Type", "")).endswith(relationship_suffix):
            matches.append(relation)
    if len(matches) != 1 or matches[0].get("TargetMode") == "External":
        return None
    target = matches[0].get("Target")
    if not isinstance(target, str):
        return None
    part = _package_target_part(source_part, target)
    if (
        part is None
        or part not in names
        or not part.startswith(expected_prefix)
        or not part.endswith(".xml")
    ):
        return None
    return part


def _pptx_visible_deck_model(
    path: Path,
) -> tuple[tuple[Any, ...], dict[tuple[Any, ...], dict[str, Any]]] | None:
    try:
        if path.stat().st_size > _DIRECT_SOURCE_MAX_BYTES:
            return None
        with zipfile.ZipFile(path) as archive:
            archive_names = archive.namelist()
            if len(archive_names) != len(set(archive_names)):
                return None
            names = set(archive_names)
            total_size = sum(item.file_size for item in archive.infolist())
            if total_size > 512 * 1024 * 1024:
                return None
            presentation = ET.fromstring(archive.read("ppt/presentation.xml"))
            relationships = ET.fromstring(
                archive.read("ppt/_rels/presentation.xml.rels")
            )
            relation_map: dict[str, ET.Element] = {}
            for relation in relationships.iter(
                "{" + _PACKAGE_REL_NAMESPACE + "}Relationship"
            ):
                relation_id = relation.get("Id")
                if not isinstance(relation_id, str) or relation_id in relation_map:
                    return None
                relation_map[relation_id] = relation
            slide_parts: list[str] = []
            for slide_id in presentation.iter(
                "{" + _PRESENTATION_NAMESPACE + "}sldId"
            ):
                relation_id = slide_id.get("{" + _OFFICE_REL_NAMESPACE + "}id")
                relation = relation_map.get(str(relation_id))
                if (
                    relation is None
                    or relation.get("TargetMode") == "External"
                    or not str(relation.get("Type", "")).endswith("/slide")
                ):
                    return None
                target = relation.get("Target")
                if not isinstance(target, str):
                    return None
                part = _package_target_part("ppt/presentation.xml", target)
                if part is None or part not in names or part in slide_parts:
                    return None
                slide_parts.append(part)
            if not slide_parts:
                return None
            roots = [ET.fromstring(archive.read(part)) for part in slide_parts]
            inherited_roots: list[tuple[ET.Element, ET.Element]] = []
            for slide_part in slide_parts:
                layout_part = _pptx_related_part(
                    archive,
                    names,
                    slide_part,
                    "/slideLayout",
                    "ppt/slideLayouts/",
                )
                if layout_part is None:
                    return None
                master_part = _pptx_related_part(
                    archive,
                    names,
                    layout_part,
                    "/slideMaster",
                    "ppt/slideMasters/",
                )
                if master_part is None:
                    return None
                inherited_roots.append(
                    (
                        ET.fromstring(archive.read(layout_part)),
                        ET.fromstring(archive.read(master_part)),
                    )
                )
    except Exception:
        return None

    deck_topology: list[Any] = []
    deck_units: dict[tuple[Any, ...], dict[str, Any]] = {}
    for slide_number, (root, inherited_pair) in enumerate(
        zip(roots, inherited_roots), 1
    ):
        model = _pptx_slide_model(root, slide_number)
        if model is None:
            return None
        topology, units = model
        layout_signature = _pptx_inherited_signature(inherited_pair[0])
        master_signature = _pptx_inherited_signature(inherited_pair[1])
        if layout_signature is None or master_signature is None:
            return None
        deck_topology.append((topology, layout_signature, master_signature))
        for unit in units:
            key = unit["key"]
            if key in deck_units:
                return None
            deck_units[key] = unit
    return tuple(deck_topology), deck_units


def _proposal_pptx_pair(
    engine: Any,
    location: str,
    document_key: str,
) -> tuple[Path, Path] | None:
    proposal_dirs: list[Path] = []
    for candidate in engine.source_root.rglob("*"):
        if (
            not candidate.is_dir()
            or candidate.is_symlink()
            or "提案" not in _normalized(candidate.name)
            or not _project_matches(engine, candidate / "__scope__", location)
        ):
            continue
        proposal_dirs.append(candidate)
    if len(proposal_dirs) != 1:
        return None
    proposal_dir = proposal_dirs[0]
    document_names = {
        _normalized(value)
        for value in _candidate_values(document_key, engine.glossary)
    }

    def matching_files(directory: Path) -> list[Path]:
        return [
            path
            for path in directory.iterdir()
            if path.is_file()
            and not path.is_symlink()
            and not path.name.startswith("~$")
            and path.suffix.casefold() == ".pptx"
            and _normalized(path.stem) in document_names
        ]

    latest = matching_files(proposal_dir)
    old_dirs = [
        path
        for path in proposal_dir.iterdir()
        if path.is_dir() and not path.is_symlink() and _normalized(path.name) == "old"
    ]
    if len(latest) != 1 or len(old_dirs) != 1:
        return None
    old = matching_files(old_dirs[0])
    if len(old) != 1 or _normalized(old[0].name) != _normalized(latest[0].name):
        return None
    return old[0], latest[0]


def _pptx_change_descriptor(unit: Mapping[str, Any]) -> str | None:
    if unit.get("kind") == "table":
        matrix = unit.get("table_matrix")
        row = unit.get("row")
        column = unit.get("column")
        if (
            isinstance(matrix, tuple)
            and isinstance(row, int)
            and isinstance(column, int)
            and 0 <= row < len(matrix)
            and matrix
            and 0 <= column < len(matrix[0])
        ):
            parts: list[str] = []
            if row > 0 and matrix[row][0]:
                parts.append(str(matrix[row][0]))
            if column > 0 and matrix[0][column]:
                parts.append(str(matrix[0][column]))
            if parts:
                return " / ".join(parts)
    slide = unit.get("slide")
    shape = unit.get("shape")
    if isinstance(slide, int) and isinstance(shape, int):
        return f"スライド{slide}の要素{shape}"
    return None


def _pptx_old_latest_visible_text_diff(
    engine: Any,
    match: re.Match[str],
) -> StructuredCandidateDecision | None:
    pair = _proposal_pptx_pair(
        engine,
        match["location"],
        match["document_key"],
    )
    if pair is None:
        return None
    old_path, latest_path = pair
    old_model = _pptx_visible_deck_model(old_path)
    latest_model = _pptx_visible_deck_model(latest_path)
    if old_model is None or latest_model is None or old_model[0] != latest_model[0]:
        return None
    old_units, latest_units = old_model[1], latest_model[1]
    if set(old_units) != set(latest_units):
        return None
    changed = [
        key
        for key in old_units
        if old_units[key]["text"] != latest_units[key]["text"]
    ]
    if len(changed) != 1:
        return None
    key = changed[0]
    descriptor = _pptx_change_descriptor(old_units[key])
    before = str(old_units[key]["text"] or "(空欄)")
    after = str(latest_units[key]["text"] or "(空欄)")
    if (
        descriptor is None
        or before == after
        or any(len(value) > 4_000 for value in (descriptor, before, after))
    ):
        return None
    answer = f"{descriptor}: 変更前={before}、変更後={after}"
    return _decision(answer, [old_path, latest_path], engine.source_root, 7)


def _python_tree(path: Path) -> ast.Module | None:
    try:
        if path.is_symlink() or path.stat().st_size > 2 * 1024 * 1024:
            return None
        return ast.parse(path.read_text(encoding="utf-8"), filename=path.as_posix())
    except (OSError, UnicodeDecodeError, SyntaxError):
        return None


def _ast_functions(tree: ast.AST, name: str) -> list[ast.FunctionDef]:
    body = tree.body if isinstance(tree, (ast.Module, ast.FunctionDef)) else []
    return [
        node
        for node in body
        if isinstance(node, ast.FunctionDef) and node.name == name
    ]


def _ast_statement_bound_names(statement: ast.stmt) -> set[str]:
    """Collect bindings in one statement without entering nested scopes."""

    result: set[str] = set()

    class BindingVisitor(ast.NodeVisitor):
        def visit_Name(self, node: ast.Name) -> None:
            if isinstance(node.ctx, (ast.Store, ast.Del)):
                result.add(node.id)

        def visit_Import(self, node: ast.Import) -> None:
            for alias in node.names:
                result.add(alias.asname or alias.name.split(".", 1)[0])

        def visit_ImportFrom(self, node: ast.ImportFrom) -> None:
            for alias in node.names:
                if alias.name != "*":
                    result.add(alias.asname or alias.name)

        def visit_FunctionDef(self, node: ast.FunctionDef) -> None:
            result.add(node.name)

        def visit_AsyncFunctionDef(self, node: ast.AsyncFunctionDef) -> None:
            result.add(node.name)

        def visit_ClassDef(self, node: ast.ClassDef) -> None:
            result.add(node.name)

        def visit_Lambda(self, node: ast.Lambda) -> None:
            return

        def visit_ExceptHandler(self, node: ast.ExceptHandler) -> None:
            if isinstance(node.name, str):
                result.add(node.name)
            for child in node.body:
                self.visit(child)

    BindingVisitor().visit(statement)
    return result


def _ast_statement_loads_name(statement: ast.stmt, name: str) -> bool:
    return any(
        isinstance(node, ast.Name)
        and isinstance(node.ctx, ast.Load)
        and node.id == name
        for node in ast.walk(statement)
    )


def _ast_imported_names(
    tree: ast.Module,
    module: str,
    function: ast.FunctionDef | None = None,
) -> dict[str, str] | None:
    """Resolve direct imports whose local names cannot later be rebound."""

    records: list[tuple[str, int, str, str]] = []

    def collect(scope: str, body: Sequence[ast.stmt]) -> bool:
        for position, node in enumerate(body):
            if (
                not isinstance(node, ast.ImportFrom)
                or node.level != 0
                or node.module != module
            ):
                continue
            for alias in node.names:
                if alias.name == "*":
                    return False
                records.append(
                    (scope, position, alias.asname or alias.name, alias.name)
                )
        return True

    if not collect("module", tree.body):
        return None
    if function is not None and not collect("function", function.body):
        return None

    result: dict[str, str] = {}
    for _, _, local, original in records:
        if local in result:
            # Multiple binding sites are execution-order dependent even when
            # they import the same source symbol.
            return None
        result[local] = original

    function_parameters: set[str] = set()
    if function is not None:
        function_parameters = {
            argument.arg
            for argument in (
                *function.args.posonlyargs,
                *function.args.args,
                *function.args.kwonlyargs,
            )
        }
    module_function_position = next(
        (
            index
            for index, statement in enumerate(tree.body)
            if statement is function
        ),
        None,
    )
    for scope, position, local, _ in records:
        body = tree.body if scope == "module" else function.body
        if local in function_parameters:
            return None
        if scope == "module" and (
            module_function_position is None or position >= module_function_position
        ):
            return None
        if scope == "function" and any(
            _ast_statement_loads_name(statement, local)
            for statement in body[:position]
        ):
            return None
        if any(
            local in _ast_statement_bound_names(statement)
            for statement in body[position + 1 :]
        ):
            return None
        if scope == "module" and function is not None and any(
            local in _ast_statement_bound_names(statement)
            for statement in function.body
        ):
            return None
    return result


def _ast_imported_module_names(
    tree: ast.Module,
    module: str,
    function: ast.FunctionDef,
) -> set[str] | None:
    records: list[tuple[int, str]] = []
    for position, statement in enumerate(tree.body):
        if not isinstance(statement, ast.Import):
            continue
        for alias in statement.names:
            if alias.name == module:
                records.append(
                    (position, alias.asname or alias.name.split(".", 1)[0])
                )
    if len(records) != 1:
        return None
    position, local = records[0]
    function_position = next(
        (
            index
            for index, statement in enumerate(tree.body)
            if statement is function
        ),
        None,
    )
    parameters = {
        argument.arg
        for argument in (
            *function.args.posonlyargs,
            *function.args.args,
            *function.args.kwonlyargs,
        )
    }
    if (
        function_position is None
        or position >= function_position
        or local in parameters
        or any(
            local in _ast_statement_bound_names(statement)
            for statement in tree.body[position + 1 :]
        )
        or any(
            local in _ast_statement_bound_names(statement)
            for statement in function.body
        )
    ):
        return None
    return {local}


def _ast_call_name(node: ast.AST) -> str | None:
    if not isinstance(node, ast.Call) or not isinstance(node.func, ast.Name):
        return None
    if any(isinstance(argument, ast.Starred) for argument in node.args):
        return None
    if any(keyword.arg is None for keyword in node.keywords):
        return None
    return node.func.id


def _ast_direct_statement_call(statement: ast.stmt) -> ast.Call | None:
    value: ast.AST | None = None
    if isinstance(statement, (ast.Assign, ast.AnnAssign)):
        value = statement.value
    elif isinstance(statement, (ast.Expr, ast.Return)):
        value = statement.value
    return value if isinstance(value, ast.Call) else None


def _ast_single_name_target(node: ast.AST) -> str | None:
    return node.id if isinstance(node, ast.Name) else None


def _ast_tuple_first_target(node: ast.AST) -> str | None:
    if not isinstance(node, (ast.Tuple, ast.List)) or len(node.elts) < 1:
        return None
    if any(not isinstance(element, ast.Name) for element in node.elts):
        return None
    return node.elts[0].id


def _ast_module_int_constants(tree: ast.Module) -> dict[str, int] | None:
    bindings: dict[str, list[int | None]] = {}
    for statement in tree.body:
        literal_values: dict[str, int] = {}
        if (
            isinstance(statement, ast.Assign)
            and isinstance(statement.value, ast.Constant)
            and isinstance(statement.value.value, int)
            and not isinstance(statement.value.value, bool)
            and all(isinstance(target, ast.Name) for target in statement.targets)
        ):
            literal_values = {
                target.id: statement.value.value for target in statement.targets
            }
        elif (
            isinstance(statement, ast.AnnAssign)
            and isinstance(statement.target, ast.Name)
            and isinstance(statement.value, ast.Constant)
            and isinstance(statement.value.value, int)
            and not isinstance(statement.value.value, bool)
        ):
            literal_values = {statement.target.id: statement.value.value}
        for name in _ast_statement_bound_names(statement):
            bindings.setdefault(name, []).append(literal_values.get(name))
    return {
        name: values[0]
        for name, values in bindings.items()
        if len(values) == 1 and values[0] is not None
    }


def _ast_int_literal(node: ast.AST, constants: Mapping[str, int]) -> int | None:
    if isinstance(node, ast.Constant) and isinstance(node.value, int) and not isinstance(node.value, bool):
        return node.value
    if isinstance(node, ast.Name):
        return constants.get(node.id)
    return None


def _ast_function_default(
    function: ast.FunctionDef,
    parameter: str,
) -> ast.AST | None:
    positional = [*function.args.posonlyargs, *function.args.args]
    names = [argument.arg for argument in positional]
    if parameter in names:
        index = names.index(parameter)
        default_offset = len(positional) - len(function.args.defaults)
        if index >= default_offset:
            return function.args.defaults[index - default_offset]
        return None
    keyword_names = [argument.arg for argument in function.args.kwonlyargs]
    if parameter in keyword_names:
        return function.args.kw_defaults[keyword_names.index(parameter)]
    return None


def _ast_effective_int_argument(
    call: ast.Call,
    function: ast.FunctionDef,
    parameter: str,
    constants: Mapping[str, int],
) -> int | None:
    positional = [*function.args.posonlyargs, *function.args.args]
    names = [argument.arg for argument in positional]
    explicit: list[ast.AST] = []
    if parameter in names:
        index = names.index(parameter)
        if len(call.args) > index:
            explicit.append(call.args[index])
    for keyword in call.keywords:
        if keyword.arg == parameter:
            explicit.append(keyword.value)
    if len(explicit) > 1:
        return None
    expression = explicit[0] if explicit else _ast_function_default(function, parameter)
    if expression is None:
        return None
    value = _ast_int_literal(expression, constants)
    return value if value is not None and 1 <= value <= 1_000_000 else None


def _ast_dtype_predicates(
    node: ast.AST,
    series_name: str,
    pandas_names: set[str],
) -> tuple[str, ...] | None:
    if not isinstance(node, ast.BoolOp) or not isinstance(node.op, ast.Or):
        return None
    predicates: list[str] = []
    for value in node.values:
        if (
            not isinstance(value, ast.Call)
            or not isinstance(value.func, ast.Attribute)
            or len(value.args) != 1
            or value.keywords
            or not isinstance(value.args[0], ast.Name)
            or value.args[0].id != series_name
            or not isinstance(value.func.value, ast.Attribute)
            or value.func.value.attr != "types"
            or not isinstance(value.func.value.value, ast.Attribute)
            or value.func.value.value.attr != "api"
            or not isinstance(value.func.value.value.value, ast.Name)
            or value.func.value.value.value.id not in pandas_names
        ):
            return None
        predicates.append(value.func.attr)
    expected = {
        "is_object_dtype",
        "is_string_dtype",
        "is_categorical_dtype",
    }
    if set(predicates) != expected or len(predicates) != len(expected):
        return None
    labels = {
        "is_object_dtype": "object",
        "is_string_dtype": "string",
        "is_categorical_dtype": "category",
    }
    return tuple(labels[predicate] for predicate in predicates)


def _ast_dropna_nunique_series(node: ast.AST) -> str | None:
    if (
        isinstance(node, ast.Call)
        and isinstance(node.func, ast.Name)
        and node.func.id == "int"
        and len(node.args) == 1
        and not node.keywords
    ):
        node = node.args[0]
    if (
        not isinstance(node, ast.Call)
        or node.args
        or node.keywords
        or not isinstance(node.func, ast.Attribute)
        or node.func.attr != "nunique"
        or not isinstance(node.func.value, ast.Call)
    ):
        return None
    dropna = node.func.value
    if (
        dropna.args
        or dropna.keywords
        or not isinstance(dropna.func, ast.Attribute)
        or dropna.func.attr != "dropna"
        or not isinstance(dropna.func.value, ast.Name)
    ):
        return None
    return dropna.func.value.id


def _ast_selection_rule(
    tree: ast.Module,
    function: ast.FunctionDef,
    call: ast.Call,
) -> tuple[tuple[str, ...], int] | None:
    positional = [*function.args.posonlyargs, *function.args.args]
    if not positional or function.decorator_list:
        return None
    frame_name = positional[0].arg
    pandas_names = _ast_imported_module_names(tree, "pandas", function)
    if pandas_names is None:
        return None
    constants = _ast_module_int_constants(tree)
    if constants is None:
        return None
    candidates: list[tuple[tuple[str, ...], str, str]] = []
    for loop in [node for node in function.body if isinstance(node, ast.For)]:
        if (
            not isinstance(loop.target, ast.Name)
            or not isinstance(loop.iter, ast.Attribute)
            or loop.iter.attr != "columns"
            or not isinstance(loop.iter.value, ast.Name)
            or loop.iter.value.id != frame_name
        ):
            continue
        column_name = loop.target.id
        series_assignments = [
            node
            for node in loop.body
            if isinstance(node, ast.Assign)
            and len(node.targets) == 1
            and isinstance(node.targets[0], ast.Name)
            and isinstance(node.value, ast.Subscript)
            and isinstance(node.value.value, ast.Name)
            and node.value.value.id == frame_name
            and isinstance(node.value.slice, ast.Name)
            and node.value.slice.id == column_name
        ]
        if len(series_assignments) != 1:
            continue
        series_name = series_assignments[0].targets[0].id
        dtype_assignments: list[tuple[str, tuple[str, ...], int]] = []
        for position, node in enumerate(loop.body):
            if (
                isinstance(node, ast.Assign)
                and len(node.targets) == 1
                and isinstance(node.targets[0], ast.Name)
            ):
                predicates = _ast_dtype_predicates(
                    node.value, series_name, pandas_names
                )
                if predicates is not None:
                    dtype_assignments.append((node.targets[0].id, predicates, position))
        if len(dtype_assignments) != 1:
            continue
        dtype_name, predicates, dtype_position = dtype_assignments[0]
        category_ifs = [
            (position, node)
            for position, node in enumerate(loop.body)
            if isinstance(node, ast.If)
            and isinstance(node.test, ast.Name)
            and node.test.id == dtype_name
        ]
        if len(category_ifs) != 1 or category_ifs[0][0] <= dtype_position:
            continue
        category_position, category_if = category_ifs[0]
        unique_assignments = [
            node
            for node in category_if.body
            if isinstance(node, ast.Assign)
            and len(node.targets) == 1
            and isinstance(node.targets[0], ast.Name)
            and _ast_dropna_nunique_series(node.value) == series_name
        ]
        if len(unique_assignments) != 1:
            continue
        unique_name = unique_assignments[0].targets[0].id
        exclusion_ifs = []
        for node in category_if.body:
            if not isinstance(node, ast.If) or not isinstance(node.test, ast.Compare):
                continue
            comparison = node.test
            if (
                len(comparison.ops) == 1
                and isinstance(comparison.ops[0], ast.GtE)
                and len(comparison.comparators) == 1
                and isinstance(comparison.left, ast.Name)
                and comparison.left.id == unique_name
                and isinstance(comparison.comparators[0], ast.Name)
                and any(isinstance(child, ast.Continue) for child in node.body)
            ):
                exclusion_ifs.append(node)
        if len(exclusion_ifs) != 1:
            continue
        limit_name = exclusion_ifs[0].test.comparators[0].id
        append_sites: list[tuple[int, str]] = []
        for position, node in enumerate(loop.body):
            if not isinstance(node, ast.Expr) or not isinstance(node.value, ast.Call):
                continue
            append = node.value
            if (
                isinstance(append.func, ast.Attribute)
                and append.func.attr == "append"
                and isinstance(append.func.value, ast.Name)
                and len(append.args) == 1
                and isinstance(append.args[0], ast.Name)
                and append.args[0].id == column_name
                and not append.keywords
            ):
                append_sites.append((position, append.func.value.id))
        if (
            len(append_sites) != 1
            or append_sites[0][0] <= category_position
        ):
            continue
        selected_name = append_sites[0][1]
        returns = [node for node in function.body if isinstance(node, ast.Return)]
        if len(returns) != 1 or not isinstance(returns[0].value, ast.Tuple):
            continue
        first_return = returns[0].value.elts[0] if returns[0].value.elts else None
        if (
            not isinstance(first_return, ast.Call)
            or first_return.args
            or first_return.keywords
            or not isinstance(first_return.func, ast.Attribute)
            or first_return.func.attr != "copy"
            or not isinstance(first_return.func.value, ast.Subscript)
            or not isinstance(first_return.func.value.value, ast.Name)
            or first_return.func.value.value.id != frame_name
            or not isinstance(first_return.func.value.slice, ast.Name)
            or first_return.func.value.slice.id != selected_name
        ):
            continue
        candidates.append((predicates, limit_name, selected_name))
    if len(candidates) != 1:
        return None
    predicates, limit_name, _ = candidates[0]
    limit = _ast_effective_int_argument(call, function, limit_name, constants)
    return (predicates, limit) if limit is not None else None


def _ast_select_dtype_include(node: ast.AST, frame_name: str) -> tuple[str, ...] | None:
    if (
        not isinstance(node, ast.Call)
        or node.args
        or node.keywords
        or not isinstance(node.func, ast.Attribute)
        or node.func.attr != "tolist"
        or not isinstance(node.func.value, ast.Attribute)
        or node.func.value.attr != "columns"
        or not isinstance(node.func.value.value, ast.Call)
    ):
        return None
    select = node.func.value.value
    if (
        not isinstance(select.func, ast.Attribute)
        or select.func.attr != "select_dtypes"
        or not isinstance(select.func.value, ast.Name)
        or select.func.value.id != frame_name
        or select.args
    ):
        return None
    includes = [keyword.value for keyword in select.keywords if keyword.arg == "include"]
    if len(includes) != 1 or not isinstance(includes[0], (ast.List, ast.Tuple)):
        return None
    values = []
    for element in includes[0].elts:
        if not isinstance(element, ast.Constant) or not isinstance(element.value, str):
            return None
        values.append(unicodedata.normalize("NFKC", element.value).casefold())
    if set(values) != {"number", "bool"} or len(values) != 2:
        return None
    return tuple(values)


def _ast_preprocessor_rule(
    tree: ast.Module,
    function: ast.FunctionDef,
    label: str,
) -> tuple[str, ...] | None:
    positional = [*function.args.posonlyargs, *function.args.args]
    if not positional or function.decorator_list:
        return None
    frame_name = positional[0].arg
    transformer_imports = _ast_imported_names(
        tree, "sklearn.compose", function
    )
    if transformer_imports is None:
        return None
    transformer_names = {
        local
        for local, original in transformer_imports.items()
        if original == "ColumnTransformer"
    }
    if len(transformer_names) != 1:
        return None
    direct_returns = [
        node
        for node in function.body
        if isinstance(node, ast.Return)
        and isinstance(node.value, ast.Call)
        and isinstance(node.value.func, ast.Name)
        and node.value.func.id in transformer_names
    ]
    all_calls = [
        node
        for node in ast.walk(function)
        if isinstance(node, ast.Call)
        and isinstance(node.func, ast.Name)
        and node.func.id in transformer_names
    ]
    if (
        len(direct_returns) != 1
        or len(all_calls) != 1
        or all_calls[0] is not direct_returns[0].value
    ):
        return None
    call = direct_returns[0].value
    transformer_values = [keyword.value for keyword in call.keywords if keyword.arg == "transformers"]
    if not transformer_values and call.args:
        transformer_values = [call.args[0]]
    if len(transformer_values) != 1 or not isinstance(transformer_values[0], (ast.List, ast.Tuple)):
        return None
    label_matches: list[str] = []
    for entry in transformer_values[0].elts:
        if not isinstance(entry, (ast.Tuple, ast.List)) or len(entry.elts) < 3:
            return None
        entry_label = entry.elts[0]
        column_value = entry.elts[2]
        if not isinstance(entry_label, ast.Constant) or not isinstance(entry_label.value, str):
            return None
        if _normalized(entry_label.value) == _normalized(label):
            if not isinstance(column_value, ast.Name):
                return None
            label_matches.append(column_value.id)
    if len(label_matches) != 1:
        return None
    categorical_name = label_matches[0]

    numeric_assignments: list[tuple[str, tuple[str, ...]]] = []
    for node in function.body:
        if (
            isinstance(node, ast.Assign)
            and len(node.targets) == 1
            and isinstance(node.targets[0], ast.Name)
        ):
            includes = _ast_select_dtype_include(node.value, frame_name)
            if includes is not None:
                numeric_assignments.append((node.targets[0].id, includes))
    if len(numeric_assignments) != 1:
        return None
    numeric_name, numeric_types = numeric_assignments[0]
    category_assignments = [
        node
        for node in function.body
        if isinstance(node, ast.Assign)
        and len(node.targets) == 1
        and isinstance(node.targets[0], ast.Name)
        and node.targets[0].id == categorical_name
        and isinstance(node.value, ast.ListComp)
    ]
    if len(category_assignments) != 1:
        return None
    comprehension = category_assignments[0].value
    if (
        len(comprehension.generators) != 1
        or not isinstance(comprehension.elt, ast.Name)
    ):
        return None
    generator = comprehension.generators[0]
    item_name = comprehension.elt.id
    if (
        generator.is_async
        or not isinstance(generator.target, ast.Name)
        or generator.target.id != item_name
        or not isinstance(generator.iter, ast.Attribute)
        or generator.iter.attr != "columns"
        or not isinstance(generator.iter.value, ast.Name)
        or generator.iter.value.id != frame_name
        or len(generator.ifs) != 1
        or not isinstance(generator.ifs[0], ast.Compare)
    ):
        return None
    comparison = generator.ifs[0]
    if (
        not isinstance(comparison.left, ast.Name)
        or comparison.left.id != item_name
        or len(comparison.ops) != 1
        or not isinstance(comparison.ops[0], ast.NotIn)
        or len(comparison.comparators) != 1
        or not isinstance(comparison.comparators[0], ast.Name)
        or comparison.comparators[0].id != numeric_name
    ):
        return None
    return numeric_types


def _readme_python_entrypoint(readme: Path, analysis_root: Path) -> Path | None:
    try:
        if readme.stat().st_size > 2 * 1024 * 1024:
            return None
        text = readme.read_text(encoding="utf-8")
    except (OSError, UnicodeDecodeError):
        return None
    run_sections = re.findall(
        r"(?ims)^##\s+Run\s*$\r?\n(.*?)(?=^##\s|\Z)",
        text,
    )
    if len(run_sections) != 1:
        return None
    commands = re.findall(
        r"(?im)^[ \t]*(?:python(?:3(?:\.\d+)?)?|[^\s`]*[\\/]python(?:\.exe)?)"
        r"[ \t]+([^\s`]+\.py)(?:[ \t]+[^\r\n]*)?$",
        run_sections[0],
    )
    candidates: set[str] = set()
    for command in commands:
        normalized = unicodedata.normalize("NFKC", command).replace("\\", "/")
        relative = PurePosixPath(normalized)
        if relative.is_absolute() or ".." in relative.parts:
            return None
        candidates.add(relative.as_posix())
    if len(candidates) != 1:
        return None
    relative = Path(next(iter(candidates)))
    entrypoint = (analysis_root / relative).resolve()
    try:
        entrypoint.relative_to(analysis_root.resolve())
    except ValueError:
        return None
    if (
        entrypoint.name != "run_train.py"
        or not entrypoint.is_file()
        or entrypoint.is_symlink()
    ):
        return None
    return entrypoint


def _analysis_project_readme(engine: Any, location: str) -> Path | None:
    matches = [
        path
        for path in _all_files(engine, {".md"})
        if _project_matches(engine, path, location)
        and _normalized(path.name) == "readme.md"
        and _normalized(path.parent.name) == "analysis_project"
    ]
    return matches[0] if len(matches) == 1 else None


def _python_categorical_rule_facts(
    engine: Any,
    location: str,
    label: str,
) -> tuple[tuple[str, ...], int, tuple[str, ...], list[Path]] | None:
    readme = _analysis_project_readme(engine, location)
    if readme is None:
        return None
    analysis_root = readme.parent
    entrypoint = _readme_python_entrypoint(readme, analysis_root)
    modeling = analysis_root / "src" / "modeling.py"
    features = analysis_root / "src" / "features.py"
    if (
        entrypoint is None
        or not modeling.is_file()
        or modeling.is_symlink()
        or not features.is_file()
        or features.is_symlink()
    ):
        return None
    run_tree = _python_tree(entrypoint)
    modeling_tree = _python_tree(modeling)
    features_tree = _python_tree(features)
    if run_tree is None or modeling_tree is None or features_tree is None:
        return None
    mains = _ast_functions(run_tree, "main")
    if len(mains) != 1 or mains[0].decorator_list:
        return None
    main = mains[0]
    run_feature_imports = _ast_imported_names(
        run_tree, "src.features", main
    )
    run_model_imports = _ast_imported_names(run_tree, "src.modeling", main)
    if run_feature_imports is None or run_model_imports is None:
        return None
    selection_sites: list[tuple[ast.Call, str, str]] = []
    selection_statements: list[ast.Assign] = []
    for node in ast.walk(main):
        if not isinstance(node, ast.Assign) or len(node.targets) != 1:
            continue
        call_name = _ast_call_name(node.value)
        selected_name = _ast_tuple_first_target(node.targets[0])
        if (
            call_name in run_feature_imports
            and selected_name is not None
            and isinstance(node.value, ast.Call)
            and node.value.args
            and isinstance(node.value.args[0], ast.Name)
        ):
            selection_sites.append(
                (node.value, selected_name, run_feature_imports[call_name])
            )
            selection_statements.append(node)
    if len(selection_sites) != 1:
        return None
    selection_call, selected_name, selection_function_name = selection_sites[0]
    selection_statement = selection_statements[0]
    selection_positions = [
        position
        for position, statement in enumerate(main.body)
        if statement is selection_statement
    ]
    if len(selection_positions) != 1:
        return None
    selection_position = selection_positions[0]
    model_sites: list[tuple[ast.Call, str]] = []
    for node in ast.walk(main):
        if not isinstance(node, ast.Call):
            continue
        call_name = _ast_call_name(node)
        if (
            call_name in run_model_imports
            and node.args
            and isinstance(node.args[0], ast.Name)
            and node.args[0].id == selected_name
        ):
            model_sites.append((node, run_model_imports[call_name]))
    if len(model_sites) != 1:
        return None
    model_positions = [
        position
        for position, statement in enumerate(main.body)
        if _ast_direct_statement_call(statement) is model_sites[0][0]
    ]
    if len(model_positions) != 1 or model_positions[0] <= selection_position:
        return None
    model_position = model_positions[0]
    if any(
        isinstance(statement, (ast.Return, ast.Raise))
        for statement in main.body[:model_position]
    ):
        return None
    if any(
        selected_name in _ast_statement_bound_names(statement)
        for statement in main.body[selection_position + 1 : model_position]
    ):
        return None
    selection_functions = _ast_functions(features_tree, selection_function_name)
    model_functions = _ast_functions(modeling_tree, model_sites[0][1])
    if len(selection_functions) != 1 or len(model_functions) != 1:
        return None
    selection_rule = _ast_selection_rule(
        features_tree,
        selection_functions[0],
        selection_call,
    )
    if selection_rule is None:
        return None

    model_function = model_functions[0]
    if model_function.decorator_list:
        return None
    model_parameters = [*model_function.args.posonlyargs, *model_function.args.args]
    model_feature_imports = _ast_imported_names(
        modeling_tree, "src.features", model_function
    )
    if not model_parameters or model_feature_imports is None:
        return None
    model_frame = model_parameters[0].arg
    preprocessor_sites: list[tuple[ast.Call, str]] = []
    for node in ast.walk(model_function):
        if not isinstance(node, ast.Call):
            continue
        call_name = _ast_call_name(node)
        if (
            call_name in model_feature_imports
            and node.args
            and isinstance(node.args[0], ast.Name)
            and node.args[0].id == model_frame
        ):
            preprocessor_sites.append((node, model_feature_imports[call_name]))
    if len(preprocessor_sites) != 1:
        return None
    preprocessor_call, preprocessor_function_name = preprocessor_sites[0]
    preprocessor_assignments = [
        (position, statement.targets[0].id)
        for position, statement in enumerate(model_function.body)
        if isinstance(statement, ast.Assign)
        and len(statement.targets) == 1
        and isinstance(statement.targets[0], ast.Name)
        and statement.value is preprocessor_call
    ]
    if len(preprocessor_assignments) != 1:
        return None
    preprocessor_position, preprocessor_name = preprocessor_assignments[0]
    model_returns = [
        (position, statement)
        for position, statement in enumerate(model_function.body)
        if isinstance(statement, ast.Return) and statement.value is not None
    ]
    if (
        len(model_returns) != 1
        or model_returns[0][0] <= preprocessor_position
        or not _ast_statement_loads_name(model_returns[0][1], preprocessor_name)
        or any(
            preprocessor_name in _ast_statement_bound_names(statement)
            for statement in model_function.body[
                preprocessor_position + 1 : model_returns[0][0]
            ]
        )
    ):
        return None
    preprocessor_functions = _ast_functions(
        features_tree, preprocessor_function_name
    )
    if len(preprocessor_functions) != 1:
        return None
    numeric_types = _ast_preprocessor_rule(
        features_tree,
        preprocessor_functions[0],
        label,
    )
    if numeric_types is None:
        return None
    predicates, limit = selection_rule
    return predicates, limit, numeric_types, [readme, entrypoint, modeling, features]


def _python_categorical_dtype_unique_rule(
    engine: Any,
    match: re.Match[str],
) -> StructuredCandidateDecision | None:
    facts = _python_categorical_rule_facts(
        engine,
        match["location"],
        match["label"],
    )
    if facts is None:
        return None
    predicates, limit, _numeric_types, paths = facts
    label = match["label"]
    dtype_text = "/".join(predicates)
    answer = (
        f"{label}は{dtype_text} dtypeを候補とする。"
        f"欠損を除いたユニーク数が{limit}未満なら{label}として採用する。"
    )
    return _decision(answer, paths, engine.source_root, 8)


def _metric_tokens(value: str) -> tuple[str, ...]:
    tokens = re.findall(r"[A-Za-z]+[0-9]*|[0-9]+|[一-龥ぁ-んァ-ンー]+", value.casefold())
    ignored = {"score", "metric", "value", "スコア", "指標", "値", "詳細値"}
    return tuple(sorted(token for token in tokens if token not in ignored))


def _metric_value_pattern(metric: str) -> re.Pattern[str]:
    pieces = [re.escape(piece) for piece in re.split(r"\s+", metric.strip()) if piece]
    label = r"\s+".join(pieces)
    return re.compile(
        label + r"(?:\s*スコア)?\s*[:=]\s*(-?\d+(?:\.\d+)?)",
        flags=re.IGNORECASE,
    )


def _report_metric_values(
    engine: Any,
    location: str,
    report_kind: str,
    metric: str,
) -> tuple[set[Decimal], list[Path]]:
    pattern = _metric_value_pattern(metric)
    values: set[Decimal] = set()
    sources: list[Path] = []
    for path in _all_files(engine, {".docx"}):
        if not _project_matches(engine, path, location):
            continue
        normalized_parts = [_normalized(part) for part in path.relative_to(engine.source_root).parts]
        if not any("報告資料" in part for part in normalized_parts):
            continue
        paragraphs = _docx_paragraphs(path)
        if not paragraphs or _normalized(report_kind) not in _normalized("\n".join(paragraphs)):
            continue
        found = {
            Decimal(capture)
            for paragraph in paragraphs
            for capture in pattern.findall(paragraph)
        }
        if found:
            values.update(found)
            sources.append(path)
    return values, sources


def _json_metric_values(value: Any, metric_tokens: tuple[str, ...]) -> set[Decimal]:
    matches: set[Decimal] = set()
    if isinstance(value, Mapping):
        for key, child in value.items():
            if _metric_tokens(str(key)) == metric_tokens and not isinstance(child, bool):
                if isinstance(child, (Decimal, int, str)):
                    try:
                        matches.add(Decimal(str(child)))
                    except Exception:
                        pass
            matches.update(_json_metric_values(child, metric_tokens))
    elif isinstance(value, list):
        for child in value:
            matches.update(_json_metric_values(child, metric_tokens))
    return matches


def _report_metric_delta(
    engine: Any,
    match: re.Match[str],
) -> StructuredCandidateDecision | None:
    report_values, reports = _report_metric_values(
        engine,
        match["location"],
        match["report_kind"],
        match["metric"],
    )
    metrics = [
        path
        for path in _analysis_files(engine, match["location"], match["container"])
        if "analysis_outputs" in path.parts and "analysis_project" not in path.parts
    ]
    if len(report_values) != 1 or not reports or len(metrics) != 1:
        return None
    try:
        document = json.loads(
            metrics[0].read_text(encoding="utf-8"),
            parse_float=Decimal,
            parse_int=Decimal,
        )
    except Exception:
        return None
    final_values = _json_metric_values(document, _metric_tokens(match["metric"]))
    if len(final_values) != 1:
        return None
    digits = int(match["digits"])
    if digits > 18:
        return None
    difference = next(iter(final_values)) - next(iter(report_values))
    quantizer = Decimal(1).scaleb(-digits)
    answer = format(difference.quantize(quantizer, rounding=ROUND_HALF_UP), f".{digits}f")
    return _decision(answer, [*reports, metrics[0]], engine.source_root, 4)


def decide_extended(engine: Any, question_id: str, question: str) -> StructuredCandidateDecision | None:
    """Return a resolved decision for one fully matched extended grammar."""

    del question_id
    match = DATE_RANGE.fullmatch(question)
    if match:
        return _date_range(engine, match)
    match = ASSIGNEE_COUNT.fullmatch(question)
    if match:
        return _assignee_count(engine, match)
    match = PROJECT_PERSON_ASSIGNMENT_ROLE.fullmatch(question)
    if match:
        return _project_person_assignment_role(engine, match)
    match = PHASE_LATEST.fullmatch(question)
    if match:
        return _phase_latest(engine, match)
    match = BUFFER_SUM.fullmatch(question)
    if match:
        return _buffer_sum(engine, match)
    match = CHECKPOINT_TASKS.fullmatch(question)
    if match:
        return _checkpoint_tasks(engine, match)
    if MISSING_ROWS_MAX.fullmatch(question):
        return _missing_rows(engine)
    match = STANDARDIZED_SHARE.fullmatch(question)
    if match:
        return _standardized_share(engine, match)
    match = INTERACTION_COLUMNS.fullmatch(question)
    if match:
        return _interaction_columns(engine, match)
    match = GB_PARAMS.fullmatch(question)
    if match:
        return _gb_params(engine, match)
    match = XLSX_CHART_SERIES_COLUMN.fullmatch(question)
    if match:
        return _xlsx_chart_series_column(engine, match)
    match = REGRESSION_PREDICTION.fullmatch(question)
    if match:
        return _regression_prediction(engine, match)
    match = NEGATIVE_CORRELATION.fullmatch(question)
    if match:
        return _negative_correlation(engine, match)
    match = HIGHLIGHT_ROWS.fullmatch(question)
    if match:
        return _highlight_rows(engine, match)
    match = BLUE_SUM.fullmatch(question)
    if match:
        return _blue_sum(engine, match)
    match = YELLOW_INTERSECTION.fullmatch(question)
    if match:
        return _yellow_intersection(engine, match)
    match = COHORT_GROUP_MEAN_ARGMAX.fullmatch(question)
    if match:
        return _cohort_group_mean_argmax(engine, match)
    match = MULTI_FILTER_MEAN_HALF_UP.fullmatch(question)
    if match:
        return _multi_filter_mean_half_up(engine, match)
    match = NOTEBOOK_HEATMAP_MIN_ABS_CORRELATION.fullmatch(question)
    if match:
        return _notebook_heatmap_min_abs_correlation(engine, match)
    match = REPORT_METRIC_DELTA.fullmatch(question)
    if match:
        return _report_metric_delta(engine, match)
    match = CONTRACT_HOURS_RATIO_TAX_DELTA.fullmatch(question)
    if match:
        return _contract_hours_ratio_tax_delta(engine, match)
    match = EXCEL_AUTOFILTER_CONDITIONS.fullmatch(question)
    if match:
        return _excel_autofilter_conditions(engine, match)
    match = PIVOT_AVERAGE_ARGMAX_CONDITIONS_AND_AGGREGATE.fullmatch(question)
    if match:
        return _pivot_average_argmax(engine, match, include_aggregate=True)
    match = PIVOT_AVERAGE_ARGMAX_CONDITIONS.fullmatch(question)
    if match:
        return _pivot_average_argmax(engine, match, include_aggregate=False)
    match = PPTX_OLD_LATEST_VISIBLE_TEXT_DIFF.fullmatch(question)
    if match:
        return _pptx_old_latest_visible_text_diff(engine, match)
    match = PYTHON_CATEGORICAL_DTYPE_UNIQUE_RULE.fullmatch(question)
    if match:
        return _python_categorical_dtype_unique_rule(engine, match)
    match = DOCX_HIGHLIGHTED_TEXT.fullmatch(question)
    if match:
        return _docx_highlighted_text(engine, match)
    match = PPTX_SHAPE_FILL_TEXT.fullmatch(question)
    if match:
        return _pptx_shape_fill_text(engine, match)
    match = ALL_PROJECT_MILESTONE_CUTOFF.fullmatch(question)
    if match:
        return _all_project_milestone_cutoff(engine, match)
    if ALL_PROJECT_PAID_GROSS_TAX_SUM.fullmatch(question):
        return _all_project_paid_gross_tax_sum(engine)

    from analysis_artifact_rules import decide_question as decide_analysis_artifact

    analysis = decide_analysis_artifact(engine, question)
    if analysis is not None:
        return analysis

    from excel_native_rules import decide_question as decide_excel_native

    excel = decide_excel_native(engine, question)
    if excel is not None:
        return excel

    from xlsx_highlight_projection_rules import (
        decide_question as decide_xlsx_highlight,
    )

    xlsx_highlight = decide_xlsx_highlight(engine, question)
    if xlsx_highlight is not None:
        return xlsx_highlight

    from xlsx_pivot_highlight_rules import (
        decide_question as decide_xlsx_pivot_highlight,
    )

    xlsx_pivot_highlight = decide_xlsx_pivot_highlight(engine, question)
    if xlsx_pivot_highlight is not None:
        return xlsx_pivot_highlight

    from xlsx_histogram_rules import decide_question as decide_xlsx_histogram

    xlsx_histogram = decide_xlsx_histogram(engine, question)
    if xlsx_histogram is not None:
        return xlsx_histogram

    from xlsx_formula_ml_rules import decide_question as decide_xlsx_formula_ml

    xlsx_formula_ml = decide_xlsx_formula_ml(engine, question)
    if xlsx_formula_ml is not None:
        return xlsx_formula_ml

    from xlsx_version_diff_rules import decide_question as decide_xlsx_version_diff

    xlsx_version_diff = decide_xlsx_version_diff(engine, question)
    if xlsx_version_diff is not None:
        return xlsx_version_diff

    from notebook_version_diff_rules import (
        decide_question as decide_notebook_version_diff,
    )

    notebook_version_diff = decide_notebook_version_diff(engine, question)
    if notebook_version_diff is not None:
        return notebook_version_diff

    from report_metric_delta_graph_rules import decide_question as decide_report_metric_delta

    report_metric_delta = decide_report_metric_delta(engine, question)
    if report_metric_delta is not None:
        return report_metric_delta

    from contract_contact_graph_rules import decide_question as decide_contract_contact

    contract_contact = decide_contract_contact(engine, question)
    if contract_contact is not None:
        return contract_contact

    from cross_document_finance_rules import (
        decide_question as decide_cross_document_finance,
    )

    cross_document_finance = decide_cross_document_finance(engine, question)
    if cross_document_finance is not None:
        return cross_document_finance

    from cross_project_portfolio_rules import (
        decide_question as decide_cross_project_portfolio,
    )

    cross_project_portfolio = decide_cross_project_portfolio(engine, question)
    if cross_project_portfolio is not None:
        return cross_project_portfolio

    from cross_project_personnel_graph_rules import (
        decide_question as decide_cross_project_personnel,
    )

    cross_project_personnel = decide_cross_project_personnel(engine, question)
    if cross_project_personnel is not None:
        return cross_project_personnel

    from model_comparison_graph_rules import decide_question as decide_model_comparison

    model_comparison = decide_model_comparison(engine, question)
    if model_comparison is not None:
        return model_comparison

    from reported_feature_correlation_graph_rules import (
        decide_question as decide_reported_feature_correlation,
    )

    reported_feature_correlation = decide_reported_feature_correlation(engine, question)
    if reported_feature_correlation is not None:
        return reported_feature_correlation

    from priority_task_owner_graph_rules import decide_question as decide_priority_task_owner

    priority_task_owner = decide_priority_task_owner(engine, question)
    if priority_task_owner is not None:
        return priority_task_owner

    from document_answerability_rules import decide_question as decide_answerability

    answerability = decide_answerability(engine, question)
    if answerability is not None:
        return answerability

    from glossary_evidence_rules import decide_question as decide_glossary_evidence

    glossary_evidence = decide_glossary_evidence(engine, question)
    if glossary_evidence is not None:
        return glossary_evidence

    from notebook_correlation_rules import decide_question as decide_notebook_correlation

    notebook_correlation = decide_notebook_correlation(engine, question)
    if notebook_correlation is not None:
        return notebook_correlation

    from notebook_axis_tick_rules import decide_question as decide_notebook_axis_tick

    notebook_axis_tick = decide_notebook_axis_tick(engine, question)
    if notebook_axis_tick is not None:
        return notebook_axis_tick

    from xlsx_role_task_graph_rules import decide_question as decide_xlsx_role_task

    xlsx_role_task = decide_xlsx_role_task(engine, question)
    if xlsx_role_task is not None:
        return xlsx_role_task

    from pptx_schedule_rules import decide_question as decide_pptx_schedule

    pptx_schedule = decide_pptx_schedule(engine, question)
    if pptx_schedule is not None:
        return pptx_schedule

    from pptx_scope_exclusion_rules import decide_question as decide_pptx_scope_exclusion

    pptx_scope_exclusion = decide_pptx_scope_exclusion(engine, question)
    if pptx_scope_exclusion is not None:
        return pptx_scope_exclusion

    from pptx_feature_legend_rules import decide_question as decide_pptx_feature_legend

    pptx_feature_legend = decide_pptx_feature_legend(engine, question)
    if pptx_feature_legend is not None:
        return pptx_feature_legend

    from pptx_revision_summary_rules import decide_question as decide_pptx_revision_summary

    pptx_revision_summary = decide_pptx_revision_summary(engine, question)
    if pptx_revision_summary is not None:
        return pptx_revision_summary

    from pdf_operational_role_rules import (
        decide_question as decide_pdf_operational_role,
    )

    pdf_operational_role = decide_pdf_operational_role(engine, question)
    if pdf_operational_role is not None:
        return pdf_operational_role

    from pdf_native_style_rules import decide_question as decide_pdf_native_style

    pdf_native_style = decide_pdf_native_style(engine, question)
    if pdf_native_style is not None:
        return pdf_native_style

    from pdf_investment_coefficient_rules import (
        decide_question as decide_pdf_investment_coefficient,
    )

    pdf_investment_coefficient = decide_pdf_investment_coefficient(engine, question)
    if pdf_investment_coefficient is not None:
        return pdf_investment_coefficient

    from pdf_highlight_trend_rules import decide_question as decide_pdf_highlight_trend

    pdf_highlight_trend = decide_pdf_highlight_trend(engine, question)
    if pdf_highlight_trend is not None:
        return pdf_highlight_trend

    from pdf_action_transition_rules import decide_question as decide_pdf_action_transition

    pdf_action_transition = decide_pdf_action_transition(engine, question)
    if pdf_action_transition is not None:
        return pdf_action_transition

    from pdf_action_content_graph_rules import decide_question as decide_pdf_action_content

    pdf_action_content = decide_pdf_action_content(engine, question)
    if pdf_action_content is not None:
        return pdf_action_content

    from docx_page_structure_rules import decide_question as decide_docx_page_structure

    docx_page_structure = decide_docx_page_structure(engine, question)
    if docx_page_structure is not None:
        return docx_page_structure

    from docx_native_style_rules import (
        decide_question as decide_docx_native_style,
    )

    docx_native_style = decide_docx_native_style(engine, question)
    if docx_native_style is not None:
        return docx_native_style

    from docx_mixed_content_rules import decide_question as decide_docx_mixed

    docx_mixed = decide_docx_mixed(engine, question)
    if docx_mixed is not None:
        return docx_mixed

    from pptx_mixed_content_rules import decide_question as decide_pptx_mixed

    pptx_mixed = decide_pptx_mixed(engine, question)
    if pptx_mixed is not None:
        return pptx_mixed

    from pptx_version_diff_rules import decide_question as decide_pptx_version_diff

    pptx_version_diff = decide_pptx_version_diff(engine, question)
    if pptx_version_diff is not None:
        return pptx_version_diff

    from pptx_spatial_rules import decide_question as decide_pptx_spatial

    pptx_spatial = decide_pptx_spatial(engine, question)
    if pptx_spatial is not None:
        return pptx_spatial

    # Existing extended rules retain precedence.  Independent source-specific
    # executors are reachable only after their deterministic graph contracts
    # have been reconstructed.
    from proposal_metric_rules import decide_extended as decide_proposal_metric

    proposal = decide_proposal_metric(engine, "graph-runtime", question)
    if proposal is not None:
        return proposal

    from pdf_visual_rules import decide_pdf_visual

    return decide_pdf_visual(engine, question)


__all__ = [
    "GRAPH_RULE_VERSION",
    "decide_extended",
    "graph_contract_for_question",
    "validate_graph_contract",
]
