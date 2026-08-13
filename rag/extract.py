"""共有ドライブ配下のファイルからテキストと非テキスト資産を抽出する.

設計方針: **情報は捨てない。分離して記録する。**

配布データには、テキスト抽出を意図的にすり抜けるよう作られた情報が含まれる
（画像として埋め込まれた統計表、テキスト層のないPDF、書式でしか表現されない強調など）。
これらを「ノイズ」として破棄すると、解ける可能性そのものを失う。

そこでこのモジュールは、扱えない情報を捨てるのではなく

  1. 実体を assets/ に取り出して保全する
  2. 索引には「ここに未処理の資産がある」というスタブを残す
  3. 何をどれだけ取りこぼしたかを台帳（coverage）に記録する

という方針をとる。後段でOCRやVLMを追加したときに、対象がすぐ分かる状態にしておく。
"""

from __future__ import annotations

import base64
import csv
import hashlib
import io
import json
import re
import unicodedata
import zipfile
from dataclasses import dataclass, asdict, field
from pathlib import Path

MAX_TABLE_ROWS = 60      # 表形式データから読む最大行数
MAX_CELL_CHARS = 200     # 1セルの最大文字数
MAX_CELL_SRC = 6000      # ノートブック1セルの最大文字数
MAX_ASSET_BYTES = 40 * 1024 * 1024

SKIP_DIRS = {".git", "__pycache__", ".ipynb_checkpoints", "node_modules"}
IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"}
IGNORE_SUFFIXES = {".pyc", ".lock"}

DATA_URI_RE = re.compile(r"data:image/(\w+);base64,([A-Za-z0-9+/=\s]+)")


@dataclass
class Section:
    """1ファイル内の 1 まとまり（ページ/スライド/シート単位）."""

    path: str
    project: str
    filename: str
    kind: str        # docx / pptx / xlsx / pdf / code / data / text / asset
    location: str
    text: str

    def to_dict(self) -> dict:
        return asdict(self)


@dataclass
class Coverage:
    """あるファイルについて、何を抽出し何を残したかの記録."""

    path: str
    status: str          # extracted / partial / deferred / failed
    detail: str = ""
    assets: list = field(default_factory=list)


class Context:
    """抽出中の共有状態（資産の保存先と台帳）."""

    def __init__(self, assets_dir: Path):
        self.assets_dir = assets_dir
        self.assets_dir.mkdir(parents=True, exist_ok=True)
        self.coverage: list = []

    def save_asset(self, data: bytes, ext: str, origin: str) -> str | None:
        """埋め込み資産をファイルとして取り出す。戻り値は assets/ 内の名前."""
        if not data or len(data) > MAX_ASSET_BYTES:
            return None
        digest = hashlib.sha1(data).hexdigest()[:12]
        name = f"{digest}.{ext}"
        target = self.assets_dir / name
        if not target.exists():
            target.write_bytes(data)
        manifest = self.assets_dir / "manifest.jsonl"
        with manifest.open("a", encoding="utf-8") as fh:
            fh.write(json.dumps(
                {"asset": name, "bytes": len(data), "origin": origin},
                ensure_ascii=False,
            ) + "\n")
        return name

    def record(self, cov: Coverage) -> None:
        self.coverage.append(cov)


def nfc(s: str) -> str:
    """macOS のファイル名は NFD なので、比較前に NFC へ揃える."""
    return unicodedata.normalize("NFC", s)


# --------------------------------------------------------------------------
# パスワード保護ファイルの復号
# --------------------------------------------------------------------------

def password_candidates(path: Path, aliases, dates) -> list:
    """規則ベースとファイル名ベースの両方からパスワード候補を作る.

    社内規定の導出形式 DA-[略号]-[YYYYMMDD]-[拡張子] に加えて、
    ファイル名そのものに埋め込まれた文字列も候補に含める
    （規定どおりでないファイルが実在するため）。
    """
    ext = path.suffix.lstrip(".").lower()
    cands = []
    for m in re.finditer(r"(?:pw|pass|password)[-_]?([0-9A-Za-z]{6,})", path.stem, re.I):
        cands.append(m.group(1))
    for m in re.finditer(r"([A-Za-z]{2,}\d{8})", path.stem):
        cands.append(m.group(1))
    for a in aliases:
        for d in dates:
            cands.append(f"DA-{a}-{d}-{ext}")
            cands.append(f"DA-{a.upper()}-{d}-{ext}")
    seen, out = set(), []
    for c in cands:
        if c not in seen:
            seen.add(c)
            out.append(c)
    return out


def try_decrypt(path: Path, candidates):
    try:
        import msoffcrypto
    except ImportError:
        return None
    for pw in candidates:
        try:
            buf = io.BytesIO()
            with open(path, "rb") as fh:
                office = msoffcrypto.OfficeFile(fh)
                office.load_key(password=pw)
                office.decrypt(buf)
            buf.seek(0)
            return buf.getvalue()
        except Exception:
            continue
    return None


def _as_stream(path: Path, decrypted):
    return io.BytesIO(decrypted) if decrypted else str(path)


# --------------------------------------------------------------------------
# 形式別パーサ  戻り値: (blocks, coverage_detail, asset_names)
# --------------------------------------------------------------------------

def parse_docx(src, ctx: Context, origin: str):
    from docx import Document

    doc = Document(src)
    parts, truncated = [], 0
    for p in doc.paragraphs:
        t = p.text.strip()
        if t:
            parts.append(t)
    for i, table in enumerate(doc.tables, 1):
        rows = []
        for r in table.rows[:MAX_TABLE_ROWS]:
            cells = []
            for c in r.cells:
                v = c.text.strip()
                if len(v) > MAX_CELL_CHARS:
                    truncated += 1
                cells.append(v[:MAX_CELL_CHARS])
            if any(cells):
                rows.append(" | ".join(cells))
        if len(table.rows) > MAX_TABLE_ROWS:
            truncated += len(table.rows) - MAX_TABLE_ROWS
        if rows:
            parts.append(f"[表{i}]\n" + "\n".join(rows))
    assets = _extract_ooxml_media(src, ctx, origin)
    if assets:
        parts.append("[埋め込み画像] " + " ".join(assets))
    detail = f"切り詰め{truncated}箇所" if truncated else ""
    return ([("本文", "\n".join(parts))] if parts else []), detail, assets


def parse_pptx(src, ctx: Context, origin: str):
    from pptx import Presentation

    prs = Presentation(src)
    out = []
    for idx, slide in enumerate(prs.slides, 1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    parts.append(t)
            if getattr(shape, "has_table", False):
                rows = []
                for r in shape.table.rows:
                    cells = [c.text.strip()[:MAX_CELL_CHARS] for c in r.cells]
                    if any(cells):
                        rows.append(" | ".join(cells))
                if rows:
                    parts.append("[表]\n" + "\n".join(rows))
        if slide.has_notes_slide:
            note = slide.notes_slide.notes_text_frame.text.strip()
            if note:
                parts.append(f"[ノート] {note}")
        if parts:
            out.append((f"スライド{idx}", "\n".join(parts)))
    assets = _extract_ooxml_media(src, ctx, origin)
    if assets:
        out.append(("埋め込み資産", "[埋め込み画像] " + " ".join(assets)))
    return out, "", assets


def parse_xlsx(src, ctx: Context, origin: str):
    from openpyxl import load_workbook

    wb = load_workbook(src, data_only=True)
    out, truncated = [], 0
    for ws in wb.worksheets:
        rows = []
        for row in ws.iter_rows(max_row=MAX_TABLE_ROWS, values_only=True):
            cells = ["" if v is None else str(v)[:MAX_CELL_CHARS] for v in row]
            if any(cells):
                rows.append(" | ".join(cells))
        if ws.max_row and ws.max_row > MAX_TABLE_ROWS:
            truncated += ws.max_row - MAX_TABLE_ROWS
        meta = []
        if ws.auto_filter and ws.auto_filter.ref:
            conds = []
            for fc in ws.auto_filter.filterColumn:
                vals = getattr(getattr(fc, "filters", None), "filter", None) or []
                conds.append(f"列{fc.colId}: {', '.join(map(str, vals))}")
            if conds:
                meta.append("[フィルタ条件] " + " / ".join(conds))
        body = "\n".join(meta + rows)
        if body.strip():
            out.append((ws.title, body))
    wb.close()
    assets = _extract_ooxml_media(src, ctx, origin)
    if assets:
        out.append(("埋め込み資産", "[埋め込み画像] " + " ".join(assets)))
    detail = f"{truncated}行を未読（先頭{MAX_TABLE_ROWS}行のみ）" if truncated else ""
    return out, detail, assets


def _extract_ooxml_media(src, ctx: Context, origin: str) -> list:
    """OOXML パッケージ内の media/ を資産として取り出す."""
    names = []
    try:
        z = zipfile.ZipFile(src if isinstance(src, str) else io.BytesIO(src.getvalue()))
    except Exception:
        return names
    for n in z.namelist():
        if "/media/" not in n:
            continue
        ext = n.rsplit(".", 1)[-1].lower()
        if ext not in {"png", "jpg", "jpeg", "gif", "emf", "wmf", "bmp"}:
            continue
        try:
            data = z.read(n)
        except Exception:
            continue
        name = ctx.save_asset(data, ext, f"{origin}::{n}")
        if name:
            names.append(name)
    return names


def parse_pdf(path: Path, ctx: Context, origin: str):
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    out, empty = [], 0
    for i, page in enumerate(reader.pages, 1):
        try:
            t = (page.extract_text() or "").strip()
        except Exception:
            t = ""
        if t:
            out.append((f"p.{i}", t))
        else:
            empty += 1
    if not out:
        # テキスト層が無いPDF。索引から消さず、未処理として登録する。
        stub = (
            f"[テキスト層なし・OCR未実施] このPDFは全{len(reader.pages)}ページとも "
            f"テキストを抽出できない画像PDFです。内容を読むにはOCRが必要です。"
        )
        return [("全体", stub)], f"画像PDF {len(reader.pages)}ページ・OCR未実施", []
    detail = f"{empty}ページがテキスト抽出不可" if empty else ""
    return out, detail, []


def parse_ipynb(path: Path, ctx: Context, origin: str):
    """ノートブックを解析する.

    埋め込み画像（markdownのdata URI・出力のimage/png）は、
    base65のままでは検索の役に立たないが、内容そのものは資料である。
    そこで実体を assets/ に取り出し、本文には参照スタブだけを残す。
    """
    nb = json.loads(path.read_text(encoding="utf-8", errors="ignore"))
    parts, assets, truncated = [], [], 0

    def pull_data_uris(text: str, where: str) -> str:
        def repl(m):
            ext = m.group(1).lower()
            try:
                data = base64.b64decode(re.sub(r"\s", "", m.group(2)))
            except Exception:
                return "[埋め込み画像(復号失敗)]"
            name = ctx.save_asset(data, ext, f"{origin}::{where}")
            if not name:
                return "[埋め込み画像(サイズ超過)]"
            assets.append(name)
            return f"[埋め込み画像 asset={name} 出典={where} ※OCR/VLM未実施]"
        return DATA_URI_RE.sub(repl, text)

    for i, cell in enumerate(nb.get("cells", []), 1):
        raw = "".join(cell.get("source", []))
        src = pull_data_uris(raw, f"cell{i}").strip()
        if len(src) > MAX_CELL_SRC:
            truncated += 1
            src = src[:MAX_CELL_SRC] + " …(以下略)"
        if src:
            parts.append(f"[cell{i} {cell.get('cell_type')}]\n{src}")
        for o in cell.get("outputs", []):
            data = o.get("data", {})
            for mime, payload in data.items():
                if mime.startswith("image/"):
                    ext = mime.split("/")[-1]
                    try:
                        blob = base64.b64decode(
                            payload if isinstance(payload, str) else "".join(payload)
                        )
                    except Exception:
                        continue
                    name = ctx.save_asset(blob, ext, f"{origin}::cell{i}出力")
                    if name:
                        assets.append(name)
                        parts.append(
                            f"[cell{i} 出力画像 asset={name} ※OCR/VLM未実施]"
                        )
            txt = "".join(o.get("text", [])) or "".join(
                data.get("text/plain", []) if isinstance(data.get("text/plain"), list)
                else [data.get("text/plain", "")]
            )
            if txt and txt.strip():
                parts.append(f"[cell{i} output]\n{txt.strip()[:3000]}")

    detail = []
    if assets:
        detail.append(f"画像{len(assets)}点を assets/ へ保全")
    if truncated:
        detail.append(f"{truncated}セルを切り詰め")
    return ([("notebook", "\n\n".join(parts))] if parts else []), " / ".join(detail), assets


def parse_table_file(path: Path, ctx: Context, origin: str):
    delim = "\t" if path.suffix.lower() == ".tsv" else ","
    total = 0
    with path.open(encoding="utf-8", errors="ignore", newline="") as fh:
        reader = csv.reader(fh, delimiter=delim)
        rows = []
        for i, row in enumerate(reader):
            total = i + 1
            if i < MAX_TABLE_ROWS:
                rows.append(" | ".join(c[:MAX_CELL_CHARS] for c in row))
    if not rows:
        return [], "", []
    head = (
        f"[列] {rows[0]}\n[全{total}行のうち先頭{len(rows) - 1}行]\n"
        + "\n".join(rows[1:])
    )
    detail = f"{total - len(rows)}行を未読" if total > len(rows) else ""
    return [("表", head)], detail, []


def parse_plain(path: Path, ctx: Context, origin: str):
    t = path.read_text(encoding="utf-8", errors="ignore").strip()
    return ([("本文", t)] if t else []), "", []


def parse_image(path: Path, ctx: Context, origin: str):
    """画像ファイル本体。索引から消えないようスタブを登録する."""
    size = ""
    try:
        from PIL import Image

        with Image.open(path) as im:
            size = f"{im.size[0]}x{im.size[1]}"
    except Exception:
        pass
    stub = (
        f"[画像ファイル {nfc(path.name)} {size}] "
        f"この画像の内容は未解析です。読むにはOCRまたは画像認識が必要です。"
    )
    return [("画像", stub)], "OCR/VLM未実施", []


# --------------------------------------------------------------------------
# 走査
# --------------------------------------------------------------------------

KIND_BY_SUFFIX = {
    ".docx": "docx", ".pptx": "pptx", ".xlsx": "xlsx", ".pdf": "pdf",
    ".ipynb": "code", ".py": "code", ".json": "code", ".toml": "code",
    ".csv": "data", ".tsv": "data",
    ".md": "text", ".txt": "text",
}


def _project_of(rel: Path) -> str:
    parts = [nfc(p) for p in rel.parts]
    for i, p in enumerate(parts):
        if p == "プロジェクト" and i + 1 < len(parts):
            return parts[i + 1]
    return ""


def _aliases_for(project: str, glossary) -> list:
    if not project or glossary is None:
        return []
    project = nfc(project)
    out = []
    for alias, canons in glossary.entries.items():
        for c in canons:
            c = nfc(c)
            if c and (c in project or project in c):
                if alias not in out:
                    out.append(alias)
    return out


DATE_RE = re.compile(r"(20\d{2})[-/年\.]?(\d{1,2})[-/月\.]?(\d{1,2})")


def _dates_in(texts) -> list:
    out = []
    for t in texts:
        for y, m, d in DATE_RE.findall(t):
            s = f"{y}{int(m):02d}{int(d):02d}"
            if s not in out:
                out.append(s)
    return out


def _is_encrypted(path: Path) -> bool:
    if path.suffix.lower() not in {".docx", ".pptx", ".xlsx"}:
        return False
    return not zipfile.is_zipfile(path)


def extract_all(share_root: Path, glossary=None, verbose=True, assets_dir=None):
    """share_root 配下を走査して (sections, coverage) を返す（2パス）."""
    ctx = Context(assets_dir or (share_root.parent / "rag" / "assets"))

    files = []
    for p in sorted(share_root.rglob("*")):
        if not p.is_file() or p.name.startswith("~$"):
            continue
        if any(d in SKIP_DIRS for d in p.parts):
            continue
        suf = p.suffix.lower()
        if suf in IGNORE_SUFFIXES:
            continue
        if suf in KIND_BY_SUFFIX or suf in IMAGE_SUFFIXES:
            files.append(p)

    sections: list = []
    locked: list = []

    def handle(path: Path, data) -> bool:
        rel = path.relative_to(share_root)
        origin = nfc(str(rel))
        suf = path.suffix.lower()
        kind = "asset" if suf in IMAGE_SUFFIXES else KIND_BY_SUFFIX[suf]
        try:
            if suf in IMAGE_SUFFIXES:
                blocks, detail, assets = parse_image(path, ctx, origin)
            elif kind == "docx":
                blocks, detail, assets = parse_docx(_as_stream(path, data), ctx, origin)
            elif kind == "pptx":
                blocks, detail, assets = parse_pptx(_as_stream(path, data), ctx, origin)
            elif kind == "xlsx":
                blocks, detail, assets = parse_xlsx(_as_stream(path, data), ctx, origin)
            elif kind == "pdf":
                blocks, detail, assets = parse_pdf(path, ctx, origin)
            elif suf == ".ipynb":
                blocks, detail, assets = parse_ipynb(path, ctx, origin)
            elif kind == "data":
                blocks, detail, assets = parse_table_file(path, ctx, origin)
            else:
                blocks, detail, assets = parse_plain(path, ctx, origin)
        except Exception as e:
            if _is_encrypted(path):
                return False
            ctx.record(Coverage(origin, "failed", f"{type(e).__name__}: {e}"))
            if verbose:
                print(f"  ! 解析失敗 {origin}: {type(e).__name__}")
            return True
        for loc, text in blocks:
            sections.append(Section(
                path=origin, project=_project_of(rel), filename=nfc(path.name),
                kind=kind, location=loc, text=text,
            ))
        status = "deferred" if kind == "asset" else ("partial" if detail else "extracted")
        ctx.record(Coverage(origin, status, detail, assets))
        return True

    for path in files:
        if not handle(path, None):
            locked.append(path)

    if verbose:
        print(f"  抽出: {len(sections)} セクション / 保護ファイル {len(locked)} 件")

    for path in locked:
        rel = path.relative_to(share_root)
        origin = nfc(str(rel))
        project = _project_of(rel)
        aliases = _aliases_for(project, glossary)
        same = [s.text for s in sections if s.project == project]
        dates = _dates_in(same) or _dates_in([s.text for s in sections])
        data = try_decrypt(path, password_candidates(path, aliases, dates))
        if data is None:
            ctx.record(Coverage(origin, "failed", "パスワード保護・復号できず"))
            if verbose:
                print(f"  ! 復号できず {origin}")
            continue
        if verbose:
            print(f"  + 復号成功 {origin}")
        handle(path, data)

    return sections, ctx.coverage


def write_coverage(coverage, path: Path) -> None:
    """抽出カバレッジ台帳を Markdown で書き出す."""
    order = {"failed": 0, "deferred": 1, "partial": 2, "extracted": 3}
    label = {
        "extracted": "抽出済み", "partial": "一部のみ",
        "deferred": "未処理", "failed": "失敗",
    }
    rows = sorted(coverage, key=lambda c: (order.get(c.status, 9), c.path))
    counts = {}
    for c in coverage:
        counts[c.status] = counts.get(c.status, 0) + 1
    lines = [
        "# 抽出カバレッジ台帳", "",
        "各ファイルについて、何を取り込み何を残したかの記録。",
        "「未処理」「一部のみ」「失敗」は、現時点で回答に使えていない情報を表す。", "",
        "| 状態 | 件数 |", "|---|---|",
    ]
    for k in ("failed", "deferred", "partial", "extracted"):
        if k in counts:
            lines.append(f"| {label[k]} | {counts[k]} |")
    lines += ["", "| 状態 | ファイル | 内容 | 保全した資産 |", "|---|---|---|---|"]
    for c in rows:
        assets = f"{len(c.assets)}点" if c.assets else ""
        lines.append(f"| {label.get(c.status, c.status)} | `{c.path}` | {c.detail} | {assets} |")
    path.write_text("\n".join(lines) + "\n", encoding="utf-8")


if __name__ == "__main__":
    import sys
    from glossary import build_glossary

    root = Path(sys.argv[1]).resolve()
    gl = build_glossary(root)
    secs, cov = extract_all(root, gl)
    print(f"合計 {len(secs)} セクション / {len(set(s.path for s in secs))} ファイル")
