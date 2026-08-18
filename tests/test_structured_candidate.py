from __future__ import annotations

import csv
import copy
import json
import sys
import tempfile
import unittest
import zipfile
from dataclasses import replace
from pathlib import Path
from types import SimpleNamespace
from unittest.mock import patch


ROOT = Path(__file__).resolve().parents[1]
RAG = ROOT / "rag"
if str(RAG) not in sys.path:
    sys.path.insert(0, str(RAG))

from glossary import Glossary  # noqa: E402
import score_candidate_rules as candidate_rules  # noqa: E402
from score_candidate_rules import (  # noqa: E402
    decide_extended,
    graph_contract_for_question,
    validate_graph_contract,
)
from structured_candidate import (  # noqa: E402
    StructuredCandidateEngine,
    _table_from_matrix,
)


def write_csv(path: Path, rows: list[list[object]]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("w", encoding="utf-8", newline="") as handle:
        writer = csv.writer(handle)
        writer.writerows(rows)


def write_marker_workbook(path: Path, marker: str) -> None:
    from openpyxl import Workbook

    path.parent.mkdir(parents=True, exist_ok=True)
    workbook = Workbook()
    workbook.active["A1"] = marker
    workbook.save(path)
    workbook.close()


def write_opaque_chart_ex_workbook(
    path: Path,
    series_label: str,
    *,
    duplicate_chart: bool = False,
    second_series: bool = False,
) -> None:
    """Write the minimal OOXML parts needed for a chart-series fixture."""

    path.parent.mkdir(parents=True, exist_ok=True)
    duplicate_anchor = ""
    if duplicate_chart:
        duplicate_anchor = (
            '<xdr:twoCellAnchor><xdr:graphicFrame><xdr:nvGraphicFramePr>'
            '<xdr:cNvPr id="3" name="グラフ 7"/>'
            '</xdr:nvGraphicFramePr><a:graphic><a:graphicData>'
            '<cx:chart r:id="rId2"/>'
            '</a:graphicData></a:graphic></xdr:graphicFrame></xdr:twoCellAnchor>'
        )
    extra_series = ""
    if second_series:
        extra_series = (
            '<cx:series><cx:tx><cx:txData><cx:v>other_field</cx:v>'
            '</cx:txData></cx:tx></cx:series>'
        )
    parts = {
        "xl/workbook.xml": (
            '<?xml version="1.0" encoding="UTF-8"?>'
            '<workbook xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" '
            'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
            '<sheets><sheet name="DataView" sheetId="42" r:id="rId5"/></sheets>'
            '</workbook>'
        ),
        "xl/_rels/workbook.xml.rels": (
            '<?xml version="1.0" encoding="UTF-8"?>'
            '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
            '<Relationship Id="rId5" '
            'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/worksheet" '
            'Target="worksheets/sheet42.xml"/>'
            '</Relationships>'
        ),
        "xl/worksheets/sheet42.xml": (
            '<?xml version="1.0" encoding="UTF-8"?>'
            '<worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" '
            'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
            '<sheetData/><drawing r:id="rId9"/></worksheet>'
        ),
        "xl/worksheets/_rels/sheet42.xml.rels": (
            '<?xml version="1.0" encoding="UTF-8"?>'
            '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
            '<Relationship Id="rId9" '
            'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/drawing" '
            'Target="../drawings/drawingZ.xml"/>'
            '</Relationships>'
        ),
        "xl/drawings/drawingZ.xml": (
            '<?xml version="1.0" encoding="UTF-8"?>'
            '<xdr:wsDr xmlns:xdr="http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing" '
            'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
            'xmlns:cx="http://schemas.microsoft.com/office/drawing/2014/chartex" '
            'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
            '<xdr:twoCellAnchor><xdr:graphicFrame><xdr:nvGraphicFramePr>'
            '<xdr:cNvPr id="2" name="グラフ 7"/>'
            '</xdr:nvGraphicFramePr><a:graphic><a:graphicData>'
            '<cx:chart r:id="rId2"/>'
            '</a:graphicData></a:graphic></xdr:graphicFrame></xdr:twoCellAnchor>'
            f'{duplicate_anchor}</xdr:wsDr>'
        ),
        "xl/drawings/_rels/drawingZ.xml.rels": (
            '<?xml version="1.0" encoding="UTF-8"?>'
            '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
            '<Relationship Id="rId2" '
            'Type="http://schemas.microsoft.com/office/2014/relationships/chartEx" '
            'Target="../charts/chartExY.xml"/>'
            '</Relationships>'
        ),
        "xl/charts/chartExY.xml": (
            '<?xml version="1.0" encoding="UTF-8"?>'
            '<cx:chartSpace xmlns:cx="http://schemas.microsoft.com/office/drawing/2014/chartex">'
            '<cx:chart><cx:plotArea><cx:plotAreaRegion>'
            f'<cx:series><cx:tx><cx:txData><cx:f>_xlchart.opaque</cx:f><cx:v>{series_label}</cx:v>'
            f'</cx:txData></cx:tx></cx:series>{extra_series}'
            '</cx:plotAreaRegion></cx:plotArea></cx:chart></cx:chartSpace>'
        ),
    }
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as archive:
        for name, value in parts.items():
            archive.writestr(name, value)


def synthetic_tabular_pivot_workbook(
    *,
    beta_raw: int = 3,
    beta_display: int | None = None,
) -> object:
    """Build an opaque Pivot object without embedding an answer fixture."""

    from openpyxl import Workbook

    workbook = Workbook()
    pivot_sheet = workbook.active
    pivot_sheet.title = "summary_z"
    raw_sheet = workbook.create_sheet("records_z")
    raw_rows = [
        ["cohort_z", "flag_q", "band_x", "measure_u", "other_v"],
        ["alpha", "on", 1, 4, 30],
        ["alpha", "on", 1, 8, 30],
        ["alpha", "on", 2, 10, 40],
        ["beta", "off", 3, beta_raw, 50],
    ]
    for row in raw_rows:
        raw_sheet.append(row)

    beta_pivot = beta_raw if beta_display is None else beta_display
    pivot_rows = [
        ["cohort_z", "flag_q", "band_x", "平均 / measure_u", "平均 / other_v"],
        ["alpha", "on", 1, 6, 30],
        [None, None, 2, 10, 40],
        ["alpha 集計", None, None, 999, 999],
        ["beta", "off", 3, beta_pivot, 50],
        ["総計", None, None, 999, 999],
    ]
    for row_number, row in enumerate(pivot_rows, 3):
        for column, value in enumerate(row, 1):
            pivot_sheet.cell(row_number, column, value)

    index = lambda value: SimpleNamespace(v=value)
    items = [
        SimpleNamespace(t="data", r=0, x=[index(0), index(0), index(0)]),
        SimpleNamespace(t="data", r=2, x=[index(1)]),
        SimpleNamespace(t="default", r=0, x=[index(0)]),
        SimpleNamespace(t="data", r=0, x=[index(1), index(1), index(2)]),
        SimpleNamespace(t="grand", r=0, x=[index(0)]),
    ]
    cache_fields = [
        SimpleNamespace(name=name)
        for name in ("cohort_z", "flag_q", "band_x", "measure_u", "other_v")
    ]
    pivot = SimpleNamespace(
        compact=False,
        compactData=False,
        outline=False,
        multipleFieldFilters=False,
        dataOnRows=False,
        pageFields=[],
        filters=[],
        colFields=[SimpleNamespace(x=-2)],
        rowFields=[SimpleNamespace(x=index) for index in (0, 1, 2)],
        dataFields=[
            SimpleNamespace(
                name="平均 / measure_u",
                fld=3,
                subtotal="average",
                showDataAs="normal",
            ),
            SimpleNamespace(
                name="平均 / other_v",
                fld=4,
                subtotal="average",
                showDataAs="normal",
            ),
        ],
        rowItems=items,
        location=SimpleNamespace(ref="A3:E8", firstDataCol=3),
        cache=SimpleNamespace(
            cacheFields=cache_fields,
            cacheSource=SimpleNamespace(
                type="worksheet",
                connectionId=None,
                consolidation=None,
                worksheetSource=SimpleNamespace(
                    sheet="records_z",
                    ref="A1:E5",
                    name=None,
                ),
            ),
        ),
    )
    pivot_sheet._pivots = [pivot]
    return workbook


def synthetic_compact_pivot_workbook(
    *,
    sales_raw: int = 150,
    sales_display: int | None = None,
    second_average: bool = False,
    broken_indent: bool = False,
) -> object:
    from openpyxl import Workbook
    from openpyxl.styles import Alignment

    workbook = Workbook()
    pivot_sheet = workbook.active
    pivot_sheet.title = "layers_q"
    raw_sheet = workbook.create_sheet("records_q")
    headers = ["status_p", "segment_q", "tier_r", "domain_s", "income_z"]
    if second_average:
        headers.append("measure_y")
    raw_rows = [
        ["No", "Female", "Single", "Human", 100],
        ["No", "Female", "Single", "Tech", 200],
        ["Yes", "Male", "Divorced", "Sales", sales_raw],
    ]
    if second_average:
        raw_rows = [row + [index * 10] for index, row in enumerate(raw_rows, 1)]
    raw_sheet.append(headers)
    for row in raw_rows:
        raw_sheet.append(row)

    display_sales = sales_raw if sales_display is None else sales_display
    pivot_rows = [
        ("No", 0, 150),
        ("Female", 1, 150),
        ("Single", 2, 150),
        ("Human", 3, 100),
        ("Tech", 3, 200),
        ("Yes", 0, display_sales),
        ("Male", 1, display_sales),
        ("Divorced", 2, display_sales),
        ("Sales", 3, display_sales),
    ]
    pivot_sheet.cell(3, 1, "row_labels")
    pivot_sheet.cell(3, 2, "平均 / income_z")
    if second_average:
        pivot_sheet.cell(3, 3, "平均 / measure_y")
    for row_number, (label, indent, value) in enumerate(pivot_rows, 4):
        pivot_sheet.cell(row_number, 1, label)
        effective_indent = 2 if broken_indent and label == "Human" else indent
        pivot_sheet.cell(row_number, 1).alignment = Alignment(indent=effective_indent)
        pivot_sheet.cell(row_number, 2, value)
        if second_average:
            pivot_sheet.cell(row_number, 3, 10)
    grand_row = 4 + len(pivot_rows)
    pivot_sheet.cell(grand_row, 1, "総計")
    pivot_sheet.cell(grand_row, 2, 999)
    if second_average:
        pivot_sheet.cell(grand_row, 3, 999)

    index = lambda value: SimpleNamespace(v=value)
    row_items = [
        SimpleNamespace(t="data", r=indent, x=[index(position)])
        for position, (_, indent, _) in enumerate(pivot_rows)
    ]
    row_items.append(SimpleNamespace(t="grand", r=0, x=[index(0)]))
    cache_fields = [SimpleNamespace(name=name) for name in headers]
    data_fields = [
        SimpleNamespace(
            name="平均 / income_z",
            fld=4,
            subtotal="average",
            showDataAs="normal",
        )
    ]
    if second_average:
        data_fields.append(
            SimpleNamespace(
                name="平均 / measure_y",
                fld=5,
                subtotal="average",
                showDataAs="normal",
            )
        )
    max_column = "F" if second_average else "E"
    pivot = SimpleNamespace(
        compact=True,
        compactData=True,
        outline=True,
        multipleFieldFilters=False,
        dataOnRows=False,
        pageFields=[],
        filters=[],
        colFields=[SimpleNamespace(x=-2)] if second_average else [],
        rowFields=[SimpleNamespace(x=value) for value in (0, 1, 2, 3)],
        dataFields=data_fields,
        rowItems=row_items,
        location=SimpleNamespace(
            ref=f"A3:{'C' if second_average else 'B'}{grand_row}",
            firstDataCol=1,
        ),
        cache=SimpleNamespace(
            cacheFields=cache_fields,
            cacheSource=SimpleNamespace(
                type="worksheet",
                connectionId=None,
                consolidation=None,
                worksheetSource=SimpleNamespace(
                    sheet="records_q",
                    ref=f"A1:{max_column}4",
                    name=None,
                ),
            ),
        ),
    )
    pivot_sheet._pivots = [pivot]
    return workbook


def write_opaque_proposal_deck(
    path: Path,
    person: str,
    *,
    duty: str = "review_z",
    style_variant: bool = False,
    extra_shape: bool = False,
    unsupported_group: bool = False,
    layout_text: str | None = None,
    master_text: str | None = None,
) -> None:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches

    path.parent.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.5))
    title.text = "opaque staffing"
    table_shape = slide.shapes.add_table(
        3,
        2,
        Inches(0.5),
        Inches(1.0),
        Inches(8),
        Inches(2.0),
    )
    table = table_shape.table
    values = (
        ("role_x", "person_y"),
        ("lead_q", "Alpha"),
        (duty, person),
    )
    for row_index, row in enumerate(values):
        for column_index, value in enumerate(row):
            cell = table.cell(row_index, column_index)
            cell.text = value
            if style_variant:
                run = cell.text_frame.paragraphs[0].runs[0]
                run.font.bold = True
                run.font.color.rgb = RGBColor(38, 88, 166)
    if extra_shape:
        extra = slide.shapes.add_textbox(
            Inches(0.5), Inches(3.5), Inches(3), Inches(0.5)
        )
        extra.text = "topology-extra"
    if unsupported_group:
        slide.shapes.add_group_shape()
    presentation.save(path)

    if layout_text is None and master_text is None:
        return

    import xml.etree.ElementTree as ET
    import zipfile

    presentation_namespace = (
        "http://schemas.openxmlformats.org/presentationml/2006/main"
    )
    drawing_namespace = (
        "http://schemas.openxmlformats.org/drawingml/2006/main"
    )

    def add_visible_shape(payload: bytes, text: str) -> bytes:
        root = ET.fromstring(payload)
        shape_tree = root.find(
            "./p:cSld/p:spTree", {"p": presentation_namespace}
        )
        if shape_tree is None:
            raise AssertionError("opaque inherited fixture has no shape tree")
        shape = ET.SubElement(
            shape_tree, "{" + presentation_namespace + "}sp"
        )
        non_visual = ET.SubElement(
            shape, "{" + presentation_namespace + "}nvSpPr"
        )
        ET.SubElement(
            non_visual,
            "{" + presentation_namespace + "}cNvPr",
            {"id": "9000", "name": "opaque inherited text"},
        )
        ET.SubElement(
            non_visual, "{" + presentation_namespace + "}cNvSpPr"
        )
        ET.SubElement(non_visual, "{" + presentation_namespace + "}nvPr")
        shape_properties = ET.SubElement(
            shape, "{" + presentation_namespace + "}spPr"
        )
        geometry = ET.SubElement(
            shape_properties,
            "{" + drawing_namespace + "}prstGeom",
            {"prst": "rect"},
        )
        ET.SubElement(geometry, "{" + drawing_namespace + "}avLst")
        text_body = ET.SubElement(
            shape, "{" + presentation_namespace + "}txBody"
        )
        ET.SubElement(text_body, "{" + drawing_namespace + "}bodyPr")
        ET.SubElement(text_body, "{" + drawing_namespace + "}lstStyle")
        paragraph = ET.SubElement(
            text_body, "{" + drawing_namespace + "}p"
        )
        run = ET.SubElement(paragraph, "{" + drawing_namespace + "}r")
        rendered = ET.SubElement(run, "{" + drawing_namespace + "}t")
        rendered.text = text
        return ET.tostring(root, encoding="utf-8", xml_declaration=True)

    rewritten = path.with_name(path.stem + "-rewrite.pptx")
    with zipfile.ZipFile(path) as source_archive:
        entries = [
            (item, source_archive.read(item.filename))
            for item in source_archive.infolist()
        ]
    with zipfile.ZipFile(rewritten, "w") as target_archive:
        for item, payload in entries:
            if (
                layout_text is not None
                and item.filename.startswith("ppt/slideLayouts/slideLayout")
                and item.filename.endswith(".xml")
            ):
                payload = add_visible_shape(payload, layout_text)
            if (
                master_text is not None
                and item.filename.startswith("ppt/slideMasters/slideMaster")
                and item.filename.endswith(".xml")
            ):
                payload = add_visible_shape(payload, master_text)
            target_archive.writestr(item, payload)
    rewritten.replace(path)


def write_opaque_categorical_project(
    project: Path,
    *,
    threshold: int = 37,
    call_limit: str | None = None,
    label: str = "catz",
    missing_dtype: bool = False,
    duplicate_label: bool = False,
    extra_entrypoint: bool = False,
) -> None:
    analysis = project / "04.analysis" / "analysis_project"
    scripts = analysis / "scripts"
    source = analysis / "src"
    scripts.mkdir(parents=True, exist_ok=True)
    source.mkdir(parents=True, exist_ok=True)
    run_lines = ["python scripts/run_train.py"]
    if extra_entrypoint:
        run_lines.append("python scripts/other.py")
        (scripts / "other.py").write_text("pass\n", encoding="utf-8")
    (analysis / "README.md").write_text(
        "# opaque analysis\n\n## Run\n\n```bash\n"
        + "\n".join(run_lines)
        + "\n```\n",
        encoding="utf-8",
    )
    argument = "" if call_limit is None else f", {call_limit}"
    (scripts / "run_train.py").write_text(
        "from src.features import choose_fields\n"
        "from src.modeling import assemble_model\n\n"
        "def main():\n"
        "    opaque_frame = object()\n"
        f"    filtered_frame, audit = choose_fields(opaque_frame{argument})\n"
        "    model = assemble_model(filtered_frame)\n"
        "    return model\n\n"
        "if __name__ == '__main__':\n"
        "    main()\n",
        encoding="utf-8",
    )
    (source / "modeling.py").write_text(
        "from src.features import make_processor\n\n"
        "def assemble_model(filtered_frame):\n"
        "    processor = make_processor(filtered_frame)\n"
        "    return processor\n",
        encoding="utf-8",
    )
    third_predicate = (
        "pd.api.types.is_datetime64_any_dtype(vector)"
        if missing_dtype
        else "pd.api.types.is_categorical_dtype(vector)"
    )
    transformer_rows = [
        '("numz", object(), numeric_fields)',
        f'("{label}", object(), categorical_fields)',
    ]
    if duplicate_label:
        transformer_rows.append(
            f'("{label.upper()}", object(), categorical_fields)'
        )
    (source / "features.py").write_text(
        "import pandas as pd\n"
        "from sklearn.compose import ColumnTransformer\n\n"
        f"UNIQUE_BOUND = {threshold}\n\n"
        "def choose_fields(frame, bound=UNIQUE_BOUND):\n"
        "    retained = []\n"
        "    rejected = []\n"
        "    for field_key in frame.columns:\n"
        "        vector = frame[field_key]\n"
        "        categorical_candidate = (\n"
        "            pd.api.types.is_object_dtype(vector)\n"
        "            or pd.api.types.is_string_dtype(vector)\n"
        f"            or {third_predicate}\n"
        "        )\n"
        "        if categorical_candidate:\n"
        "            distinct = int(vector.dropna().nunique())\n"
        "            if distinct >= bound:\n"
        "                rejected.append(field_key)\n"
        "                continue\n"
        "        retained.append(field_key)\n"
        "    return frame[retained].copy(), {\"rejected\": rejected}\n\n"
        "def make_processor(frame):\n"
        "    numeric_fields = frame.select_dtypes(include=[\"number\", \"bool\"]).columns.tolist()\n"
        "    categorical_fields = [field_key for field_key in frame.columns if field_key not in numeric_fields]\n"
        "    return ColumnTransformer(transformers=[\n        "
        + ",\n        ".join(transformer_rows)
        + "\n    ])\n",
        encoding="utf-8",
    )


class StructuredCandidateTest(unittest.TestCase):
    def setUp(self) -> None:
        self.temporary = tempfile.TemporaryDirectory()
        self.root = Path(self.temporary.name)
        self.glossary = Glossary()

    def tearDown(self) -> None:
        self.temporary.cleanup()

    def test_declared_alias_duplicate_and_ordinal_label_resolve_exact_list(self) -> None:
        self.glossary.add("ORG", "正式組織")
        self.glossary.add("PL", "plan")
        rows = [
            ["TaskID", "Phase"],
            ["T01", "1. Setup"],
            ["T02", "2. Ready"],
            ["T03", "2. Ready"],
        ]
        write_csv(self.root / "正式組織" / "primary" / "plan.csv", rows)
        write_csv(self.root / "正式組織" / "copy" / "plan.csv", rows)
        question = (
            "ORGのPLにおいて、PhaseがReadyに一致するTaskIDを"
            "すべて挙げてください。"
        )
        decision = StructuredCandidateEngine(self.root, self.glossary).decide(
            "opaque-list", question
        )
        self.assertEqual(decision.status, "resolved")
        self.assertIsNotNone(decision.result)
        self.assertEqual(decision.result.answer, "T02、T03")
        self.assertEqual(len(decision.result.source_paths), 2)

    def test_legal_form_normalization_is_generic(self) -> None:
        write_csv(
            self.root / "株式会社北極研究連盟" / "rows.csv",
            [["TaskID", "State"], ["A1", "Open"], ["A2", "Closed"]],
        )
        question = (
            "北極研究連盟のrows.csvにおいて、StateがOpenに一致する"
            "TaskIDをすべて挙げてください。"
        )
        decision = StructuredCandidateEngine(self.root, self.glossary).decide(
            "opaque-legal", question
        )
        self.assertEqual(decision.status, "resolved")
        self.assertEqual(decision.result.answer, "A1")

    def test_different_tables_with_same_scope_hold_instead_of_picking_one(self) -> None:
        write_csv(
            self.root / "組織A" / "one" / "rows.csv",
            [["TaskID", "State"], ["A1", "Open"]],
        )
        write_csv(
            self.root / "組織A" / "two" / "rows.csv",
            [["TaskID", "State"], ["B1", "Open"]],
        )
        question = (
            "組織Aのrows.csvにおいて、StateがOpenに一致する"
            "TaskIDをすべて挙げてください。"
        )
        decision = StructuredCandidateEngine(self.root, self.glossary).decide(
            "ambiguous-source", question
        )
        self.assertEqual(decision.status, "hold")
        self.assertEqual(decision.reason, "table_not_unique")

    def test_compound_emits_terminal_identifier_output(self) -> None:
        write_csv(
            self.root / "組織B" / "dataset.csv",
            [
                ["RowID", "Category", "Amount", "Metric"],
                ["R1", "Blue", 200, 10],
                ["R2", "Blue", 300, 20],
                ["R3", "Blue", 400, 20],
                ["R4", "Red", 500, 99],
            ],
        )
        question = (
            "組織Bのdataset.csvにおいて、CategoryがBlueであり、かつ"
            "Amountが100より大きいデータを抽出し、Metricの平均値を計算して"
            "ください。その平均値に最も近いMetricのRowIDをすべて答えてください。"
        )
        decision = StructuredCandidateEngine(self.root, self.glossary).decide(
            "opaque-compound", question
        )
        self.assertEqual(decision.status, "resolved")
        self.assertEqual(decision.result.answer, "R2、R3")
        self.assertEqual(decision.result.operation_count, 6)
        self.assertEqual(decision.result.output_count, 2)

    def test_unsupported_question_never_reads_a_source_as_an_answer(self) -> None:
        write_csv(
            self.root / "組織C" / "rows.csv",
            [["TaskID", "State"], ["A1", "Open"]],
        )
        decision = StructuredCandidateEngine(self.root, self.glossary).decide(
            "unsupported", "組織Cのrows.csvを要約してください。"
        )
        self.assertEqual(decision.status, "unsupported")
        self.assertIsNone(decision.result)

    def test_xlsx_chart_series_column_uses_ooxml_relationships_and_fails_closed(self) -> None:
        from question_graph_runtime import build_graph_plan

        source = (
            self.root
            / "opaque_project"
            / "03.データ"
            / "opaque_book.xlsx"
        )
        write_opaque_chart_ex_workbook(source, "field_alpha")
        question = (
            "opaque_projectのopaque_book.xlsxのDataViewにあるグラフ7は"
            "どのカラムを可視化したものですか。"
        )
        contract = graph_contract_for_question(question)
        self.assertIsNotNone(contract)
        self.assertEqual(contract["graph_rule_version"], "1.4")
        self.assertEqual(contract["rule_id"], "xlsx_chart_series_column")
        self.assertEqual(contract["scope"]["sheet"], "DataView")
        self.assertEqual(contract["scope"]["chart_index"], 7)
        self.assertEqual(
            [node["operator"] for node in contract["operation_graph"]["nodes"]],
            ["retrieve", "select", "select", "resolve", "verify", "project"],
        )
        self.assertEqual(
            contract["requested_output"]["answer_shape"],
            {"container": "scalar", "value_type": "identifier", "unit": None},
        )
        self.assertTrue(validate_graph_contract(question, contract))

        plan = build_graph_plan("opaque-chart", question)
        engine = StructuredCandidateEngine(self.root, self.glossary)
        first = engine.decide_from_graph("opaque-chart", question, plan)
        self.assertEqual(first.status, "resolved")
        self.assertEqual(first.result.answer, "field_alpha")
        self.assertEqual(first.result.operation_count, 6)
        first_sha = first.result.source_sha256

        write_opaque_chart_ex_workbook(source, "field_beta")
        changed = engine.decide_from_graph("opaque-chart-changed", question, plan)
        self.assertEqual(changed.status, "resolved")
        self.assertEqual(changed.result.answer, "field_beta")
        self.assertNotEqual(changed.result.source_sha256, first_sha)

        write_opaque_chart_ex_workbook(
            source,
            "field_beta",
            duplicate_chart=True,
        )
        duplicate = engine.decide_from_graph("opaque-chart-duplicate", question, plan)
        self.assertEqual(duplicate.status, "hold")

        write_opaque_chart_ex_workbook(
            source,
            "field_beta",
            second_series=True,
        )
        multiple = engine.decide_from_graph("opaque-chart-multiple", question, plan)
        self.assertEqual(multiple.status, "hold")

    def test_header_detection_allows_title_row_and_optional_unrelated_blanks(self) -> None:
        table = _table_from_matrix(
            self.root / "book.xlsx",
            "a" * 64,
            "Sheet1",
            [
                ["Project schedule", None, None],
                ["TaskID", "Phase", "Note"],
                ["T01", "Ready", None],
            ],
            ["TaskID", "Phase"],
        )
        self.assertIsNotNone(table)
        self.assertEqual(table.headers, ("TaskID", "Phase", "Note"))
        self.assertEqual(table.rows[0], ("T01", "Ready", ""))

    def test_project_assignment_role_uses_one_complete_scoped_roster(self) -> None:
        from openpyxl import Workbook
        from question_graph_runtime import build_graph_plan

        source_root = self.root / "opaque-assignment"
        glossary = Glossary()
        glossary.add("NEBULA", "opaque_project_nebula")
        plan_path = (
            source_root
            / "プロジェクト"
            / "opaque_project_nebula"
            / "02.計画"
            / "schedule.xlsx"
        )

        def write_roster(rows: list[tuple[object, object]]) -> None:
            plan_path.parent.mkdir(parents=True, exist_ok=True)
            workbook = Workbook()
            sheet = workbook.active
            sheet.title = "team_roster"
            sheet.append(["役割一覧"])
            sheet.append(["役割", "氏名"])
            for row in rows:
                sheet.append(row)
            workbook.save(plan_path)
            workbook.close()

        write_roster(
            [
                ("role_alpha", "person_alpha"),
                ("role_beta", "person_beta"),
            ]
        )
        positive_question = (
            "NEBULA案件において、person_alphaさんは"
            "どの役割としてアサインされていますか。"
        )
        absent_question = positive_question.replace("person_alpha", "person_missing")
        contract = graph_contract_for_question(positive_question)
        self.assertIsNotNone(contract)
        self.assertEqual(contract["graph_rule_version"], "1.4")
        self.assertEqual(contract["rule_id"], "project_person_assignment_role")
        self.assertEqual(
            [node["operator"] for node in contract["operation_graph"]["nodes"]],
            [
                "retrieve",
                "select_authoritative",
                "verify_complete",
                "filter",
                "boolean_test",
                "project",
            ],
        )
        self.assertEqual(
            contract["requested_output"]["answer_shape"],
            {"container": "scalar", "value_type": "string", "unit": None},
        )
        self.assertTrue(validate_graph_contract(positive_question, contract))
        self.assertIsNone(graph_contract_for_question(positive_question + "補足"))

        graph_plan = build_graph_plan("opaque-assignment-plan", positive_question)
        self.assertEqual(graph_plan.strict_status, "pass")
        self.assertEqual(
            graph_plan.strict_reasons,
            ("extended_graph_certified",),
        )
        engine = StructuredCandidateEngine(source_root, glossary)
        first = engine.decide_from_graph(
            "opaque-positive",
            positive_question,
            graph_plan,
        )
        self.assertEqual(first.status, "resolved")
        self.assertEqual(first.result.answer, "role_alpha")
        self.assertEqual(first.result.operation_count, 6)
        first_hash = first.result.source_sha256

        absent_plan = build_graph_plan("opaque-absent-plan", absent_question)
        absent = engine.decide_from_graph(
            "opaque-absent",
            absent_question,
            absent_plan,
        )
        self.assertEqual(absent.status, "resolved")
        self.assertEqual(absent.result.answer, "アサインされていない")

        # Source mutation, not a question-specific exception, changes both a
        # positive role and the previously absent member's membership state.
        write_roster(
            [
                ("role_delta", "person_alpha"),
                ("role_beta", "person_beta"),
                ("role_gamma", "person_missing"),
            ]
        )
        changed = engine.decide_from_graph(
            "opaque-changed",
            positive_question,
            graph_plan,
        )
        self.assertEqual(changed.status, "resolved")
        self.assertEqual(changed.result.answer, "role_delta")
        self.assertNotEqual(changed.result.source_sha256, first_hash)
        now_present = engine.decide_from_graph(
            "opaque-now-present",
            absent_question,
            absent_plan,
        )
        self.assertEqual(now_present.status, "resolved")
        self.assertEqual(now_present.result.answer, "role_gamma")

    def test_project_assignment_role_fails_closed_without_unique_complete_scope(self) -> None:
        from openpyxl import Workbook
        from question_graph_runtime import build_graph_plan

        question = (
            "opaque_scope案件において、person_missingさんは"
            "どの役割としてアサインされていますか。"
        )
        plan = build_graph_plan("opaque-coverage-plan", question)

        def write_book(
            source_root: Path,
            canonical: str,
            rows: list[tuple[object, object]],
            *,
            duplicate_table: bool = False,
        ) -> None:
            target = (
                source_root
                / "プロジェクト"
                / canonical
                / "02.計画"
                / "schedule.xlsx"
            )
            target.parent.mkdir(parents=True, exist_ok=True)
            workbook = Workbook()
            for index in range(2 if duplicate_table else 1):
                sheet = workbook.active if index == 0 else workbook.create_sheet()
                sheet.title = f"team_roster_{index}"
                sheet.append(["役割一覧"])
                sheet.append(["役割", "氏名"])
                for row in rows:
                    sheet.append(row)
            workbook.save(target)
            workbook.close()

        incomplete_root = self.root / "opaque-incomplete"
        write_book(
            incomplete_root,
            "opaque_scope",
            [("role_alpha", "person_alpha"), ("role_beta", None)],
        )
        incomplete = StructuredCandidateEngine(
            incomplete_root,
            Glossary(),
        ).decide_from_graph("opaque-incomplete", question, plan)
        self.assertEqual(incomplete.status, "hold")

        duplicate_root = self.root / "opaque-duplicate"
        write_book(
            duplicate_root,
            "opaque_scope",
            [("role_alpha", "person_alpha")],
            duplicate_table=True,
        )
        duplicate = StructuredCandidateEngine(
            duplicate_root,
            Glossary(),
        ).decide_from_graph("opaque-duplicate", question, plan)
        self.assertEqual(duplicate.status, "hold")

        ambiguous_root = self.root / "opaque-ambiguous"
        write_book(
            ambiguous_root,
            "株式会社opaque_scope",
            [("role_alpha", "person_alpha")],
        )
        write_book(
            ambiguous_root,
            "合同会社opaque_scope",
            [("role_beta", "person_beta")],
        )
        ambiguous = StructuredCandidateEngine(
            ambiguous_root,
            Glossary(),
        ).decide_from_graph("opaque-ambiguous", question, plan)
        self.assertEqual(ambiguous.status, "hold")

    def test_opaque_cohort_group_mean_argmax_is_source_driven(self) -> None:
        source = self.root / "銀河診療所" / "03.データ" / "observations.csv"
        rows = [
            ["flag", "Gender", "Age", "marker_z"],
            ["on", "Female", 21, 4],
            ["on", "Female", 37, 7],
            ["on", "Female", 37, 9],
            ["off", "Female", 21, 99],
            ["on", "Male", 52, 100],
        ]
        write_csv(source, rows)
        question = (
            "銀河診療所のプロジェクトデータ（observations.csv）において、"
            "flag=onの女性の中で、marker_zの平均値が最も高い年齢は何歳ですか。"
        )
        contract = graph_contract_for_question(question)
        self.assertIsNotNone(contract)
        self.assertEqual(contract["rule_id"], "cohort_group_mean_argmax")
        self.assertEqual(contract["requested_output"]["answer_shape"]["unit"], "歳")
        self.assertTrue(validate_graph_contract(question, contract))

        engine = StructuredCandidateEngine(self.root, self.glossary)
        first = engine.decide("ignored-one", question)
        self.assertEqual(first.status, "resolved")
        self.assertEqual(first.result.answer, "37歳")

        rows[1][-1] = 20
        write_csv(source, rows)
        second = engine.decide("ignored-two", question)
        self.assertEqual(second.status, "resolved")
        self.assertEqual(second.result.answer, "21歳")

    def test_opaque_multi_filter_mean_uses_half_up_after_mean(self) -> None:
        from openpyxl import Workbook

        source = self.root / "流星信用" / "03.データ" / "train.xlsx"
        source.parent.mkdir(parents=True, exist_ok=True)

        def write_book(values: tuple[int, int]) -> None:
            workbook = Workbook()
            sheet = workbook.active
            sheet.append(["tenure", "tier", "channel", "principal_x"])
            sheet.append(["long", "K2", "web", values[0]])
            sheet.append(["long", "K2", "web", values[1]])
            sheet.append(["long", "K3", "web", 999])
            workbook.save(source)
            workbook.close()

        write_book((10, 11))
        question = (
            "流星信用の分析対象データにおいて、tenure=long、tier=K2、"
            "channel=webに該当するprincipal_xの平均を算出してください。"
            "四捨五入して整数値で出してください。"
        )
        contract = graph_contract_for_question(question)
        self.assertIsNotNone(contract)
        self.assertEqual(contract["rule_id"], "multi_filter_mean_half_up")
        self.assertEqual(
            [node["operator"] for node in contract["operation_graph"]["nodes"]],
            ["retrieve", "filter", "filter", "filter", "mean", "calculate"],
        )
        self.assertTrue(validate_graph_contract(question, contract))

        engine = StructuredCandidateEngine(self.root, self.glossary)
        first = engine.decide("ignored-one", question)
        self.assertEqual(first.status, "resolved")
        self.assertEqual(first.result.answer, "11")

        write_book((20, 21))
        second = engine.decide("ignored-two", question)
        self.assertEqual(second.status, "resolved")
        self.assertEqual(second.result.answer, "21")

    def test_notebook_heatmap_top_n_controls_source_recomputed_correlation(self) -> None:
        project = self.root / "彗星評価" / "04.分析" / "analysis_project"
        data = project / "data" / "source.csv"
        notebook = project / "notebooks" / "viz.ipynb"
        write_csv(
            data,
            [
                ["strong", "medium", "weak", "outside", "outcome"],
                [2, 1, 1, 1, 1],
                [4, 2, 1, -1, 2],
                [6, 3, 2, 1, 3],
                [8, 4, 1, -1, 4],
                [10, 6, 2, 1, 5],
                [12, 5, 3, -1, 6],
            ],
        )

        def write_notebook(top_n: int) -> None:
            notebook.parent.mkdir(parents=True, exist_ok=True)
            source = (
                'csv_rel = Path("data/source.csv")\n'
                'target_name = "outcome"\n'
                "features = [c for c in numeric.columns if c != target_name]\n"
                "ranked = numeric[features].corrwith(numeric[target_name]).abs().sort_values(ascending=False)\n"
                f"visible = ranked.head({top_n}).index.tolist()\n"
                "matrix = numeric[visible].corr()\n"
                "sns.heatmap(matrix)\n"
            )
            notebook.write_text(
                json.dumps(
                    {
                        "cells": [{"cell_type": "code", "source": source.splitlines(True)}],
                        "metadata": {},
                        "nbformat": 4,
                        "nbformat_minor": 5,
                    }
                ),
                encoding="utf-8",
            )

        write_notebook(3)
        question = (
            "彗星評価の viz.ipynb にある特徴量相関ヒートマップの図で"
            "可視化されている特徴量のうち、outcomeとの相関係数の絶対値が"
            "最も小さい特徴量名を答えてください。"
        )
        contract = graph_contract_for_question(question)
        self.assertIsNotNone(contract)
        self.assertEqual(contract["rule_id"], "notebook_heatmap_min_abs_correlation")
        self.assertTrue(validate_graph_contract(question, contract))

        engine = StructuredCandidateEngine(self.root, self.glossary)
        first = engine.decide("ignored-one", question)
        self.assertEqual(first.status, "resolved")
        self.assertEqual(first.result.answer, "weak")

        write_notebook(2)
        second = engine.decide("ignored-two", question)
        self.assertEqual(second.status, "resolved")
        self.assertEqual(second.result.answer, "medium")

        mismatched_target = question.replace("outcomeとの相関係数", "outsideとの相関係数")
        rejected = engine.decide("ignored-three", mismatched_target)
        self.assertEqual(rejected.status, "unsupported")
        self.assertIsNone(rejected.result)

    def test_report_and_json_metric_delta_rounds_only_final_decimal(self) -> None:
        from docx import Document

        project = self.root / "星雲病院"
        report = project / "05.会議" / "報告資料" / "途中報告資料.docx"
        metrics = project / "04.分析" / "analysis_outputs" / "summary.json"
        report.parent.mkdir(parents=True, exist_ok=True)
        metrics.parent.mkdir(parents=True, exist_ok=True)
        document = Document()
        document.add_paragraph("これは途中報告資料である。")
        document.add_paragraph("Omega Zeta = 0.123456789")
        document.save(report)
        metrics.write_text(
            json.dumps({"evaluation": {"zeta_omega": "0.223456789"}}),
            encoding="utf-8",
        )
        question = (
            "星雲病院案件において、途中報告資料に記載されたOmega Zetaスコアの詳細値と、"
            "最終分析出力summary.jsonに記録されているOmega Zetaスコアの詳細値を用いて、"
            "改善幅を小数第4位まで答えてください。"
        )
        contract = graph_contract_for_question(question)
        self.assertIsNotNone(contract)
        self.assertEqual(contract["rule_id"], "report_metrics_decimal_delta")
        self.assertEqual(
            contract["requested_output"]["display_precision"],
            {"mode": "decimal_places", "digits": 4},
        )
        self.assertTrue(validate_graph_contract(question, contract))
        self.assertIsNone(graph_contract_for_question(question + "補足"))

        engine = StructuredCandidateEngine(self.root, self.glossary)
        first = engine.decide("ignored-one", question)
        self.assertEqual(first.status, "resolved")
        self.assertEqual(first.result.answer, "0.1000")

        metrics.write_text(
            json.dumps({"evaluation": {"zeta_omega": "0.223506789"}}),
            encoding="utf-8",
        )
        second = engine.decide("ignored-two", question)
        self.assertEqual(second.status, "resolved")
        self.assertEqual(second.result.answer, "0.1001")

    def test_contract_hours_ratio_tax_delta_is_source_driven_and_exact(self) -> None:
        from types import SimpleNamespace

        from docx import Document

        source = self.root / "銀河医療" / "01.契約" / "契約書.docx"
        source.parent.mkdir(parents=True, exist_ok=True)

        def write_contract(unit: int, hours: int, *, inconsistent: bool = False) -> None:
            expected_ex_tax = unit * hours + (1 if inconsistent else 0)
            tax = expected_ex_tax // 10
            expected_with_tax = expected_ex_tax + tax
            document = Document()
            document.add_paragraph(
                "本契約の料金モデルはtime_and_materialsとし、"
                "実績工数に基づく事後精算（月次精算）とする。"
            )
            document.add_paragraph(f"時間単価は{unit:,}円（消費税別）とする。")
            document.add_paragraph(f"想定総工数は{hours}時間とする。")
            document.add_paragraph(
                f"見込金額は、税抜{expected_ex_tax:,}円、"
                f"消費税{tax:,}円、税込{expected_with_tax:,}円とする。"
            )
            document.add_paragraph(
                "最終請求額は、実績工数に時間単価を乗じ、"
                "これに消費税を加算した金額とする。"
            )
            document.save(source)

        question = (
            "銀河医療の契約条件において、仮に実績工数が見込工数の5分の2だった場合、"
            "最終請求金額（税込）は見込金額（税込）よりいくら少なくなりますか。"
        )
        contract = graph_contract_for_question(question)
        self.assertIsNotNone(contract)
        self.assertEqual(contract["graph_rule_version"], "1.4")
        self.assertEqual(contract["rule_id"], "contract_hours_ratio_tax_delta")
        self.assertEqual(contract["scope"]["container"], "契約書.docx")
        self.assertEqual(
            [node["operator"] for node in contract["operation_graph"]["nodes"]],
            ["retrieve", "verify", "calculate", "calculate"],
        )
        self.assertEqual(
            contract["requested_output"]["answer_shape"],
            {"container": "scalar", "value_type": "integer", "unit": "円"},
        )
        self.assertTrue(validate_graph_contract(question, contract))
        self.assertIsNone(graph_contract_for_question(question + "補足"))
        self.assertIsNone(
            graph_contract_for_question(question.replace("5分の2", "0分の2"))
        )

        write_contract(12_000, 50)
        engine = StructuredCandidateEngine(self.root, self.glossary)
        plan = SimpleNamespace(branch_intents=())
        first = engine.decide_from_graph("opaque-one", question, plan)
        self.assertEqual(first.status, "resolved")
        self.assertEqual(first.result.answer, "396,000円")
        self.assertEqual(first.result.operation_count, 4)
        self.assertEqual(first.result.output_count, 1)
        first_hash = first.result.source_sha256

        write_contract(18_000, 50)
        second = engine.decide_from_graph("opaque-two", question, plan)
        self.assertEqual(second.status, "resolved")
        self.assertEqual(second.result.answer, "594,000円")
        self.assertNotEqual(second.result.source_sha256, first_hash)

        write_contract(18_000, 50, inconsistent=True)
        self.assertIsNone(decide_extended(engine, "ignored-inconsistent", question))

        # No rounding convention is declared, so a fractional-yen result must
        # not be guessed even when every source amount is internally valid.
        write_contract(101, 1)
        self.assertIsNone(decide_extended(engine, "ignored-rounding", question))

        write_contract(18_000, 50)
        duplicate = self.root / "銀河医療" / "archive" / "契約書.docx"
        duplicate.parent.mkdir(parents=True, exist_ok=True)
        duplicate.write_bytes(source.read_bytes())
        self.assertIsNone(decide_extended(engine, "ignored-duplicate", question))

    def test_excel_autofilter_conditions_use_relative_columns_and_source_values(self) -> None:
        from types import SimpleNamespace

        from openpyxl import Workbook
        from openpyxl.worksheet.filters import (
            CustomFilter,
            CustomFilters,
            FilterColumn,
            Filters,
        )

        source = self.root / "流星人材" / "03.データ" / "opaque.xlsx"
        source.parent.mkdir(parents=True, exist_ok=True)

        def write_book(
            left_value: str,
            right_value: str,
            *,
            reverse: bool = False,
            custom: bool = False,
        ) -> None:
            workbook = Workbook()
            sheet = workbook.active
            sheet.title = "records"
            for offset, header in enumerate(("segment_z", "region_q", "flag_x")):
                sheet.cell(row=3, column=3 + offset, value=header)
            for row_number, values in enumerate(
                (("amber", "north", "7"), ("violet", "south", "9")),
                4,
            ):
                for offset, value in enumerate(values):
                    sheet.cell(row=row_number, column=3 + offset, value=value)
            sheet.auto_filter.ref = "C3:E5"
            if custom:
                columns = [
                    FilterColumn(
                        colId=0,
                        customFilters=CustomFilters(
                            customFilter=[CustomFilter(operator="equal", val=left_value)]
                        ),
                    )
                ]
            else:
                columns = [
                    FilterColumn(colId=0, filters=Filters(filter=[left_value])),
                    FilterColumn(colId=2, filters=Filters(filter=[right_value])),
                ]
                if reverse:
                    columns.reverse()
            sheet.auto_filter.filterColumn = columns
            workbook.save(source)
            workbook.close()

        question = (
            "流星人材のopaque.xlsxにおいて、recordsシートで"
            "フィルターで抽出されている条件を教えてください。"
        )
        contract = graph_contract_for_question(question)
        self.assertIsNotNone(contract)
        self.assertEqual(contract["graph_rule_version"], "1.4")
        self.assertEqual(contract["rule_id"], "excel_autofilter_conditions")
        self.assertEqual(contract["scope"]["sheet"], "records")
        self.assertEqual(
            [node["operator"] for node in contract["operation_graph"]["nodes"]],
            ["retrieve", "verify", "project", "list"],
        )
        self.assertEqual(
            contract["requested_output"]["answer_shape"],
            {"container": "key_value", "value_type": "string", "unit": None},
        )
        self.assertTrue(validate_graph_contract(question, contract))
        self.assertIsNone(graph_contract_for_question(question + "補足"))

        write_book("amber", "7")
        engine = StructuredCandidateEngine(self.root, self.glossary)
        plan = SimpleNamespace(branch_intents=())
        first = engine.decide_from_graph("opaque-one", question, plan)
        self.assertEqual(first.status, "resolved")
        self.assertEqual(first.result.answer, "segment_z=amber、flag_x=7")
        self.assertEqual(first.result.operation_count, 4)
        self.assertEqual(first.result.output_count, 1)
        first_hash = first.result.source_sha256

        write_book("violet", "9", reverse=True)
        second = engine.decide_from_graph("opaque-two", question, plan)
        self.assertEqual(second.status, "resolved")
        self.assertEqual(second.result.answer, "segment_z=violet、flag_x=9")
        self.assertNotEqual(second.result.source_sha256, first_hash)

        missing_sheet = question.replace("recordsシート", "missingシート")
        self.assertIsNone(decide_extended(engine, "ignored-sheet", missing_sheet))

        write_book("violet", "9", custom=True)
        self.assertIsNone(decide_extended(engine, "ignored-custom", question))

    def test_pivot_average_graph_contracts_are_typed_full_question_rules(self) -> None:
        aggregate_question = (
            "星雲診療のopaque.xlsx内の PivotTable で集計されている表から、"
            "measure_uの平均が最も高いものの抽出条件と集計内容を答えてください。"
        )
        conditions_question = (
            "軌道人材のopaque.xlsxのlayers_qシートにおいて、"
            "平均報酬が最も高い層の抽出条件を教えてください。"
        )
        expected_operators = [
            "retrieve",
            "verify",
            "project",
            "group",
            "mean",
            "verify_complete",
            "argmax_all",
            "project",
        ]
        aggregate = graph_contract_for_question(aggregate_question)
        conditions = graph_contract_for_question(conditions_question)
        self.assertIsNotNone(aggregate)
        self.assertIsNotNone(conditions)
        self.assertEqual(aggregate["graph_rule_version"], "1.4")
        self.assertEqual(
            aggregate["rule_id"],
            "pivot_average_argmax_conditions_and_aggregate",
        )
        self.assertEqual(
            conditions["rule_id"],
            "pivot_average_argmax_conditions",
        )
        self.assertEqual(conditions["scope"]["sheet"], "layers_q")
        self.assertEqual(
            [
                node["operator"]
                for node in aggregate["operation_graph"]["nodes"]
            ],
            expected_operators,
        )
        self.assertEqual(
            aggregate["requested_output"]["answer_shape"],
            {"container": "key_value", "value_type": "string", "unit": None},
        )
        self.assertEqual(
            aggregate["requested_output"]["cardinality"],
            "multiple",
        )
        self.assertTrue(validate_graph_contract(aggregate_question, aggregate))
        self.assertTrue(validate_graph_contract(conditions_question, conditions))
        tampered = json.loads(json.dumps(conditions))
        tampered["scope"]["sheet"] = "other"
        self.assertFalse(validate_graph_contract(conditions_question, tampered))
        self.assertIsNone(graph_contract_for_question(aggregate_question + "追記"))
        self.assertIsNone(graph_contract_for_question(conditions_question + "追記"))

    def test_tabular_pivot_recomputes_all_leaf_means_and_rejects_stale_view(self) -> None:
        source = self.root / "星雲診療" / "opaque.xlsx"
        write_marker_workbook(source, "first")
        question = (
            "星雲診療のopaque.xlsx内の PivotTable で集計されている表から、"
            "measure_uの平均が最も高いものの抽出条件と集計内容を答えてください。"
        )
        engine = StructuredCandidateEngine(self.root, self.glossary)
        with patch.object(
            candidate_rules,
            "_load_pivot_workbook",
            side_effect=lambda _: synthetic_tabular_pivot_workbook(),
        ):
            first = engine.decide_from_graph("opaque-one", question, None)
        self.assertEqual(first.status, "resolved")
        self.assertEqual(
            first.result.answer,
            "cohort_z=alpha、flag_q=on、band_x=2で抽出されたデータに対する"
            "平均 / measure_u",
        )
        self.assertEqual(first.result.operation_count, 8)
        self.assertEqual(first.result.output_count, 1)
        first_hash = first.result.source_sha256

        # Mutating both raw records and every corresponding leaf changes the
        # winner without changing the question or the generic rule.
        write_marker_workbook(source, "second")
        with patch.object(
            candidate_rules,
            "_load_pivot_workbook",
            side_effect=lambda _: synthetic_tabular_pivot_workbook(beta_raw=30),
        ):
            changed = engine.decide_from_graph("opaque-two", question, None)
        self.assertEqual(changed.status, "resolved")
        self.assertEqual(
            changed.result.answer,
            "cohort_z=beta、flag_q=off、band_x=3で抽出されたデータに対する"
            "平均 / measure_u",
        )
        self.assertNotEqual(changed.result.source_sha256, first_hash)

        # A refreshed raw table with a stale Pivot leaf is not answerable.
        with patch.object(
            candidate_rules,
            "_load_pivot_workbook",
            side_effect=lambda _: synthetic_tabular_pivot_workbook(
                beta_raw=30,
                beta_display=3,
            ),
        ):
            self.assertIsNone(decide_extended(engine, "ignored-stale", question))

        duplicate = self.root / "星雲診療" / "archive" / "opaque.xlsx"
        duplicate.parent.mkdir(parents=True, exist_ok=True)
        duplicate.write_bytes(source.read_bytes())
        self.assertIsNone(decide_extended(engine, "ignored-duplicate", question))

    def test_compact_pivot_uses_deepest_hierarchy_and_single_metric_fallback(self) -> None:
        source = self.root / "軌道人材" / "opaque.xlsx"
        write_marker_workbook(source, "compact-one")
        question = (
            "軌道人材のopaque.xlsxのlayers_qシートにおいて、"
            "平均報酬が最も高い層の抽出条件を教えてください。"
        )
        engine = StructuredCandidateEngine(self.root, self.glossary)
        with patch.object(
            candidate_rules,
            "_load_pivot_workbook",
            side_effect=lambda _: synthetic_compact_pivot_workbook(),
        ):
            first = engine.decide_from_graph("opaque-one", question, None)
        self.assertEqual(first.status, "resolved")
        self.assertEqual(
            first.result.answer,
            "status_p=No、segment_q=Female、tier_r=Single、domain_s=Tech",
        )
        self.assertEqual(first.result.operation_count, 8)
        self.assertEqual(first.result.output_count, 1)

        # Intermediate compact-layout totals are deliberately larger than the
        # leaf values; only the complete depth-four hierarchy can win.
        with patch.object(
            candidate_rules,
            "_load_pivot_workbook",
            side_effect=lambda _: synthetic_compact_pivot_workbook(sales_raw=250),
        ):
            changed = engine.decide_from_graph("opaque-two", question, None)
        self.assertEqual(changed.status, "resolved")
        self.assertEqual(
            changed.result.answer,
            "status_p=Yes、segment_q=Male、tier_r=Divorced、domain_s=Sales",
        )

        # The Japanese surface label can bind to a single source measure, but
        # never to one of two equally valid average measures.
        with patch.object(
            candidate_rules,
            "_load_pivot_workbook",
            side_effect=lambda _: synthetic_compact_pivot_workbook(
                second_average=True
            ),
        ):
            self.assertIsNone(decide_extended(engine, "ignored-ambiguous", question))

        # Layout metadata and visible indentation must describe the same tree.
        with patch.object(
            candidate_rules,
            "_load_pivot_workbook",
            side_effect=lambda _: synthetic_compact_pivot_workbook(
                broken_indent=True
            ),
        ):
            self.assertIsNone(decide_extended(engine, "ignored-indent", question))

        missing_sheet = question.replace("layers_qシート", "missingシート")
        self.assertIsNone(decide_extended(engine, "ignored-sheet", missing_sheet))

    def test_pptx_diff_and_python_rule_have_full_typed_graph_contracts(self) -> None:
        diff_question = (
            "銀河アセットのplan_xについて、oldフォルダ内の旧版と"
            "提案フォルダ直下の最新版を比較し、変更された箇所を"
            "変更前と変更後で答えてください。"
        )
        dtype_question = (
            "星系医療の分析コードにおいて、CATZは dtype と"
            "ユニーク数の条件でどのように判定していますか。"
        )
        diff_contract = graph_contract_for_question(diff_question)
        dtype_contract = graph_contract_for_question(dtype_question)
        self.assertIsNotNone(diff_contract)
        self.assertIsNotNone(dtype_contract)
        self.assertEqual(diff_contract["graph_rule_version"], "1.4")
        self.assertEqual(
            diff_contract["rule_id"],
            "pptx_old_latest_visible_text_diff",
        )
        self.assertEqual(
            diff_contract["requested_output"]["answer_shape"],
            {"container": "list", "value_type": "change_record", "unit": None},
        )
        self.assertEqual(diff_contract["requested_output"]["cardinality"], "all")
        self.assertEqual(
            [
                node["operator"]
                for node in diff_contract["operation_graph"]["nodes"]
            ],
            [
                "retrieve",
                "retrieve",
                "verify",
                "project",
                "compare",
                "filter",
                "list",
            ],
        )
        self.assertEqual(
            dtype_contract["rule_id"],
            "python_categorical_dtype_unique_rule",
        )
        self.assertEqual(
            dtype_contract["requested_output"]["answer_shape"],
            {"container": "scalar", "value_type": "string", "unit": None},
        )
        self.assertEqual(dtype_contract["scope"]["classification_label"], "CATZ")
        self.assertTrue(validate_graph_contract(diff_question, diff_contract))
        self.assertTrue(validate_graph_contract(dtype_question, dtype_contract))
        self.assertIsNone(graph_contract_for_question(diff_question + "追記"))
        self.assertIsNone(graph_contract_for_question(dtype_question + "追記"))

    def test_pptx_old_latest_diff_uses_visible_text_and_exact_topology(self) -> None:
        proposal = self.root / "株式会社銀河アセット" / "00.提案"
        old = proposal / "old" / "plan_x.pptx"
        latest = proposal / "plan_x.pptx"
        question = (
            "銀河アセットのplan_xについて、oldフォルダ内の旧版と"
            "提案フォルダ直下の最新版を比較し、変更された箇所を"
            "変更前と変更後で答えてください。"
        )
        write_opaque_proposal_deck(
            old,
            "Omega",
            layout_text="layout-stable",
            master_text="master-stable",
        )
        write_opaque_proposal_deck(
            latest,
            "Sigma",
            style_variant=True,
            layout_text="layout-stable",
            master_text="master-stable",
        )
        engine = StructuredCandidateEngine(self.root, self.glossary)
        first = engine.decide_from_graph("opaque-one", question, None)
        self.assertEqual(first.status, "resolved")
        self.assertEqual(
            first.result.answer,
            "review_z / person_y: 変更前=Omega、変更後=Sigma",
        )
        self.assertEqual(first.result.operation_count, 7)
        self.assertEqual(first.result.output_count, 1)

        # Opaque text mutation changes the source-derived record.  Font/style
        # serialization changes remain outside the visible-text comparison.
        write_opaque_proposal_deck(old, "Theta")
        write_opaque_proposal_deck(latest, "Lambda", style_variant=True)
        changed = engine.decide_from_graph("opaque-two", question, None)
        self.assertEqual(changed.status, "resolved")
        self.assertEqual(
            changed.result.answer,
            "review_z / person_y: 変更前=Theta、変更後=Lambda",
        )

        # Zero or multiple visible diffs are not silently summarized.
        write_opaque_proposal_deck(old, "Same")
        write_opaque_proposal_deck(latest, "Same", style_variant=True)
        self.assertIsNone(decide_extended(engine, "ignored-zero", question))
        write_opaque_proposal_deck(old, "Before", duty="duty_old")
        write_opaque_proposal_deck(latest, "After", duty="duty_new")
        self.assertIsNone(decide_extended(engine, "ignored-multiple", question))

        write_opaque_proposal_deck(old, "Before")
        write_opaque_proposal_deck(latest, "After", extra_shape=True)
        self.assertIsNone(decide_extended(engine, "ignored-topology", question))
        write_opaque_proposal_deck(old, "Before", unsupported_group=True)
        write_opaque_proposal_deck(latest, "After", unsupported_group=True)
        self.assertIsNone(decide_extended(engine, "ignored-object", question))

        # A single slide-text change is still unsafe when referenced inherited
        # visible text changed independently; neither change may be omitted.
        write_opaque_proposal_deck(
            old, "Before", layout_text="layout-old", master_text="master-stable"
        )
        write_opaque_proposal_deck(
            latest, "After", layout_text="layout-new", master_text="master-stable"
        )
        self.assertIsNone(decide_extended(engine, "ignored-layout", question))
        write_opaque_proposal_deck(
            old, "Before", layout_text="layout-stable", master_text="master-old"
        )
        write_opaque_proposal_deck(
            latest, "After", layout_text="layout-stable", master_text="master-new"
        )
        self.assertIsNone(decide_extended(engine, "ignored-master", question))

    def test_python_categorical_rule_traces_runtime_default_and_fails_closed(self) -> None:
        project = self.root / "星系医療"
        question = (
            "星系医療の分析コードにおいて、CATZは dtype と"
            "ユニーク数の条件でどのように判定していますか。"
        )
        write_opaque_categorical_project(project, threshold=37)
        engine = StructuredCandidateEngine(self.root, self.glossary)
        first = engine.decide_from_graph("opaque-one", question, None)
        self.assertEqual(first.status, "resolved")
        self.assertEqual(
            first.result.answer,
            "CATZはobject/string/category dtypeを候補とする。"
            "欠損を除いたユニーク数が37未満ならCATZとして採用する。",
        )
        self.assertEqual(first.result.operation_count, 8)
        self.assertEqual(first.result.output_count, 1)
        self.assertEqual(len(first.result.source_paths), 4)

        # A literal runtime override is authoritative over the function
        # default; changing either source value is reflected without fixtures.
        write_opaque_categorical_project(project, threshold=91, call_limit="23")
        changed = engine.decide_from_graph("opaque-two", question, None)
        self.assertEqual(changed.status, "resolved")
        self.assertIn("ユニーク数が23未満", changed.result.answer)
        self.assertNotIn("91未満", changed.result.answer)

        write_opaque_categorical_project(project, call_limit="runtime_limit()")
        self.assertIsNone(decide_extended(engine, "ignored-dynamic", question))
        write_opaque_categorical_project(project, call_limit="12, bound=13")
        self.assertIsNone(decide_extended(engine, "ignored-multiple-limit", question))
        write_opaque_categorical_project(project, missing_dtype=True)
        self.assertIsNone(decide_extended(engine, "ignored-dtype", question))
        write_opaque_categorical_project(project, duplicate_label=True)
        self.assertIsNone(decide_extended(engine, "ignored-label", question))
        write_opaque_categorical_project(project, extra_entrypoint=True)
        self.assertIsNone(decide_extended(engine, "ignored-entry", question))

        # Calls in dead control flow, imported-name rebinding, and a dynamic
        # reassignment after a literal default must never be certified as the
        # effective runtime pipeline.
        write_opaque_categorical_project(project)
        entrypoint = (
            project
            / "04.analysis"
            / "analysis_project"
            / "scripts"
            / "run_train.py"
        )
        entry_text = entrypoint.read_text(encoding="utf-8").replace(
            "    opaque_frame = object()\n"
            "    filtered_frame, audit = choose_fields(opaque_frame)\n"
            "    model = assemble_model(filtered_frame)\n"
            "    return model\n",
            "    if False:\n"
            "        opaque_frame = object()\n"
            "        filtered_frame, audit = choose_fields(opaque_frame)\n"
            "        model = assemble_model(filtered_frame)\n"
            "    return None\n",
        )
        entrypoint.write_text(entry_text, encoding="utf-8")
        self.assertIsNone(decide_extended(engine, "ignored-dead", question))

        write_opaque_categorical_project(project)
        entry_text = entrypoint.read_text(encoding="utf-8").replace(
            "from src.modeling import assemble_model\n\n",
            "from src.modeling import assemble_model\n"
            "choose_fields = lambda frame, *args, **kwargs: (frame, {})\n\n",
        )
        entrypoint.write_text(entry_text, encoding="utf-8")
        self.assertIsNone(decide_extended(engine, "ignored-rebound", question))

        write_opaque_categorical_project(project)
        features = (
            project / "04.analysis" / "analysis_project" / "src" / "features.py"
        )
        feature_text = features.read_text(encoding="utf-8").replace(
            "UNIQUE_BOUND = 37\n\n",
            "UNIQUE_BOUND = 37\nUNIQUE_BOUND = runtime_bound()\n\n",
        )
        features.write_text(feature_text, encoding="utf-8")
        self.assertIsNone(decide_extended(engine, "ignored-reassigned", question))

        write_opaque_categorical_project(project)
        feature_text = features.read_text(encoding="utf-8").replace(
            "import pandas as pd\n",
            "import pandas as pd\npd = runtime_dtype_provider()\n",
        )
        features.write_text(feature_text, encoding="utf-8")
        self.assertIsNone(decide_extended(engine, "ignored-dtype-binding", question))

    def test_docx_highlight_text_uses_document_identity_and_run_style(self) -> None:
        from types import SimpleNamespace

        from docx import Document
        from docx.enum.text import WD_COLOR_INDEX

        canonical = "恒星航行研究所"
        self.glossary.add("QX9", canonical)
        material_dir = self.root / canonical / "05.会議" / "報告資料"
        target = material_dir / "opaque-a.docx"

        def write_material(
            path: Path,
            key: str,
            pieces: tuple[str, ...],
            *,
            mention: str | None = None,
        ) -> None:
            path.parent.mkdir(parents=True, exist_ok=True)
            document = Document()
            document.add_paragraph("分析進捗報告書")
            document.add_paragraph(f"チェックポイント: {key}")
            paragraph = document.add_paragraph()
            for piece in pieces:
                run = paragraph.add_run(piece)
                run.font.highlight_color = WD_COLOR_INDEX.YELLOW
            ignored = document.add_paragraph().add_run("非対象")
            ignored.font.highlight_color = WD_COLOR_INDEX.BRIGHT_GREEN
            if mention:
                document.add_paragraph(f"次回は {mention} を予定する。")
            document.save(path)

        write_material(target, "M73", ("総計", ":", " 123 ", "arb"))
        write_material(
            material_dir / "opaque-decoy.docx",
            "M72",
            ("デコイ",),
            mention="M91",
        )
        # The key matches, but a meeting record is not the requested material.
        write_material(
            self.root / canonical / "05.会議" / "会議録" / "minutes.docx",
            "M73",
            ("議事録デコイ",),
        )
        question = (
            "QX9のM73資料（docx）において、黄色でハイライトされている部分を"
            "すべて抜き出してください。"
        )
        contract = graph_contract_for_question(question)
        self.assertIsNotNone(contract)
        self.assertEqual(contract["graph_rule_version"], "1.4")
        self.assertEqual(contract["rule_id"], "docx_highlighted_text_projection")
        self.assertEqual(contract["scope"]["document_key"], "M73")
        self.assertEqual(contract["scope"]["color"], "yellow")
        self.assertEqual(
            [node["operator"] for node in contract["operation_graph"]["nodes"]],
            ["retrieve", "match", "verify", "filter", "group", "project", "list"],
        )
        self.assertEqual(
            contract["requested_output"]["answer_shape"],
            {"container": "list", "value_type": "string", "unit": None},
        )
        self.assertTrue(validate_graph_contract(question, contract))
        self.assertIsNone(graph_contract_for_question(question + "補足"))
        self.assertIsNone(
            graph_contract_for_question(question.replace("黄色", "黄金色"))
        )

        engine = StructuredCandidateEngine(self.root, self.glossary)
        plan = SimpleNamespace(branch_intents=())
        first = engine.decide_from_graph("opaque-one", question, plan)
        self.assertEqual(first.status, "resolved")
        self.assertEqual(first.result.answer, "総計: 123 arb")
        self.assertEqual(first.result.operation_count, 7)
        self.assertEqual(first.result.output_count, 1)
        first_hash = first.result.source_sha256

        # Run fragmentation is not semantic, while source text mutation is.
        write_material(target, "M73", ("総計:", " 987 arb"))
        second = engine.decide_from_graph("opaque-two", question, plan)
        self.assertEqual(second.status, "resolved")
        self.assertEqual(second.result.answer, "総計: 987 arb")
        self.assertNotEqual(second.result.source_sha256, first_hash)

        # A body mention must not masquerade as a document identity label.
        missing_identity = question.replace("M73資料", "M91資料")
        self.assertIsNone(decide_extended(engine, "ignored-mention", missing_identity))

        duplicate = material_dir / "duplicate" / "opaque-b.docx"
        duplicate.parent.mkdir(parents=True, exist_ok=True)
        duplicate.write_bytes(target.read_bytes())
        self.assertIsNone(decide_extended(engine, "ignored-duplicate", question))

    def test_pptx_shape_fill_text_uses_slide_order_and_source_rgb(self) -> None:
        from types import SimpleNamespace

        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches

        project = self.root / "深空人材機構"
        source = project / "00.提案" / "設計書.pptx"
        source.parent.mkdir(parents=True, exist_ok=True)

        def write_deck(texts: tuple[str, ...], rgb: str) -> None:
            presentation = Presentation()
            blank = presentation.slide_layouts[6]
            for _ in range(3):
                presentation.slides.add_slide(blank)
            slide = presentation.slides[2]
            for index, text in enumerate(texts):
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.CHEVRON,
                    Inches(0.5 + 2.1 * index),
                    Inches(1.0),
                    Inches(1.9),
                    Inches(0.8),
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(*bytes.fromhex(rgb))
                shape.line.fill.background()
                shape.text = text
            neutral = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.5),
                Inches(2.2),
                Inches(1.9),
                Inches(0.8),
            )
            neutral.fill.solid()
            neutral.fill.fore_color.rgb = RGBColor(64, 64, 64)
            neutral.text = "neutral"
            presentation.save(source)

        question = (
            "深空人材機構の設計書P3において、赤で強調されている箇所の文字列を"
            "抜き出してください。"
        )
        contract = graph_contract_for_question(question)
        self.assertIsNotNone(contract)
        self.assertEqual(contract["graph_rule_version"], "1.4")
        self.assertEqual(contract["rule_id"], "pptx_shape_fill_text_projection")
        self.assertEqual(contract["scope"]["slide"], 3)
        self.assertEqual(contract["scope"]["color"], "red")
        self.assertEqual(
            [node["operator"] for node in contract["operation_graph"]["nodes"]],
            ["retrieve", "verify", "select", "filter", "project", "list"],
        )
        self.assertEqual(
            contract["requested_output"]["answer_shape"],
            {"container": "list", "value_type": "string", "unit": None},
        )
        self.assertTrue(validate_graph_contract(question, contract))
        self.assertIsNone(graph_contract_for_question(question + "補足"))
        self.assertIsNone(
            graph_contract_for_question(question.replace("設計書P3", "設計書P0"))
        )
        self.assertIsNone(
            graph_contract_for_question(question.replace("赤で", "えんじで"))
        )

        write_deck(("Theta",), "A23B2C")
        engine = StructuredCandidateEngine(self.root, self.glossary)
        plan = SimpleNamespace(branch_intents=())
        first = engine.decide_from_graph("opaque-one", question, plan)
        self.assertEqual(first.status, "resolved")
        self.assertEqual(first.result.answer, "Theta")
        self.assertEqual(first.result.operation_count, 6)
        self.assertEqual(first.result.output_count, 1)
        first_hash = first.result.source_sha256

        # Text, count and the exact RGB can change without changing the rule.
        write_deck(("Iota", "Kappa"), "B2473A")
        second = engine.decide_from_graph("opaque-two", question, plan)
        self.assertEqual(second.status, "resolved")
        self.assertEqual(second.result.answer, "Iota、Kappa")
        self.assertNotEqual(second.result.source_sha256, first_hash)

        # Brown-orange is outside the conservative red hue cone.
        write_deck(("not-red",), "C06020")
        self.assertIsNone(decide_extended(engine, "ignored-orange", question))

        write_deck(("restored",), "A23B2C")
        duplicate = project / "archive" / "設計書.pptx"
        duplicate.parent.mkdir(parents=True, exist_ok=True)
        duplicate.write_bytes(source.read_bytes())
        self.assertIsNone(decide_extended(engine, "ignored-duplicate", question))

    def test_all_project_milestone_cutoff_is_complete_ordered_and_metamorphic(self) -> None:
        from datetime import datetime

        from openpyxl import Workbook

        def write_plan(
            canonical: str,
            events: list[tuple[str, str]],
            *,
            revision: str = "",
        ) -> Path:
            target = (
                self.root
                / "プロジェクト"
                / canonical
                / "02.計画"
                / f"工程表{revision}.xlsx"
            )
            target.parent.mkdir(parents=True, exist_ok=True)
            workbook = Workbook()
            sheet = workbook.active
            sheet.title = "WBS"
            sheet.append(["番号", "イベント", "開始日", "終了日"])
            sheet.append(
                [
                    "P1",
                    "星門審査資料作成",
                    datetime.fromisoformat("2031-01-01"),
                    datetime.fromisoformat("2031-01-02"),
                ]
            )
            for index, (milestone, when) in enumerate(events, 1):
                event_date = datetime.fromisoformat(when)
                sheet.append(
                    [f"E{index}", f"{milestone}実施", event_date, event_date]
                )
            workbook.save(target)
            workbook.close()
            return target

        projects = (
            ("惑星機構甲", "ORB9", [("星門審査", "2032-05-08")]),
            ("惑星機構乙", "AX2", [("星門審査", "2032-04-30")]),
            ("惑星機構丙", "NOV7", [("軌道評議", "2032-05-09")]),
            (
                "惑星機構丁",
                "DUAL4",
                [("星門審査", "2032-05-07"), ("軌道評議", "2032-04-01")],
            ),
        )
        for canonical, alias, events in projects:
            self.glossary.add(alias, canonical, primary=True)
            write_plan(canonical, events)
        # A byte-distinct revision with the same extracted event date is not
        # an ambiguity: both sources independently assert the same milestone.
        write_plan(
            "惑星機構乙",
            [("星門審査", "2032-04-30")],
            revision="_rev2",
        )

        question = (
            "星門審査または軌道評議が2032年5月10日以前に実施された案件を、"
            "主略称ですべて挙げてください。"
        )
        contract = graph_contract_for_question(question)
        self.assertIsNotNone(contract)
        self.assertEqual(
            contract["rule_id"],
            "all_project_milestone_cutoff_primary_alias",
        )
        self.assertEqual(
            contract["requested_output"]["answer_shape"],
            {"container": "list", "value_type": "identifier", "unit": None},
        )
        self.assertTrue(validate_graph_contract(question, contract))

        engine = StructuredCandidateEngine(self.root, self.glossary)
        first = engine.decide("ignored-one", question)
        self.assertEqual(first.status, "resolved")
        # ORB9 is declared before AX2 even though its source date is later.
        # With no explicit sort in the question, the source glossary's primary
        # declaration order is the authoritative within-branch order.
        self.assertEqual(first.result.answer, "ORB9、AX2、DUAL4、NOV7")
        self.assertEqual(first.result.output_count, 1)
        self.assertEqual(len(first.result.source_paths), 5)

        wording_variant = (
            "星門審査もしくは軌道評議が2032/05/10までに開催された案件を"
            "全て主略称にて答えてください。"
        )
        variant_contract = graph_contract_for_question(wording_variant)
        self.assertIsNotNone(variant_contract)
        self.assertTrue(validate_graph_contract(wording_variant, variant_contract))
        variant = engine.decide("ignored-two", wording_variant)
        self.assertEqual(variant.status, "resolved")
        self.assertEqual(variant.result.answer, "ORB9、AX2、DUAL4、NOV7")

        # OR-branch order is semantic.  Reversing the question branches moves
        # the right-only group first, and the dual-match project follows that
        # first branch using its branch-local source date.
        reversed_question = (
            "軌道評議または星門審査が2032年5月10日以前に実施された案件を、"
            "主略称ですべて挙げてください。"
        )
        reversed_decision = engine.decide("ignored-reversed", reversed_question)
        self.assertEqual(reversed_decision.status, "resolved")
        self.assertEqual(
            reversed_decision.result.answer,
            "NOV7、DUAL4、ORB9、AX2",
        )

        # A date mutation does not override the source glossary declaration
        # order when the question itself contains no explicit sort request.
        write_plan("惑星機構甲", [("星門審査", "2032-04-01")])
        changed = engine.decide("ignored-three", question)
        self.assertEqual(changed.status, "resolved")
        self.assertEqual(changed.result.answer, "ORB9、AX2、DUAL4、NOV7")

    def test_all_project_milestone_cutoff_fails_closed_on_incomplete_project(self) -> None:
        from datetime import datetime

        from openpyxl import Workbook

        canonical = "衛星機構壱"
        self.glossary.add("SAT1", canonical, primary=True)
        plan = self.root / "プロジェクト" / canonical / "02.計画" / "WBS.xlsx"
        plan.parent.mkdir(parents=True, exist_ok=True)
        workbook = Workbook()
        sheet = workbook.active
        sheet.append(["イベント", "開始日", "終了日"])
        sheet.append(
            [
                "星門審査実施",
                datetime.fromisoformat("2032-03-01"),
                datetime.fromisoformat("2032-03-02"),
            ]
        )
        workbook.save(plan)
        workbook.close()
        question = (
            "星門審査または軌道評議が2032-05-10以前に実施された案件を、"
            "主略称ですべて挙げてください。"
        )
        contract = graph_contract_for_question(question)
        self.assertIsNotNone(contract)
        engine = StructuredCandidateEngine(self.root, self.glossary)
        self.assertIsNone(decide_extended(engine, "ignored", question))

        self.assertIsNone(
            graph_contract_for_question(
                question.replace("2032-05-10", "2032-02-30")
            )
        )

    def test_all_project_paid_gross_tax_sum_is_complete_and_metamorphic(self) -> None:
        from datetime import datetime

        from docx import Document
        from openpyxl import Workbook

        source_root = self.root / "opaque-paid-case"
        glossary = Glossary()

        def write_project(
            canonical: str,
            alias: str,
            *,
            paid_gross: int,
            contract_gross: int,
            due: str,
            revision: str = "",
            status: str = "完了",
            include_amount: bool = True,
            duplicate_sheet: bool = False,
        ) -> None:
            self.assertEqual(contract_gross % 11, 0)
            glossary.add(alias, canonical, primary=True)
            project = source_root / "プロジェクト" / canonical
            contract = project / "01.契約" / "契約書.docx"
            contract.parent.mkdir(parents=True, exist_ok=True)
            document = Document()
            document.add_paragraph("本契約は固定価格契約とする。")
            document.add_paragraph("消費税率: 10%")
            table = document.add_table(rows=2, cols=7)
            headers = (
                "payment_key",
                "amount_net",
                "tax",
                "amount_gross",
                "condition",
                "due_date",
                "note",
            )
            for index, value in enumerate(headers):
                table.rows[0].cells[index].text = value
            tax = contract_gross // 11
            net = contract_gross - tax
            values = (
                "phase_q",
                f"{net:,}円",
                f"{tax:,}円",
                f"{contract_gross:,}円",
                "opaque acceptance",
                due,
                "source terms",
            )
            for index, value in enumerate(values):
                table.rows[1].cells[index].text = value
            document.save(contract)

            plan = project / "02.計画" / f"工程表{revision}.xlsx"
            plan.parent.mkdir(parents=True, exist_ok=True)
            workbook = Workbook()
            sheet = workbook.active
            sheet.title = "ledger_q"
            sheet.append(["kind_q", "event_q", "detail_q", "date_q", "state_q"])
            detail = (
                f"税込金額 {paid_gross:,}円" if include_amount else "opaque transfer"
            )
            row = ["支払", "phase_q", detail, datetime.fromisoformat(due), status]
            sheet.append(row)
            if duplicate_sheet:
                copy = workbook.create_sheet("ledger_copy_q")
                copy.append(["kind_q", "event_q", "detail_q", "date_q", "state_q"])
                copy.append(row)
            workbook.save(plan)
            workbook.close()

        gross_alpha = 121_000
        gross_beta = 242_000
        gross_gamma = 363_000
        write_project(
            "opaque_project_alpha",
            "OPA1",
            paid_gross=gross_alpha,
            contract_gross=110_000,
            due="2034-01-11",
        )
        # A completed plan row may omit the amount only when its date joins to
        # one exact contract installment.
        write_project(
            "opaque_project_beta",
            "OPB2",
            paid_gross=gross_beta,
            contract_gross=gross_beta,
            due="2034-02-12",
            include_amount=False,
            duplicate_sheet=True,
        )
        # Explicit revision semantics select r2; the pending r1 must not veto
        # or contribute to the current payment ledger.
        write_project(
            "opaque_project_gamma",
            "OPG3",
            paid_gross=gross_gamma,
            contract_gross=330_000,
            due="2034-03-13",
            revision="_r1",
            status="未着手",
        )
        write_project(
            "opaque_project_gamma",
            "OPG3",
            paid_gross=gross_gamma,
            contract_gross=330_000,
            due="2034-03-13",
            revision="_r2",
        )

        question = (
            "全案件で支払った税込金額をもとに、"
            "消費税額の総額を計算してください。"
        )
        contract = graph_contract_for_question(question)
        self.assertIsNotNone(contract)
        self.assertEqual(contract["graph_rule_version"], "1.4")
        self.assertEqual(contract["rule_id"], "all_project_paid_gross_tax_sum")
        self.assertEqual(
            [node["operator"] for node in contract["operation_graph"]["nodes"]],
            [
                "enumerate_projects",
                "retrieve",
                "decrypt",
                "select_authoritative",
                "match",
                "verify_complete",
                "deduplicate",
                "retrieve_actual",
                "verify_arithmetic",
                "override_estimate",
                "calculate",
                "sum",
            ],
        )
        self.assertEqual(
            contract["requested_output"]["answer_shape"],
            {"container": "scalar", "value_type": "integer", "unit": "円"},
        )
        self.assertTrue(validate_graph_contract(question, contract))
        self.assertIsNone(graph_contract_for_question(question + "追記"))

        engine = StructuredCandidateEngine(source_root, glossary)
        graph_plan = SimpleNamespace(branch_intents=())
        first = engine.decide_from_graph("opaque-run-one", question, graph_plan)
        expected = sum((gross_alpha, gross_beta, gross_gamma)) // 11
        self.assertEqual(first.status, "resolved")
        self.assertEqual(first.result.answer, f"{expected:,}円")
        self.assertEqual(first.result.operation_count, 12)
        self.assertEqual(first.result.output_count, 1)

        wording_variant = (
            "すべての案件において支払済みの税込金額をもとに、"
            "消費税額の総額を算出してください。"
        )
        variant = engine.decide_from_graph(
            "opaque-run-variant",
            wording_variant,
            graph_plan,
        )
        self.assertEqual(variant.status, "resolved")
        self.assertEqual(variant.result.answer, f"{expected:,}円")

        changed_gross = gross_alpha + 1_100
        write_project(
            "opaque_project_alpha",
            "OPA1",
            paid_gross=changed_gross,
            contract_gross=110_000,
            due="2034-01-11",
        )
        changed = engine.decide_from_graph("opaque-run-two", question, graph_plan)
        changed_expected = sum((changed_gross, gross_beta, gross_gamma)) // 11
        self.assertEqual(changed.status, "resolved")
        self.assertEqual(changed.result.answer, f"{changed_expected:,}円")
        self.assertNotEqual(changed.result.source_sha256, first.result.source_sha256)

    def test_all_project_paid_gross_tax_sum_overrides_tm_estimate_with_actual(self) -> None:
        from datetime import datetime
        from decimal import Decimal

        from docx import Document
        from openpyxl import Workbook
        from pptx import Presentation
        from pptx.util import Inches

        source_root = self.root / "opaque-tm-actual"
        glossary = Glossary()
        rate = 20_000

        def write_contract_and_plan(
            canonical: str,
            alias: str,
            *,
            estimate_gross: int,
            due: str,
            status: str = "完了",
        ) -> Path:
            glossary.add(alias, canonical, primary=True)
            project = source_root / "プロジェクト" / canonical
            contract = project / "01.契約" / "契約書.docx"
            contract.parent.mkdir(parents=True, exist_ok=True)
            document = Document()
            document.add_paragraph(
                "契約形態: time_and_materials。報酬は実績工数に基づき精算する。"
            )
            document.add_paragraph(f"時間単価:{rate:,}円")
            document.add_paragraph("消費税率:10%")
            tax = estimate_gross // 11
            net = estimate_gross - tax
            table = document.add_table(rows=1, cols=4)
            for index, value in enumerate(
                (f"{net:,}円", f"{tax:,}円", f"{estimate_gross:,}円", due)
            ):
                table.rows[0].cells[index].text = value
            document.save(contract)

            plan = project / "02.計画" / "工程表.xlsx"
            plan.parent.mkdir(parents=True, exist_ok=True)
            workbook = Workbook()
            sheet = workbook.active
            sheet.append(["kind", "detail", "due", "state"])
            sheet.append(
                [
                    "支払",
                    f"最終精算 税込{estimate_gross:,}円(見積)",
                    datetime.fromisoformat(due),
                    status,
                ]
            )
            workbook.save(plan)
            workbook.close()
            return project

        def write_report(
            project: Path,
            *,
            actual_hours: str,
            disclaimed: bool = False,
            old: bool = False,
        ) -> None:
            hours = Decimal(actual_hours)
            net = int(hours * rate)
            tax = net // 10
            gross = net + tax
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.5), Inches(9), Inches(6)
            )
            note = (
                "\n上記は例示。実際の請求書は工数記録を優先する。"
                if disclaimed
                else ""
            )
            box.text = (
                "費用・請求(確定事項)\n"
                f"実績工数:{actual_hours}時間\n"
                f"時間単価:¥{rate:,}\n"
                f"最終請求金額(税抜):¥{net:,}\n"
                f"消費税額:¥{tax:,}\n"
                f"最終請求金額(税込):¥{gross:,}{note}"
            )
            report_dir = project / "06.報告書"
            if old:
                report_dir /= "old"
            report_dir.mkdir(parents=True, exist_ok=True)
            presentation.save(report_dir / "opaque_最終報告.pptx")

        project_actual = write_contract_and_plan(
            "opaque_actual",
            "OA1",
            estimate_gross=132_000,
            due="2036-01-10",
        )
        write_report(project_actual, actual_hours="7.5")
        # A stale report cannot compete with the unique current report.
        write_report(project_actual, actual_hours="99", old=True)

        project_estimate = write_contract_and_plan(
            "opaque_estimate_only",
            "OE2",
            estimate_gross=220_000,
            due="2036-02-11",
        )
        write_report(project_estimate, actual_hours="10", disclaimed=True)

        question = (
            "全案件で支払った税込金額をもとに、"
            "消費税額の総額を計算してください。"
        )
        engine = StructuredCandidateEngine(source_root, glossary)
        first = engine.decide_from_graph(
            "opaque-actual-one",
            question,
            SimpleNamespace(branch_intents=()),
        )
        expected = int(Decimal("7.5") * rate) // 10 + 220_000 // 11
        self.assertEqual(first.status, "resolved")
        self.assertEqual(first.result.answer, f"{expected:,}円")
        self.assertEqual(first.result.operation_count, 12)
        self.assertEqual(first.result.output_count, 1)

        write_report(project_actual, actual_hours="8")
        changed = engine.decide_from_graph(
            "opaque-actual-two",
            question,
            SimpleNamespace(branch_intents=()),
        )
        changed_expected = 8 * rate // 10 + 220_000 // 11
        self.assertEqual(changed.status, "resolved")
        self.assertEqual(changed.result.answer, f"{changed_expected:,}円")
        self.assertNotEqual(first.result.source_sha256, changed.result.source_sha256)

    def test_all_project_paid_gross_tax_sum_fails_closed_without_every_payment(self) -> None:
        from datetime import datetime

        from docx import Document
        from openpyxl import Workbook

        question = (
            "全案件で支払った税込金額をもとに、"
            "消費税額の総額を計算してください。"
        )

        def make_case(
            case_name: str,
            rows_by_project: tuple[tuple[str, list[list[object]]], ...],
            *,
            contract_due: str = "2035-01-10",
        ) -> tuple[StructuredCandidateEngine, Glossary]:
            source_root = self.root / case_name
            glossary = Glossary()
            for ordinal, (canonical, rows) in enumerate(rows_by_project, 1):
                glossary.add(f"PX{ordinal}", canonical, primary=True)
                project = source_root / "プロジェクト" / canonical
                contract = project / "01.契約" / "契約書.docx"
                contract.parent.mkdir(parents=True, exist_ok=True)
                document = Document()
                document.add_paragraph("本契約は固定価格契約とする。")
                document.add_paragraph("消費税率: 10%")
                table = document.add_table(rows=1, cols=5)
                values = (
                    "phase_x",
                    "110,000円",
                    "11,000円",
                    "121,000円",
                    contract_due,
                )
                for index, value in enumerate(values):
                    table.rows[0].cells[index].text = value
                document.save(contract)
                plan = project / "02.計画" / "工程表.xlsx"
                plan.parent.mkdir(parents=True, exist_ok=True)
                workbook = Workbook()
                sheet = workbook.active
                sheet.append(["kind_x", "detail_x", "date_x", "state_x"])
                for row in rows:
                    sheet.append(row)
                workbook.save(plan)
                workbook.close()
            return StructuredCandidateEngine(source_root, glossary), glossary

        completed_row = [
            "支払",
            "税込金額 121,000円",
            datetime.fromisoformat("2035-01-10"),
            "完了",
        ]
        pending_row = [
            "支払",
            "税込金額 121,000円",
            datetime.fromisoformat("2035-01-10"),
            "未着手",
        ]
        engine, _ = make_case(
            "pending-project",
            (("opaque_one", [completed_row]), ("opaque_two", [pending_row])),
        )
        self.assertIsNone(decide_extended(engine, "ignored-pending", question))
        wrapped = engine.decide_from_graph(
            "opaque-pending",
            question,
            SimpleNamespace(branch_intents=()),
        )
        self.assertEqual(wrapped.status, "hold")

        marker_only = [
            "支払関連",
            "opaque marker only",
            datetime.fromisoformat("2035-02-20"),
            "完了",
        ]
        engine, _ = make_case(
            "marker-only",
            (("opaque_marker", [marker_only]),),
        )
        self.assertIsNone(decide_extended(engine, "ignored-marker", question))

        conflict_rows = [
            completed_row,
            [
                "支払",
                "税込金額 122,100円",
                datetime.fromisoformat("2035-01-10"),
                "完了",
            ],
        ]
        engine, _ = make_case(
            "conflicting-payment",
            (("opaque_conflict", conflict_rows),),
        )
        self.assertIsNone(decide_extended(engine, "ignored-conflict", question))

        fractional_tax = [
            "支払",
            "税込金額 121,001円",
            datetime.fromisoformat("2035-01-10"),
            "完了",
        ]
        engine, _ = make_case(
            "fractional-yen",
            (("opaque_fraction", [fractional_tax]),),
        )
        self.assertIsNone(decide_extended(engine, "ignored-rounding", question))

    def test_all_project_paid_gross_tax_sum_fails_closed_on_tm_actual_ambiguity(self) -> None:
        from datetime import datetime

        from docx import Document
        from openpyxl import Workbook
        from pptx import Presentation
        from pptx.util import Inches

        question = (
            "全案件で支払った税込金額をもとに、"
            "消費税額の総額を計算してください。"
        )
        rate = 10_000
        estimate_gross = 121_000
        due = "2037-03-12"

        def make_tm_case(
            case_name: str,
            *,
            status: str,
            claims: tuple[tuple[str, int, int, int], ...],
            duplicate_report: bool = False,
        ) -> StructuredCandidateEngine:
            source_root = self.root / case_name
            canonical = "opaque_tm"
            glossary = Glossary()
            glossary.add("OTM1", canonical, primary=True)
            project = source_root / "プロジェクト" / canonical
            contract = project / "01.契約" / "契約書.docx"
            contract.parent.mkdir(parents=True, exist_ok=True)
            document = Document()
            document.add_paragraph(
                "time_and_materials契約とし、実績工数に基づき精算する。"
            )
            document.add_paragraph(f"時間単価:{rate:,}円")
            document.add_paragraph("消費税率:10%")
            table = document.add_table(rows=1, cols=4)
            for index, value in enumerate(
                ("110,000円", "11,000円", "121,000円", due)
            ):
                table.rows[0].cells[index].text = value
            document.save(contract)

            plan = project / "02.計画" / "工程表.xlsx"
            plan.parent.mkdir(parents=True, exist_ok=True)
            workbook = Workbook()
            sheet = workbook.active
            sheet.append(["kind", "detail", "due", "state"])
            sheet.append(
                [
                    "支払",
                    f"税込金額 {estimate_gross:,}円(見積)",
                    datetime.fromisoformat(due),
                    status,
                ]
            )
            workbook.save(plan)
            workbook.close()

            def save_report(name: str) -> None:
                presentation = Presentation()
                for hours, net, tax, gross in claims:
                    slide = presentation.slides.add_slide(
                        presentation.slide_layouts[6]
                    )
                    box = slide.shapes.add_textbox(
                        Inches(0.5), Inches(0.5), Inches(9), Inches(6)
                    )
                    box.text = (
                        "費用・請求(確定事項)\n"
                        f"実績工数:{hours}時間\n"
                        f"時間単価:¥{rate:,}\n"
                        f"最終請求金額(税抜):¥{net:,}\n"
                        f"消費税額:¥{tax:,}\n"
                        f"最終請求金額(税込):¥{gross:,}"
                    )
                report = project / "06.報告書" / name
                report.parent.mkdir(parents=True, exist_ok=True)
                presentation.save(report)

            save_report("opaque_最終報告.pptx")
            if duplicate_report:
                save_report("opaque_copy_最終報告.pptx")
            return StructuredCandidateEngine(source_root, glossary)

        pending = make_tm_case(
            "tm-pending-proof",
            status="未着手",
            claims=(("12", 120_000, 12_000, 132_000),),
        )
        self.assertIsNone(decide_extended(pending, "ignored-pending-tm", question))

        conflicting = make_tm_case(
            "tm-conflicting-actuals",
            status="完了",
            claims=(
                ("12", 120_000, 12_000, 132_000),
                ("13", 130_000, 13_000, 143_000),
            ),
        )
        self.assertIsNone(decide_extended(conflicting, "ignored-conflict-tm", question))

        inconsistent = make_tm_case(
            "tm-inconsistent-arithmetic",
            status="完了",
            claims=(("12", 121_000, 12_100, 133_100),),
        )
        self.assertIsNone(decide_extended(inconsistent, "ignored-math-tm", question))

        duplicate = make_tm_case(
            "tm-duplicate-final",
            status="完了",
            claims=(("12", 120_000, 12_000, 132_000),),
            duplicate_report=True,
        )
        self.assertIsNone(decide_extended(duplicate, "ignored-duplicate-tm", question))

    def test_native_office_extended_rules_cannot_bypass_live_graph_plan(self) -> None:
        engine = StructuredCandidateEngine(self.root, self.glossary)
        questions = (
            "site_alphaとの契約書において、太字で記載されている部分を抽出してください。",
            "site_alphaのdata_alpha.xlsxにおいて、Sheet1の黄色にハイライトされたセルの"
            "抽出条件と集計内容を答えてください。",
            "site_alphaのdeck_alpha.pptxにおいて、この案件にかかる金額の提示が"
            "まとまっているのは何ページですか。",
            "site_alphaのdeck_v1.pptxからdeck_v2.pptxに修正されたもののうち、"
            "案件遂行に関連する変更を挙げてください。",
            "folder_alphaにあるmap_alphaにおいて、person_alphaさんから見て"
            "右側に座っている人の名前をすべて挙げてください。",
        )
        for question in questions:
            with self.subTest(question=question):
                decision = engine.decide("legacy-direct", question)
                self.assertEqual(decision.status, "hold")
                self.assertEqual(decision.reason, "extended_graph_plan_required")

    def test_native_office_extended_rules_require_resolved_graph_branch(self) -> None:
        from question_graph_runtime import build_graph_plan

        engine = StructuredCandidateEngine(self.root, self.glossary)
        questions = (
            "site_alphaとの契約書において、太字で記載されている部分を抽出してください。",
            "site_alphaのdata_alpha.xlsxにおいて、Sheet1の黄色にハイライトされたセルの"
            "抽出条件と集計内容を答えてください。",
            "site_alphaのdeck_alpha.pptxにおいて、この案件にかかる金額の提示が"
            "まとまっているのは何ページですか。",
            "site_alphaのdeck_v1.pptxからdeck_v2.pptxに修正されたもののうち、"
            "案件遂行に関連する変更を挙げてください。",
            "folder_alphaにあるmap_alphaにおいて、person_alphaさんから見て"
            "右側に座っている人の名前をすべて挙げてください。",
        )
        for question in questions:
            with self.subTest(question=question):
                plan = build_graph_plan("branch-status", question)
                branches = copy.deepcopy(plan.branch_intents)
                branches[0]["status"] = "hold"
                tampered = replace(plan, branch_intents=branches)
                decision = engine.decide_from_graph(
                    "branch-status", question, tampered
                )
                self.assertEqual(decision.status, "hold")
                self.assertEqual(
                    decision.reason, "extended_graph_plan_not_certified"
                )


if __name__ == "__main__":
    unittest.main()
