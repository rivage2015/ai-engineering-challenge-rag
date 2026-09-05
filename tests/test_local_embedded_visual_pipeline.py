from __future__ import annotations

import base64
import copy
import hashlib
import importlib.util
import json
import sqlite3
import sys
import tempfile
import unittest
import zipfile
from pathlib import Path
from unittest import mock


ROOT = Path(__file__).resolve().parents[1]
SCRIPTS = ROOT / "scripts"
ENGINE = ROOT / "distribution" / "macos-local-memory" / "engine"
sys.path.insert(0, str(SCRIPTS))

import adapt_layer1_to_local_memory as adapter  # noqa: E402
import build_cross_document_semantic_graph as graph_builder  # noqa: E402
import build_search_units as search_units  # noqa: E402
import local_image_ocr  # noqa: E402
import local_visual_observation  # noqa: E402
import probe_intermediate_records as records  # noqa: E402
import validate_intermediate_records_streaming  # noqa: E402
import validate_search_units_streaming  # noqa: E402


def load_engine(name: str):
    path = ENGINE / f"{name}.py"
    spec = importlib.util.spec_from_file_location(
        f"embedded_visual_{name}", path
    )
    assert spec is not None and spec.loader is not None
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


RUN_AT = "2031-04-01T00:00:00+00:00"
OCR_FRAME = "source_orientation_1_top_left_normalized_1000"
PROVISIONAL_MARKER = "[暫定読取]"

# A complete, tiny PNG.  The production image readers are mocked, but keeping a
# real raster payload makes the container/source-digest boundary representative.
PNG_BYTES = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR4"
    "2mP8/x8AAusB9Y9ZK7sAAAAASUVORK5CYII="
)


def high_ocr_observation(image_path: Path, text: str) -> dict[str, object]:
    raw_sha256 = hashlib.sha256(image_path.read_bytes()).hexdigest()
    bbox = [100, 100, 700, 80]
    supporters = [
        {
            "pass": "apple_vision_primary",
            "engine": "apple_vision",
            "independence_group": "apple_vision",
            "line_id": "vision-1",
            "raw_text": text,
            "bbox": bbox,
            "bbox_coordinate_system": OCR_FRAME,
            "confidence": 0.95,
        },
        {
            "pass": "tesseract_psm3",
            "engine": "tesseract",
            "independence_group": "tesseract",
            "line_id": "tesseract-1",
            "raw_text": text,
            "bbox": bbox,
            "bbox_coordinate_system": OCR_FRAME,
            "confidence": 0.90,
        },
    ]
    line = {
        "text": text,
        "bbox": bbox,
        "bbox_coordinate_system": OCR_FRAME,
        "overlap": 1.0,
        "primary_confidence": 0.95,
        "audit_confidence": 0.90,
        "agreement_type": "independent_agreement",
        "quality_tier": "high",
        "provenance": {
            "primary_pass": "apple_vision_primary",
            "audit_pass": "tesseract_psm3",
            "primary_engine": "apple_vision",
            "audit_engine": "tesseract",
            "primary_independence_group": "apple_vision",
            "audit_independence_group": "tesseract",
            "primary_line_id": "vision-1",
            "audit_line_id": "tesseract-1",
            "primary_bbox_coordinate_system": OCR_FRAME,
            "audit_bbox_coordinate_system": OCR_FRAME,
            "comparison_coordinate_system": OCR_FRAME,
            "supporters": supporters,
        },
    }
    dimensions = {"width_px": 1, "height_px": 1}
    return {
        "input_sha256": raw_sha256,
        "source_dimensions": dimensions,
        "dimensions": dimensions,
        "image_format": "PNG",
        "orientation": 1,
        "orientation_source": "imageio_canonicalizer",
        "canonicalization": {
            "status": "completed",
            "canonical_dimensions": dimensions,
            "canonical_orientation": 1,
        },
        "ocr_input_sha256": raw_sha256,
        "ocr_input_dimensions": dimensions,
        "ocr_input_orientation": 1,
        "coordinate_frame_policy": "canonical_orientation_1",
        "engines": {
            "apple_vision_primary": {"status": "completed"},
            "tesseract_psm3": {"status": "completed"},
        },
        "independent_engines": True,
        "consensus_lines": [line],
        "read_lines": [line],
        "unlocated_transcript": None,
        "unresolved_count": 0,
    }


def provisional_visual_observation(
    image_path: Path,
    *,
    expected_input_sha256: str | None = None,
) -> dict[str, object]:
    image_sha256 = hashlib.sha256(image_path.read_bytes()).hexdigest()
    if expected_input_sha256 is not None:
        assert expected_input_sha256 == image_sha256
    observation = {
        "visible_objects": [
            {
                "object_id": "o1",
                "kind": "chart",
                "description": "架空の棒グラフ",
            }
        ],
        "explicit_labels": [
            {"label_id": "l1", "text": "幻影計画"},
            {"label_id": "l2", "text": "担当者 視覚太郎"},
        ],
        "explicit_relations": [
            {"source_ref": "l2", "relation": "labels", "target_ref": "o1"}
        ],
        "labeled_values": [
            {
                "value_id": "v1",
                "label_text": "2031年",
                "series_label": "稼働数",
                "value_text": "12",
                "unit_text": "回",
                "value_status": "exact_label",
                "unclear_reason": "",
            }
        ],
        "warnings": [],
    }
    text = "\n".join(
        [
            f"{PROVISIONAL_MARKER} object o1: chart / 架空の棒グラフ",
            f"{PROVISIONAL_MARKER} label l1: 幻影計画",
            f"{PROVISIONAL_MARKER} label l2: 担当者 視覚太郎",
            f"{PROVISIONAL_MARKER} value v1: 2031年 / 稼働数 = 12 回 / exact_label",
        ]
    )
    return {
        "schema_version": "0.1",
        "record_type": "local_visual_observation",
        "observation_type": "whole_image_literal_visual_observation",
        "status": "provisional",
        "quality_tier": "provisional",
        "provisional_marker": PROVISIONAL_MARKER,
        "text": text,
        "observation": observation,
        "question_independent": True,
        "model": "gemma4:12b",
        "model_digest": "a" * 64,
        "prompt_sha256": local_visual_observation.VISUAL_OBSERVATION_PROMPT_SHA256,
        "input_image_sha256": image_sha256,
        "model_output_sha256": "b" * 64,
        "runner": local_visual_observation.VISUAL_OBSERVATION_RUNNER,
        "runner_version": local_visual_observation.VISUAL_OBSERVATION_VERSION,
        "host": "127.0.0.1",
        "port": 11434,
        "temperature": 0,
        "strict_json": True,
        "external_network_used": False,
        "downloads_performed": False,
    }


def write_jsonl(path: Path, values: list[dict[str, object]]) -> None:
    path.write_text(
        "".join(records.canonical_json(value) + "\n" for value in values),
        encoding="utf-8",
    )


def materialize_intermediate(
    probe: records.Probe,
    source_root: Path,
    output: Path,
) -> None:
    output.mkdir()
    write_jsonl(output / "documents.jsonl", probe.documents)
    # SearchUnit grouping depends on extraction order, so do not sort here.
    write_jsonl(output / "evidence.jsonl", probe.evidence)
    write_jsonl(output / "relations.jsonl", probe.relations)
    document = probe.documents[0]
    relative_path = document["source"]["relative_path"]
    evidence_path = output / "evidence.jsonl"
    state = {
        "state_version": "1",
        "build_status": "complete",
        "source_root": str(source_root.resolve()),
        "extractor": probe.extractor,
        "extractor_version": probe.extractor_version,
        "run_at": probe.run_at,
        "input_paths": [relative_path],
        "entries": {
            relative_path: {
                "document_id": document["document_id"],
                "relative_path": relative_path,
                "source_sha256": document["source"]["sha256"],
                "status": document["extraction"]["status"],
                "shards": {
                    "evidence": {
                        "relative_path": evidence_path.name,
                        "sha256": records.digest_file(evidence_path),
                        "size_bytes": evidence_path.stat().st_size,
                        "record_count": len(probe.evidence),
                    }
                },
            }
        },
        "totals": {
            "documents": len(probe.documents),
            "evidence": len(probe.evidence),
            "relations": len(probe.relations),
        },
    }
    (output / "build-state.json").write_text(
        records.canonical_json(state) + "\n", encoding="utf-8"
    )


class LocalEmbeddedVisualPipelineTests(unittest.TestCase):
    def setUp(self) -> None:
        self.temporary = tempfile.TemporaryDirectory(prefix="aiec-embedded-visual-")
        self.work = Path(self.temporary.name)
        self.source_root = self.work / "source"
        self.source_root.mkdir()

    def tearDown(self) -> None:
        self.temporary.cleanup()

    def _extract_office_helper(self) -> tuple[records.Probe, Path]:
        source = self.source_root / "fictional-report.docx"
        source.write_bytes(b"fictional OOXML container boundary")
        probe = records.Probe(self.source_root, RUN_AT, None, diagnostic=False)
        document = probe.add_document(source, "fictional-safe-office-helper")
        location = {
            "source_member": "word/media/image1.png",
            "object_index": 3,
        }
        content_ref = f"{source.name}::word/media/image1.png"
        visual_origin = probe._embedded_visual_origin(
            PNG_BYTES,
            document,
            location_prefix=location,
            source_name="word/media/image1.png",
            visual_origin_kind="office_embedded_image",
        )
        parent = probe.add_evidence(
            document["document_id"],
            "image",
            location,
            records.content(content_ref=content_ref, mime_type="image/png"),
            ordinal=3,
            native_properties={
                "embedded_sha256": hashlib.sha256(PNG_BYTES).hexdigest(),
                "size_bytes": len(PNG_BYTES),
                "visual_origin": {
                    **visual_origin,
                },
            },
        )
        probe.contain_document(document["document_id"], parent["evidence_id"])
        projected = probe._project_embedded_image_bytes(
            PNG_BYTES,
            document,
            parent_id=parent["evidence_id"],
            location_prefix=location,
            content_ref=content_ref,
            source_name="word/media/image1.png",
            visual_origin=visual_origin,
        )
        self.assertEqual(projected, 2)
        probe.finalize_document()
        return probe, source

    def _extract_notebook(self) -> tuple[records.Probe, Path]:
        source = self.source_root / "fictional-analysis.ipynb"
        notebook = {
            "cells": [
                {
                    "cell_type": "code",
                    "execution_count": None,
                    "metadata": {},
                    "source": ["# fictional embedded chart"],
                    "outputs": [
                        {
                            "output_type": "display_data",
                            "metadata": {},
                            "data": {
                                "image/png": base64.b64encode(PNG_BYTES).decode("ascii")
                            },
                        }
                    ],
                }
            ],
            "metadata": {},
            "nbformat": 4,
            "nbformat_minor": 5,
        }
        source.write_text(json.dumps(notebook), encoding="utf-8")
        probe = records.Probe(self.source_root, RUN_AT, None, diagnostic=False)
        probe.extract(source)
        return probe, source

    def _assert_pipeline(
        self,
        probe: records.Probe,
        source: Path,
        *,
        expected_kind: str,
        expected_locator: dict[str, object],
    ) -> None:
        self.assertEqual(len(probe.documents), 1)
        source_sha256 = records.digest_file(source)
        images = [item for item in probe.evidence if item["evidence_type"] == "image"]
        ocr_lines = [
            item for item in probe.evidence if item["evidence_type"] == "ocr_line"
        ]
        visual = [
            item
            for item in probe.evidence
            if item.get("provenance", {}).get("extraction_method")
            == "local_vlm_visual_observation_provisional"
        ]
        self.assertEqual(len(images), 1)
        self.assertEqual(len(ocr_lines), 1)
        self.assertEqual(len(visual), 1)

        expected_ocr_tier = (
            "provisional"
            if expected_kind in {
                "office_embedded_image",
                "notebook_embedded_image",
            }
            else "high"
        )
        self.assertEqual(
            expected_ocr_tier,
            ocr_lines[0]["native_properties"]["quality_tier"],
        )
        if expected_kind in {
            "office_embedded_image",
            "notebook_embedded_image",
        }:
            self.assertEqual(
                "display_transform_unresolved",
                ocr_lines[0]["native_properties"]["agreement_type"],
            )
            self.assertEqual(
                "independent_agreement",
                ocr_lines[0]["native_properties"][
                    "embedded_source_agreement_type"
                ],
            )
            self.assertEqual(
                "adaptive_local_ocr_provisional",
                ocr_lines[0]["provenance"]["extraction_method"],
            )
            self.assertFalse(
                ocr_lines[0]["native_properties"][
                    "display_transform_resolved"
                ]
            )

        parent = images[0]
        for projected in (*ocr_lines, *visual):
            self.assertEqual(projected["parent_evidence_id"], parent["evidence_id"])
            origin = projected["native_properties"]["visual_origin"]
            self.assertEqual(origin["kind"], expected_kind)
            self.assertEqual(origin["source_relative_path"], source.name)
            self.assertEqual(origin["source_sha256"], source_sha256)
            self.assertEqual(origin["source_location"], expected_locator)
            self.assertEqual(
                origin["materialization"]["embedded_sha256"],
                hashlib.sha256(PNG_BYTES).hexdigest(),
            )
            self.assertFalse(origin["materialization"]["external_network_used"])

        intermediate = self.work / f"intermediate-{expected_kind}"
        search_output = self.work / f"search-{expected_kind}"
        semantic_output = self.work / f"semantic-{expected_kind}"
        materialize_intermediate(probe, self.source_root, intermediate)
        search_units.build(intermediate, search_output, 500)
        adapter.adapt(
            intermediate,
            self.source_root.resolve(),
            semantic_output,
            search_output,
        )

        packets = [
            json.loads(line)
            for line in (search_output / "search_units.jsonl")
            .read_text(encoding="utf-8")
            .splitlines()
            if line
        ]
        image_packets = [
            item for item in packets if item["unit_type"] == "image_text_packet"
        ]
        self.assertEqual(len(image_packets), 1)
        self.assertEqual(image_packets[0]["context"]["container_kind"], expected_kind)
        self.assertEqual(
            expected_ocr_tier,
            image_packets[0]["context"]["quality_tier"],
        )
        for key, value in expected_locator.items():
            if key == "locator_text":
                # The packet replaces free-form locator text with its audited
                # grouping contract; structured container coordinates remain.
                continue
            self.assertEqual(image_packets[0]["locator"][key], value)
        provisional_visual_units = [
            item
            for item in packets
            if item["unit_type"] == "text_chunk"
            and item.get("context", {}).get("quality_tier") == "provisional"
        ]
        self.assertEqual(len(provisional_visual_units), 1)
        self.assertEqual(
            provisional_visual_units[0]["context"]["container_kind"],
            expected_kind,
        )
        self.assertEqual(
            provisional_visual_units[0]["context"]["provisional_marker"],
            PROVISIONAL_MARKER,
        )
        self.assertTrue(all(
            line.startswith(PROVISIONAL_MARKER + " ")
            for line in provisional_visual_units[0]["text"]["search_text"].splitlines()
        ))

        semantic = [
            json.loads(line)
            for line in (semantic_output / "semantic-evidence.jsonl")
            .read_text(encoding="utf-8")
            .splitlines()
            if line
        ]
        visual_semantic = [
            item
            for item in semantic
            if item["extraction_method"]
            == "local_vlm_visual_observation_provisional"
        ]
        self.assertEqual(len(visual_semantic), 1)
        projected = visual_semantic[0]
        self.assertEqual(projected["source"]["relative_path"], source.name)
        self.assertEqual(projected["source"]["sha256"], source_sha256)
        self.assertEqual(projected["quality_tier"], "provisional")
        self.assertEqual(projected["provisional_marker"], PROVISIONAL_MARKER)
        self.assertIn("幻影計画", projected["observed_text"])
        self.assertTrue(
            all(
                line.startswith(PROVISIONAL_MARKER + " ")
                for line in projected["observed_text"].splitlines()
            )
        )
        # Child observations have their own object index and locator text.  The
        # stable parent container keys survive the projection, while the full
        # original location was checked above in visual_origin.source_location.
        for key, value in expected_locator.items():
            if key in {"object_index", "locator_text"}:
                continue
            self.assertEqual(projected["locator"][key], value)
        self.assertEqual(
            projected["locator"]["locator_text"],
            (
                expected_locator.get("locator_text", "") + ";"
                if expected_locator.get("locator_text") else ""
            ) + "visual_observation=whole_image",
        )
        if "object_index" in expected_locator:
            self.assertEqual(
                projected["locator"]["image_object_index"],
                expected_locator["object_index"],
            )

        # Searchability and verified-graph eligibility are separate gates.  The
        # provisional VLM record reaches Semantic Evidence, but the graph
        # quality gate must quarantine it from verified nodes and edges.
        self.assertEqual(
            graph_builder._quality_disposition(
                projected, projected["observed_text"]
            ),
            "excluded_provisional",
        )
        if expected_kind in {
            "office_embedded_image",
            "notebook_embedded_image",
        }:
            projected_ocr = next(
                item for item in semantic
                if item["extraction_method"]
                == "adaptive_local_ocr_provisional"
            )
            self.assertEqual("provisional", projected_ocr["quality_tier"])
            self.assertEqual(
                "excluded_provisional",
                graph_builder._quality_disposition(
                    projected_ocr, projected_ocr["observed_text"]
                ),
            )

    def test_office_embedded_image_preserves_member_and_source_lineage(self) -> None:
        with (
            mock.patch.object(
                local_image_ocr,
                "extract",
                side_effect=lambda path: high_ocr_observation(
                    path, "埋め込み画像の確定OCR"
                ),
            ),
            mock.patch.object(
                local_visual_observation,
                "observe_path",
                side_effect=provisional_visual_observation,
            ),
        ):
            probe, source = self._extract_office_helper()
        self._assert_pipeline(
            probe,
            source,
            expected_kind="office_embedded_image",
            expected_locator={
                "source_member": "word/media/image1.png",
                "object_index": 3,
            },
        )

    def test_notebook_embedded_image_preserves_cell_and_source_lineage(self) -> None:
        with (
            mock.patch.object(
                local_image_ocr,
                "extract",
                side_effect=lambda path: high_ocr_observation(
                    path, "Notebook画像の確定OCR"
                ),
            ),
            mock.patch.object(
                local_visual_observation,
                "observe_path",
                side_effect=provisional_visual_observation,
            ),
        ):
            probe, source = self._extract_notebook()
        self._assert_pipeline(
            probe,
            source,
            expected_kind="notebook_embedded_image",
            expected_locator={
                "notebook_cell_index": 1,
                "object_index": 1,
                "locator_text": "cell=1;output=1;output-image=1",
            },
        )

    def test_notebook_referenced_markdown_attachment_is_visually_read(self) -> None:
        source = self.source_root / "attachment.ipynb"
        notebook = {
            "cells": [{
                "cell_type": "markdown",
                "metadata": {},
                "source": ["![chart](attachment:chart.png)"],
                "attachments": {
                    "chart.png": {
                        "image/png": base64.b64encode(PNG_BYTES).decode("ascii")
                    },
                    "unused.png": {
                        "image/png": base64.b64encode(PNG_BYTES).decode("ascii")
                    },
                },
            }],
            "metadata": {},
            "nbformat": 4,
            "nbformat_minor": 5,
        }
        source.write_text(json.dumps(notebook), encoding="utf-8")
        with (
            mock.patch.object(
                local_image_ocr,
                "extract",
                side_effect=lambda path: high_ocr_observation(path, "添付画像の確定OCR"),
            ),
            mock.patch.object(
                local_visual_observation,
                "observe_path",
                side_effect=provisional_visual_observation,
            ),
        ):
            probe = records.Probe(self.source_root, RUN_AT, None, diagnostic=False)
            probe.extract(source)
        images = [item for item in probe.evidence if item["evidence_type"] == "image"]
        self.assertEqual(1, len(images))
        self.assertEqual("chart.png", images[0]["native_properties"]["attachment_name"])
        self.assertIn("attachment-image=1", images[0]["location"]["locator_text"])

    def test_notebook_code_data_uri_is_preserved_and_not_treated_as_displayed(self) -> None:
        encoded = base64.b64encode(PNG_BYTES).decode("ascii")
        source_text = f'payload = "data:image/png;base64,{encoded}"\n'
        source = self.source_root / "literal-data-uri.ipynb"
        source.write_text(
            json.dumps({
                "cells": [{
                    "cell_type": "code",
                    "execution_count": None,
                    "metadata": {},
                    "source": [source_text],
                    "outputs": [],
                }],
                "metadata": {},
                "nbformat": 4,
                "nbformat_minor": 5,
            }),
            encoding="utf-8",
        )
        probe = records.Probe(self.source_root, RUN_AT, None, diagnostic=False)
        probe.extract(source)
        notebook_cell = next(
            item for item in probe.evidence
            if item["evidence_type"] == "notebook_cell"
        )
        self.assertEqual(source_text, notebook_cell["content"]["raw_text"])
        self.assertFalse(any(
            item["evidence_type"] in {"image", "ocr_line"}
            for item in probe.evidence
        ))

    def test_ooxml_relationship_walker_excludes_orphans_and_binds_xlsx_cell(self) -> None:
        docx_path = self.work / "referenced.docx"
        with zipfile.ZipFile(docx_path, "w") as archive:
            archive.writestr("_rels/.rels", '''<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rIdOffice" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/></Relationships>''')
            archive.writestr("word/document.xml", '''<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:v="urn:schemas-microsoft-com:vml"><w:body name="rIdUnused"><a:blip r:embed="rIdImage"/><w:object r:id="rIdFakeImage"/><v:imagedata r:id="rIdVmlImage"/></w:body></w:document>''')
            archive.writestr("word/_rels/document.xml.rels", '''<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rIdImage" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image" Target="media/used.png"/><Relationship Id="rIdFakeImage" Type="http://example.invalid/not-an-image" Target="media/fake.png"/><Relationship Id="rIdVmlImage" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image" Target="media/vml.png"/><Relationship Id="rIdUnused" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image" Target="media/unreferenced.png"/></Relationships>''')
            archive.writestr("word/media/used.png", PNG_BYTES)
            archive.writestr("word/media/fake.png", PNG_BYTES)
            archive.writestr("word/media/vml.png", PNG_BYTES)
            archive.writestr("word/media/unreferenced.png", PNG_BYTES)
            archive.writestr("word/media/orphan.png", PNG_BYTES)
        with zipfile.ZipFile(docx_path) as archive:
            placements = records.referenced_ooxml_media(
                archive, media_prefixes=("word/media/",)
            )
        self.assertEqual(
            ["word/media/used.png", "word/media/vml.png"],
            [item["member"] for item in placements],
        )

        xlsx_path = self.work / "placed.xlsx"
        with zipfile.ZipFile(xlsx_path, "w") as archive:
            archive.writestr("_rels/.rels", '''<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rIdOffice" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="xl/workbook.xml"/></Relationships>''')
            archive.writestr("xl/workbook.xml", '''<workbook xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"><sheets><sheet name="担当表" sheetId="1" r:id="rIdSheet"/></sheets></workbook>''')
            archive.writestr("xl/_rels/workbook.xml.rels", '''<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rIdSheet" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/worksheet" Target="worksheets/sheet1.xml"/></Relationships>''')
            archive.writestr("xl/worksheets/sheet1.xml", '''<worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"><drawing r:id="rIdDrawing"/></worksheet>''')
            archive.writestr("xl/worksheets/_rels/sheet1.xml.rels", '''<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rIdDrawing" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/drawing" Target="../drawings/drawing1.xml"/></Relationships>''')
            archive.writestr("xl/drawings/drawing1.xml", '''<xdr:wsDr xmlns:xdr="http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"><xdr:oneCellAnchor><xdr:from><xdr:col>2</xdr:col><xdr:row>4</xdr:row></xdr:from><xdr:pic><a:blip r:embed="rIdImage"/></xdr:pic></xdr:oneCellAnchor><xdr:twoCellAnchor><xdr:from><xdr:col>1</xdr:col><xdr:row>2</xdr:row></xdr:from><xdr:graphicFrame><a:graphic><a:graphicData uri="http://schemas.openxmlformats.org/drawingml/2006/chart"><c:chart r:id="rIdChart"/></a:graphicData></a:graphic></xdr:graphicFrame></xdr:twoCellAnchor></xdr:wsDr>''')
            archive.writestr("xl/drawings/_rels/drawing1.xml.rels", '''<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rIdImage" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image" Target="../media/image1.png"/><Relationship Id="rIdChart" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/chart" Target="../charts/chart1.xml"/></Relationships>''')
            archive.writestr("xl/media/image1.png", PNG_BYTES)
            archive.writestr("xl/charts/chart1.xml", '''<c:chartSpace xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart"><c:chart><c:plotArea><c:barChart><c:ser><c:tx><c:strRef><c:strCache><c:pt idx="0"><c:v>受付件数</c:v></c:pt></c:strCache></c:strRef></c:tx><c:cat><c:strRef><c:f>集計!$A$2:$A$3</c:f><c:strCache><c:pt idx="0"><c:v>1月</c:v></c:pt><c:pt idx="1"><c:v>2月</c:v></c:pt></c:strCache></c:strRef></c:cat><c:val><c:numRef><c:f>集計!$B$2:$B$3</c:f><c:numCache><c:pt idx="0"><c:v>10</c:v></c:pt><c:pt idx="1"><c:v>20</c:v></c:pt></c:numCache></c:numRef></c:val></c:ser></c:barChart></c:plotArea></c:chart></c:chartSpace>''')
            archive.writestr("xl/charts/orphan.xml", '''<c:chartSpace xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart"><c:chart><c:title><c:v>表示されない図表</c:v></c:title></c:chart></c:chartSpace>''')
        with zipfile.ZipFile(xlsx_path) as archive:
            placements = records.referenced_ooxml_media(
                archive, media_prefixes=("xl/media/",)
            )
        self.assertEqual(1, len(placements))
        self.assertEqual("担当表", placements[0]["sheet_name"])
        self.assertEqual("C5", placements[0]["cell"])
        with zipfile.ZipFile(xlsx_path) as archive:
            chart_placements = records.referenced_ooxml_charts(archive)
            chart_payload = records._ooxml_chart_payload(
                archive.read(chart_placements[0]["member"]),
                chart_placements[0]["member"],
            )
        self.assertEqual(1, len(chart_placements))
        self.assertEqual("xl/charts/chart1.xml", chart_placements[0]["member"])
        self.assertEqual("担当表", chart_placements[0]["sheet_name"])
        self.assertEqual("B3", chart_placements[0]["cell"])
        self.assertEqual(
            [{"category": "1月", "value": "10"}, {"category": "2月", "value": "20"}],
            chart_payload["series"][0]["points"],
        )

    def test_xlsx_drawing_cells_bind_only_namespaced_relationship_uses(self) -> None:
        source = self.work / "drawing-cells.xlsx"
        with zipfile.ZipFile(source, "w") as archive:
            archive.writestr(
                "xl/drawings/drawing1.xml",
                '''<xdr:wsDr xmlns:xdr="http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"><xdr:oneCellAnchor name="rIdImage"><xdr:from><xdr:col>0</xdr:col><xdr:row>0</xdr:row></xdr:from><xdr:sp/></xdr:oneCellAnchor><xdr:oneCellAnchor><xdr:from><xdr:col>3</xdr:col><xdr:row>6</xdr:row></xdr:from><xdr:pic><a:blip r:embed="rIdImage"/><a:blip r:embed="rIdImage"/></xdr:pic></xdr:oneCellAnchor></xdr:wsDr>''',
            )
        with zipfile.ZipFile(source) as archive:
            cells = records._xlsx_drawing_cells(
                archive,
                "xl/drawings/drawing1.xml",
                "rIdImage",
            )
        self.assertEqual(["D7", "D7"], cells)

    def test_verified_ooxml_paths_reject_unknown_namespaces_carriers_and_types(self) -> None:
        wrong_rels = self.work / "wrong-rels.docx"
        with zipfile.ZipFile(wrong_rels, "w") as archive:
            archive.writestr(
                "_rels/.rels",
                '<evil:Relationships xmlns:evil="https://invalid.example/rels"/>',
            )
        with zipfile.ZipFile(wrong_rels) as archive:
            with self.assertRaisesRegex(
                ValueError, "ooxml_relationship_root_invalid"
            ):
                records._ooxml_relationships(archive)

        fake_use = self.work / "fake-use.docx"
        with zipfile.ZipFile(fake_use, "w") as archive:
            archive.writestr("_rels/.rels", '''<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rIdOffice" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/></Relationships>''')
            archive.writestr("word/document.xml", '''<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"><w:body><w:bogus r:id="rIdOfficial"/><a:blip r:embed="rIdEvilType"/></w:body></w:document>''')
            archive.writestr("word/_rels/document.xml.rels", '''<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rIdOfficial" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image" Target="media/official.png"/><Relationship Id="rIdEvilType" Type="https://invalid.example/image" Target="media/evil.png"/></Relationships>''')
            archive.writestr("word/media/official.png", PNG_BYTES)
            archive.writestr("word/media/evil.png", PNG_BYTES)
        with zipfile.ZipFile(fake_use) as archive:
            self.assertEqual([], records.referenced_ooxml_media(
                archive, media_prefixes=("word/media/",)
            ))

        hidden_custom = self.work / "hidden-custom.xlsx"
        with zipfile.ZipFile(hidden_custom, "w") as archive:
            archive.writestr("_rels/.rels", '''<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rIdOffice" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="xl/custom/item.xml"/></Relationships>''')
            archive.writestr("xl/custom/item.xml", '''<x:item xmlns:x="https://invalid.example/custom" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"><a:blip r:embed="rIdImage"/><c:chart r:id="rIdChart"/></x:item>''')
            archive.writestr("xl/custom/_rels/item.xml.rels", '''<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rIdImage" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image" Target="../media/poison.png"/><Relationship Id="rIdChart" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/chart" Target="../charts/poison.xml"/></Relationships>''')
            archive.writestr("xl/media/poison.png", PNG_BYTES)
            archive.writestr("xl/charts/poison.xml", '''<c:chartSpace xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart"><c:chart><c:plotArea><c:barChart><c:ser><c:cat><c:strCache><c:pt idx="0"><c:v>poison</c:v></c:pt></c:strCache></c:cat><c:val><c:numCache><c:pt idx="0"><c:v>99</c:v></c:pt></c:numCache></c:val></c:ser></c:barChart></c:plotArea></c:chart></c:chartSpace>''')
        with zipfile.ZipFile(hidden_custom) as archive:
            self.assertEqual([], records.referenced_ooxml_media(
                archive, media_prefixes=("xl/media/",)
            ))
            self.assertEqual([], records.referenced_ooxml_charts(archive))

        with self.assertRaisesRegex(ValueError, "ooxml_chart_root_invalid"):
            records._ooxml_chart_payload(
                b'<c:chartSpace xmlns:c="https://invalid.example/chart"><c:chart><c:ser><c:pt idx="0"><c:v>forged</c:v></c:pt></c:ser></c:chart></c:chartSpace>',
                "xl/charts/chart1.xml",
            )

        wrong_diagram = self.work / "wrong-diagram.pptx"
        with zipfile.ZipFile(wrong_diagram, "w") as archive:
            archive.writestr("_rels/.rels", '''<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rIdOffice" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="ppt/presentation.xml"/></Relationships>''')
            archive.writestr("ppt/presentation.xml", '''<p:presentation xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"><p:sldIdLst><p:sldId id="256" r:id="rIdSlide"/></p:sldIdLst></p:presentation>''')
            archive.writestr("ppt/_rels/presentation.xml.rels", '''<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rIdSlide" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="slides/slide1.xml"/></Relationships>''')
            archive.writestr("ppt/slides/slide1.xml", '''<p:sld xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:dgm="http://schemas.openxmlformats.org/drawingml/2006/diagram" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"><p:cSld><p:spTree><p:graphicFrame><a:graphic><a:graphicData uri="http://schemas.openxmlformats.org/drawingml/2006/diagram"><dgm:relIds r:dm="rIdDiagram"/></a:graphicData></a:graphic></p:graphicFrame></p:spTree></p:cSld></p:sld>''')
            archive.writestr("ppt/slides/_rels/slide1.xml.rels", '''<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rIdDiagram" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/diagramData" Target="../diagrams/data1.xml"/></Relationships>''')
            archive.writestr("ppt/diagrams/data1.xml", '''<evil:dataModel xmlns:evil="https://invalid.example/diagram"><evil:pt modelId="n1"><evil:t>forged</evil:t></evil:pt></evil:dataModel>''')
        with zipfile.ZipFile(wrong_diagram) as archive:
            with self.assertRaisesRegex(
                ValueError, "ooxml_diagram_data_root_invalid"
            ):
                records.referenced_pptx_diagrams(archive)

    def test_pptx_direct_visuals_are_projected_and_inherited_artwork_is_held(self) -> None:
        pptx_path = self.work / "visual-roles.pptx"
        with zipfile.ZipFile(pptx_path, "w") as archive:
            archive.writestr("_rels/.rels", '''<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rIdOffice" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="ppt/presentation.xml"/></Relationships>''')
            archive.writestr("ppt/presentation.xml", '''<p:presentation xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"><p:sldIdLst><p:sldId id="256" r:id="rIdSlide"/></p:sldIdLst></p:presentation>''')
            archive.writestr("ppt/_rels/presentation.xml.rels", '''<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rIdSlide" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="slides/slide1.xml"/></Relationships>''')
            archive.writestr("ppt/slides/slide1.xml", '''<p:sld xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:dgm="http://schemas.openxmlformats.org/drawingml/2006/diagram" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"><p:cSld><p:bg><p:bgPr><a:blipFill><a:blip r:embed="rIdBackground"/></a:blipFill></p:bgPr></p:bg><p:spTree><p:pic><p:blipFill><a:blip r:embed="rIdDirectPicture"/></p:blipFill></p:pic><p:sp><p:spPr><a:blipFill><a:blip r:embed="rIdShapeFill"/></a:blipFill></p:spPr></p:sp><p:graphicFrame><a:graphic><a:graphicData uri="http://schemas.openxmlformats.org/drawingml/2006/diagram"><dgm:relIds r:dm="rIdDiagram" r:lo="rIdDiagramLayout"/></a:graphicData></a:graphic></p:graphicFrame></p:spTree></p:cSld></p:sld>''')
            archive.writestr("ppt/slides/_rels/slide1.xml.rels", '''<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rIdBackground" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image" Target="../media/background.png"/><Relationship Id="rIdDirectPicture" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image" Target="../media/direct.png"/><Relationship Id="rIdShapeFill" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image" Target="../media/fill.png"/><Relationship Id="rIdDiagram" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/diagramData" Target="../diagrams/data1.xml"/><Relationship Id="rIdDiagramLayout" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/diagramLayout" Target="../diagrams/layout1.xml"/><Relationship Id="rIdLayout" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/></Relationships>''')
            archive.writestr("ppt/diagrams/data1.xml", '''<dgm:dataModel xmlns:dgm="http://schemas.openxmlformats.org/drawingml/2006/diagram" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"><dgm:ptLst><dgm:pt modelId="n1"><dgm:t><a:t>受付</a:t></dgm:t></dgm:pt><dgm:pt modelId="n2"><dgm:t><a:t>案内</a:t></dgm:t></dgm:pt></dgm:ptLst><dgm:cxnLst><dgm:cxn modelId="e1" srcId="n1" destId="n2" type="parOf"/></dgm:cxnLst></dgm:dataModel>''')
            archive.writestr("ppt/diagrams/layout1.xml", '''<dgm:layoutDef xmlns:dgm="http://schemas.openxmlformats.org/drawingml/2006/diagram" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"><dgm:pt modelId="fake"><a:t>本文ではないレイアウト文字</a:t></dgm:pt></dgm:layoutDef>''')
            archive.writestr("ppt/slideLayouts/slideLayout1.xml", '''<p:sldLayout xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"><p:cSld><p:spTree><p:pic><p:blipFill><a:blip r:embed="rIdLayoutPicture"/></p:blipFill></p:pic></p:spTree></p:cSld></p:sldLayout>''')
            archive.writestr("ppt/slideLayouts/_rels/slideLayout1.xml.rels", '''<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rIdLayoutPicture" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image" Target="../media/layout.png"/><Relationship Id="rIdMaster" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster" Target="../slideMasters/slideMaster1.xml"/></Relationships>''')
            archive.writestr("ppt/slideMasters/slideMaster1.xml", '''<p:sldMaster xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"><p:cSld><p:spTree><p:pic><p:blipFill><a:blip r:embed="rIdMasterPicture"/></p:blipFill></p:pic></p:spTree></p:cSld></p:sldMaster>''')
            archive.writestr("ppt/slideMasters/_rels/slideMaster1.xml.rels", '''<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rIdMasterPicture" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image" Target="../media/master.png"/></Relationships>''')
            for name in ("background.png", "direct.png", "fill.png", "layout.png", "master.png"):
                archive.writestr(f"ppt/media/{name}", PNG_BYTES)

        with zipfile.ZipFile(pptx_path) as archive:
            placements = records.referenced_ooxml_media(
                archive, media_prefixes=("ppt/media/",)
            )
            self.assertEqual(
                2,
                records._count_unresolved_pptx_inherited_media(archive),
            )
        self.assertEqual(
            ["background", "picture", "shape_fill"],
            [item["usage_kind"] for item in placements],
        )
        self.assertTrue(all(item["slide_number"] == 1 for item in placements))
        self.assertTrue(all(
            item["pptx_display_scope"] == "direct_slide"
            for item in placements
        ))
        with zipfile.ZipFile(pptx_path) as archive:
            diagrams = records.referenced_pptx_diagrams(archive)
        self.assertEqual(1, len(diagrams))
        self.assertEqual(["受付", "案内"], [
            item["text"] for item in diagrams[0]["points"]
        ])
        self.assertEqual(
            {"modelId": "e1", "srcId": "n1", "destId": "n2", "type": "parOf"},
            diagrams[0]["connections"][0],
        )

        source = self.source_root / "visual-roles.pptx"
        source.write_bytes(pptx_path.read_bytes())
        probe = records.Probe(self.source_root, RUN_AT, None, diagnostic=False)
        document = probe.add_document(source, "test-pptx-visual-role-projector")
        with (
            zipfile.ZipFile(source) as archive,
            mock.patch.object(
                probe, "_project_embedded_image_bytes", return_value=1
            ),
        ):
            projected = probe._project_ooxml_referenced_media(
                archive,
                document,
                media_prefixes=("ppt/media/",),
                skip_direct_pptx_pictures=True,
            )
        self.assertEqual(2, projected)
        images = [
            item for item in probe.evidence if item["evidence_type"] == "image"
        ]
        self.assertEqual(2, len(images))
        self.assertNotIn(
            "ppt/media/direct.png",
            {item["location"]["source_member"] for item in images},
        )
        self.assertNotIn(
            "ppt/media/layout.png",
            {item["location"]["source_member"] for item in images},
        )
        self.assertNotIn(
            "ppt/media/master.png",
            {item["location"]["source_member"] for item in images},
        )
        self.assertTrue(all(item["location"]["slide_number"] == 1 for item in images))
        self.assertEqual("partial", document["extraction"]["status"])
        self.assertTrue(any(
            "effective slide visibility is unresolved" in warning
            for warning in document["extraction"]["warnings"]
        ))

    def test_ooxml_verified_visuals_ignore_extlst_and_nested_binding_poison(self) -> None:
        pptx_path = self.work / "extlst-poison.pptx"
        with zipfile.ZipFile(pptx_path, "w") as archive:
            archive.writestr("_rels/.rels", '''<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rIdOffice" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="ppt/presentation.xml"/></Relationships>''')
            archive.writestr("ppt/presentation.xml", '''<p:presentation xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"><p:sldIdLst><p:sldId id="256" r:id="rIdSlide"/></p:sldIdLst><p:extLst><p:ext uri="poison"><p:sldId id="999" r:id="rIdPoisonSlide"/></p:ext></p:extLst></p:presentation>''')
            archive.writestr("ppt/_rels/presentation.xml.rels", '''<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rIdSlide" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="slides/slide1.xml"/><Relationship Id="rIdPoisonSlide" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="slides/poison.xml"/></Relationships>''')
            archive.writestr("ppt/slides/slide1.xml", '''<p:sld xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart" xmlns:dgm="http://schemas.openxmlformats.org/drawingml/2006/diagram" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"><p:cSld><p:spTree/></p:cSld><p:extLst><p:ext uri="poison"><p:graphicFrame><a:graphic><a:graphicData uri="http://schemas.openxmlformats.org/drawingml/2006/chart"><c:chart r:id="rIdPoisonChart"/></a:graphicData></a:graphic></p:graphicFrame><p:graphicFrame><a:graphic><a:graphicData uri="http://schemas.openxmlformats.org/drawingml/2006/diagram"><dgm:relIds r:dm="rIdPoisonDiagram"/></a:graphicData></a:graphic></p:graphicFrame></p:ext></p:extLst></p:sld>''')
            archive.writestr("ppt/slides/_rels/slide1.xml.rels", '''<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rIdPoisonChart" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/chart" Target="../charts/poison.xml"/><Relationship Id="rIdPoisonDiagram" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/diagramData" Target="../diagrams/poison.xml"/></Relationships>''')
            archive.writestr("ppt/slides/poison.xml", '''<p:sld xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"><p:cSld><p:spTree><p:graphicFrame><a:graphic><a:graphicData uri="http://schemas.openxmlformats.org/drawingml/2006/chart"><c:chart r:id="rIdNestedChart"/></a:graphicData></a:graphic></p:graphicFrame></p:spTree></p:cSld></p:sld>''')
            archive.writestr("ppt/slides/_rels/poison.xml.rels", '''<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rIdNestedChart" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/chart" Target="../charts/nested.xml"/></Relationships>''')
            for member in ("ppt/charts/poison.xml", "ppt/charts/nested.xml"):
                archive.writestr(member, '''<c:chartSpace xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart"><c:chart/></c:chartSpace>''')
            archive.writestr("ppt/diagrams/poison.xml", '''<dgm:dataModel xmlns:dgm="http://schemas.openxmlformats.org/drawingml/2006/diagram" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"><dgm:ptLst><dgm:pt modelId="poison"><dgm:t><a:t>poison</a:t></dgm:t></dgm:pt></dgm:ptLst></dgm:dataModel>''')
        with zipfile.ZipFile(pptx_path) as archive:
            self.assertEqual([], records.referenced_ooxml_charts(archive))
            self.assertEqual([], records.referenced_pptx_diagrams(archive))

        xlsx_path = self.work / "xlsx-extlst-poison.xlsx"
        with zipfile.ZipFile(xlsx_path, "w") as archive:
            archive.writestr("_rels/.rels", '''<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rIdOffice" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="xl/workbook.xml"/></Relationships>''')
            archive.writestr("xl/workbook.xml", '''<workbook xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"><sheets><sheet name="Visible" sheetId="1" r:id="rIdSheet"/></sheets><extLst><ext uri="poison"><sheet name="Poison" sheetId="2" r:id="rIdPoisonSheet"/></ext></extLst></workbook>''')
            archive.writestr("xl/_rels/workbook.xml.rels", '''<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rIdSheet" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/worksheet" Target="worksheets/sheet1.xml"/><Relationship Id="rIdPoisonSheet" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/worksheet" Target="worksheets/poison.xml"/></Relationships>''')
            archive.writestr("xl/worksheets/sheet1.xml", '''<worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"><drawing r:id="rIdDrawing"/></worksheet>''')
            archive.writestr("xl/worksheets/_rels/sheet1.xml.rels", '''<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rIdDrawing" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/drawing" Target="../drawings/drawing1.xml"/></Relationships>''')
            archive.writestr("xl/drawings/drawing1.xml", '''<xdr:wsDr xmlns:xdr="http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"><xdr:oneCellAnchor><xdr:from><xdr:col>0</xdr:col><xdr:row>0</xdr:row></xdr:from><xdr:extLst><xdr:ext><xdr:graphicFrame><a:graphic><a:graphicData uri="http://schemas.openxmlformats.org/drawingml/2006/chart"><c:chart r:id="rIdPoisonChart"/></a:graphicData></a:graphic></xdr:graphicFrame></xdr:ext></xdr:extLst></xdr:oneCellAnchor></xdr:wsDr>''')
            archive.writestr("xl/drawings/_rels/drawing1.xml.rels", '''<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rIdPoisonChart" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/chart" Target="../charts/poison.xml"/></Relationships>''')
            archive.writestr("xl/worksheets/poison.xml", '''<worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"><drawing r:id="rIdPoisonDrawing"/></worksheet>''')
            archive.writestr("xl/worksheets/_rels/poison.xml.rels", '''<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rIdPoisonDrawing" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/drawing" Target="../drawings/poison.xml"/></Relationships>''')
            archive.writestr("xl/drawings/poison.xml", '''<xdr:wsDr xmlns:xdr="http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"><xdr:oneCellAnchor><xdr:from><xdr:col>1</xdr:col><xdr:row>1</xdr:row></xdr:from><xdr:graphicFrame><a:graphic><a:graphicData uri="http://schemas.openxmlformats.org/drawingml/2006/chart"><c:chart r:id="rIdNestedChart"/></a:graphicData></a:graphic></xdr:graphicFrame></xdr:oneCellAnchor></xdr:wsDr>''')
            archive.writestr("xl/drawings/_rels/poison.xml.rels", '''<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rIdNestedChart" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/chart" Target="../charts/nested.xml"/></Relationships>''')
            for member in ("xl/charts/poison.xml", "xl/charts/nested.xml"):
                archive.writestr(member, '''<c:chartSpace xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart"><c:chart/></c:chartSpace>''')
        with zipfile.ZipFile(xlsx_path) as archive:
            self.assertEqual([], records.referenced_ooxml_charts(archive))

    def test_ooxml_direct_slide_and_sheet_bindings_are_unique(self) -> None:
        duplicate_slides = self.work / "duplicate-slides.pptx"
        with zipfile.ZipFile(duplicate_slides, "w") as archive:
            archive.writestr("ppt/presentation.xml", '''<p:presentation xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"><p:sldIdLst><p:sldId id="256" r:id="rIdSlide"/><p:sldId id="257" r:id="rIdSlide"/></p:sldIdLst></p:presentation>''')
            archive.writestr("ppt/_rels/presentation.xml.rels", '''<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rIdSlide" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="slides/slide1.xml"/></Relationships>''')
            archive.writestr("ppt/slides/slide1.xml", '''<p:sld xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:cSld><p:spTree/></p:cSld></p:sld>''')
        with zipfile.ZipFile(duplicate_slides) as archive:
            with self.assertRaisesRegex(ValueError, "ooxml_slide_binding_ambiguous"):
                records._pptx_slide_context(
                    archive, records._ooxml_relationships(archive)
                )

        duplicate_sheets = self.work / "duplicate-sheets.xlsx"
        with zipfile.ZipFile(duplicate_sheets, "w") as archive:
            archive.writestr("xl/workbook.xml", '''<workbook xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"><sheets><sheet name="One" sheetId="1" r:id="rIdOne"/><sheet name="Two" sheetId="2" r:id="rIdTwo"/></sheets></workbook>''')
            archive.writestr("xl/_rels/workbook.xml.rels", '''<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rIdOne" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/worksheet" Target="worksheets/sheet1.xml"/><Relationship Id="rIdTwo" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/worksheet" Target="worksheets/sheet1.xml"/></Relationships>''')
            archive.writestr("xl/worksheets/sheet1.xml", '''<worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main"/>''')
        with zipfile.ZipFile(duplicate_sheets) as archive:
            with self.assertRaisesRegex(ValueError, "ooxml_sheet_binding_invalid"):
                records._xlsx_sheet_context(
                    archive, records._ooxml_relationships(archive)
                )

    def test_chart_cache_pairs_only_identical_unique_point_indexes(self) -> None:
        def payload(category_cache: str, value_cache: str) -> dict[str, object]:
            raw = f'''<c:chartSpace xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart"><c:chart><c:plotArea><c:barChart><c:ser><c:cat>{category_cache}</c:cat><c:val>{value_cache}</c:val></c:ser></c:barChart></c:plotArea></c:chart></c:chartSpace>'''.encode()
            return records._ooxml_chart_payload(raw, "xl/charts/chart1.xml")

        mismatched = payload(
            '<c:strCache><c:pt idx="0"><c:v>A</c:v></c:pt><c:pt idx="2"><c:v>C</c:v></c:pt></c:strCache>',
            '<c:numCache><c:pt idx="0"><c:v>10</c:v></c:pt><c:pt idx="1"><c:v>20</c:v></c:pt></c:numCache>',
        )["series"][0]
        self.assertEqual("unresolved_cache_index_mismatch", mismatched["cache_status"])
        self.assertEqual([], mismatched["points"])

        duplicate = payload(
            '<c:strCache><c:pt idx="0"><c:v>A</c:v></c:pt><c:pt idx="0"><c:v>B</c:v></c:pt></c:strCache>',
            '<c:numCache><c:pt idx="0"><c:v>10</c:v></c:pt><c:pt idx="1"><c:v>20</c:v></c:pt></c:numCache>',
        )["series"][0]
        self.assertEqual("unresolved_cache_structure", duplicate["cache_status"])
        self.assertEqual([], duplicate["points"])

        multi_level = payload(
            '<c:multiLvlStrCache><c:lvl><c:pt idx="0"><c:v>東</c:v></c:pt></c:lvl><c:lvl><c:pt idx="0"><c:v>1月</c:v></c:pt></c:lvl></c:multiLvlStrCache>',
            '<c:numCache><c:pt idx="0"><c:v>10</c:v></c:pt><c:pt idx="1"><c:v>20</c:v></c:pt></c:numCache>',
        )["series"][0]
        self.assertEqual("unresolved_cache_structure", multi_level["cache_status"])
        self.assertEqual([], multi_level["points"])

    def test_native_pptx_chart_cache_reaches_validated_semantic_evidence(self) -> None:
        from pptx import Presentation
        from pptx.chart.data import ChartData
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.util import Inches

        source = self.source_root / "monthly-chart.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        chart_data = ChartData()
        chart_data.categories = ["2031", "2032"]
        chart_data.add_series("受付件数", (12, 18))
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(1), Inches(1), Inches(4), Inches(3),
            chart_data,
        ).chart
        chart.has_title = True
        chart.chart_title.text_frame.text = "月別受付件数"
        presentation.save(source)

        with zipfile.ZipFile(source) as archive:
            placements = records.referenced_ooxml_charts(archive)
        self.assertEqual(1, len(placements))
        self.assertEqual(1, placements[0]["slide_number"])
        self.assertTrue(placements[0]["member"].startswith("ppt/charts/"))

        probe = records.Probe(self.source_root, RUN_AT, None, diagnostic=False)
        probe.extract(source)
        chart_evidence = [
            item for item in probe.evidence
            if item["evidence_type"] in {"chart", "chart_series"}
        ]
        self.assertEqual(2, len(chart_evidence))
        self.assertTrue(all(
            item["provenance"]["extraction_method"]
            == "verified_ooxml_chart_cache"
            for item in chart_evidence
        ))
        series = next(
            item for item in chart_evidence
            if item["evidence_type"] == "chart_series"
        )
        self.assertIn("2031: 12", series["content"]["raw_text"])
        self.assertIn("2032: 18", series["content"]["raw_text"])

        intermediate = self.work / "chart-intermediate"
        search_output = self.work / "chart-search"
        semantic_output = self.work / "chart-semantic"
        materialize_intermediate(probe, self.source_root, intermediate)
        state = search_units.build(intermediate, search_output, 500)
        self.assertEqual(
            {"chart_series": 1, "chart_summary": 1},
            state["counts_by_type"],
        )
        validate_search_units_streaming.validate(search_output, intermediate)
        adapter.adapt(
            intermediate,
            self.source_root.resolve(),
            semantic_output,
            search_output,
        )
        semantic = [
            json.loads(line)
            for line in (semantic_output / "semantic-evidence.jsonl")
            .read_text(encoding="utf-8")
            .splitlines()
        ]
        projected_series = [
            item for item in semantic
            if item.get("adapter", {}).get("unit_type") == "chart_series"
        ]
        self.assertEqual(1, len(projected_series))
        self.assertIn("2032: 18", projected_series[0]["observed_text"])

    def test_native_xlsx_chart_uses_canonical_drawing_placement(self) -> None:
        from openpyxl import Workbook
        from openpyxl.chart import BarChart, Reference

        source = self.source_root / "native-chart.xlsx"
        workbook = Workbook()
        sheet = workbook.active
        sheet.title = "集計"
        sheet.append(["年", "受付件数"])
        sheet.append([2031, 12])
        sheet.append([2032, 18])
        chart = BarChart()
        chart.add_data(Reference(sheet, min_col=2, min_row=1, max_row=3), titles_from_data=True)
        chart.set_categories(Reference(sheet, min_col=1, min_row=2, max_row=3))
        sheet.add_chart(chart, "D2")
        workbook.save(source)

        with zipfile.ZipFile(source) as archive:
            placements = records.referenced_ooxml_charts(archive)
        self.assertEqual(1, len(placements))
        self.assertEqual("集計", placements[0]["sheet_name"])
        self.assertEqual("D2", placements[0]["cell"])

        probe = records.Probe(self.source_root, RUN_AT, None, diagnostic=False)
        probe.extract(source)
        chart_evidence = [
            item for item in probe.evidence
            if item["evidence_type"] in {"chart", "chart_series"}
        ]
        self.assertEqual(2, len(chart_evidence))
        self.assertTrue(all(
            item["provenance"]["extraction_method"]
            == "verified_ooxml_chart_cache"
            for item in chart_evidence
        ))
        self.assertTrue(all(
            item["location"].get("sheet_name") == "集計"
            and item["location"].get("cell") == "D2"
            for item in chart_evidence
        ))

    def test_pptx_smartart_text_and_raw_connections_reach_search(self) -> None:
        from pptx import Presentation

        base = self.work / "smartart-base.pptx"
        source = self.source_root / "smartart.pptx"
        presentation = Presentation()
        presentation.slides.add_slide(presentation.slide_layouts[6])
        presentation.save(base)
        with zipfile.ZipFile(base) as existing, zipfile.ZipFile(source, "w") as output:
            for info in existing.infolist():
                raw = existing.read(info.filename)
                if info.filename == "ppt/slides/slide1.xml":
                    value = raw.decode("utf-8").replace(
                        "</p:spTree>",
                        '<p:graphicFrame><p:nvGraphicFramePr><p:cNvPr id="2" name="SmartArt 1"/><p:cNvGraphicFramePr/><p:nvPr/></p:nvGraphicFramePr><p:xfrm><a:off xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" x="0" y="0"/><a:ext xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" cx="1" cy="1"/></p:xfrm><a:graphic xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"><a:graphicData uri="http://schemas.openxmlformats.org/drawingml/2006/diagram"><dgm:relIds xmlns:dgm="http://schemas.openxmlformats.org/drawingml/2006/diagram" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" r:dm="rIdSmart"/></a:graphicData></a:graphic></p:graphicFrame></p:spTree>',
                    ).replace(
                        "</p:sld>",
                        '<p:extLst><p:ext uri="poison"><dgm:relIds xmlns:dgm="http://schemas.openxmlformats.org/drawingml/2006/diagram" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" r:dm="rIdPoisonSmart"/></p:ext></p:extLst></p:sld>',
                    )
                    raw = value.encode("utf-8")
                elif info.filename == "ppt/slides/_rels/slide1.xml.rels":
                    value = raw.decode("utf-8").replace(
                        "</Relationships>",
                        '<Relationship Id="rIdSmart" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/diagramData" Target="../diagrams/data1.xml"/><Relationship Id="rIdPoisonSmart" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/diagramData" Target="../diagrams/poison.xml"/></Relationships>',
                    )
                    raw = value.encode("utf-8")
                output.writestr(info, raw)
            output.writestr(
                "ppt/diagrams/data1.xml",
                '''<dgm:dataModel xmlns:dgm="http://schemas.openxmlformats.org/drawingml/2006/diagram" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"><dgm:ptLst><dgm:pt modelId="n1"><dgm:t><a:t>受付</a:t></dgm:t></dgm:pt><dgm:pt modelId="n2"><dgm:t><a:t>案内</a:t></dgm:t></dgm:pt></dgm:ptLst><dgm:cxnLst><dgm:cxn modelId="e1" srcId="n1" destId="n2" type="parOf"/><dgm:cxn modelId="e2" srcId="n1" destId="n2" type="sibTrans"/></dgm:cxnLst></dgm:dataModel>''',
            )
            output.writestr(
                "ppt/diagrams/poison.xml",
                '''<dgm:dataModel xmlns:dgm="http://schemas.openxmlformats.org/drawingml/2006/diagram" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"><dgm:ptLst><dgm:pt modelId="poison"><dgm:t><a:t>extLst poison</a:t></dgm:t></dgm:pt></dgm:ptLst></dgm:dataModel>''',
            )

        probe = records.Probe(self.source_root, RUN_AT, None, diagnostic=False)
        probe.extract(source)
        smartart_text = [
            item["content"].get("raw_text")
            for item in probe.evidence
            if item.get("native_properties", {}).get("smartart_model_id")
        ]
        self.assertEqual(["受付", "案内"], smartart_text)
        self.assertNotIn("extLst poison", smartart_text)
        connection_text = [
            item for item in probe.evidence
            if "smartart_connection" in item.get("native_properties", {})
        ]
        self.assertEqual(2, len(connection_text))
        diagram_relations = [
            item for item in probe.relations
            if item["relation_type"] == "diagram_connection"
        ]
        self.assertEqual(1, len(diagram_relations))
        self.assertEqual(
            ["parOf", "sibTrans"],
            [
                item["type"]
                for item in diagram_relations[0]["properties"]["raw_connections"]
            ],
        )
        self.assertFalse(
            diagram_relations[0]["properties"]["semantic_interpretation_performed"]
        )

        intermediate = self.work / "smartart-intermediate"
        search_output = self.work / "smartart-search"
        materialize_intermediate(probe, self.source_root, intermediate)
        validate_intermediate_records_streaming.validate(
            intermediate,
            self.source_root,
            published_schema=(
                importlib.util.find_spec("jsonschema") is not None
            ),
        )
        search_units.build(intermediate, search_output, 500)
        validate_search_units_streaming.validate(search_output, intermediate)
        searchable = (search_output / "search_units.jsonl").read_text(
            encoding="utf-8"
        )
        self.assertIn("受付", searchable)
        self.assertIn("案内", searchable)
        self.assertIn("原形式type=parOf", searchable)
        self.assertIn("原形式type=sibTrans", searchable)

        inventory = self.work / "smartart-inventory.jsonl"
        source_stat = source.stat()
        write_jsonl(inventory, [{
            "relative_path": source.name,
            "kind": "file",
            "size_bytes": source_stat.st_size,
            "mtime_ns": source_stat.st_mtime_ns,
            "sha256": hashlib.sha256(source.read_bytes()).hexdigest(),
            "read_status": "observed",
        }])
        adaptive_output = self.work / "smartart-adaptive"
        adaptive_builder = load_engine("build_adaptive_semantic_graph")
        adaptive_validator = load_engine("validate_adaptive_semantic_graph")
        adaptive_builder.build(
            self.source_root, inventory, adaptive_output, SCRIPTS
        )
        validation = adaptive_validator.validate(
            adaptive_output, self.source_root, inventory
        )
        adaptive_relations_path = (
            adaptive_output / "layer1-intermediate" / "relations.jsonl"
        )
        adaptive_relations = [
            json.loads(line)
            for line in adaptive_relations_path.read_text(
                encoding="utf-8"
            ).splitlines()
            if line.strip()
        ]
        adaptive_diagram_relations = [
            item for item in adaptive_relations
            if item["relation_type"] == "diagram_connection"
        ]
        self.assertEqual(1, len(adaptive_diagram_relations))
        self.assertEqual(
            len(adaptive_relations), validation["structural_relations"]
        )

        def change_provenance(item: dict[str, object]) -> None:
            item["provenance"]["rule_or_model"] = "unverified inference"

        def change_source_member(item: dict[str, object]) -> None:
            item["properties"]["source_member"] = "ppt/diagrams/other.xml"

        def change_slide(item: dict[str, object]) -> None:
            item["properties"]["slide_number"] = 2

        def change_raw_connection(item: dict[str, object]) -> None:
            item["properties"]["raw_connections"][0]["type"] = "forged"

        def change_endpoints(item: dict[str, object]) -> None:
            item["from_ref"], item["to_ref"] = (
                item["to_ref"], item["from_ref"]
            )

        def change_support(item: dict[str, object]) -> None:
            item["supporting_evidence_ids"] = list(reversed(
                item["supporting_evidence_ids"]
            ))

        def remove_support(item: dict[str, object]) -> None:
            item["supporting_evidence_ids"] = (
                item["supporting_evidence_ids"][:-1]
            )

        tamper_cases = {
            "provenance": change_provenance,
            "source_member": change_source_member,
            "slide_number": change_slide,
            "raw_connections": change_raw_connection,
            "endpoints": change_endpoints,
            "supporting_evidence": change_support,
            "missing_supporting_evidence": remove_support,
        }
        diagram_relation_index = next(
            index for index, item in enumerate(adaptive_relations)
            if item["relation_type"] == "diagram_connection"
        )
        for label, tamper in tamper_cases.items():
            with self.subTest(adaptive_relation_tamper=label):
                changed = copy.deepcopy(adaptive_relations)
                tamper(changed[diagram_relation_index])
                write_jsonl(adaptive_relations_path, changed)
                with self.assertRaisesRegex(
                    ValueError, "native_structural_relations_mismatch"
                ):
                    adaptive_validator.validate(
                        adaptive_output, self.source_root, inventory
                    )
        write_jsonl(adaptive_relations_path, adaptive_relations)
        adaptive_validator.validate(
            adaptive_output, self.source_root, inventory
        )

        index_builder = load_engine("build_local_semantic_index")
        semantic_evidence = [
            json.loads(line)
            for line in (
                adaptive_output / "semantic-evidence.jsonl"
            ).read_text(encoding="utf-8").splitlines()
            if line.strip()
        ]
        semantic_evidence_by_id = {
            item["evidence_id"]: item for item in semantic_evidence
        }
        index_builder._validate_attested_smartart_connection(
            adaptive_diagram_relations[0]["relation_id"],
            adaptive_diagram_relations[0],
            semantic_evidence_by_id,
        )
        for label, tamper in tamper_cases.items():
            with self.subTest(graph_relation_tamper=label):
                changed_relation = copy.deepcopy(
                    adaptive_diagram_relations[0]
                )
                tamper(changed_relation)
                with self.assertRaisesRegex(ValueError, "graph_smartart_"):
                    index_builder._validate_attested_smartart_connection(
                        changed_relation["relation_id"],
                        changed_relation,
                        semantic_evidence_by_id,
                    )
        changed_evidence = copy.deepcopy(semantic_evidence_by_id)
        changed_source = changed_evidence[
            adaptive_diagram_relations[0]["from_ref"]["record_id"]
        ]
        changed_source["locator"]["object_id"] = "forged-model-id"
        changed_source["locator"]["locator_text"] = (
            changed_source["locator"]["locator_text"].rsplit(
                "point=", 1
            )[0]
            + "point=forged-model-id"
        )
        with self.assertRaisesRegex(ValueError, "graph_smartart_"):
            index_builder._validate_attested_smartart_connection(
                adaptive_diagram_relations[0]["relation_id"],
                adaptive_diagram_relations[0],
                changed_evidence,
            )

        security_output = self.work / "smartart-security"
        security_output.mkdir()
        security_builder = load_engine("content_security_gate")
        security_validator = load_engine("validate_content_security_gate")
        security_builder.build(
            adaptive_output / "semantic-evidence.jsonl",
            adaptive_output / "semantic-documents.jsonl",
            security_output,
            created_at=RUN_AT,
        )
        security_validator.validate(
            adaptive_output / "semantic-evidence.jsonl",
            adaptive_output / "semantic-documents.jsonl",
            security_output,
        )
        semantic_documents = [
            json.loads(line)
            for line in (
                adaptive_output / "semantic-documents.jsonl"
            ).read_text(encoding="utf-8").splitlines()
            if line.strip()
        ]
        safe_evidence = [
            json.loads(line)
            for line in (
                security_output / "safe-answer-evidence.jsonl"
            ).read_text(encoding="utf-8").splitlines()
            if line.strip()
        ]
        lineage_relations = [
            json.loads(line)
            for line in (
                adaptive_output / "semantic-lineage-relations.jsonl"
            ).read_text(encoding="utf-8").splitlines()
            if line.strip()
        ]
        all_relations = [*adaptive_relations, *lineage_relations]
        graph_context = {
            "lineage_context": {
                "output_dir": adaptive_output,
                "source_root": self.source_root,
                "inventory": inventory,
            },
            "security_context": {"gate_dir": security_output},
        }

        def insert_evidence(
            connection: sqlite3.Connection,
            evidence_records: list[dict[str, object]],
        ) -> None:
            for item in evidence_records:
                observed_text = str(item["observed_text"])
                connection.execute(
                    "INSERT INTO evidence VALUES (?, ?, ?, ?, ?, ?, ?, ?)",
                    (
                        item["evidence_id"],
                        item["document_id"],
                        item["source"]["relative_path"],
                        records.canonical_json(item["locator"]),
                        observed_text,
                        observed_text,
                        0,
                        hashlib.sha256(
                            observed_text.encode("utf-8")
                        ).hexdigest(),
                    ),
                )
            connection.commit()

        connection = sqlite3.connect(":memory:")
        try:
            index_builder.initialize(connection)
            insert_evidence(connection, safe_evidence)
            graph_report = index_builder.project_verified_structural_graph(
                connection,
                semantic_documents,
                safe_evidence,
                all_relations,
                **graph_context,
            )
            smartart_edges = connection.execute(
                "SELECT relation_id, from_node_id, to_node_id, basis_rule, "
                "properties_json FROM graph_edges "
                "WHERE relation_type = 'diagram_connection'"
            ).fetchall()
            self.assertEqual(1, len(smartart_edges))
            self.assertEqual(
                adaptive_diagram_relations[0]["relation_id"],
                smartart_edges[0][0],
            )
            self.assertEqual(
                "native SmartArt srcId/destId connection",
                smartart_edges[0][3],
            )
            self.assertEqual(
                adaptive_diagram_relations[0]["properties"],
                json.loads(smartart_edges[0][4]),
            )
            self.assertIn(
                adaptive_diagram_relations[0]["relation_id"],
                graph_report["security_partition"]["promoted_relation_ids"],
            )
        finally:
            connection.close()

        self_claimed = copy.deepcopy(adaptive_diagram_relations[0])
        self_claimed["provenance"]["generated_by"] = "self-claimed-smartart"
        self_claimed["relation_id"] = index_builder._stable_relation_id(
            self_claimed
        )
        held_connection = sqlite3.connect(":memory:")
        try:
            index_builder.initialize(held_connection)
            insert_evidence(held_connection, safe_evidence)
            held_report = index_builder.project_verified_structural_graph(
                held_connection,
                semantic_documents,
                safe_evidence,
                [*all_relations, self_claimed],
                **graph_context,
            )
            self.assertIn(
                self_claimed["relation_id"],
                held_report["skipped_relations"]["not_explicit"],
            )
            self.assertIsNone(held_connection.execute(
                "SELECT relation_id FROM graph_edges WHERE relation_id = ?",
                (self_claimed["relation_id"],),
            ).fetchone())
        finally:
            held_connection.close()

        mislabeled_containment = copy.deepcopy(next(
            item for item in adaptive_relations
            if item["relation_type"] == "contains"
            and item["from_ref"]["record_type"] == "document"
        ))
        mislabeled_containment["provenance"]["rule_or_model"] = (
            "native SmartArt srcId/destId connection"
        )
        mislabeled_containment["relation_id"] = (
            index_builder._stable_relation_id(mislabeled_containment)
        )
        boundary_connection = sqlite3.connect(":memory:")
        try:
            index_builder.initialize(boundary_connection)
            insert_evidence(boundary_connection, safe_evidence)
            boundary_report = index_builder.project_verified_structural_graph(
                boundary_connection,
                semantic_documents,
                safe_evidence,
                [mislabeled_containment],
            )
            self.assertIn(
                mislabeled_containment["relation_id"],
                boundary_report["skipped_relations"]["not_explicit"],
            )
            self.assertEqual(0, boundary_report["edge_count"])
        finally:
            boundary_connection.close()

    def test_embedded_visual_document_budget_stops_before_second_image(self) -> None:
        source = self.source_root / "budget.docx"
        source.write_bytes(b"container")
        probe = records.Probe(self.source_root, RUN_AT, None, diagnostic=False)
        document = probe.add_document(source, "test")
        with (
            mock.patch.object(records, "MAX_EMBEDDED_VISUALS_PER_DOCUMENT", 1),
            mock.patch.object(probe, "_project_local_image_evidence", return_value=1) as project,
        ):
            first = probe._project_embedded_image_bytes(
                PNG_BYTES, document, parent_id="ev_parent_1",
                location_prefix={"object_index": 1}, content_ref="first",
                source_name="first.png",
            )
            second = probe._project_embedded_image_bytes(
                PNG_BYTES, document, parent_id="ev_parent_2",
                location_prefix={"object_index": 2}, content_ref="second",
                source_name="second.png",
            )
        self.assertEqual((1, 0), (first, second))
        self.assertEqual(1, project.call_count)
        self.assertEqual("partial", document["extraction"]["status"])


if __name__ == "__main__":
    unittest.main()
